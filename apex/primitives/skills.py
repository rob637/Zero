"""
Skills — Local compute primitives that produce creative outputs.

Skills differ from Connectors: connectors talk to external services,
skills CREATE things locally using libraries + LLM intelligence.

Each skill is a Primitive. The planner sees them like any other tool.
"""

import json
import logging
import os
import platform
import subprocess
from datetime import datetime
from pathlib import Path
from typing import Any, Callable, Dict, List, Optional

from .base import Primitive, StepResult, get_data_index

logger = logging.getLogger(__name__)


def _output_dir() -> Path:
    """Get or create the output directory for skill artifacts."""
    if platform.system() == "Windows":
        base = Path.home() / "Documents" / "Ziggy"
    else:
        base = Path.home() / "Ziggy"
    base.mkdir(parents=True, exist_ok=True)
    return base


# ============================================================
#  PHOTO BOOK SKILL
# ============================================================

class PhotoBookSkill(Primitive):
    """Create beautiful PDF photo books from local photos.

    Flow:
    1. Search the local file index for photos matching a query
    2. Read EXIF metadata (location, date, camera)
    3. Ask LLM to generate captions for each photo
    4. Lay out photos + captions into a PDF photo book
    5. Return the PDF path for the user to open/print
    """

    def __init__(self, llm_complete: Optional[Callable] = None):
        self._llm = llm_complete

    @property
    def name(self) -> str:
        return "PHOTO_BOOK"

    def get_operations(self) -> Dict[str, str]:
        return {
            "create": "Create a PDF photo book from local photos. Search by theme, location, date, or event. AI generates captions. Output is a printable PDF.",
            "preview": "Preview what photos would be included in a photo book without generating it.",
        }

    def get_param_schema(self) -> Dict[str, Dict[str, Any]]:
        return {
            "create": {
                "query": {"type": "str", "required": True, "description": "What photos to include: a theme, location, date range, or event (e.g. 'Africa vacation', 'summer 2024', 'family beach photos')"},
                "title": {"type": "str", "required": False, "description": "Book title. Auto-generated from query if not provided."},
                "max_photos": {"type": "int", "required": False, "description": "Maximum photos to include (default: 20)"},
                "layout": {"type": "str", "required": False, "description": "Layout style: 'classic' (1 photo/page), 'collage' (2-3 photos/page), 'magazine' (mixed). Default: classic"},
            },
            "preview": {
                "query": {"type": "str", "required": True, "description": "Search query for photos"},
                "max_photos": {"type": "int", "required": False, "description": "Maximum photos to preview (default: 20)"},
            },
        }

    async def execute(self, operation: str, params: Dict[str, Any]) -> StepResult:
        if operation == "create":
            return await self._create_book(params)
        elif operation == "preview":
            return await self._preview(params)
        return StepResult(False, error=f"Unknown operation: {operation}")

    async def _find_photos(self, query: str, max_photos: int = 20) -> List[Dict[str, Any]]:
        """Search the local file index for photos matching the query."""
        index = get_data_index()
        if not index:
            return []

        # Search using FTS
        results = index.search(query, source="local_files", limit=max_photos * 3)

        # Filter to image files only
        image_exts = {".jpg", ".jpeg", ".png", ".tiff", ".heic", ".heif", ".webp"}
        photos = []
        for obj in results:
            path = obj.get("source_id", "") or obj.get("raw", {}).get("path", "")
            ext = os.path.splitext(path)[1].lower()
            if ext in image_exts and os.path.isfile(path):
                photos.append({
                    "path": path,
                    "name": os.path.basename(path),
                    "body": obj.get("body", ""),
                    "timestamp": obj.get("timestamp", ""),
                })
                if len(photos) >= max_photos:
                    break

        return photos

    async def _generate_captions(self, photos: List[Dict]) -> List[str]:
        """Ask LLM to generate captions for photos based on metadata."""
        if not self._llm or not photos:
            return [p.get("name", "") for p in photos]

        photo_descriptions = []
        for i, p in enumerate(photos):
            photo_descriptions.append(f"Photo {i+1}: {p['name']}\nMetadata: {p.get('body', 'No metadata')}")

        prompt = f"""Generate short, engaging photo captions for a photo book.

{chr(10).join(photo_descriptions)}

For each photo, write a 1-2 sentence caption that would look great in a printed photo book.
Use the metadata (location, date, camera) to add context when available.
Be warm and descriptive, not generic.

Return a JSON array of strings, one caption per photo. Return ONLY the JSON array."""

        try:
            response = await self._llm(prompt)
            response = response.strip()
            if response.startswith("```"):
                response = response.split("\n", 1)[1] if "\n" in response else response[3:]
                response = response.rsplit("```", 1)[0]
            captions = json.loads(response)
            if isinstance(captions, list) and len(captions) == len(photos):
                return captions
        except Exception as e:
            logger.warning(f"Caption generation failed: {e}")

        # Fallback: use filename
        return [p.get("name", f"Photo {i+1}") for i, p in enumerate(photos)]

    async def _create_book(self, params: Dict[str, Any]) -> StepResult:
        """Create a PDF photo book."""
        query = params.get("query", "")
        if not query:
            return StepResult(False, error="Please specify what photos to include (e.g. 'Africa vacation', 'summer 2024')")

        title = params.get("title", "")
        max_photos = min(params.get("max_photos", 20), 100)
        layout = params.get("layout", "classic")

        # 1. Find photos
        photos = await self._find_photos(query, max_photos)
        if not photos:
            return StepResult(False, error=f"No photos found matching '{query}'. Make sure your Pictures folder is indexed.")

        # 2. Generate captions
        captions = await self._generate_captions(photos)

        # 3. Generate title if not provided
        if not title:
            title = query.title()

        # 4. Build PDF
        try:
            output_path = self._build_pdf(photos, captions, title, layout)
        except ImportError:
            return StepResult(False, error="PDF generation requires the fpdf2 library. Run: pip install fpdf2")
        except Exception as e:
            return StepResult(False, error=f"Failed to create photo book: {e}")

        # 5. Try to open the file
        try:
            if platform.system() == "Windows":
                os.startfile(output_path)
            elif platform.system() == "Darwin":
                subprocess.Popen(["open", output_path])
            else:
                subprocess.Popen(["xdg-open", output_path])
        except Exception:
            pass

        return StepResult(True, data={
            "file": str(output_path),
            "title": title,
            "photo_count": len(photos),
            "message": f"Photo book '{title}' created with {len(photos)} photos → {output_path}"
        })

    def _build_pdf(
        self,
        photos: List[Dict],
        captions: List[str],
        title: str,
        layout: str,
    ) -> str:
        """Generate the PDF photo book."""
        from fpdf import FPDF

        pdf = FPDF(orientation="P", unit="mm", format="A4")
        pdf.set_auto_page_break(auto=False)

        # --- Title Page ---
        pdf.add_page()
        pdf.set_fill_color(20, 20, 20)
        pdf.rect(0, 0, 210, 297, "F")
        pdf.set_text_color(255, 255, 255)
        pdf.set_font("Helvetica", "B", 36)
        pdf.set_y(120)
        pdf.cell(0, 20, title, align="C", new_x="LMARGIN", new_y="NEXT")
        pdf.set_font("Helvetica", "", 14)
        pdf.set_text_color(180, 180, 180)
        pdf.cell(0, 10, datetime.now().strftime("%B %Y"), align="C", new_x="LMARGIN", new_y="NEXT")
        pdf.cell(0, 10, f"{len(photos)} photos", align="C")

        # --- Photo Pages ---
        if layout == "collage":
            self._layout_collage(pdf, photos, captions)
        else:
            self._layout_classic(pdf, photos, captions)

        # Save
        out_dir = _output_dir()
        safe_title = "".join(c if c.isalnum() or c in " -_" else "" for c in title).strip()
        filename = f"PhotoBook - {safe_title or 'Untitled'}.pdf"
        output_path = str(out_dir / filename)
        pdf.output(output_path)
        return output_path

    def _layout_classic(self, pdf, photos, captions):
        """One photo per page with caption below."""
        for i, (photo, caption) in enumerate(zip(photos, captions)):
            pdf.add_page()
            pdf.set_fill_color(255, 255, 255)
            pdf.rect(0, 0, 210, 297, "F")

            # Photo — centered, max width 180mm, max height 220mm
            try:
                from PIL import Image
                with Image.open(photo["path"]) as img:
                    w, h = img.size
                    aspect = w / h

                    max_w, max_h = 180, 210
                    if aspect > max_w / max_h:
                        img_w = max_w
                        img_h = max_w / aspect
                    else:
                        img_h = max_h
                        img_w = max_h * aspect

                    x = (210 - img_w) / 2
                    y = 15
                    pdf.image(photo["path"], x=x, y=y, w=img_w, h=img_h)

                    # Caption below photo
                    caption_y = y + img_h + 8
            except Exception:
                caption_y = 230

            pdf.set_y(caption_y)
            pdf.set_text_color(60, 60, 60)
            pdf.set_font("Helvetica", "I", 11)
            pdf.multi_cell(0, 6, caption, align="C")

            # Page number
            pdf.set_y(285)
            pdf.set_text_color(160, 160, 160)
            pdf.set_font("Helvetica", "", 8)
            pdf.cell(0, 5, str(i + 1), align="C")

    def _layout_collage(self, pdf, photos, captions):
        """Two photos per page with captions."""
        pairs = list(zip(photos, captions))
        for i in range(0, len(pairs), 2):
            pdf.add_page()
            pdf.set_fill_color(255, 255, 255)
            pdf.rect(0, 0, 210, 297, "F")

            for slot, idx in enumerate([i, i + 1]):
                if idx >= len(pairs):
                    break
                photo, caption = pairs[idx]
                y_start = 10 + slot * 140

                try:
                    from PIL import Image
                    with Image.open(photo["path"]) as img:
                        w, h = img.size
                        aspect = w / h
                        max_w, max_h = 180, 115
                        if aspect > max_w / max_h:
                            img_w = max_w
                            img_h = max_w / aspect
                        else:
                            img_h = max_h
                            img_w = max_h * aspect

                        x = (210 - img_w) / 2
                        pdf.image(photo["path"], x=x, y=y_start, w=img_w, h=img_h)
                        caption_y = y_start + img_h + 3
                except Exception:
                    caption_y = y_start + 115

                pdf.set_y(caption_y)
                pdf.set_text_color(60, 60, 60)
                pdf.set_font("Helvetica", "I", 10)
                pdf.multi_cell(0, 5, caption, align="C")

            # Page number
            pdf.set_y(285)
            pdf.set_text_color(160, 160, 160)
            pdf.set_font("Helvetica", "", 8)
            pdf.cell(0, 5, str(i // 2 + 1), align="C")

    async def _preview(self, params: Dict[str, Any]) -> StepResult:
        """Preview what photos would be in the book."""
        query = params.get("query", "")
        if not query:
            return StepResult(False, error="Please specify a search query")

        max_photos = min(params.get("max_photos", 20), 100)
        photos = await self._find_photos(query, max_photos)

        if not photos:
            return StepResult(False, error=f"No photos found matching '{query}'")

        return StepResult(True, data={
            "query": query,
            "photo_count": len(photos),
            "photos": [{"name": p["name"], "path": p["path"], "metadata": p.get("body", "")} for p in photos],
            "message": f"Found {len(photos)} photos matching '{query}'. Use photo_book create to generate the PDF."
        })


# ============================================================
#  REPORT SKILL
# ============================================================

class ReportSkill(Primitive):
    """Generate PDF/DOCX reports from data and natural language descriptions.

    Pull data from any connected service, generate charts and narrative,
    output as a professional PDF report.
    """

    def __init__(self, llm_complete: Optional[Callable] = None):
        self._llm = llm_complete

    @property
    def name(self) -> str:
        return "REPORT"

    def get_operations(self) -> Dict[str, str]:
        return {
            "create": "Generate a PDF report. Provide data, a title, and optionally section descriptions. AI writes the narrative and creates charts.",
        }

    def get_param_schema(self) -> Dict[str, Dict[str, Any]]:
        return {
            "create": {
                "title": {"type": "str", "required": True, "description": "Report title"},
                "sections": {"type": "list", "required": True, "description": "List of section dicts with 'heading' and 'content' (text, data table, or chart spec)"},
                "data": {"type": "dict", "required": False, "description": "Supporting data (tables, metrics) to include"},
            },
        }

    async def execute(self, operation: str, params: Dict[str, Any]) -> StepResult:
        if operation != "create":
            return StepResult(False, error=f"Unknown operation: {operation}")

        title = params.get("title", "Report")
        sections = params.get("sections", [])

        if not sections:
            return StepResult(False, error="Provide at least one section with 'heading' and 'content'")

        try:
            from fpdf import FPDF
        except ImportError:
            return StepResult(False, error="Report generation requires fpdf2. Run: pip install fpdf2")

        try:
            pdf = FPDF(orientation="P", unit="mm", format="A4")
            pdf.set_auto_page_break(auto=True, margin=20)

            # Title page
            pdf.add_page()
            pdf.set_font("Helvetica", "B", 28)
            pdf.set_y(100)
            pdf.cell(0, 15, title, align="C", new_x="LMARGIN", new_y="NEXT")
            pdf.set_font("Helvetica", "", 12)
            pdf.set_text_color(120, 120, 120)
            pdf.cell(0, 10, datetime.now().strftime("%B %d, %Y"), align="C")

            # Content pages
            for section in sections:
                pdf.add_page()
                heading = section.get("heading", "")
                content = section.get("content", "")

                pdf.set_text_color(0, 0, 0)
                pdf.set_font("Helvetica", "B", 18)
                pdf.cell(0, 12, heading, new_x="LMARGIN", new_y="NEXT")
                pdf.ln(4)

                pdf.set_font("Helvetica", "", 11)
                if isinstance(content, str):
                    pdf.multi_cell(0, 6, content)
                elif isinstance(content, list):
                    # Render as table
                    if content and isinstance(content[0], dict):
                        headers = list(content[0].keys())
                        col_w = 190 / len(headers)
                        pdf.set_font("Helvetica", "B", 10)
                        for h in headers:
                            pdf.cell(col_w, 8, str(h), border=1)
                        pdf.ln()
                        pdf.set_font("Helvetica", "", 9)
                        for row in content[:50]:
                            for h in headers:
                                pdf.cell(col_w, 7, str(row.get(h, "")), border=1)
                            pdf.ln()

            # Save
            out_dir = _output_dir()
            safe = "".join(c if c.isalnum() or c in " -_" else "" for c in title).strip()
            path = str(out_dir / f"Report - {safe or 'Untitled'}.pdf")
            pdf.output(path)

            try:
                if platform.system() == "Windows":
                    os.startfile(path)
                elif platform.system() == "Darwin":
                    subprocess.Popen(["open", path])
            except Exception:
                pass

            return StepResult(True, data={
                "file": path,
                "title": title,
                "sections": len(sections),
                "message": f"Report '{title}' generated → {path}"
            })
        except Exception as e:
            return StepResult(False, error=f"Report generation failed: {e}")


# ============================================================
#  DATA VIZ SKILL
# ============================================================

class DataVizSkill(Primitive):
    """Create charts and visualizations from data."""

    def __init__(self, llm_complete: Optional[Callable] = None):
        self._llm = llm_complete

    @property
    def name(self) -> str:
        return "DATA_VIZ"

    def get_operations(self) -> Dict[str, str]:
        return {
            "chart": "Create a chart (bar, line, pie, scatter) from data. AI picks the best visualization if not specified.",
        }

    def get_param_schema(self) -> Dict[str, Dict[str, Any]]:
        return {
            "chart": {
                "data": {"type": "list", "required": True, "description": "Data to visualize — list of dicts or list of numbers"},
                "chart_type": {"type": "str", "required": False, "description": "Chart type: bar, line, pie, scatter, area. Auto-detected if omitted."},
                "title": {"type": "str", "required": False, "description": "Chart title"},
                "x_field": {"type": "str", "required": False, "description": "Field name for X axis"},
                "y_field": {"type": "str", "required": False, "description": "Field name for Y axis"},
            },
        }

    async def execute(self, operation: str, params: Dict[str, Any]) -> StepResult:
        if operation != "chart":
            return StepResult(False, error=f"Unknown operation: {operation}")

        try:
            import matplotlib
            matplotlib.use("Agg")
            import matplotlib.pyplot as plt
        except ImportError:
            return StepResult(False, error="Charts require matplotlib. Run: pip install matplotlib")

        data = params.get("data", [])
        if not data:
            return StepResult(False, error="No data provided")

        chart_type = params.get("chart_type", "bar")
        title = params.get("title", "Chart")
        x_field = params.get("x_field")
        y_field = params.get("y_field")

        try:
            fig, ax = plt.subplots(figsize=(10, 6))

            if isinstance(data[0], dict):
                keys = list(data[0].keys())
                x_key = x_field or keys[0]
                y_key = y_field or (keys[1] if len(keys) > 1 else keys[0])
                x_vals = [d.get(x_key, "") for d in data]
                y_vals = [float(d.get(y_key, 0)) for d in data]
            else:
                x_vals = list(range(len(data)))
                y_vals = [float(v) for v in data]
                x_key, y_key = "Index", "Value"

            if chart_type == "pie":
                ax.pie(y_vals, labels=x_vals, autopct="%1.1f%%")
            elif chart_type == "line":
                ax.plot(x_vals, y_vals, marker="o")
                ax.set_xlabel(x_key)
                ax.set_ylabel(y_key)
            elif chart_type == "scatter":
                ax.scatter(x_vals, y_vals)
                ax.set_xlabel(x_key)
                ax.set_ylabel(y_key)
            else:
                ax.bar(range(len(x_vals)), y_vals, tick_label=[str(x) for x in x_vals])
                ax.set_xlabel(x_key)
                ax.set_ylabel(y_key)
                plt.xticks(rotation=45, ha="right")

            ax.set_title(title)
            plt.tight_layout()

            out_dir = _output_dir()
            safe = "".join(c if c.isalnum() or c in " -_" else "" for c in title).strip()
            path = str(out_dir / f"Chart - {safe or 'chart'}.png")
            fig.savefig(path, dpi=150)
            plt.close(fig)

            try:
                if platform.system() == "Windows":
                    os.startfile(path)
                elif platform.system() == "Darwin":
                    subprocess.Popen(["open", path])
            except Exception:
                pass

            return StepResult(True, data={
                "file": path,
                "title": title,
                "chart_type": chart_type,
                "data_points": len(data),
                "message": f"Chart '{title}' saved → {path}"
            })
        except Exception as e:
            return StepResult(False, error=f"Chart creation failed: {e}")


# ============================================================
#  FILE CONVERTER SKILL
# ============================================================

class FileConverterSkill(Primitive):
    """Convert files between formats."""

    @property
    def name(self) -> str:
        return "CONVERT"

    def get_operations(self) -> Dict[str, str]:
        return {
            "convert": "Convert a file to a different format. Supports: images (resize, format change), CSV↔XLSX, text→PDF.",
        }

    def get_param_schema(self) -> Dict[str, Dict[str, Any]]:
        return {
            "convert": {
                "input_path": {"type": "str", "required": True, "description": "Path to the input file"},
                "output_format": {"type": "str", "required": True, "description": "Target format: pdf, png, jpg, xlsx, csv, txt"},
                "resize": {"type": "str", "required": False, "description": "For images: resize to WxH (e.g. '800x600' or '50%')"},
            },
        }

    async def execute(self, operation: str, params: Dict[str, Any]) -> StepResult:
        if operation != "convert":
            return StepResult(False, error=f"Unknown operation: {operation}")

        input_path = params.get("input_path", "")
        output_format = params.get("output_format", "").lower().strip(".")
        resize = params.get("resize")

        if not input_path or not os.path.isfile(input_path):
            return StepResult(False, error=f"File not found: {input_path}")
        if not output_format:
            return StepResult(False, error="Specify output_format (e.g. 'pdf', 'png', 'csv')")

        name = os.path.splitext(os.path.basename(input_path))[0]
        out_dir = _output_dir()
        output_path = str(out_dir / f"{name}.{output_format}")
        ext = os.path.splitext(input_path)[1].lower()

        try:
            # Image conversions
            if ext in (".jpg", ".jpeg", ".png", ".bmp", ".tiff", ".webp", ".gif"):
                from PIL import Image
                with Image.open(input_path) as img:
                    if resize:
                        if "%" in resize:
                            pct = float(resize.strip("%")) / 100
                            img = img.resize((int(img.width * pct), int(img.height * pct)))
                        elif "x" in resize.lower():
                            w, h = resize.lower().split("x")
                            img = img.resize((int(w), int(h)))
                    if output_format == "pdf":
                        rgb = img.convert("RGB")
                        rgb.save(output_path)
                    else:
                        img.save(output_path)
                return StepResult(True, data={"file": output_path, "message": f"Converted → {output_path}"})

            # CSV → XLSX
            if ext == ".csv" and output_format == "xlsx":
                import csv
                from io import StringIO
                with open(input_path, "r", encoding="utf-8") as f:
                    reader = list(csv.reader(f))
                # Use fpdf2 or openpyxl if available
                try:
                    import openpyxl
                    wb = openpyxl.Workbook()
                    ws = wb.active
                    for row in reader:
                        ws.append(row)
                    wb.save(output_path)
                    return StepResult(True, data={"file": output_path, "message": f"Converted → {output_path}"})
                except ImportError:
                    return StepResult(False, error="XLSX conversion requires openpyxl. Run: pip install openpyxl")

            # Text/MD → PDF
            if ext in (".txt", ".md", ".log") and output_format == "pdf":
                from fpdf import FPDF
                pdf = FPDF()
                pdf.add_page()
                pdf.set_font("Helvetica", "", 11)
                with open(input_path, "r", encoding="utf-8", errors="replace") as f:
                    for line in f:
                        pdf.multi_cell(0, 6, line.rstrip())
                pdf.output(output_path)
                return StepResult(True, data={"file": output_path, "message": f"Converted → {output_path}"})

            return StepResult(False, error=f"Unsupported conversion: {ext} → .{output_format}")
        except Exception as e:
            return StepResult(False, error=f"Conversion failed: {e}")


# ============================================================
#  EXPENSE REPORT SKILL
# ============================================================

class ExpenseReportSkill(Primitive):
    """Generate expense reports from CSV/data."""

    def __init__(self, llm_complete: Optional[Callable] = None):
        self._llm = llm_complete

    @property
    def name(self) -> str:
        return "EXPENSE_REPORT"

    def get_operations(self) -> Dict[str, str]:
        return {
            "create": "Create an expense report PDF from a CSV file or list of expenses. Categories are auto-detected by AI.",
        }

    def get_param_schema(self) -> Dict[str, Dict[str, Any]]:
        return {
            "create": {
                "source": {"type": "str", "required": False, "description": "Path to CSV/spreadsheet file with expenses"},
                "expenses": {"type": "list", "required": False, "description": "List of expense dicts with 'date', 'description', 'amount', optionally 'category'"},
                "title": {"type": "str", "required": False, "description": "Report title (default: 'Expense Report')"},
                "period": {"type": "str", "required": False, "description": "Reporting period (e.g. 'March 2026')"},
            },
        }

    async def execute(self, operation: str, params: Dict[str, Any]) -> StepResult:
        if operation != "create":
            return StepResult(False, error=f"Unknown operation: {operation}")

        try:
            from fpdf import FPDF
        except ImportError:
            return StepResult(False, error="Requires fpdf2. Run: pip install fpdf2")

        expenses = params.get("expenses", [])
        source = params.get("source", "")
        title = params.get("title", "Expense Report")
        period = params.get("period", datetime.now().strftime("%B %Y"))

        # Load from CSV if provided
        if source and os.path.isfile(source):
            import csv
            with open(source, "r", encoding="utf-8") as f:
                reader = csv.DictReader(f)
                for row in reader:
                    expenses.append({
                        "date": row.get("date", row.get("Date", "")),
                        "description": row.get("description", row.get("Description", row.get("vendor", row.get("Vendor", "")))),
                        "amount": row.get("amount", row.get("Amount", "0")),
                        "category": row.get("category", row.get("Category", "")),
                    })

        if not expenses:
            return StepResult(False, error="No expenses provided. Pass a CSV file path or a list of expense dicts.")

        # Calculate totals
        total = sum(float(e.get("amount", 0)) for e in expenses)
        categories = {}
        for e in expenses:
            cat = e.get("category", "Uncategorized")
            categories[cat] = categories.get(cat, 0) + float(e.get("amount", 0))

        # Build PDF
        pdf = FPDF()
        pdf.add_page()
        pdf.set_font("Helvetica", "B", 22)
        pdf.cell(0, 12, title, new_x="LMARGIN", new_y="NEXT")
        pdf.set_font("Helvetica", "", 12)
        pdf.set_text_color(100, 100, 100)
        pdf.cell(0, 8, f"Period: {period}", new_x="LMARGIN", new_y="NEXT")
        pdf.cell(0, 8, f"Total: ${total:,.2f}", new_x="LMARGIN", new_y="NEXT")
        pdf.ln(5)

        # Category summary
        pdf.set_text_color(0, 0, 0)
        pdf.set_font("Helvetica", "B", 14)
        pdf.cell(0, 10, "Summary by Category", new_x="LMARGIN", new_y="NEXT")
        pdf.set_font("Helvetica", "", 11)
        for cat, amt in sorted(categories.items(), key=lambda x: -x[1]):
            pdf.cell(120, 7, cat, border="B")
            pdf.cell(0, 7, f"${amt:,.2f}", border="B", align="R", new_x="LMARGIN", new_y="NEXT")
        pdf.ln(5)

        # Line items
        pdf.set_font("Helvetica", "B", 14)
        pdf.cell(0, 10, "Line Items", new_x="LMARGIN", new_y="NEXT")
        pdf.set_font("Helvetica", "B", 9)
        pdf.cell(30, 7, "Date", border=1)
        pdf.cell(80, 7, "Description", border=1)
        pdf.cell(40, 7, "Category", border=1)
        pdf.cell(30, 7, "Amount", border=1, align="R", new_x="LMARGIN", new_y="NEXT")
        pdf.set_font("Helvetica", "", 9)
        for e in expenses:
            pdf.cell(30, 6, str(e.get("date", "")), border=1)
            pdf.cell(80, 6, str(e.get("description", ""))[:45], border=1)
            pdf.cell(40, 6, str(e.get("category", "")), border=1)
            pdf.cell(30, 6, f"${float(e.get('amount', 0)):,.2f}", border=1, align="R", new_x="LMARGIN", new_y="NEXT")

        out_dir = _output_dir()
        safe = "".join(c if c.isalnum() or c in " -_" else "" for c in title).strip()
        path = str(out_dir / f"{safe or 'Expense Report'} - {period}.pdf")
        pdf.output(path)

        try:
            if platform.system() == "Windows":
                os.startfile(path)
        except Exception:
            pass

        return StepResult(True, data={
            "file": path,
            "total": total,
            "expense_count": len(expenses),
            "categories": categories,
            "message": f"Expense report generated: {len(expenses)} items, ${total:,.2f} total → {path}"
        })


# ============================================================
#  PRESENTATION BUILDER SKILL
# ============================================================

class PresentationBuilderSkill(Primitive):
    """Generate PPTX slide decks locally from a topic or outline.

    AI generates slide content, python-pptx renders it.
    No cloud API needed — works fully offline.
    """

    def __init__(self, llm_complete: Optional[Callable] = None):
        self._llm = llm_complete

    @property
    def name(self) -> str:
        return "SLIDE_DECK"

    def get_operations(self) -> Dict[str, str]:
        return {
            "create": "Generate a PowerPoint slide deck from a topic, outline, or data. AI writes the content for each slide.",
        }

    def get_param_schema(self) -> Dict[str, Dict[str, Any]]:
        return {
            "create": {
                "topic": {"type": "str", "required": True, "description": "Presentation topic or title (e.g. 'Q1 Sales Review', 'Project Kickoff')"},
                "slides": {"type": "int", "required": False, "description": "Number of slides (default: 8)"},
                "outline": {"type": "list", "required": False, "description": "Optional list of slide titles/topics to cover"},
                "audience": {"type": "str", "required": False, "description": "Target audience (e.g. 'executive team', 'engineering', 'customers')"},
            },
        }

    async def execute(self, operation: str, params: Dict[str, Any]) -> StepResult:
        if operation != "create":
            return StepResult(False, error=f"Unknown operation: {operation}")

        try:
            from pptx import Presentation
            from pptx.util import Inches, Pt
            from pptx.enum.text import PP_ALIGN
        except ImportError:
            return StepResult(False, error="Requires python-pptx. Run: pip install python-pptx")

        topic = params.get("topic", "")
        if not topic:
            return StepResult(False, error="Please provide a presentation topic")

        num_slides = min(params.get("slides", 8), 30)
        outline = params.get("outline", [])
        audience = params.get("audience", "general")

        # Generate slide content via LLM
        slide_data = await self._generate_slides(topic, num_slides, outline, audience)
        if not slide_data:
            return StepResult(False, error="Failed to generate slide content")

        # Build PPTX
        try:
            prs = Presentation()
            prs.slide_width = Inches(13.333)
            prs.slide_height = Inches(7.5)

            for i, slide_info in enumerate(slide_data):
                slide_title = slide_info.get("title", f"Slide {i+1}")
                bullets = slide_info.get("bullets", [])
                notes = slide_info.get("notes", "")

                if i == 0:
                    # Title slide
                    layout = prs.slide_layouts[0]
                    slide = prs.slides.add_slide(layout)
                    slide.shapes.title.text = slide_title
                    if slide.placeholders[1]:
                        slide.placeholders[1].text = slide_info.get("subtitle", "")
                else:
                    # Content slide
                    layout = prs.slide_layouts[1]
                    slide = prs.slides.add_slide(layout)
                    slide.shapes.title.text = slide_title
                    body = slide.placeholders[1]
                    tf = body.text_frame
                    tf.clear()
                    for j, bullet in enumerate(bullets):
                        if j == 0:
                            tf.text = bullet
                        else:
                            p = tf.add_paragraph()
                            p.text = bullet
                            p.level = 0

                # Speaker notes
                if notes:
                    slide.notes_slide.notes_text_frame.text = notes

            out_dir = _output_dir()
            safe = "".join(c if c.isalnum() or c in " -_" else "" for c in topic).strip()
            path = str(out_dir / f"{safe or 'Presentation'}.pptx")
            prs.save(path)

            _try_open(path)

            return StepResult(True, data={
                "file": path,
                "title": topic,
                "slide_count": len(slide_data),
                "message": f"Presentation '{topic}' created with {len(slide_data)} slides → {path}"
            })
        except Exception as e:
            return StepResult(False, error=f"PPTX generation failed: {e}")

    async def _generate_slides(self, topic, num_slides, outline, audience):
        if not self._llm:
            # Fallback without LLM
            slides = [{"title": topic, "subtitle": datetime.now().strftime("%B %Y"), "bullets": [], "notes": ""}]
            for i in range(1, num_slides):
                title = outline[i-1] if i-1 < len(outline) else f"Section {i}"
                slides.append({"title": title, "bullets": ["Content here"], "notes": ""})
            return slides

        outline_text = ""
        if outline:
            outline_text = f"\nThe slides should cover these topics in order:\n" + "\n".join(f"- {t}" for t in outline)

        prompt = f"""Create a {num_slides}-slide presentation about: {topic}
Audience: {audience}{outline_text}

Return a JSON array of slide objects. Each object has:
- "title": slide title
- "subtitle": (first slide only) subtitle text
- "bullets": list of 3-5 bullet point strings (not for the title slide)
- "notes": 1-2 sentences of speaker notes

Slide 1 should be the title slide. Last slide should be a summary or Q&A.
Make content specific, actionable, and professional.
Return ONLY the JSON array."""

        try:
            response = await self._llm(prompt)
            response = response.strip()
            if response.startswith("```"):
                response = response.split("\n", 1)[1] if "\n" in response else response[3:]
                response = response.rsplit("```", 1)[0]
            return json.loads(response)
        except Exception as e:
            logger.warning(f"Slide generation failed: {e}")
            return None


# ============================================================
#  INVOICE SKILL
# ============================================================

class InvoiceSkill(Primitive):
    """Generate professional PDF invoices."""

    def __init__(self, llm_complete: Optional[Callable] = None):
        self._llm = llm_complete

    @property
    def name(self) -> str:
        return "INVOICE"

    def get_operations(self) -> Dict[str, str]:
        return {
            "create": "Generate a professional PDF invoice with line items, tax, and totals.",
        }

    def get_param_schema(self) -> Dict[str, Dict[str, Any]]:
        return {
            "create": {
                "client": {"type": "str", "required": True, "description": "Client/company name"},
                "items": {"type": "list", "required": True, "description": "List of dicts with 'description', 'quantity', 'unit_price'"},
                "invoice_number": {"type": "str", "required": False, "description": "Invoice number (auto-generated if omitted)"},
                "from_name": {"type": "str", "required": False, "description": "Your name/company name"},
                "due_date": {"type": "str", "required": False, "description": "Payment due date"},
                "tax_rate": {"type": "float", "required": False, "description": "Tax rate as percentage (e.g. 8.5 for 8.5%)"},
                "notes": {"type": "str", "required": False, "description": "Additional notes or payment instructions"},
                "currency": {"type": "str", "required": False, "description": "Currency symbol (default: $)"},
            },
        }

    async def execute(self, operation: str, params: Dict[str, Any]) -> StepResult:
        if operation != "create":
            return StepResult(False, error=f"Unknown operation: {operation}")

        try:
            from fpdf import FPDF
        except ImportError:
            return StepResult(False, error="Requires fpdf2. Run: pip install fpdf2")

        client = params.get("client", "")
        items = params.get("items", [])
        if not client or not items:
            return StepResult(False, error="Provide 'client' name and 'items' list")

        inv_num = params.get("invoice_number", f"INV-{datetime.now().strftime('%Y%m%d-%H%M')}")
        from_name = params.get("from_name", "")
        due_date = params.get("due_date", "")
        tax_rate = float(params.get("tax_rate", 0))
        notes = params.get("notes", "")
        currency = params.get("currency", "$")

        # Calculate totals
        subtotal = 0
        for item in items:
            qty = float(item.get("quantity", 1))
            price = float(item.get("unit_price", 0))
            item["_total"] = qty * price
            subtotal += item["_total"]

        tax = subtotal * (tax_rate / 100) if tax_rate else 0
        total = subtotal + tax

        try:
            pdf = FPDF()
            pdf.add_page()

            # Header
            pdf.set_font("Helvetica", "B", 28)
            pdf.set_text_color(40, 40, 40)
            pdf.cell(0, 15, "INVOICE", new_x="LMARGIN", new_y="NEXT")

            # Invoice details
            pdf.set_font("Helvetica", "", 10)
            pdf.set_text_color(100, 100, 100)
            pdf.cell(100, 6, f"Invoice #: {inv_num}", new_x="LMARGIN", new_y="NEXT")
            pdf.cell(100, 6, f"Date: {datetime.now().strftime('%B %d, %Y')}", new_x="LMARGIN", new_y="NEXT")
            if due_date:
                pdf.cell(100, 6, f"Due: {due_date}", new_x="LMARGIN", new_y="NEXT")
            pdf.ln(5)

            # From / To
            pdf.set_text_color(0, 0, 0)
            pdf.set_font("Helvetica", "B", 10)
            if from_name:
                pdf.cell(95, 6, "From:", new_x="RIGHT")
            else:
                pdf.cell(95, 6, "", new_x="RIGHT")
            pdf.cell(95, 6, "Bill To:", new_x="LMARGIN", new_y="NEXT")
            pdf.set_font("Helvetica", "", 10)
            if from_name:
                pdf.cell(95, 6, from_name, new_x="RIGHT")
            else:
                pdf.cell(95, 6, "", new_x="RIGHT")
            pdf.cell(95, 6, client, new_x="LMARGIN", new_y="NEXT")
            pdf.ln(10)

            # Line items table header
            pdf.set_fill_color(45, 55, 72)
            pdf.set_text_color(255, 255, 255)
            pdf.set_font("Helvetica", "B", 10)
            pdf.cell(90, 8, "  Description", border=0, fill=True)
            pdf.cell(25, 8, "Qty", border=0, fill=True, align="C")
            pdf.cell(35, 8, "Unit Price", border=0, fill=True, align="R")
            pdf.cell(35, 8, "Total  ", border=0, fill=True, align="R", new_x="LMARGIN", new_y="NEXT")

            # Line items
            pdf.set_text_color(40, 40, 40)
            pdf.set_font("Helvetica", "", 10)
            for i, item in enumerate(items):
                bg = (248, 249, 250) if i % 2 == 0 else (255, 255, 255)
                pdf.set_fill_color(*bg)
                desc = str(item.get("description", ""))[:50]
                qty = float(item.get("quantity", 1))
                price = float(item.get("unit_price", 0))
                line_total = item["_total"]
                pdf.cell(90, 7, f"  {desc}", fill=True)
                pdf.cell(25, 7, f"{qty:g}", fill=True, align="C")
                pdf.cell(35, 7, f"{currency}{price:,.2f}", fill=True, align="R")
                pdf.cell(35, 7, f"{currency}{line_total:,.2f}  ", fill=True, align="R", new_x="LMARGIN", new_y="NEXT")

            pdf.ln(3)

            # Totals
            pdf.set_font("Helvetica", "", 11)
            x_label = 115
            x_val = 150
            pdf.set_x(x_label)
            pdf.cell(35, 7, "Subtotal:", align="R")
            pdf.cell(35, 7, f"{currency}{subtotal:,.2f}  ", align="R", new_x="LMARGIN", new_y="NEXT")
            if tax_rate:
                pdf.set_x(x_label)
                pdf.cell(35, 7, f"Tax ({tax_rate}%):", align="R")
                pdf.cell(35, 7, f"{currency}{tax:,.2f}  ", align="R", new_x="LMARGIN", new_y="NEXT")

            pdf.set_font("Helvetica", "B", 13)
            pdf.set_x(x_label)
            pdf.cell(35, 10, "Total:", align="R")
            pdf.cell(35, 10, f"{currency}{total:,.2f}  ", align="R", new_x="LMARGIN", new_y="NEXT")

            # Notes
            if notes:
                pdf.ln(10)
                pdf.set_font("Helvetica", "B", 10)
                pdf.set_text_color(100, 100, 100)
                pdf.cell(0, 6, "Notes:", new_x="LMARGIN", new_y="NEXT")
                pdf.set_font("Helvetica", "", 9)
                pdf.multi_cell(0, 5, notes)

            out_dir = _output_dir()
            safe = "".join(c if c.isalnum() or c in " -_" else "" for c in client).strip()
            path = str(out_dir / f"Invoice {inv_num} - {safe}.pdf")
            pdf.output(path)

            _try_open(path)

            return StepResult(True, data={
                "file": path,
                "invoice_number": inv_num,
                "client": client,
                "subtotal": subtotal,
                "tax": tax,
                "total": total,
                "items": len(items),
                "message": f"Invoice {inv_num} for {client}: {currency}{total:,.2f} → {path}"
            })
        except Exception as e:
            return StepResult(False, error=f"Invoice generation failed: {e}")


# ============================================================
#  MEETING PREP SKILL
# ============================================================

class MeetingPrepSkill(Primitive):
    """Generate meeting preparation documents.

    Pulls calendar events, related emails, and tasks to create
    a briefing document for an upcoming meeting.
    """

    def __init__(self, llm_complete: Optional[Callable] = None):
        self._llm = llm_complete

    @property
    def name(self) -> str:
        return "MEETING_PREP"

    def get_operations(self) -> Dict[str, str]:
        return {
            "create": "Generate a meeting prep document. Searches indexed emails, calendar, and tasks for context, then creates a briefing PDF.",
        }

    def get_param_schema(self) -> Dict[str, Dict[str, Any]]:
        return {
            "create": {
                "meeting": {"type": "str", "required": True, "description": "Meeting title, topic, or attendee name to prep for"},
                "context": {"type": "str", "required": False, "description": "Additional context or specific questions to address"},
                "include_emails": {"type": "bool", "required": False, "description": "Search indexed emails for related threads (default: true)"},
                "include_tasks": {"type": "bool", "required": False, "description": "Include related open tasks (default: true)"},
            },
        }

    async def execute(self, operation: str, params: Dict[str, Any]) -> StepResult:
        if operation != "create":
            return StepResult(False, error=f"Unknown operation: {operation}")

        try:
            from fpdf import FPDF
        except ImportError:
            return StepResult(False, error="Requires fpdf2. Run: pip install fpdf2")

        meeting = params.get("meeting", "")
        if not meeting:
            return StepResult(False, error="Specify the meeting title or topic")

        context = params.get("context", "")
        include_emails = params.get("include_emails", True)
        include_tasks = params.get("include_tasks", True)

        # Gather context from index
        gathered = await self._gather_context(meeting, include_emails, include_tasks)

        # Generate briefing via LLM
        briefing = await self._generate_briefing(meeting, context, gathered)

        # Build PDF
        try:
            pdf = FPDF()
            pdf.add_page()
            pdf.set_auto_page_break(auto=True, margin=20)

            # Header
            pdf.set_font("Helvetica", "B", 22)
            pdf.set_text_color(30, 30, 30)
            pdf.cell(0, 12, "Meeting Prep", new_x="LMARGIN", new_y="NEXT")
            pdf.set_font("Helvetica", "B", 14)
            pdf.set_text_color(60, 60, 60)
            pdf.cell(0, 8, meeting, new_x="LMARGIN", new_y="NEXT")
            pdf.set_font("Helvetica", "", 10)
            pdf.set_text_color(120, 120, 120)
            pdf.cell(0, 6, f"Prepared: {datetime.now().strftime('%B %d, %Y at %I:%M %p')}", new_x="LMARGIN", new_y="NEXT")
            pdf.ln(8)

            # Briefing sections
            pdf.set_text_color(0, 0, 0)
            for section in briefing:
                pdf.set_font("Helvetica", "B", 13)
                pdf.cell(0, 9, section.get("heading", ""), new_x="LMARGIN", new_y="NEXT")
                pdf.set_font("Helvetica", "", 10)
                content = section.get("content", "")
                if isinstance(content, list):
                    for item in content:
                        pdf.cell(5, 6, "•")
                        pdf.multi_cell(0, 6, f" {item}")
                else:
                    pdf.multi_cell(0, 6, content)
                pdf.ln(4)

            out_dir = _output_dir()
            safe = "".join(c if c.isalnum() or c in " -_" else "" for c in meeting).strip()
            path = str(out_dir / f"Meeting Prep - {safe or 'Meeting'}.pdf")
            pdf.output(path)

            _try_open(path)

            return StepResult(True, data={
                "file": path,
                "meeting": meeting,
                "sections": len(briefing),
                "context_items": len(gathered),
                "message": f"Meeting prep for '{meeting}' → {path}"
            })
        except Exception as e:
            return StepResult(False, error=f"Meeting prep failed: {e}")

    async def _gather_context(self, meeting: str, include_emails: bool, include_tasks: bool) -> List[Dict]:
        """Search the data index for related content."""
        items = []
        index = get_data_index()
        if not index:
            return items

        # Search for related emails
        if include_emails:
            try:
                results = index.search(meeting, source="gmail", limit=10)
                results += index.search(meeting, source="outlook", limit=10)
                for r in results[:10]:
                    items.append({
                        "type": "email",
                        "title": r.get("title", ""),
                        "body": (r.get("body", "") or "")[:300],
                        "timestamp": r.get("timestamp", ""),
                    })
            except Exception:
                pass

        # Search for related tasks
        if include_tasks:
            try:
                results = index.search(meeting, source="todoist", limit=5)
                results += index.search(meeting, source="microsoft_todo", limit=5)
                results += index.search(meeting, source="jira", limit=5)
                for r in results[:8]:
                    items.append({
                        "type": "task",
                        "title": r.get("title", ""),
                        "body": (r.get("body", "") or "")[:200],
                    })
            except Exception:
                pass

        # Search local files for related docs
        try:
            results = index.search(meeting, source="local_files", limit=5)
            for r in results[:5]:
                items.append({
                    "type": "file",
                    "title": r.get("title", ""),
                    "body": (r.get("body", "") or "")[:200],
                })
        except Exception:
            pass

        return items

    async def _generate_briefing(self, meeting: str, context: str, gathered: List[Dict]) -> List[Dict]:
        """Generate the meeting briefing via LLM."""
        if not self._llm:
            sections = [{"heading": "Meeting Overview", "content": f"Briefing for: {meeting}"}]
            if gathered:
                sections.append({"heading": "Related Items", "content": [g["title"] for g in gathered]})
            return sections

        gathered_text = ""
        if gathered:
            gathered_text = "\n\nRelated context found:\n"
            for g in gathered:
                gathered_text += f"- [{g['type']}] {g['title']}: {g.get('body', '')[:150]}\n"

        context_text = f"\nAdditional context: {context}" if context else ""

        prompt = f"""Create a concise meeting prep briefing for: {meeting}{context_text}{gathered_text}

Return a JSON array of sections. Each section has:
- "heading": section title
- "content": either a string paragraph or a list of bullet point strings

Include these sections:
1. Meeting Objective — what this meeting aims to accomplish
2. Key Context — relevant background from related emails/files/tasks
3. Discussion Points — 3-5 topics to cover
4. Potential Questions — questions that may come up
5. Action Items to Raise — things to propose or follow up on

Be specific and actionable based on the context provided. Keep it to 1 page.
Return ONLY the JSON array."""

        try:
            response = await self._llm(prompt)
            response = response.strip()
            if response.startswith("```"):
                response = response.split("\n", 1)[1] if "\n" in response else response[3:]
                response = response.rsplit("```", 1)[0]
            return json.loads(response)
        except Exception as e:
            logger.warning(f"Briefing generation failed: {e}")
            return [{"heading": "Meeting Overview", "content": f"Briefing for: {meeting}"}]


# ============================================================
#  TRAVEL ITINERARY SKILL
# ============================================================

class TravelItinerarySkill(Primitive):
    """Generate travel itinerary PDFs.

    Combines calendar events, flight/hotel info, and local recommendations
    into a printable day-by-day itinerary.
    """

    def __init__(self, llm_complete: Optional[Callable] = None):
        self._llm = llm_complete

    @property
    def name(self) -> str:
        return "ITINERARY"

    def get_operations(self) -> Dict[str, str]:
        return {
            "create": "Generate a travel itinerary PDF. Provide destination, dates, and preferences. AI fills in activities, dining, and logistics.",
        }

    def get_param_schema(self) -> Dict[str, Dict[str, Any]]:
        return {
            "create": {
                "destination": {"type": "str", "required": True, "description": "Travel destination (e.g. 'Tokyo', 'Paris, France')"},
                "dates": {"type": "str", "required": True, "description": "Travel dates (e.g. 'June 15-22, 2026')"},
                "travelers": {"type": "int", "required": False, "description": "Number of travelers (default: 1)"},
                "interests": {"type": "str", "required": False, "description": "Interests/preferences (e.g. 'food, history, hiking')"},
                "flights": {"type": "list", "required": False, "description": "Flight details if known (list of dicts with airline, number, time)"},
                "hotel": {"type": "str", "required": False, "description": "Hotel name and address if booked"},
                "budget": {"type": "str", "required": False, "description": "Budget level: budget, moderate, luxury"},
            },
        }

    async def execute(self, operation: str, params: Dict[str, Any]) -> StepResult:
        if operation != "create":
            return StepResult(False, error=f"Unknown operation: {operation}")

        try:
            from fpdf import FPDF
        except ImportError:
            return StepResult(False, error="Requires fpdf2. Run: pip install fpdf2")

        destination = params.get("destination", "")
        dates = params.get("dates", "")
        if not destination or not dates:
            return StepResult(False, error="Provide 'destination' and 'dates'")

        interests = params.get("interests", "sightseeing, food, culture")
        travelers = params.get("travelers", 1)
        flights = params.get("flights", [])
        hotel = params.get("hotel", "")
        budget = params.get("budget", "moderate")

        # Generate itinerary via LLM
        itinerary = await self._generate_itinerary(
            destination, dates, interests, travelers, flights, hotel, budget
        )

        if not itinerary:
            return StepResult(False, error="Failed to generate itinerary")

        try:
            pdf = FPDF()
            pdf.set_auto_page_break(auto=True, margin=20)

            # Title page
            pdf.add_page()
            pdf.set_fill_color(30, 58, 95)
            pdf.rect(0, 0, 210, 297, "F")
            pdf.set_text_color(255, 255, 255)
            pdf.set_font("Helvetica", "B", 32)
            pdf.set_y(100)
            pdf.cell(0, 16, destination, align="C", new_x="LMARGIN", new_y="NEXT")
            pdf.set_font("Helvetica", "", 16)
            pdf.set_text_color(200, 210, 230)
            pdf.cell(0, 10, dates, align="C", new_x="LMARGIN", new_y="NEXT")
            pdf.ln(5)
            pdf.set_font("Helvetica", "", 12)
            details = []
            if travelers > 1:
                details.append(f"{travelers} travelers")
            details.append(budget.title())
            pdf.cell(0, 8, " · ".join(details), align="C")

            # Logistics page
            if flights or hotel:
                pdf.add_page()
                pdf.set_text_color(0, 0, 0)
                pdf.set_font("Helvetica", "B", 18)
                pdf.cell(0, 12, "Travel Logistics", new_x="LMARGIN", new_y="NEXT")
                pdf.ln(4)

                if flights:
                    pdf.set_font("Helvetica", "B", 12)
                    pdf.cell(0, 8, "Flights", new_x="LMARGIN", new_y="NEXT")
                    pdf.set_font("Helvetica", "", 10)
                    for f in flights:
                        line = f"{f.get('airline', '')} {f.get('number', '')} — {f.get('time', '')}"
                        pdf.cell(0, 6, line.strip(), new_x="LMARGIN", new_y="NEXT")
                    pdf.ln(4)

                if hotel:
                    pdf.set_font("Helvetica", "B", 12)
                    pdf.cell(0, 8, "Accommodation", new_x="LMARGIN", new_y="NEXT")
                    pdf.set_font("Helvetica", "", 10)
                    pdf.multi_cell(0, 6, hotel)

            # Day-by-day pages
            for day in itinerary:
                pdf.add_page()
                pdf.set_text_color(30, 58, 95)
                pdf.set_font("Helvetica", "B", 18)
                pdf.cell(0, 12, day.get("day", ""), new_x="LMARGIN", new_y="NEXT")
                pdf.set_text_color(100, 100, 100)
                pdf.set_font("Helvetica", "I", 10)
                pdf.cell(0, 6, day.get("theme", ""), new_x="LMARGIN", new_y="NEXT")
                pdf.ln(4)

                pdf.set_text_color(0, 0, 0)
                for activity in day.get("activities", []):
                    pdf.set_font("Helvetica", "B", 10)
                    time_str = activity.get("time", "")
                    name_str = activity.get("name", "")
                    pdf.cell(25, 7, time_str)
                    pdf.cell(0, 7, name_str, new_x="LMARGIN", new_y="NEXT")
                    desc = activity.get("description", "")
                    if desc:
                        pdf.set_font("Helvetica", "", 9)
                        pdf.set_text_color(80, 80, 80)
                        pdf.set_x(35)
                        pdf.multi_cell(0, 5, desc)
                        pdf.set_text_color(0, 0, 0)
                    pdf.ln(2)

                # Tips
                tips = day.get("tips", "")
                if tips:
                    pdf.ln(3)
                    pdf.set_font("Helvetica", "I", 9)
                    pdf.set_text_color(100, 100, 100)
                    pdf.multi_cell(0, 5, f"Tip: {tips}")

            out_dir = _output_dir()
            safe = "".join(c if c.isalnum() or c in " -_" else "" for c in destination).strip()
            path = str(out_dir / f"Itinerary - {safe or 'Trip'}.pdf")
            pdf.output(path)

            _try_open(path)

            return StepResult(True, data={
                "file": path,
                "destination": destination,
                "dates": dates,
                "days": len(itinerary),
                "message": f"Travel itinerary for {destination} ({dates}) → {path}"
            })
        except Exception as e:
            return StepResult(False, error=f"Itinerary generation failed: {e}")

    async def _generate_itinerary(self, destination, dates, interests, travelers, flights, hotel, budget):
        if not self._llm:
            return [{"day": "Day 1", "theme": "Arrival", "activities": [{"time": "Morning", "name": f"Arrive in {destination}", "description": ""}], "tips": ""}]

        prompt = f"""Create a day-by-day travel itinerary for {destination}.
Dates: {dates}
Travelers: {travelers}
Interests: {interests}
Budget: {budget}
{"Hotel: " + hotel if hotel else ""}

Return a JSON array with one object per day. Each day has:
- "day": "Day 1 — Monday, June 15"
- "theme": short theme for the day (e.g. "Historic Old Town")
- "activities": list of activities, each with "time" (e.g. "9:00 AM"), "name", "description" (1 sentence)
- "tips": one practical tip for the day

Include 4-6 activities per day: morning, lunch, afternoon, dinner, evening.
Include specific restaurant/cafe names and real attractions.
Return ONLY the JSON array."""

        try:
            response = await self._llm(prompt)
            response = response.strip()
            if response.startswith("```"):
                response = response.split("\n", 1)[1] if "\n" in response else response[3:]
                response = response.rsplit("```", 1)[0]
            return json.loads(response)
        except Exception as e:
            logger.warning(f"Itinerary generation failed: {e}")
            return None


# ============================================================
#  SOCIAL MEDIA KIT SKILL
# ============================================================

class SocialMediaKitSkill(Primitive):
    """Generate social media post packages.

    Creates optimized images (resized for each platform) and
    AI-generated captions with hashtags for multiple platforms.
    """

    def __init__(self, llm_complete: Optional[Callable] = None):
        self._llm = llm_complete

    @property
    def name(self) -> str:
        return "SOCIAL_KIT"

    def get_operations(self) -> Dict[str, str]:
        return {
            "create": "Create a social media kit: resize an image for multiple platforms and generate platform-specific captions with hashtags.",
        }

    def get_param_schema(self) -> Dict[str, Dict[str, Any]]:
        return {
            "create": {
                "image": {"type": "str", "required": False, "description": "Path to source image. If omitted, creates text-only posts."},
                "topic": {"type": "str", "required": True, "description": "What the post is about"},
                "platforms": {"type": "list", "required": False, "description": "Target platforms (default: ['instagram', 'twitter', 'linkedin', 'facebook']). Options: instagram, twitter, linkedin, facebook, tiktok, pinterest"},
                "tone": {"type": "str", "required": False, "description": "Tone: professional, casual, humorous, inspirational (default: professional)"},
                "hashtag_count": {"type": "int", "required": False, "description": "Number of hashtags per post (default: 5)"},
            },
        }

    # Platform image sizes (width x height in pixels)
    PLATFORM_SIZES = {
        "instagram": (1080, 1080),
        "twitter": (1200, 675),
        "linkedin": (1200, 627),
        "facebook": (1200, 630),
        "tiktok": (1080, 1920),
        "pinterest": (1000, 1500),
    }

    async def execute(self, operation: str, params: Dict[str, Any]) -> StepResult:
        if operation != "create":
            return StepResult(False, error=f"Unknown operation: {operation}")

        topic = params.get("topic", "")
        if not topic:
            return StepResult(False, error="Specify a topic for the social media post")

        image_path = params.get("image", "")
        platforms = params.get("platforms", ["instagram", "twitter", "linkedin", "facebook"])
        tone = params.get("tone", "professional")
        hashtag_count = params.get("hashtag_count", 5)

        out_dir = _output_dir() / "social_kit"
        out_dir.mkdir(parents=True, exist_ok=True)

        results = {"platforms": {}, "files": []}

        # Resize images for each platform
        if image_path and os.path.isfile(image_path):
            try:
                from PIL import Image
                for plat in platforms:
                    size = self.PLATFORM_SIZES.get(plat, (1200, 630))
                    with Image.open(image_path) as img:
                        resized = self._smart_resize(img, size)
                        safe_topic = "".join(c if c.isalnum() or c == " " else "" for c in topic)[:30].strip()
                        filename = f"{safe_topic} - {plat}.jpg"
                        save_path = str(out_dir / filename)
                        resized.save(save_path, "JPEG", quality=92)
                        results["files"].append(save_path)
                        if plat not in results["platforms"]:
                            results["platforms"][plat] = {}
                        results["platforms"][plat]["image"] = save_path
            except ImportError:
                pass
            except Exception as e:
                logger.warning(f"Image resize failed: {e}")

        # Generate captions
        captions = await self._generate_captions(topic, platforms, tone, hashtag_count)
        for plat, caption in captions.items():
            if plat not in results["platforms"]:
                results["platforms"][plat] = {}
            results["platforms"][plat]["caption"] = caption

        # Save captions to text file
        caption_file = str(out_dir / "captions.txt")
        with open(caption_file, "w", encoding="utf-8") as f:
            for plat in platforms:
                info = results["platforms"].get(plat, {})
                f.write(f"{'='*40}\n")
                f.write(f"{plat.upper()}\n")
                f.write(f"{'='*40}\n")
                f.write(info.get("caption", "") + "\n\n")
        results["files"].append(caption_file)

        _try_open(str(out_dir))

        return StepResult(True, data={
            **results,
            "topic": topic,
            "platform_count": len(platforms),
            "message": f"Social media kit for '{topic}': {len(platforms)} platforms, {len(results['files'])} files → {out_dir}"
        })

    def _smart_resize(self, img, target_size):
        """Resize with center crop to fill target aspect ratio."""
        from PIL import Image
        target_w, target_h = target_size
        target_ratio = target_w / target_h
        img_ratio = img.width / img.height

        if img_ratio > target_ratio:
            # Image is wider — crop sides
            new_w = int(img.height * target_ratio)
            left = (img.width - new_w) // 2
            img = img.crop((left, 0, left + new_w, img.height))
        elif img_ratio < target_ratio:
            # Image is taller — crop top/bottom
            new_h = int(img.width / target_ratio)
            top = (img.height - new_h) // 2
            img = img.crop((0, top, img.width, top + new_h))

        return img.resize(target_size, Image.LANCZOS)

    async def _generate_captions(self, topic, platforms, tone, hashtag_count):
        if not self._llm:
            return {p: f"Check out {topic}! #{'#'.join(topic.split()[:3])}" for p in platforms}

        platforms_str = ", ".join(platforms)
        prompt = f"""Write social media captions for each platform about: {topic}

Platforms: {platforms_str}
Tone: {tone}
Include {hashtag_count} relevant hashtags per post.

Platform guidelines:
- Instagram: engaging, emoji-rich, up to 2200 chars, hashtags at the end
- Twitter/X: concise, under 280 chars total including hashtags
- LinkedIn: professional, thought-leadership, 1-2 paragraphs
- Facebook: conversational, medium length, question to drive engagement
- TikTok: trendy, Gen-Z friendly, short and punchy
- Pinterest: descriptive, keyword-rich for search

Return a JSON object where keys are platform names and values are the caption strings.
Return ONLY the JSON object."""

        try:
            response = await self._llm(prompt)
            response = response.strip()
            if response.startswith("```"):
                response = response.split("\n", 1)[1] if "\n" in response else response[3:]
                response = response.rsplit("```", 1)[0]
            return json.loads(response)
        except Exception as e:
            logger.warning(f"Caption generation failed: {e}")
            return {p: topic for p in platforms}


# ============================================================
#  HELPER
# ============================================================

def _try_open(path: str):
    """Try to open a file with the system default application."""
    try:
        if platform.system() == "Windows":
            os.startfile(path)
        elif platform.system() == "Darwin":
            subprocess.Popen(["open", path])
        else:
            subprocess.Popen(["xdg-open", path])
    except Exception:
        pass
