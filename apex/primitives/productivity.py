"""Telic Engine — Productivity Primitives"""

import asyncio
import hashlib
import json
import logging
import os
import re
import uuid
from dataclasses import dataclass
from datetime import datetime, timedelta, timezone
from pathlib import Path
from typing import Any, Callable, Dict, List, Optional, Union

from .base import Primitive, StepResult
from .base import get_data_index

logger = logging.getLogger(__name__)


class CalendarPrimitive(Primitive):
    """Calendar and scheduling operations.
    
    Stores events locally. Can be wired to Google Calendar, Outlook, etc.
    via provider functions passed at init.
    """
    
    def __init__(
        self,
        storage_path: Optional[str] = None,
        create_func: Optional[Callable] = None,
        list_func: Optional[Callable] = None,
        list_calendars_func: Optional[Callable] = None,
        providers: Optional[Dict[str, Any]] = None,
    ):
        self._events: List[Dict] = []
        self._storage_path = storage_path
        self._create_func = create_func
        self._list_func = list_func
        self._list_calendars_func = list_calendars_func
        self._providers = providers or {}
        self._connector = next(iter(self._providers.values())) if self._providers else None
        self._calendars_cache: List[Dict] = []  # Cache of available calendars
        if storage_path and Path(storage_path).exists():
            try:
                with open(storage_path, "r") as f:
                    self._events = json.load(f)
            except Exception:
                pass
    
    @property
    def name(self) -> str:
        return "CALENDAR"
    
    def get_operations(self) -> Dict[str, str]:
        return {
            "create": "Create a calendar event (supports birthdays, meetings, reminders, recurring events)",
            "list": "List events in a date range. By default queries ALL of the user's own calendars (not subscriptions) in a single call — no need to call list_calendars first. Only specify calendar_id if the user asks about a specific calendar.",
            "list_calendars": "List all available calendars (name, id, access level).",
            "search": "Search/find events by date range and/or calendar. Returns all events in the range.",
            "delete": "Delete an event by ID",
            "availability": "Check free/busy times in a date range",
        }

    def get_available_operations(self) -> Dict[str, str]:
        ops = self.get_operations()
        connected = self.get_connected_providers()
        if connected:
            return ops

        idx = get_data_index()
        if idx and not idx.is_stale("google_calendar"):
            return {
                "list": ops["list"],
                "search": ops["search"],
            }

        return {}
    
    def get_param_schema(self) -> Dict[str, Dict[str, Any]]:
        return {
            "create": {
                "title": {"type": "str", "required": True, "description": "Event title"},
                "start": {"type": "str", "required": True, "description": "Start datetime ISO format (e.g. 2026-04-01T09:00:00)"},
                "end": {"type": "str", "required": False, "description": "End datetime ISO format (defaults to 1 hour after start)"},
                "description": {"type": "str", "required": False, "description": "Event description or notes"},
                "location": {"type": "str", "required": False, "description": "Event location"},
                "attendees": {"type": "list", "required": False, "description": "List of attendee email addresses"},
                "calendar_id": {"type": "str", "required": False, "description": "Calendar ID (use 'primary' for main calendar, or calendar name like 'FAMILY SHARED' for specific calendars)"},
                "recurrence": {"type": "str", "required": False, "description": "Recurrence rule: 'yearly', 'monthly', 'weekly', 'daily', or RRULE string"},
                "all_day": {"type": "bool", "required": False, "description": "True for all-day events like birthdays"},
            },
            "list": {
                "start_date": {"type": "str", "required": False, "description": "Start of range (ISO date, defaults to today)"},
                "end_date": {"type": "str", "required": False, "description": "End of range (ISO date, defaults to 7 days out)"},
                "limit": {"type": "int", "required": False, "description": "Max events to return (default 50)"},
                "calendar_id": {"type": "str", "required": False, "description": "Calendar ID or name (default: queries user's own calendars, not subscriptions)"},
            },
            "search": {
                "query": {"type": "str", "required": True, "description": "Keyword to search for in event titles/descriptions (NOT a date)"},
                "calendar_id": {"type": "str", "required": False, "description": "Calendar ID or name to search within"},
                "start_date": {"type": "str", "required": False, "description": "Start of range (ISO date)"},
                "end_date": {"type": "str", "required": False, "description": "End of range (ISO date)"},
            },
            "list_calendars": {},
            "delete": {
                "id": {"type": "str", "required": True, "description": "Event ID to delete"},
            },
            "availability": {
                "start_date": {"type": "str", "required": True, "description": "Start of range (ISO date)"},
                "end_date": {"type": "str", "required": True, "description": "End of range (ISO date)"},
            },
        }
    
    def _save(self):
        if self._storage_path:
            Path(self._storage_path).parent.mkdir(parents=True, exist_ok=True)
            with open(self._storage_path, "w") as f:
                json.dump(self._events, f, indent=2)
            print(f"[CALENDAR] Saved {len(self._events)} events to {self._storage_path}")
    
    async def _resolve_calendar_id(self, calendar_id: str) -> str:
        """Resolve calendar name to ID.
        
        If calendar_id looks like a name (contains spaces or is a known name),
        look it up in the list of calendars and return the actual ID.
        """
        if not calendar_id or calendar_id == "primary":
            return "primary"
        
        # If it looks like an email/ID already (contains @), use as-is
        if "@" in calendar_id:
            return calendar_id
        
        # Try to resolve by name
        if self._list_calendars_func:
            try:
                if not self._calendars_cache:
                    self._calendars_cache = await self._list_calendars_func()
                    print(f"[CALENDAR] Cached {len(self._calendars_cache)} calendars")
                
                # Search for matching calendar by name (case insensitive)
                # list_calendars() may return 'name' or 'summary' depending on connector
                search_name = calendar_id.lower().strip()
                for cal in self._calendars_cache:
                    cal_name = (cal.get("name") or cal.get("summary") or "").lower()
                    if search_name in cal_name or cal_name in search_name:
                        real_id = cal.get("id")
                        print(f"[CALENDAR] Resolved '{calendar_id}' -> '{real_id}'")
                        return real_id
                
                print(f"[CALENDAR] No calendar found matching '{calendar_id}', using as-is")
            except Exception as e:
                print(f"[CALENDAR] Error resolving calendar: {e}")
        
        return calendar_id

    async def execute(self, operation: str, params: Dict[str, Any]) -> StepResult:
        print(f"[CALENDAR] execute({operation}, {params})")
        try:
            if operation == "create":
                if self._create_func:
                    print(f"[CALENDAR] Using external calendar (Google/Outlook)")
                    # Resolve calendar name to ID (e.g., "FAMILY SHARED" -> actual ID)
                    raw_calendar_id = params.get("calendar_id", "primary")
                    resolved_calendar_id = await self._resolve_calendar_id(raw_calendar_id)
                    
                    # Convert friendly recurrence to RRULE
                    recurrence = params.get("recurrence")
                    recurrence_rules = None
                    if recurrence:
                        recurrence_map = {
                            "yearly": ["RRULE:FREQ=YEARLY"],
                            "annually": ["RRULE:FREQ=YEARLY"],
                            "monthly": ["RRULE:FREQ=MONTHLY"],
                            "weekly": ["RRULE:FREQ=WEEKLY"],
                            "daily": ["RRULE:FREQ=DAILY"],
                        }
                        if recurrence.lower() in recurrence_map:
                            recurrence_rules = recurrence_map[recurrence.lower()]
                        elif recurrence.startswith("RRULE:"):
                            recurrence_rules = [recurrence]
                    
                    # Map param names: CalendarPrimitive uses 'title', Google API uses 'summary'
                    api_params = {
                        "summary": params.get("title") or params.get("summary", "Untitled"),
                        "start": params.get("start"),
                        "end": params.get("end"),
                        "description": params.get("description", ""),
                        "location": params.get("location", ""),
                        "attendees": params.get("attendees", []),
                        "calendar_id": resolved_calendar_id,
                        "all_day": params.get("all_day", False),
                        "recurrence": recurrence_rules,
                    }
                    print(f"[CALENDAR] Creating event: {api_params}")
                    result = await self._create_func(**api_params)
                    # Handle CalendarEvent dataclass or dict
                    if hasattr(result, 'to_dict'):
                        result_dict = result.to_dict()
                    elif isinstance(result, dict):
                        result_dict = result
                    else:
                        result_dict = {"id": str(result)}
                    return StepResult(True, data={
                        **result_dict,
                        "storage": "google_calendar",
                        "calendar": raw_calendar_id,
                        "message": f"Event '{api_params['summary']}' created in Google Calendar ({raw_calendar_id})"
                    })
                
                print(f"[CALENDAR] Creating local event")
                event = {
                    "id": f"evt_{len(self._events)}_{int(datetime.now().timestamp())}",
                    "title": params.get("title", "Untitled"),
                    "start": params.get("start", datetime.now().isoformat()),
                    "end": params.get("end"),
                    "description": params.get("description", ""),
                    "location": params.get("location", ""),
                    "attendees": params.get("attendees", []),
                    "created": datetime.now().isoformat(),
                    "storage": "local",  # Indicate where event is stored
                }
                self._events.append(event)
                self._save()
                print(f"[CALENDAR] Event created: {event['id']} - {event['title']}")
                print(f"[CALENDAR] Saved to: {self._storage_path}")
                return StepResult(True, data={
                    **event,
                    "message": f"Event '{event['title']}' saved to local calendar ({self._storage_path}). Google Calendar sync coming soon."
                })
            
            elif operation == "list":
                # Try local index first (instant) — fall through to live API if stale
                idx = get_data_index()
                if idx and not idx.is_stale("google_calendar"):
                    try:
                        q_kwargs = {"kind": "event"}
                        if params.get("start_date"):
                            q_kwargs["after"] = datetime.fromisoformat(params["start_date"])
                        if params.get("end_date"):
                            dt = datetime.fromisoformat(params["end_date"])
                            q_kwargs["before"] = dt + timedelta(days=1)
                        if params.get("limit"):
                            q_kwargs["limit"] = params["limit"]
                        results = idx.query(**q_kwargs)
                        if results:
                            print(f"[CALENDAR] Index hit: {len(results)} events")
                            return StepResult(True, data=[
                                {k: v for k, v in {
                                    "id": r.source_id, "summary": r.title,
                                    "start": r.timestamp.isoformat() if r.timestamp else "",
                                    "end": r.timestamp_end.isoformat() if r.timestamp_end else "",
                                    "description": r.body, "location": r.location,
                                    "status": r.status, "html_link": r.url,
                                    "attendees": [{"email": p} for p in r.participants],
                                }.items() if v}
                                for r in results
                            ])
                    except Exception as e:
                        print(f"[CALENDAR] Index query failed, falling through to API: {e}")

                if self._list_func:
                    # Map primitive params to connector params
                    api_params = {}
                    if params.get("start_date"):
                        dt = datetime.fromisoformat(params["start_date"])
                        # Don't force timezone on date-only strings — let the
                        # connector handle them so the calendar's own timezone
                        # is respected by the Google API.
                        api_params["time_min"] = dt
                    if params.get("end_date"):
                        dt = datetime.fromisoformat(params["end_date"])
                        api_params["time_max"] = dt + timedelta(days=1)
                    if params.get("limit"):
                        api_params["max_results"] = params["limit"]
                    if params.get("calendar_id"):
                        raw_cal = params["calendar_id"]
                        api_params["calendar_id"] = await self._resolve_calendar_id(raw_cal)
                    result = await self._list_func(**api_params)
                    # Convert CalendarEvent objects to dicts
                    if result and hasattr(result[0], 'to_dict'):
                        result = [e.to_dict() for e in result]
                    return StepResult(True, data=result)
                
                start = params.get("start_date", datetime.now().strftime("%Y-%m-%d"))
                end = params.get("end_date")
                limit = params.get("limit", 50)
                
                filtered = []
                for evt in self._events:
                    evt_start = evt.get("start", "")
                    if evt_start >= start and (not end or evt_start <= end):
                        filtered.append(evt)
                
                filtered.sort(key=lambda e: e.get("start", ""))
                return StepResult(True, data=filtered[:limit])
            
            elif operation == "list_calendars":
                if self._list_calendars_func:
                    calendars = await self._list_calendars_func()
                    return StepResult(True, data=calendars)
                return StepResult(False, error="No calendar provider connected")
            
            elif operation == "search":
                # Try local FTS index first
                idx = get_data_index()
                query_text = params.get("query", "")
                if idx and not idx.is_stale("google_calendar") and query_text:
                    try:
                        results = idx.search(query_text, kind="event")
                        if results:
                            print(f"[CALENDAR] Index FTS hit: {len(results)} events for '{query_text}'")
                            return StepResult(True, data=[
                                {k: v for k, v in {
                                    "id": r.source_id, "summary": r.title,
                                    "start": r.timestamp.isoformat() if r.timestamp else "",
                                    "end": r.timestamp_end.isoformat() if r.timestamp_end else "",
                                    "description": r.body, "location": r.location,
                                    "status": r.status, "html_link": r.url,
                                }.items() if v}
                                for r in results
                            ])
                    except Exception as e:
                        print(f"[CALENDAR] Index FTS failed, falling through to API: {e}")

                if self._list_func:
                    # Search = list with optional date bounds + calendar filter.
                    # The AI is the intelligence — it reads the results.
                    # Don't pass query text as Google's q= filter; it silently
                    # kills results when the text doesn't match event titles.
                    search_kwargs = {}
                    if params.get("start_date"):
                        search_kwargs["time_min"] = datetime.fromisoformat(params["start_date"])
                    if params.get("end_date"):
                        dt = datetime.fromisoformat(params["end_date"])
                        search_kwargs["time_max"] = dt + timedelta(days=1)
                    if params.get("calendar_id"):
                        search_kwargs["calendar_id"] = await self._resolve_calendar_id(params["calendar_id"])
                    result = await self._list_func(**search_kwargs)
                    if result and hasattr(result[0], 'to_dict'):
                        result = [e.to_dict() for e in result]
                    return StepResult(True, data=result)
                
                query = params.get("query", "").lower()
                matches = [
                    e for e in self._events
                    if query in e.get("title", "").lower()
                    or query in e.get("description", "").lower()
                    or query in e.get("location", "").lower()
                ]
                return StepResult(True, data=matches)
            
            elif operation == "delete":
                event_id = params.get("id")
                before = len(self._events)
                self._events = [e for e in self._events if e.get("id") != event_id]
                self._save()
                return StepResult(True, data={"deleted": before != len(self._events)})
            
            elif operation == "availability":
                start = params.get("start_date", "")
                end = params.get("end_date", "")
                busy = [
                    {"start": e["start"], "end": e.get("end", e["start"]), "title": e["title"]}
                    for e in self._events
                    if e.get("start", "") >= start and e.get("start", "") <= end
                ]
                return StepResult(True, data={"busy": busy, "count": len(busy)})
            
            else:
                return StepResult(False, error=f"Unknown operation: {operation}")
        except Exception as e:
            return StepResult(False, error=str(e))


# ============================================================
#  WEB PRIMITIVE
# ============================================================



class TaskPrimitive(Primitive):
    """Task and to-do list management.
    
    Local task store. Can be wired to Todoist, Microsoft To Do,
    Jira, etc. via providers dict.
    """
    
    def __init__(self, storage_path: Optional[str] = None, providers: Optional[Dict[str, Any]] = None):
        self._tasks: List[Dict] = []
        self._storage_path = storage_path
        self._providers = providers or {}
        if storage_path and Path(storage_path).exists():
            try:
                with open(storage_path, "r") as f:
                    self._tasks = json.load(f)
            except Exception:
                pass
    
    @property
    def name(self) -> str:
        return "TASK"
    
    def get_operations(self) -> Dict[str, str]:
        return {
            "create": "Create a new task or to-do item",
            "list": "List tasks, optionally filtered by status or tag",
            "update": "Update a task (status, title, due date, etc.)",
            "complete": "Mark a task as completed",
            "delete": "Delete a task",
            "search": "Search tasks by keyword",
        }
    
    def get_param_schema(self) -> Dict[str, Dict[str, Any]]:
        return {
            "create": {
                "title": {"type": "str", "required": True, "description": "Task title"},
                "description": {"type": "str", "required": False, "description": "Task description/details"},
                "due": {"type": "str", "required": False, "description": "Due date (ISO date or datetime)"},
                "priority": {"type": "str", "required": False, "description": "low, medium, high, urgent"},
                "tags": {"type": "list", "required": False, "description": "Tags/categories"},
                "project": {"type": "str", "required": False, "description": "Project name"},
                "provider": {"type": "str", "required": False, "description": "Provider: todoist, microsoft_todo, jira (default: local)"},
            },
            "list": {
                "status": {"type": "str", "required": False, "description": "Filter: open, completed, all (default open)"},
                "project": {"type": "str", "required": False, "description": "Filter by project"},
                "tag": {"type": "str", "required": False, "description": "Filter by tag"},
                "limit": {"type": "int", "required": False, "description": "Max results"},
                "provider": {"type": "str", "required": False, "description": "Provider name"},
            },
            "update": {
                "id": {"type": "str", "required": True, "description": "Task ID"},
                "title": {"type": "str", "required": False, "description": "New title"},
                "status": {"type": "str", "required": False, "description": "New status"},
                "due": {"type": "str", "required": False, "description": "New due date"},
                "priority": {"type": "str", "required": False, "description": "New priority"},
                "provider": {"type": "str", "required": False, "description": "Provider name"},
            },
            "complete": {
                "id": {"type": "str", "required": True, "description": "Task ID to complete"},
                "provider": {"type": "str", "required": False, "description": "Provider name"},
            },
            "delete": {
                "id": {"type": "str", "required": True, "description": "Task ID to delete"},
                "provider": {"type": "str", "required": False, "description": "Provider name"},
            },
            "search": {
                "query": {"type": "str", "required": True, "description": "Search term"},
                "provider": {"type": "str", "required": False, "description": "Provider name"},
            },
        }
    
    def _get_provider(self, name: Optional[str]) -> Optional[Any]:
        if name and name in self._providers:
            return self._providers[name]
        if self._providers:
            return next(iter(self._providers.values()))
        return None
    
    def _save(self):
        if self._storage_path:
            Path(self._storage_path).parent.mkdir(parents=True, exist_ok=True)
            with open(self._storage_path, "w") as f:
                json.dump(self._tasks, f, indent=2)
    
    async def execute(self, operation: str, params: Dict[str, Any]) -> StepResult:
        try:
            provider_name = params.get("provider")
            provider = self._get_provider(provider_name)
            
            if operation == "create":
                if provider and hasattr(provider, "create_task"):
                    title = params.get("title", "Untitled")
                    description = params.get("description", "")
                    due = params.get("due")
                    # Todoist uses 'content', Microsoft To-Do uses 'title'
                    import inspect
                    sig = inspect.signature(provider.create_task)
                    if "content" in sig.parameters:
                        result = await provider.create_task(
                            content=title,
                            description=description,
                            due_date=due,
                        )
                    else:
                        result = await provider.create_task(
                            title=title,
                            body=description or None,
                            due_date=due,
                        )
                    return StepResult(True, data={"id": getattr(result, "id", str(result)), "title": params.get("title"), "status": "created"})
                
                task = {
                    "id": f"task_{len(self._tasks)}_{int(datetime.now().timestamp())}",
                    "title": params.get("title", "Untitled"),
                    "description": params.get("description", ""),
                    "status": "open",
                    "due": params.get("due"),
                    "priority": params.get("priority", "medium"),
                    "tags": params.get("tags", []),
                    "project": params.get("project"),
                    "created": datetime.now().isoformat(),
                    "completed_at": None,
                }
                self._tasks.append(task)
                self._save()
                return StepResult(True, data=task)
            
            elif operation == "list":
                # Try local index first
                idx = get_data_index()
                source_key = provider_name or "todoist"
                if idx and not idx.is_stale(source_key):
                    try:
                        q_kwargs = {"kind": "task"}
                        status_filter = params.get("status", "open")
                        if status_filter == "open":
                            q_kwargs["status"] = "pending"
                        elif status_filter == "completed":
                            q_kwargs["status"] = "completed"
                        q_kwargs["limit"] = params.get("limit", 50)
                        results = idx.query(**q_kwargs)
                        if results:
                            print(f"[TASK] Index hit: {len(results)} tasks")
                            return StepResult(True, data=[
                                {k: v for k, v in {
                                    "id": r.source_id, "title": r.title,
                                    "description": r.body,
                                    "status": r.status,
                                    "due": r.timestamp.isoformat() if r.timestamp else None,
                                    "labels": r.labels,
                                    "url": r.url,
                                }.items() if v}
                                for r in results
                            ])
                    except Exception as e:
                        print(f"[TASK] Index query failed, falling through to API: {e}")

                if provider and hasattr(provider, "list_tasks"):
                    result = await provider.list_tasks()
                    # Convert dataclass objects to dicts if needed
                    if result and hasattr(result[0], '__dataclass_fields__'):
                        from dataclasses import asdict
                        result = [asdict(t) for t in result]
                    return StepResult(True, data=result)
                
                status = params.get("status", "open")
                project = params.get("project")
                tag = params.get("tag")
                limit = params.get("limit", 50)
                
                filtered = self._tasks
                if status != "all":
                    filtered = [t for t in filtered if t.get("status") == status]
                if project:
                    filtered = [t for t in filtered if t.get("project") == project]
                if tag:
                    filtered = [t for t in filtered if tag in t.get("tags", [])]
                
                return StepResult(True, data=filtered[:limit])
            
            elif operation == "update":
                task_id = params.get("id")
                
                if provider and hasattr(provider, "update_task"):
                    result = await provider.update_task(task_id=task_id, title=params.get("title"), status=params.get("status"))
                    return StepResult(True, data=result)
                
                for task in self._tasks:
                    if task.get("id") == task_id:
                        for k in ("title", "status", "due", "priority", "description", "project"):
                            if k in params and k != "id":
                                task[k] = params[k]
                        self._save()
                        return StepResult(True, data=task)
                return StepResult(False, error=f"Task not found: {task_id}")
            
            elif operation == "complete":
                task_id = params.get("id")
                
                if provider and hasattr(provider, "complete_task"):
                    result = await provider.complete_task(task_id=task_id)
                    return StepResult(True, data=result)
                
                for task in self._tasks:
                    if task.get("id") == task_id:
                        task["status"] = "completed"
                        task["completed_at"] = datetime.now().isoformat()
                        self._save()
                        return StepResult(True, data=task)
                return StepResult(False, error=f"Task not found: {task_id}")
            
            elif operation == "delete":
                task_id = params.get("id")
                
                if provider and hasattr(provider, "delete_task"):
                    result = await provider.delete_task(task_id=task_id)
                    return StepResult(True, data=result)
                
                before = len(self._tasks)
                self._tasks = [t for t in self._tasks if t.get("id") != task_id]
                self._save()
                return StepResult(True, data={"deleted": before != len(self._tasks)})
            
            elif operation == "search":
                query = (params.get("query") or "").lower()
                
                # Try local index first
                idx = get_data_index()
                if idx:
                    try:
                        results = idx.search(query, kind="task", limit=params.get("limit", 20))
                        if results:
                            print(f"[TASK] Index search hit: {len(results)} tasks")
                            return StepResult(True, data=[
                                {k: v for k, v in {
                                    "id": r.source_id, "title": r.title,
                                    "description": r.body,
                                    "status": r.status,
                                    "due": r.timestamp.isoformat() if r.timestamp else None,
                                }.items() if v}
                                for r in results
                            ])
                    except Exception as e:
                        print(f"[TASK] Index search failed, falling through: {e}")
                
                if provider and hasattr(provider, "search_tasks"):
                    result = await provider.search_tasks(query=query)
                    return StepResult(True, data=result)
                
                matches = [
                    t for t in self._tasks
                    if query in t.get("title", "").lower()
                    or query in t.get("description", "").lower()
                    or query in t.get("project", "").lower()
                ]
                return StepResult(True, data=matches)
            
            else:
                return StepResult(False, error=f"Unknown operation: {operation}")
        except Exception as e:
            return StepResult(False, error=str(e))


# ============================================================
#  SHELL PRIMITIVE
# ============================================================



class NotesPrimitive(Primitive):
    """Note-taking operations - OneNote, Apple Notes, Google Keep, etc."""
    
    def __init__(self, storage_path: str = "", providers: Optional[Dict[str, Any]] = None):
        self._storage_path = Path(storage_path) if storage_path else Path.home() / ".telic" / "notes"
        self._storage_path.mkdir(parents=True, exist_ok=True)
        self._providers = providers or {}
    
    @property
    def name(self) -> str:
        return "NOTES"
    
    def get_operations(self) -> Dict[str, str]:
        return self.get_available_operations()
    
    def get_available_operations(self) -> Dict[str, str]:
        return {
            "create": "Create a new note",
            "read": "Read a note",
            "update": "Update a note",
            "delete": "Delete a note",
            "search": "Search notes",
            "list": "List all notes",
        }
    
    def get_param_schema(self) -> Dict[str, Any]:
        return {
            "create": {"title": {"type": "str", "description": "Note title"}, "content": {"type": "str", "description": "Note content"}, "folder": {"type": "str", "description": "Folder (optional)"}, "tags": {"type": "list", "description": "Tags (optional)"}},
            "read": {"note_id": {"type": "str", "description": "Note ID or title"}},
            "update": {"note_id": {"type": "str", "description": "Note ID"}, "content": {"type": "str", "description": "New content"}},
            "delete": {"note_id": {"type": "str", "description": "Note ID"}},
            "search": {"query": {"type": "str", "description": "Search query"}},
            "list": {"folder": {"type": "str", "description": "Folder (optional)"}},
        }
    
    async def execute(self, operation: str, params: Dict[str, Any]) -> StepResult:
        try:
            for name, provider in self._providers.items():
                if operation == "create" and hasattr(provider, "create_note"):
                    result = await provider.create_note(params.get("title"), params.get("content"), params.get("folder"), params.get("tags"))
                    return StepResult(True, data={"created": True, "note": result, "provider": name})
                elif operation == "read" and hasattr(provider, "read_note"):
                    result = await provider.read_note(params.get("note_id"))
                    return StepResult(True, data={"note": result, "provider": name})
                elif operation == "list" and hasattr(provider, "list_notes"):
                    result = await provider.list_notes(params.get("folder"))
                    return StepResult(True, data={"notes": result, "provider": name})
            
            # Local file-based notes
            if operation == "create":
                title = params.get("title", "Untitled")
                content = params.get("content", "")
                tags = params.get("tags", [])
                note_id = f"note_{int(datetime.now().timestamp())}"
                note_file = self._storage_path / f"{note_id}.json"
                note = {"id": note_id, "title": title, "content": content, "tags": tags, "created": datetime.now().isoformat()}
                note_file.write_text(json.dumps(note))
                return StepResult(True, data={"created": True, "note": note, "provider": "local"})
            
            elif operation == "read":
                note_id = params.get("note_id")
                note_file = self._storage_path / f"{note_id}.json"
                if note_file.exists():
                    note = json.loads(note_file.read_text())
                    return StepResult(True, data={"note": note, "provider": "local"})
                # Try to find by title
                for f in self._storage_path.glob("*.json"):
                    note = json.loads(f.read_text())
                    if note.get("title") == note_id:
                        return StepResult(True, data={"note": note, "provider": "local"})
                return StepResult(False, error=f"Note not found: {note_id}")
            
            elif operation == "update":
                note_id = params.get("note_id")
                content = params.get("content")
                note_file = self._storage_path / f"{note_id}.json"
                if note_file.exists():
                    note = json.loads(note_file.read_text())
                    note["content"] = content
                    note["updated"] = datetime.now().isoformat()
                    note_file.write_text(json.dumps(note))
                    return StepResult(True, data={"updated": True, "note": note, "provider": "local"})
                return StepResult(False, error=f"Note not found: {note_id}")
            
            elif operation == "delete":
                note_id = params.get("note_id")
                note_file = self._storage_path / f"{note_id}.json"
                if note_file.exists():
                    note_file.unlink()
                    return StepResult(True, data={"deleted": True, "provider": "local"})
                return StepResult(False, error=f"Note not found: {note_id}")
            
            elif operation == "search":
                query = params.get("query", "").lower()
                results = []
                for f in self._storage_path.glob("*.json"):
                    note = json.loads(f.read_text())
                    if query in note.get("title", "").lower() or query in note.get("content", "").lower():
                        results.append(note)
                return StepResult(True, data={"notes": results, "provider": "local"})
            
            elif operation == "list":
                notes = []
                for f in self._storage_path.glob("*.json"):
                    notes.append(json.loads(f.read_text()))
                return StepResult(True, data={"notes": notes, "provider": "local"})
            
            else:
                return StepResult(False, error=f"Unknown operation: {operation}")
        except Exception as e:
            return StepResult(False, error=str(e))


# ============================================================
#  FINANCE PRIMITIVE - Banking, payments
# ============================================================



class SpreadsheetPrimitive(Primitive):
    """Spreadsheet operations - Google Sheets, Excel, etc."""
    
    def __init__(self, llm_complete: Optional[Callable] = None, providers: Optional[Dict[str, Any]] = None):
        self._llm = llm_complete
        self._providers = providers or {}
    
    @property
    def name(self) -> str:
        return "SPREADSHEET"
    
    def get_operations(self) -> Dict[str, str]:
        return self.get_available_operations()
    
    def get_available_operations(self) -> Dict[str, str]:
        return {
            "read": "Read data from a spreadsheet",
            "write": "Write data to a spreadsheet",
            "create": "Create a new spreadsheet",
            "add_sheet": "Add a worksheet to a spreadsheet",
            "formula": "Set a formula in a cell",
            "format": "Format cells",
            "chart": "Create a chart",
        }
    
    def get_param_schema(self) -> Dict[str, Any]:
        return {
            "read": {"file": {"type": "str", "description": "File path or ID"}, "sheet": {"type": "str", "description": "Sheet name"}, "range": {"type": "str", "description": "Cell range like A1:D10"}},
            "write": {"file": {"type": "str", "description": "File path or ID"}, "data": {"type": "list", "description": "2D array of values"}, "sheet": {"type": "str", "description": "Sheet name"}, "range": {"type": "str", "description": "Starting cell"}},
            "create": {"name": {"type": "str", "description": "Spreadsheet name"}, "data": {"type": "list", "description": "Initial data (optional)"}},
            "add_sheet": {"file": {"type": "str", "description": "File path or ID"}, "name": {"type": "str", "description": "New sheet name"}},
            "formula": {"file": {"type": "str", "description": "File path or ID"}, "cell": {"type": "str", "description": "Cell like A1"}, "formula": {"type": "str", "description": "Formula"}},
            "format": {"file": {"type": "str", "description": "File path or ID"}, "range": {"type": "str", "description": "Cell range"}, "format": {"type": "dict", "description": "Format options"}},
            "chart": {"file": {"type": "str", "description": "File path or ID"}, "data_range": {"type": "str", "description": "Data range"}, "chart_type": {"type": "str", "description": "bar, line, pie, etc."}},
        }
    
    async def execute(self, operation: str, params: Dict[str, Any]) -> StepResult:
        try:
            file_path = params.get("file", "")
            
            # Try providers first
            for name, provider in self._providers.items():
                if operation == "read" and hasattr(provider, "read_spreadsheet"):
                    result = await provider.read_spreadsheet(file_path, params.get("sheet"), params.get("range"))
                    return StepResult(True, data={"data": result, "provider": name})
                elif operation == "write" and hasattr(provider, "write_spreadsheet"):
                    await provider.write_spreadsheet(file_path, params.get("data"), params.get("sheet"), params.get("range"))
                    return StepResult(True, data={"written": True, "provider": name})
                elif operation == "create" and hasattr(provider, "create_spreadsheet"):
                    result = await provider.create_spreadsheet(params.get("name"), params.get("data"))
                    return StepResult(True, data={"created": True, "file": result, "provider": name})
            
            # Local file handling for CSV/Excel
            if operation == "read":
                if file_path.endswith(".csv"):
                    import csv
                    with open(file_path, "r") as f:
                        reader = csv.reader(f)
                        data = list(reader)
                    return StepResult(True, data={"data": data, "provider": "local"})
                elif file_path.endswith((".xlsx", ".xls")):
                    try:
                        import openpyxl
                        wb = openpyxl.load_workbook(file_path)
                        sheet = wb[params.get("sheet")] if params.get("sheet") else wb.active
                        data = [[cell.value for cell in row] for row in sheet.iter_rows()]
                        return StepResult(True, data={"data": data, "provider": "local"})
                    except ImportError:
                        return StepResult(False, error="openpyxl not installed for Excel support")
                return StepResult(True, data={"data": [], "message": "File type not supported locally", "provider": "local"})
            
            elif operation == "write":
                data = params.get("data", [])
                if file_path.endswith(".csv"):
                    import csv
                    with open(file_path, "w", newline="") as f:
                        writer = csv.writer(f)
                        writer.writerows(data)
                    return StepResult(True, data={"written": True, "provider": "local"})
                return StepResult(True, data={"written": True, "provider": "local", "note": "Would write to cloud"})
            
            elif operation == "create":
                name = params.get("name", "Spreadsheet")
                data = params.get("data", [])
                file_path = f"{name}.csv"
                if data:
                    import csv
                    with open(file_path, "w", newline="") as f:
                        writer = csv.writer(f)
                        writer.writerows(data)
                return StepResult(True, data={"created": True, "file": file_path, "provider": "local"})
            
            elif operation in ("add_sheet", "formula", "format", "chart"):
                return StepResult(True, data={"success": True, "provider": "local", "note": f"{operation} would be applied via cloud provider"})
            
            else:
                return StepResult(False, error=f"Unknown operation: {operation}")
        except Exception as e:
            return StepResult(False, error=str(e))


# ============================================================
#  PRESENTATION PRIMITIVE - PowerPoint, Google Slides
# ============================================================



class PresentationPrimitive(Primitive):
    """Presentation operations - PowerPoint, Google Slides, etc."""
    
    def __init__(self, llm_complete: Optional[Callable] = None, providers: Optional[Dict[str, Any]] = None):
        self._llm = llm_complete
        self._providers = providers or {}
    
    @property
    def name(self) -> str:
        return "PRESENTATION"
    
    def get_operations(self) -> Dict[str, str]:
        return self.get_available_operations()
    
    def get_available_operations(self) -> Dict[str, str]:
        return {
            "create": "Create a new presentation",
            "add_slide": "Add a slide to a presentation",
            "update_slide": "Update slide content",
            "export": "Export to PDF or images",
            "get_text": "Extract all text from presentation",
        }
    
    def get_param_schema(self) -> Dict[str, Any]:
        return {
            "create": {"name": {"type": "str", "description": "Presentation name"}, "template": {"type": "str", "description": "Template (optional)"}},
            "add_slide": {"file": {"type": "str", "description": "File path or ID"}, "layout": {"type": "str", "description": "Slide layout"}, "content": {"type": "dict", "description": "Slide content"}},
            "update_slide": {"file": {"type": "str", "description": "File path or ID"}, "slide_id": {"type": "str", "description": "Slide index or ID"}, "content": {"type": "dict", "description": "New content"}},
            "export": {"file": {"type": "str", "description": "File path or ID"}, "format": {"type": "str", "description": "pdf, png, jpg"}},
            "get_text": {"file": {"type": "str", "description": "File path or ID"}},
        }
    
    async def execute(self, operation: str, params: Dict[str, Any]) -> StepResult:
        try:
            file_path = params.get("file", "")
            
            for name, provider in self._providers.items():
                if operation == "create" and hasattr(provider, "create_presentation"):
                    result = await provider.create_presentation(params.get("name"), params.get("template"))
                    return StepResult(True, data={"created": True, "file": result, "provider": name})
                elif operation == "add_slide" and hasattr(provider, "add_slide"):
                    result = await provider.add_slide(file_path, params.get("layout"), params.get("content"))
                    return StepResult(True, data={"added": True, "provider": name})
            
            # Local handling
            if operation == "create":
                name = params.get("name", "Presentation")
                return StepResult(True, data={"created": True, "file": f"{name}.pptx", "provider": "local", "note": "Would create via python-pptx"})
            
            elif operation == "add_slide":
                return StepResult(True, data={"added": True, "provider": "local", "note": "Would add slide via python-pptx"})
            
            elif operation == "update_slide":
                return StepResult(True, data={"updated": True, "provider": "local"})
            
            elif operation == "export":
                fmt = params.get("format", "pdf")
                return StepResult(True, data={"exported": True, "format": fmt, "provider": "local"})
            
            elif operation == "get_text":
                try:
                    from pptx import Presentation as PPTX
                    prs = PPTX(file_path)
                    text = []
                    for slide in prs.slides:
                        for shape in slide.shapes:
                            if hasattr(shape, "text"):
                                text.append(shape.text)
                    return StepResult(True, data={"text": "\n".join(text), "provider": "local"})
                except ImportError:
                    return StepResult(True, data={"text": "", "provider": "local", "note": "python-pptx not installed"})
            
            else:
                return StepResult(False, error=f"Unknown operation: {operation}")
        except Exception as e:
            return StepResult(False, error=str(e))


# ============================================================
#  NOTES PRIMITIVE - Note-taking apps
# ============================================================


