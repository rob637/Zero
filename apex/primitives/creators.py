"""
Creator Primitives — Generic rendering plumbing.

These primitives don't decide WHAT to create — the LLM does that.
They just render structured content into output formats.

- PDF: structured sections → PDF file
- SLIDES: slide array → PPTX file
- CHART: data + chart type → PNG image
"""

import logging
import os
import platform
import subprocess
from datetime import datetime
from pathlib import Path
from typing import Any, Callable, Dict, List, Optional

from .base import Primitive, StepResult

logger = logging.getLogger(__name__)


def _output_dir() -> Path:
    """Get or create the output directory for generated artifacts."""
    if platform.system() == "Windows":
        base = Path.home() / "Documents" / "Ziggy"
    else:
        base = Path.home() / "Ziggy"
    base.mkdir(parents=True, exist_ok=True)
    return base


def _try_open(path: str):
    """Open a file with the system default application."""
    try:
        if platform.system() == "Windows":
            os.startfile(path)
        elif platform.system() == "Darwin":
            subprocess.Popen(["open", path])
        else:
            subprocess.Popen(["xdg-open", path])
    except Exception:
        pass


# ============================================================
#  PDF PRIMITIVE — Structured content → PDF
# ============================================================

class PdfPrimitive(Primitive):
    """Render structured content into a PDF file.

    The LLM decides the content and layout. This primitive
    just draws it. Supports: headings, text, tables, images,
    columns, spacers, page breaks, full-bleed images.

    Section types the LLM can use:
      heading  — {type, text, size?, color?, align?}
      text     — {type, text, size?, color?, bold?, italic?, align?}
      table    — {type, headers, rows, col_widths?}
      image    — {type, path, width?, height?, align?, caption?}
      spacer   — {type, height}
      page     — {type}  (page break)
      hr       — {type}  (horizontal rule)
      columns  — {type, widths, left, right}  (two-column text)
    """

    @property
    def name(self) -> str:
        return "PDF"

    def get_operations(self) -> Dict[str, str]:
        return {
            "create": "Generate a PDF from structured content. Pass a title and an array of sections. Each section has a 'type' (heading, text, table, image, spacer, page, hr, columns) and type-specific fields. The AI decides what sections to include.",
        }

    def get_param_schema(self) -> Dict[str, Dict[str, Any]]:
        return {
            "create": {
                "sections": {"type": "list", "required": True, "description": "Array of section objects. Each has 'type' (heading/text/table/image/spacer/page/hr/columns) plus type-specific fields."},
                "filename": {"type": "str", "required": False, "description": "Output filename (without .pdf). Defaults to title or 'document'."},
                "title": {"type": "str", "required": False, "description": "Document title (added as first heading if provided)"},
                "page_size": {"type": "str", "required": False, "description": "Page size: letter, a4, letter-landscape, a4-landscape (default: letter)"},
                "margin": {"type": "float", "required": False, "description": "Page margin in mm (default: 15)"},
            },
        }

    async def execute(self, operation: str, params: Dict[str, Any]) -> StepResult:
        if operation != "create":
            return StepResult(False, error=f"Unknown operation: {operation}")

        try:
            from fpdf import FPDF
        except ImportError:
            return StepResult(False, error="Requires fpdf2. Run: pip install fpdf2")

        sections = params.get("sections", [])
        if not sections:
            return StepResult(False, error="Provide 'sections' array with content to render")

        title = params.get("title", "")
        filename = params.get("filename", title or "document")
        page_size = params.get("page_size", "letter")
        margin = float(params.get("margin", 15))

        # Parse page size
        orientation = "P"
        fmt = "letter"
        ps = page_size.lower()
        if "landscape" in ps:
            orientation = "L"
        if "a4" in ps:
            fmt = "A4"

        try:
            pdf = FPDF(orientation=orientation, unit="mm", format=fmt)
            pdf.set_auto_page_break(auto=True, margin=margin)
            pdf.set_margins(margin, margin, margin)
            pdf.add_page()

            # Add title as first heading if provided
            if title:
                pdf.set_font("Helvetica", "B", 24)
                pdf.set_text_color(30, 30, 30)
                pdf.cell(0, 14, title, new_x="LMARGIN", new_y="NEXT")
                pdf.ln(4)

            # Render each section
            for section in sections:
                self._render_section(pdf, section, margin)

            # Save
            out_dir = _output_dir()
            safe = "".join(c if c.isalnum() or c in " -_" else "" for c in filename).strip()
            path = str(out_dir / f"{safe or 'document'}.pdf")
            pdf.output(path)

            _try_open(path)

            return StepResult(True, data={
                "file": path,
                "sections": len(sections),
                "pages": pdf.page,
                "message": f"PDF created: {len(sections)} sections, {pdf.page} pages → {path}"
            })
        except Exception as e:
            return StepResult(False, error=f"PDF creation failed: {e}")

    def _render_section(self, pdf, section: Dict, margin: float):
        """Render a single section onto the PDF."""
        sec_type = section.get("type", "text")

        if sec_type == "page":
            pdf.add_page()
            return

        if sec_type == "spacer":
            pdf.ln(section.get("height", 8))
            return

        if sec_type == "hr":
            y = pdf.get_y()
            pdf.set_draw_color(200, 200, 200)
            pdf.line(margin, y, pdf.w - margin, y)
            pdf.ln(4)
            return

        # Set text color
        color = section.get("color", "")
        if color and color.startswith("#") and len(color) == 7:
            r = int(color[1:3], 16)
            g = int(color[3:5], 16)
            b = int(color[5:7], 16)
            pdf.set_text_color(r, g, b)
        else:
            pdf.set_text_color(0, 0, 0)

        align = section.get("align", "L").upper()[0] if section.get("align") else "L"

        if sec_type == "heading":
            size = section.get("size", 18)
            pdf.set_font("Helvetica", "B", size)
            pdf.cell(0, size * 0.6, section.get("text", ""), align=align, new_x="LMARGIN", new_y="NEXT")
            pdf.ln(3)

        elif sec_type == "text":
            size = section.get("size", 10)
            style = ""
            if section.get("bold"):
                style += "B"
            if section.get("italic"):
                style += "I"
            pdf.set_font("Helvetica", style, size)
            pdf.multi_cell(0, size * 0.5, section.get("text", ""), align=align)
            pdf.ln(2)

        elif sec_type == "table":
            self._render_table(pdf, section, margin)

        elif sec_type == "image":
            self._render_image(pdf, section, margin)

        elif sec_type == "columns":
            self._render_columns(pdf, section, margin)

    def _render_table(self, pdf, section: Dict, margin: float):
        """Render a table with headers and rows."""
        headers = section.get("headers", [])
        rows = section.get("rows", [])
        if not headers and not rows:
            return

        col_count = len(headers) if headers else (len(rows[0]) if rows else 0)
        if col_count == 0:
            return

        # Calculate column widths
        col_widths = section.get("col_widths")
        if not col_widths:
            available = pdf.w - 2 * margin
            col_widths = [available / col_count] * col_count

        row_height = 7

        # Header row
        if headers:
            pdf.set_fill_color(45, 55, 72)
            pdf.set_text_color(255, 255, 255)
            pdf.set_font("Helvetica", "B", 9)
            for i, h in enumerate(headers):
                w = col_widths[i] if i < len(col_widths) else col_widths[-1]
                pdf.cell(w, row_height, f" {str(h)}", fill=True)
            pdf.ln(row_height)

        # Data rows
        pdf.set_text_color(40, 40, 40)
        pdf.set_font("Helvetica", "", 9)
        for ri, row in enumerate(rows):
            bg = (248, 249, 250) if ri % 2 == 0 else (255, 255, 255)
            pdf.set_fill_color(*bg)
            cells = row if isinstance(row, (list, tuple)) else list(row.values())
            for i, cell in enumerate(cells):
                w = col_widths[i] if i < len(col_widths) else col_widths[-1]
                pdf.cell(w, row_height, f" {str(cell)}", fill=True)
            pdf.ln(row_height)

        pdf.ln(4)

    def _render_image(self, pdf, section: Dict, margin: float):
        """Render an image, optionally with caption."""
        path = section.get("path", "")
        if not path or not os.path.isfile(path):
            return

        img_w = section.get("width", 0)
        img_h = section.get("height", 0)
        align = section.get("align", "left").lower()

        # Default width if not specified
        if not img_w and not img_h:
            img_w = pdf.w - 2 * margin

        # Calculate x position for alignment
        x = margin
        if align == "center" and img_w:
            x = (pdf.w - img_w) / 2
        elif align == "right" and img_w:
            x = pdf.w - margin - img_w

        kwargs = {"x": x, "name": path}
        if img_w:
            kwargs["w"] = img_w
        if img_h:
            kwargs["h"] = img_h

        try:
            pdf.image(**kwargs)
        except Exception as e:
            logger.warning(f"Failed to embed image {path}: {e}")
            return

        # Caption
        caption = section.get("caption", "")
        if caption:
            pdf.set_font("Helvetica", "I", 9)
            pdf.set_text_color(100, 100, 100)
            pdf.cell(0, 5, caption, align="C" if align == "center" else "L", new_x="LMARGIN", new_y="NEXT")
            pdf.set_text_color(0, 0, 0)

        pdf.ln(4)

    def _render_columns(self, pdf, section: Dict, margin: float):
        """Render two-column text layout."""
        left_text = section.get("left", "")
        right_text = section.get("right", "")
        widths = section.get("widths", [50, 50])

        available = pdf.w - 2 * margin
        left_w = available * (widths[0] / 100)
        right_w = available * (widths[1] / 100) if len(widths) > 1 else available - left_w

        y_start = pdf.get_y()
        pdf.set_font("Helvetica", "", 10)

        # Left column
        pdf.set_x(margin)
        pdf.multi_cell(left_w - 2, 5, left_text if isinstance(left_text, str) else str(left_text))
        y_after_left = pdf.get_y()

        # Right column
        pdf.set_y(y_start)
        pdf.set_x(margin + left_w)
        pdf.multi_cell(right_w, 5, right_text if isinstance(right_text, str) else str(right_text))
        y_after_right = pdf.get_y()

        pdf.set_y(max(y_after_left, y_after_right))
        pdf.ln(4)


# ============================================================
#  SLIDES PRIMITIVE — Slide array → PPTX
# ============================================================

class SlidesPrimitive(Primitive):
    """Render a slide array into a PowerPoint file.

    The LLM decides slide content. This primitive just renders.
    Each slide has: title, bullets or body text, optional image, notes.

    Slide layouts the LLM can use:
      title   — title slide with subtitle
      content — title + bullet points
      image   — title + image + optional caption
      blank   — empty slide (for full-bleed images)
    """

    @property
    def name(self) -> str:
        return "SLIDES"

    def get_operations(self) -> Dict[str, str]:
        return {
            "create": "Generate a PowerPoint file from a slide array. Each slide has 'title', 'bullets' (list of strings), optional 'notes', optional 'layout' (title/content/image/blank). The AI decides the content.",
        }

    def get_param_schema(self) -> Dict[str, Dict[str, Any]]:
        return {
            "create": {
                "slides": {"type": "list", "required": True, "description": "Array of slide objects. Each has 'title', optional 'subtitle' (title slide only), optional 'bullets' (list), optional 'notes', optional 'image' (file path), optional 'layout' (title/content/image/blank)."},
                "filename": {"type": "str", "required": False, "description": "Output filename (without .pptx)"},
            },
        }

    async def execute(self, operation: str, params: Dict[str, Any]) -> StepResult:
        if operation != "create":
            return StepResult(False, error=f"Unknown operation: {operation}")

        try:
            from pptx import Presentation
            from pptx.util import Inches, Pt
        except ImportError:
            return StepResult(False, error="Requires python-pptx. Run: pip install python-pptx")

        slides = params.get("slides", [])
        if not slides:
            return StepResult(False, error="Provide 'slides' array")

        filename = params.get("filename", "Presentation")

        try:
            prs = Presentation()
            prs.slide_width = Inches(13.333)
            prs.slide_height = Inches(7.5)

            for i, slide_data in enumerate(slides):
                layout_name = slide_data.get("layout", "title" if i == 0 else "content")
                title = slide_data.get("title", f"Slide {i+1}")
                subtitle = slide_data.get("subtitle", "")
                bullets = slide_data.get("bullets", [])
                notes = slide_data.get("notes", "")
                image_path = slide_data.get("image", "")

                if layout_name == "title":
                    layout = prs.slide_layouts[0]
                    slide = prs.slides.add_slide(layout)
                    slide.shapes.title.text = title
                    if len(slide.placeholders) > 1 and subtitle:
                        slide.placeholders[1].text = subtitle

                elif layout_name == "blank":
                    layout = prs.slide_layouts[6]
                    slide = prs.slides.add_slide(layout)
                    if image_path and os.path.isfile(image_path):
                        slide.shapes.add_picture(
                            image_path, Inches(0), Inches(0),
                            prs.slide_width, prs.slide_height
                        )

                else:
                    # content or image layout
                    layout = prs.slide_layouts[1]
                    slide = prs.slides.add_slide(layout)
                    slide.shapes.title.text = title

                    if bullets and len(slide.placeholders) > 1:
                        tf = slide.placeholders[1].text_frame
                        tf.clear()
                        for j, bullet in enumerate(bullets):
                            if j == 0:
                                tf.text = bullet
                            else:
                                p = tf.add_paragraph()
                                p.text = bullet
                                p.level = 0

                    if image_path and os.path.isfile(image_path):
                        try:
                            slide.shapes.add_picture(
                                image_path, Inches(7), Inches(1.5),
                                width=Inches(5.5)
                            )
                        except Exception:
                            pass

                # Speaker notes
                if notes:
                    slide.notes_slide.notes_text_frame.text = notes

            out_dir = _output_dir()
            safe = "".join(c if c.isalnum() or c in " -_" else "" for c in filename).strip()
            path = str(out_dir / f"{safe or 'Presentation'}.pptx")
            prs.save(path)

            _try_open(path)

            return StepResult(True, data={
                "file": path,
                "slide_count": len(slides),
                "message": f"Presentation created: {len(slides)} slides → {path}"
            })
        except Exception as e:
            return StepResult(False, error=f"PPTX creation failed: {e}")


# ============================================================
#  CHART PRIMITIVE — Data + type → image
# ============================================================

class ChartPrimitive(Primitive):
    """Render data into chart images.

    Takes data and chart type, produces a PNG. Pure rendering —
    the LLM decides what to chart, this draws it.
    """

    @property
    def name(self) -> str:
        return "CHART"

    def get_operations(self) -> Dict[str, str]:
        return {
            "create": "Create a chart image from data. Pass data (list of dicts or numbers), chart_type (bar/line/pie/scatter/area), title, and axis fields.",
        }

    def get_param_schema(self) -> Dict[str, Dict[str, Any]]:
        return {
            "create": {
                "data": {"type": "list", "required": True, "description": "Data to chart — list of dicts or list of numbers"},
                "chart_type": {"type": "str", "required": False, "description": "Chart type: bar, line, pie, scatter, area (default: bar)"},
                "title": {"type": "str", "required": False, "description": "Chart title"},
                "x_field": {"type": "str", "required": False, "description": "Field name for X axis (for dict data)"},
                "y_field": {"type": "str", "required": False, "description": "Field name for Y axis (for dict data)"},
                "filename": {"type": "str", "required": False, "description": "Output filename (without .png)"},
            },
        }

    async def execute(self, operation: str, params: Dict[str, Any]) -> StepResult:
        if operation != "create":
            return StepResult(False, error=f"Unknown operation: {operation}")

        try:
            import matplotlib
            matplotlib.use("Agg")
            import matplotlib.pyplot as plt
        except ImportError:
            return StepResult(False, error="Requires matplotlib. Run: pip install matplotlib")

        data = params.get("data", [])
        if not data:
            return StepResult(False, error="No data provided")

        chart_type = params.get("chart_type", "bar")
        title = params.get("title", "Chart")
        x_field = params.get("x_field")
        y_field = params.get("y_field")
        filename = params.get("filename", title)

        try:
            fig, ax = plt.subplots(figsize=(10, 6))

            if isinstance(data[0], dict):
                keys = list(data[0].keys())
                x_key = x_field or keys[0]
                y_key = y_field or (keys[1] if len(keys) > 1 else keys[0])
                x_vals = [d.get(x_key, "") for d in data]
                y_vals = [float(d.get(y_key, 0)) for d in data]
            else:
                x_vals = list(range(len(data)))
                y_vals = [float(v) for v in data]
                x_key, y_key = "Index", "Value"

            if chart_type == "pie":
                ax.pie(y_vals, labels=x_vals, autopct="%1.1f%%")
            elif chart_type == "line":
                ax.plot(x_vals, y_vals, marker="o")
                ax.set_xlabel(x_key)
                ax.set_ylabel(y_key)
            elif chart_type == "scatter":
                ax.scatter(x_vals, y_vals)
                ax.set_xlabel(x_key)
                ax.set_ylabel(y_key)
            elif chart_type == "area":
                ax.fill_between(range(len(y_vals)), y_vals, alpha=0.4)
                ax.plot(range(len(y_vals)), y_vals)
                if isinstance(x_vals[0], str):
                    ax.set_xticks(range(len(x_vals)))
                    ax.set_xticklabels(x_vals, rotation=45, ha="right")
                ax.set_ylabel(y_key)
            else:  # bar
                ax.bar(range(len(x_vals)), y_vals, tick_label=[str(x) for x in x_vals])
                ax.set_xlabel(x_key)
                ax.set_ylabel(y_key)
                plt.xticks(rotation=45, ha="right")

            ax.set_title(title)
            plt.tight_layout()

            out_dir = _output_dir()
            safe = "".join(c if c.isalnum() or c in " -_" else "" for c in filename).strip()
            path = str(out_dir / f"{safe or 'chart'}.png")
            fig.savefig(path, dpi=150)
            plt.close(fig)

            _try_open(path)

            return StepResult(True, data={
                "file": path,
                "title": title,
                "chart_type": chart_type,
                "data_points": len(data),
                "message": f"Chart '{title}' → {path}"
            })
        except Exception as e:
            return StepResult(False, error=f"Chart creation failed: {e}")
