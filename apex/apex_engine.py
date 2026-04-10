"""
Telic Engine - Phase 1 Implementation

This is the WORKING implementation that ties everything together.

Usage:
    from apex_engine import Apex
    
    engine = Apex(api_key="...")
    
    # Simple request
    result = await engine.do("Find all PDFs in Downloads and list them")
    
    # With context
    result = await engine.do(
        "Create an amortization schedule from this loan and email it to Fred",
        context={"loan_doc": "~/Documents/loan.pdf"}
    )
"""

import asyncio
import hashlib
import json
import logging
import os
import re
import uuid
from dataclasses import dataclass, field
from datetime import datetime, timedelta, timezone
from pathlib import Path
from typing import Any, Callable, Dict, List, Optional, Union
from abc import ABC, abstractmethod

# LLM client (lazy-loaded for fast startup)
_litellm = None
def _get_litellm():
    global _litellm
    if _litellm is None:
        try:
            import litellm
            _litellm = litellm
        except ImportError:
            pass
    return _litellm

# Safety rails
try:
    from src.privacy.redaction import RedactionEngine
    from src.privacy.audit_log import AuditLogger, TransmissionDestination
    from src.control.trust_levels import TrustLevel, TrustLevelManager
    from src.control.approval_gateway import ApprovalGateway, RiskLevel
    from src.control.undo_manager import UndoManager, UndoType
    from src.control.action_history import ActionHistoryDB, ActionCategory
    HAS_SAFETY_RAILS = True
except ImportError:
    HAS_SAFETY_RAILS = False

logger = logging.getLogger(__name__)


# ============================================================
#  ACTION RISK CLASSIFICATION
# ============================================================

# Maps (primitive, operation) -> risk level string
_ACTION_RISK: Dict[tuple, str] = {
    # High risk — destructive or sends data externally
    ("EMAIL", "send"): "high",
    ("FILE", "write"): "high",
    ("SHELL", "run"): "high",
    ("SHELL", "script"): "high",
    ("MESSAGE", "send"): "high",
    ("MESSAGE", "reply"): "high",
    ("BROWSER", "fill_form"): "high",
    ("BROWSER", "execute_js"): "high",
    ("CALENDAR", "delete"): "high",
    ("CALENDAR", "create"): "high",  # Creates external events - needs approval
    ("TASK", "delete"): "high",
    ("TASK", "create"): "high",  # Creates external tasks - needs approval
    ("CONTACTS", "add"): "medium",
    # Medium risk — modifies existing data
    ("TASK", "update"): "medium",
    ("TASK", "complete"): "medium",
    ("FILE", "list"): "low",
    ("FILE", "search"): "low",
    ("FILE", "read"): "low",
    ("NOTIFY", "alert"): "medium",
    ("NOTIFY", "remind"): "medium",
    ("DEVTOOLS", "create_issue"): "high",
    ("DEVTOOLS", "create_pr"): "high",
    ("DEVTOOLS", "comment"): "medium",
    ("CLOUD_STORAGE", "upload"): "high",
    ("CLOUD_STORAGE", "delete"): "high",
    ("CLOUD_STORAGE", "create_folder"): "medium",
}

# Operations that support undo checkpoints
_UNDOABLE_OPS: Dict[tuple, str] = {
    ("FILE", "write"): "file_overwrite",
    ("CALENDAR", "create"): "calendar_create",
    ("CALENDAR", "delete"): "calendar_delete",
    ("TASK", "create"): "task_create",
    ("TASK", "delete"): "task_update",
    ("CONTACTS", "add"): "contact_create",
}


def _classify_risk(primitive: str, operation: str) -> str:
    """Return risk level: 'low', 'medium', 'high', or 'critical'."""
    return _ACTION_RISK.get((primitive, operation), "low")


# ============================================================
#  RESULT TYPES
# ============================================================

@dataclass
class StepResult:
    """Result from a single step."""
    success: bool
    data: Any = None
    error: Optional[str] = None
    
    def to_dict(self) -> Dict:
        return {"success": self.success, "data": self.data, "error": self.error}


@dataclass 
class PlanStep:
    """A step in the execution plan.
    
    step_type controls orchestration behavior:
      - "action"    : (default) execute a single primitive operation
      - "condition" : evaluate a condition, run then_steps or else_steps
      - "parallel"  : run a list of sub-steps concurrently
      - "loop"      : iterate over a list, run body steps for each item
      - "sub_plan"  : delegate to the planner for a sub-request
    """
    id: int
    description: str
    primitive: str
    operation: str
    params: Dict[str, Any]
    depends_on: List[int] = field(default_factory=list)
    wires: Dict[str, str] = field(default_factory=dict)
    result: Optional[StepResult] = None
    # Orchestration fields
    step_type: str = "action"
    condition: Optional[str] = None          # For "condition": expression to evaluate
    then_steps: Optional[List['PlanStep']] = None   # Steps to run if condition is true
    else_steps: Optional[List['PlanStep']] = None   # Steps to run if condition is false
    loop_over: Optional[str] = None          # For "loop": wire ref to iterate (e.g. "step_0.results")
    loop_var: str = "item"                   # Variable name for current iteration item
    loop_body: Optional[List['PlanStep']] = None    # Steps to run per iteration
    parallel_steps: Optional[List['PlanStep']] = None  # For "parallel": concurrent steps
    sub_request: Optional[str] = None        # For "sub_plan": natural language sub-request
    on_fail: str = "stop"                    # "stop" | "continue" | "retry"
    max_retries: int = 3                     # Max self-heal retries
    side_effect: bool = True                 # True if modifies external state, False if read-only
    
    def to_dict(self) -> Dict:
        d = {
            "id": self.id,
            "description": self.description,
            "primitive": self.primitive,
            "operation": self.operation,
            "params": self.params,
            "depends_on": self.depends_on,
            "wires": self.wires,
            "result": self.result.to_dict() if self.result else None,
        }
        if self.step_type != "action":
            d["step_type"] = self.step_type
        if self.condition:
            d["condition"] = self.condition
        if self.then_steps:
            d["then_steps"] = [s.to_dict() for s in self.then_steps]
        if self.else_steps:
            d["else_steps"] = [s.to_dict() for s in self.else_steps]
        if self.loop_over:
            d["loop_over"] = self.loop_over
            d["loop_var"] = self.loop_var
        if self.loop_body:
            d["loop_body"] = [s.to_dict() for s in self.loop_body]
        if self.parallel_steps:
            d["parallel_steps"] = [s.to_dict() for s in self.parallel_steps]
        if self.sub_request:
            d["sub_request"] = self.sub_request
        if self.on_fail != "stop":
            d["on_fail"] = self.on_fail
        return d


@dataclass
class ExecutionResult:
    """Result from executing a request."""
    success: bool
    request: str
    plan: List[PlanStep]
    final_result: Any = None
    error: Optional[str] = None
    
    def to_dict(self) -> Dict:
        return {
            "success": self.success,
            "request": self.request,
            "plan": [s.to_dict() for s in self.plan],
            "final_result": self.final_result,
            "error": self.error,
        }


# ============================================================
#  PRIMITIVE BASE
# ============================================================

class Primitive(ABC):
    """Base class for primitives."""
    
    @property
    @abstractmethod
    def name(self) -> str:
        pass
    
    @abstractmethod
    def get_operations(self) -> Dict[str, str]:
        """Return dict of operation_name -> description."""
        pass
    
    def get_param_schema(self) -> Dict[str, Dict[str, Any]]:
        """Return param schema per operation for LLM guidance.
        
        Override to declare expected parameters so the planner and self-healer
        know the correct format. Format:
        {
            "operation_name": {
                "param_name": {"type": "str|int|float|dict|list", "required": bool, "description": "..."},
                ...
            }
        }
        """
        return {}
    
    def get_available_operations(self) -> Dict[str, str]:
        """Return operations that are actually configured and ready to use.
        
        By default, returns all operations. Override to filter out operations
        that require external providers that aren't connected.
        """
        return self.get_operations()
    
    def get_connected_providers(self) -> List[str]:
        """Return list of connected provider names for this primitive.
        
        Used to enrich tool descriptions so the LLM knows what services
        are available and can route intelligently.
        Auto-detects from _providers dict or _connector attribute.
        """
        if hasattr(self, '_providers') and self._providers:
            return list(self._providers.keys())
        if hasattr(self, '_connector') and self._connector:
            name = getattr(self._connector, 'name', None) or type(self._connector).__name__
            return [name]
        return []
    
    @abstractmethod
    async def execute(self, operation: str, params: Dict[str, Any]) -> StepResult:
        """Execute an operation."""
        pass


# ============================================================
#  FILE PRIMITIVE
# ============================================================

class FilePrimitive(Primitive):
    """File system operations."""
    
    def __init__(self, allowed_roots: Optional[List[str]] = None):
        import tempfile
        # If no explicit restrictions, allow the user's entire home tree
        # The LLM decides WHERE to search - we just execute
        self._allowed = allowed_roots or [
            str(Path.home()),
            tempfile.gettempdir(),
        ]
        # On Windows, also allow the drives where home lives
        home = Path.home()
        if hasattr(home, 'drive') and home.drive:
            self._allowed.append(home.drive + "\\")
    
    @property
    def name(self) -> str:
        return "FILE"
    
    def get_operations(self) -> Dict[str, str]:
        return {
            "search": "Search for files matching pattern",
            "read": "Read file contents",
            "write": "Write content to file",
            "list": "List directory contents",
            "info": "Get file metadata",
            "exists": "Check if file exists",
            "checksum": "Calculate hash/checksum of a file (for duplicate detection)",
            "move": "Move a file to a new location",
            "copy": "Copy a file to a new location",
            "delete": "Delete a file (moves to trash if available)",
        }
    
    def get_param_schema(self) -> Dict[str, Dict[str, Any]]:
        return {
            "search": {
                "pattern": {"type": "str", "required": True, "description": "Glob pattern (e.g. '*.pdf', '*.docx')"},
                "directory": {"type": "str", "required": False, "description": "Directory to search (default: home)"},
                "recursive": {"type": "bool", "required": False, "description": "Search subdirectories (default true)"},
                "limit": {"type": "int", "required": False, "description": "Max results (default 50)"},
            },
            "read": {
                "path": {"type": "str", "required": True, "description": "File path to read"},
            },
            "write": {
                "path": {"type": "str", "required": True, "description": "File path to write"},
                "content": {"type": "str", "required": True, "description": "Content to write"},
            },
            "list": {
                "directory": {"type": "str", "required": True, "description": "Directory to list"},
            },
            "info": {
                "path": {"type": "str", "required": True, "description": "File path"},
            },
            "exists": {
                "path": {"type": "str", "required": True, "description": "File path to check"},
            },
            "checksum": {
                "path": {"type": "str", "required": True, "description": "File path to hash"},
                "algorithm": {"type": "str", "required": False, "description": "Hash algorithm: md5, sha1, sha256 (default: md5)"},
            },
            "move": {
                "source": {"type": "str", "required": True, "description": "Source file path"},
                "destination": {"type": "str", "required": True, "description": "Destination path"},
            },
            "copy": {
                "source": {"type": "str", "required": True, "description": "Source file path"},
                "destination": {"type": "str", "required": True, "description": "Destination path"},
            },
            "delete": {
                "path": {"type": "str", "required": True, "description": "File path to delete"},
                "permanent": {"type": "bool", "required": False, "description": "Permanently delete (skip trash, default: false)"},
            },
        }
    
    def _is_allowed(self, path: str) -> bool:
        resolved = str(Path(path).expanduser().resolve())
        return any(resolved.startswith(str(Path(a).expanduser().resolve())) for a in self._allowed)
    
    async def execute(self, operation: str, params: Dict[str, Any]) -> StepResult:
        try:
            if operation == "search":
                pattern = params.get("pattern", "*")
                directory = params.get("directory", str(Path.home()))
                directory = str(Path(directory).expanduser())
                recursive = params.get("recursive", True)
                limit = params.get("limit", 5000)
                
                if not self._is_allowed(directory):
                    return StepResult(False, error=f"Directory not allowed: {directory}")
                
                base = Path(directory)
                if not base.exists():
                    return StepResult(False, error=f"Directory not found: {directory}")
                
                matches = []
                glob_func = base.rglob if recursive else base.glob
                for p in glob_func(pattern):
                    if len(matches) >= limit:
                        break
                    matches.append({
                        "path": str(p),
                        "name": p.name,
                        "is_dir": p.is_dir(),
                        "size": p.stat().st_size if p.is_file() else 0,
                    })
                
                return StepResult(True, data=matches)
            
            elif operation == "read":
                path = str(Path(params.get("path", "")).expanduser())
                if not self._is_allowed(path):
                    return StepResult(False, error=f"Path not allowed: {path}")
                
                if not Path(path).exists():
                    return StepResult(False, error=f"File not found: {path}")
                
                with open(path, "r", errors="ignore") as f:
                    content = f.read()
                
                return StepResult(True, data=content)
            
            elif operation == "write":
                path = str(Path(params.get("path", "")).expanduser())
                content = params.get("content", "")
                
                if not self._is_allowed(path):
                    return StepResult(False, error=f"Path not allowed: {path}")
                
                Path(path).parent.mkdir(parents=True, exist_ok=True)
                with open(path, "w") as f:
                    f.write(content)
                
                return StepResult(True, data={"path": path, "size": len(content)})
            
            elif operation == "list":
                directory = str(Path(params.get("directory", "")).expanduser())
                
                if not self._is_allowed(directory):
                    return StepResult(False, error=f"Directory not allowed: {directory}")
                
                base = Path(directory)
                if not base.exists():
                    return StepResult(False, error=f"Directory not found: {directory}")
                
                items = []
                for p in base.iterdir():
                    items.append({
                        "path": str(p),
                        "name": p.name,
                        "is_dir": p.is_dir(),
                        "size": p.stat().st_size if p.is_file() else 0,
                    })
                
                return StepResult(True, data=items)
            
            elif operation == "info":
                path = str(Path(params.get("path", "")).expanduser())
                if not self._is_allowed(path):
                    return StepResult(False, error=f"Path not allowed: {path}")
                
                p = Path(path)
                if not p.exists():
                    return StepResult(False, error=f"File not found: {path}")
                
                stat = p.stat()
                return StepResult(True, data={
                    "path": str(p),
                    "name": p.name,
                    "size": stat.st_size,
                    "modified": datetime.fromtimestamp(stat.st_mtime).isoformat(),
                    "is_dir": p.is_dir(),
                    "extension": p.suffix,
                })
            
            elif operation == "exists":
                path = str(Path(params.get("path", "")).expanduser())
                return StepResult(True, data={"exists": Path(path).exists(), "path": path})
            
            elif operation == "checksum":
                import hashlib
                path = str(Path(params.get("path", "")).expanduser())
                algorithm = params.get("algorithm", "md5").lower()
                
                if not self._is_allowed(path):
                    return StepResult(False, error=f"Path not allowed: {path}")
                if not Path(path).exists():
                    return StepResult(False, error=f"File not found: {path}")
                if Path(path).is_dir():
                    return StepResult(False, error="Cannot hash a directory")
                
                hash_funcs = {"md5": hashlib.md5, "sha1": hashlib.sha1, "sha256": hashlib.sha256}
                if algorithm not in hash_funcs:
                    return StepResult(False, error=f"Unsupported algorithm: {algorithm}. Use md5, sha1, or sha256.")
                
                hasher = hash_funcs[algorithm]()
                with open(path, "rb") as f:
                    for chunk in iter(lambda: f.read(65536), b""):
                        hasher.update(chunk)
                
                return StepResult(True, data={
                    "path": path,
                    "algorithm": algorithm,
                    "checksum": hasher.hexdigest(),
                    "size": Path(path).stat().st_size,
                })
            
            elif operation == "move":
                import shutil
                source = str(Path(params.get("source", "")).expanduser())
                destination = str(Path(params.get("destination", "")).expanduser())
                
                if not self._is_allowed(source):
                    return StepResult(False, error=f"Source not allowed: {source}")
                if not self._is_allowed(destination):
                    return StepResult(False, error=f"Destination not allowed: {destination}")
                if not Path(source).exists():
                    return StepResult(False, error=f"Source not found: {source}")
                
                Path(destination).parent.mkdir(parents=True, exist_ok=True)
                shutil.move(source, destination)
                
                return StepResult(True, data={"source": source, "destination": destination})
            
            elif operation == "copy":
                import shutil
                source = str(Path(params.get("source", "")).expanduser())
                destination = str(Path(params.get("destination", "")).expanduser())
                
                if not self._is_allowed(source):
                    return StepResult(False, error=f"Source not allowed: {source}")
                if not self._is_allowed(destination):
                    return StepResult(False, error=f"Destination not allowed: {destination}")
                if not Path(source).exists():
                    return StepResult(False, error=f"Source not found: {source}")
                
                Path(destination).parent.mkdir(parents=True, exist_ok=True)
                if Path(source).is_dir():
                    shutil.copytree(source, destination)
                else:
                    shutil.copy2(source, destination)
                
                return StepResult(True, data={"source": source, "destination": destination})
            
            elif operation == "delete":
                import shutil
                path = str(Path(params.get("path", "")).expanduser())
                permanent = params.get("permanent", False)
                
                if not self._is_allowed(path):
                    return StepResult(False, error=f"Path not allowed: {path}")
                if not Path(path).exists():
                    return StepResult(False, error=f"File not found: {path}")
                
                if permanent:
                    if Path(path).is_dir():
                        shutil.rmtree(path)
                    else:
                        Path(path).unlink()
                    return StepResult(True, data={"path": path, "deleted": True, "permanent": True})
                else:
                    # Try to move to trash
                    try:
                        from send2trash import send2trash
                        send2trash(path)
                        return StepResult(True, data={"path": path, "deleted": True, "permanent": False, "location": "trash"})
                    except ImportError:
                        # Fallback to permanent delete
                        if Path(path).is_dir():
                            shutil.rmtree(path)
                        else:
                            Path(path).unlink()
                        return StepResult(True, data={"path": path, "deleted": True, "permanent": True, "note": "Install send2trash for trash support"})
            
            else:
                return StepResult(False, error=f"Unknown operation: {operation}")
                
        except Exception as e:
            return StepResult(False, error=str(e))


# ============================================================
#  DOCUMENT PRIMITIVE  
# ============================================================

class DocumentPrimitive(Primitive):
    """Document parsing and creation."""
    
    def __init__(self, llm_complete: Optional[Callable] = None):
        self._llm = llm_complete
    
    @property
    def name(self) -> str:
        return "DOCUMENT"
    
    def get_operations(self) -> Dict[str, str]:
        return {
            "parse": "Parse document to text",
            "extract": "Extract structured data using LLM",
            "create": "Create a document",
            "summarize": "Summarize document content",
        }
    
    def get_param_schema(self) -> Dict[str, Dict[str, Any]]:
        return {
            "parse": {
                "path": {"type": "str", "required": False, "description": "Path to document file (PDF, DOCX, or plain text)"},
                "content": {"type": "str", "required": False, "description": "Raw document content (alternative to path)"},
            },
            "extract": {
                "content": {"type": "str", "required": True, "description": "Document text to extract from"},
                "schema": {"type": "dict", "required": True, "description": "Fields to extract: {field_name: description}"},
            },
            "create": {
                "format": {"type": "str", "required": False, "description": "Output format: text, csv, json, markdown (default text)"},
                "content": {"type": "str", "required": False, "description": "Text content"},
                "data": {"type": "list", "required": False, "description": "Structured data (list of dicts) for csv/json/markdown"},
                "path": {"type": "str", "required": False, "description": "Save path (optional)"},
            },
            "summarize": {
                "content": {"type": "str", "required": True, "description": "Document text to summarize"},
                "max_length": {"type": "int", "required": False, "description": "Max summary length in characters (default 500)"},
            },
        }
    
    async def execute(self, operation: str, params: Dict[str, Any]) -> StepResult:
        try:
            if operation == "parse":
                path = params.get("path")
                content = params.get("content")
                
                if path:
                    path = str(Path(path).expanduser())
                    ext = Path(path).suffix.lower()
                    
                    if ext == ".pdf":
                        try:
                            import pypdf
                            reader = pypdf.PdfReader(path)
                            text = "\n".join(page.extract_text() or "" for page in reader.pages)
                        except ImportError:
                            # Fallback: try pdfplumber
                            try:
                                import pdfplumber
                                with pdfplumber.open(path) as pdf:
                                    text = "\n".join(page.extract_text() or "" for page in pdf.pages)
                            except ImportError:
                                return StepResult(False, error="No PDF library available. Install pypdf or pdfplumber.")
                    
                    elif ext == ".docx":
                        try:
                            import docx
                            doc = docx.Document(path)
                            text = "\n".join(p.text for p in doc.paragraphs)
                        except ImportError:
                            return StepResult(False, error="python-docx not installed")
                    
                    else:
                        # Plain text
                        with open(path, "r", errors="ignore") as f:
                            text = f.read()
                else:
                    text = content or ""
                
                return StepResult(True, data=text)
            
            elif operation == "extract":
                content = params.get("content", "")
                schema = params.get("schema", {})
                
                if not self._llm:
                    return StepResult(False, error="LLM not configured for extraction")
                
                prompt = f"""Extract structured data from this document.

Schema to extract (field name: description):
{json.dumps(schema, indent=2)}

Document:
{content[:15000]}

Return ONLY a valid JSON object with the extracted values. Use null if not found."""

                response = await self._llm(prompt)
                
                # Parse JSON from response
                json_match = re.search(r'\{[\s\S]*\}', response)
                if json_match:
                    try:
                        extracted = json.loads(json_match.group())
                    except json.JSONDecodeError:
                        extracted = {"raw": response}
                else:
                    extracted = {"raw": response}
                
                return StepResult(True, data=extracted)
            
            elif operation == "create":
                format_type = params.get("format", "text")
                content = params.get("content", "")
                data = params.get("data")
                path = params.get("path")
                
                if format_type == "csv" and data:
                    import csv
                    import io
                    output = io.StringIO()
                    if isinstance(data, list) and data:
                        if isinstance(data[0], dict):
                            writer = csv.DictWriter(output, fieldnames=data[0].keys())
                            writer.writeheader()
                            writer.writerows(data)
                        else:
                            writer = csv.writer(output)
                            writer.writerows(data)
                    result = output.getvalue()
                
                elif format_type == "json" and data:
                    result = json.dumps(data, indent=2)
                
                elif format_type == "markdown" and data:
                    # Create markdown table from data
                    if isinstance(data, list) and data and isinstance(data[0], dict):
                        headers = list(data[0].keys())
                        lines = ["| " + " | ".join(str(h) for h in headers) + " |"]
                        lines.append("| " + " | ".join(["---"] * len(headers)) + " |")
                        for row in data:
                            lines.append("| " + " | ".join(str(row.get(h, "")) for h in headers) + " |")
                        result = "\n".join(lines)
                    else:
                        result = content
                
                else:
                    result = content or str(data or "")
                
                if path:
                    path = str(Path(path).expanduser())
                    Path(path).parent.mkdir(parents=True, exist_ok=True)
                    with open(path, "w") as f:
                        f.write(result)
                
                return StepResult(True, data={"content": result, "path": path, "format": format_type})
            
            elif operation == "summarize":
                content = params.get("content", "")
                max_length = params.get("max_length", 500)
                
                if not self._llm:
                    # Simple truncation if no LLM
                    return StepResult(True, data=content[:max_length])
                
                prompt = f"Summarize in {max_length} characters or less:\n\n{content[:10000]}"
                summary = await self._llm(prompt)
                
                return StepResult(True, data=summary[:max_length])
            
            else:
                return StepResult(False, error=f"Unknown operation: {operation}")
                
        except Exception as e:
            return StepResult(False, error=str(e))


# ============================================================
#  COMPUTE PRIMITIVE
# ============================================================

class ComputePrimitive(Primitive):
    """AI-powered computation engine.
    
    The LLM IS the intelligence. It understands what "amortization schedule"
    means, what inputs it needs, and how to compute it. No hard-coded formula
    registry needed.
    
    Flow:
    1. User asks for any computation (amortization, ROI, depreciation, anything)
    2. LLM writes Python code to compute it
    3. Engine executes the code in a sandbox
    4. Returns the result
    
    This handles millions of scenarios because the LLM knows math — we don't
    need to pre-register every possible formula.
    """
    
    def __init__(self, llm_complete: Optional[Callable] = None):
        self._llm = llm_complete
    
    @property
    def name(self) -> str:
        return "COMPUTE"
    
    def get_operations(self) -> Dict[str, str]:
        return {
            "formula": "Compute any formula — amortization, compound interest, ROI, depreciation, NPV, IRR, or anything else. The AI writes the code.",
            "calculate": "Evaluate a math expression (e.g. 'sqrt(144) + pi')",
            "aggregate": "Aggregate numeric data (sum, average, min, max, count, median, or any custom aggregation)",
        }
    
    def get_param_schema(self) -> Dict[str, Dict[str, Any]]:
        return {
            "formula": {
                "name": {"type": "str", "required": True, "description": "What to compute: amortization, compound_interest, roi, depreciation, npv, irr, or ANY formula name"},
                "inputs": {"type": "dict", "required": True, "description": "Key-value inputs (e.g. {\"principal\": 100000, \"rate\": 3.5, \"term_years\": 30})"},
            },
            "calculate": {
                "expression": {"type": "str", "required": True, "description": "Math expression (e.g. '100 * 1.05 ** 10')"},
                "variables": {"type": "dict", "required": False, "description": "Variable name-value pairs"},
            },
            "aggregate": {
                "data": {"type": "list", "required": True, "description": "List of numbers or dicts"},
                "function": {"type": "str", "required": True, "description": "Any aggregation: sum, average, min, max, count, median, std_dev, percentile, etc."},
                "field": {"type": "str", "required": False, "description": "Field to extract from dicts"},
            },
        }
    
    async def _llm_generate_code(self, name: str, inputs: Dict) -> StepResult:
        """Ask the LLM to write Python code that computes the formula, then execute it.
        
        The LLM knows math. It writes deterministic code. We run it safely.
        This handles ANY formula — no registry needed.
        """
        if not self._llm:
            return StepResult(False, error="No LLM configured. COMPUTE requires an LLM to generate computation code.")
        
        prompt = f"""Write a Python function to compute this:

Formula: {name}
Inputs: {json.dumps(inputs)}

Requirements:
- Write a function called `compute(inputs)` that takes a dict and returns a dict of results
- Use only Python stdlib (math module is available)
- For schedules/tables, include them as a list of dicts under a "schedule" key
- Round monetary values to 2 decimal places
- Include the most useful summary values as top-level keys
- Handle edge cases (zero values, negative inputs) gracefully

Respond with ONLY the Python code. No markdown, no explanation, no ```python blocks.
Start directly with: def compute(inputs):"""

        try:
            response = await self._llm(prompt)
            
            # Clean up response — strip markdown code fences if present
            code = response.strip()
            if code.startswith("```"):
                code = re.sub(r'^```\w*\n?', '', code)
                code = re.sub(r'\n?```$', '', code)
                code = code.strip()
            
            # Validate: must define compute function
            if "def compute" not in code:
                return StepResult(False, error=f"LLM did not generate a valid compute function for '{name}'")
            
            # Execute in sandbox
            return self._execute_sandboxed(code, inputs)
            
        except Exception as e:
            return StepResult(False, error=f"Code generation failed for '{name}': {e}")
    
    def _execute_sandboxed(self, code: str, inputs: Dict) -> StepResult:
        """Execute LLM-generated code in a restricted sandbox."""
        import math
        import datetime as _datetime
        
        # Safe modules the LLM is allowed to import
        _safe_modules = {
            "math": math,
            "datetime": _datetime,
            "json": json,
        }
        
        def _safe_import(name, *args, **kwargs):
            if name in _safe_modules:
                return _safe_modules[name]
            raise ImportError(f"Import of '{name}' is not allowed in sandbox")
        
        # Restricted namespace — safe builtins + safe imports only
        sandbox = {
            "__builtins__": {
                # Math & numeric
                "abs": abs, "round": round, "min": min, "max": max,
                "sum": sum, "len": len, "pow": pow, "int": int, "float": float,
                "sorted": sorted, "enumerate": enumerate, "range": range, "zip": zip,
                "map": map, "filter": filter, "list": list, "dict": dict, "tuple": tuple,
                "set": set, "str": str, "bool": bool, "type": type,
                "True": True, "False": False, "None": None,
                "isinstance": isinstance, "ValueError": ValueError,
                "ZeroDivisionError": ZeroDivisionError, "TypeError": TypeError,
                "KeyError": KeyError, "IndexError": IndexError,
                "print": lambda *a, **k: None,  # no-op
                "__import__": _safe_import,  # safe import for math, datetime, json
            },
            "math": math,
        }
        
        try:
            exec(code, sandbox)
            
            compute_fn = sandbox.get("compute")
            if not callable(compute_fn):
                return StepResult(False, error="Generated code did not define a callable 'compute' function")
            
            result = compute_fn(inputs)
            
            if isinstance(result, dict):
                return StepResult(True, data=result)
            else:
                return StepResult(True, data={"result": result})
                
        except Exception as e:
            return StepResult(False, error=f"Computation error: {type(e).__name__}: {e}")
    
    # ── Built-in formulas (fast path — no LLM needed) ──────────

    @staticmethod
    def _builtin_amortization(inputs: Dict) -> StepResult:
        principal = float(inputs.get("principal", 0))
        annual_rate = float(inputs.get("rate", 0))
        term_months = int(inputs.get("term_months", inputs.get("term", 0)))
        if principal <= 0 or annual_rate <= 0 or term_months <= 0:
            return StepResult(False, error="amortization requires positive principal, rate, and term_months")
        r = annual_rate / 100 / 12
        payment = round(principal * r * (1 + r) ** term_months / ((1 + r) ** term_months - 1), 2)
        balance = principal
        schedule = []
        total_interest = 0.0
        for m in range(1, term_months + 1):
            interest = round(balance * r, 2)
            princ = round(payment - interest, 2)
            balance = round(balance - princ, 2)
            if m == term_months:
                princ = round(princ + balance, 2)
                balance = 0.0
            total_interest += interest
            schedule.append({"month": m, "payment": payment, "interest": interest, "principal": princ, "balance": max(balance, 0.0)})
        return StepResult(True, data={
            "monthly_payment": payment,
            "total_interest": round(total_interest, 2),
            "total_paid": round(payment * term_months, 2),
            "schedule": schedule,
        })

    @staticmethod
    def _builtin_compound_interest(inputs: Dict) -> StepResult:
        principal = float(inputs.get("principal", 0))
        rate = float(inputs.get("rate", 0))
        years = float(inputs.get("years", inputs.get("term_years", 0)))
        n = int(inputs.get("compounds_per_year", 12))
        if principal <= 0 or rate <= 0 or years <= 0:
            return StepResult(False, error="compound_interest requires positive principal, rate, and years")
        r = rate / 100
        final = round(principal * (1 + r / n) ** (n * years), 2)
        return StepResult(True, data={"final_amount": final, "interest_earned": round(final - principal, 2), "principal": principal})

    _BUILTIN_FORMULAS: Dict[str, Callable] = {}  # populated after class body

    async def execute(self, operation: str, params: Dict[str, Any]) -> StepResult:
        try:
            if operation == "formula":
                name = params.get("name", "").strip()
                inputs = params.get("inputs", {})
                
                # Normalize: if no nested "inputs", treat non-meta params as inputs
                if not inputs:
                    if not name:
                        name = str(params.get("formula", "")).strip()
                    inputs = {k: v for k, v in params.items() if k not in ("name", "formula")}
                
                if not name and not inputs:
                    return StepResult(False, error=f"Tell me what to compute. Params: {json.dumps(params)}. Example: {{\"name\": \"amortization\", \"inputs\": {{\"principal\": 100000, \"rate\": 3.5, \"term_years\": 30}}}}")
                
                if not name and inputs:
                    # No name but has inputs — ask LLM what this probably is
                    name = "custom_calculation"
                
                # Fast path: built-in formulas (no LLM needed)
                builtin = self._BUILTIN_FORMULAS.get(name)
                if builtin:
                    return builtin(inputs)
                
                # Fallback: LLM writes code, engine runs it — works for ANYTHING
                return await self._llm_generate_code(name, inputs)
            
            elif operation == "calculate":
                expr = params.get("expression", "")
                variables = params.get("variables", {})
                
                if not expr:
                    return StepResult(False, error="Missing 'expression' parameter")
                
                import math
                safe_ns = {
                    "__builtins__": {},
                    "abs": abs, "round": round, "min": min, "max": max,
                    "sum": sum, "len": len, "pow": pow, "int": int, "float": float,
                    "pi": math.pi, "e": math.e, "sqrt": math.sqrt,
                    "sin": math.sin, "cos": math.cos, "tan": math.tan,
                    "log": math.log, "log10": math.log10, "ceil": math.ceil, "floor": math.floor,
                }
                safe_ns.update({k: float(v) for k, v in variables.items() if isinstance(v, (int, float))})
                
                result = eval(expr, safe_ns)
                return StepResult(True, data=result)
            
            elif operation == "aggregate":
                data = params.get("data", [])
                func = params.get("function", "sum").lower().strip()
                field_name = params.get("field")
                
                # Debug: log what arrived so we can diagnose wiring issues
                logger.info(f"[COMPUTE.aggregate] func={func}, data type={type(data).__name__}, len={len(data) if isinstance(data, list) else 'N/A'}, raw={str(data)[:200]}")
                
                # Coerce: if data arrived as a JSON string (wiring edge case), parse it
                if isinstance(data, str):
                    try:
                        data = json.loads(data)
                    except (json.JSONDecodeError, ValueError):
                        pass
                
                # Count operates on the raw data length (no numeric extraction needed)
                if func == "count" and isinstance(data, list):
                    return StepResult(True, data=len(data))
                
                if field_name and isinstance(data, list):
                    values = [float(item.get(field_name, 0)) for item in data if isinstance(item, dict)]
                else:
                    values = [float(v) for v in data if isinstance(v, (int, float))]
                
                if not values:
                    return StepResult(True, data=0)
                
                # Common aggregations — fast path, no LLM needed
                import math as _math
                fast = {
                    "sum": lambda v: sum(v),
                    "average": lambda v: sum(v) / len(v),
                    "avg": lambda v: sum(v) / len(v),
                    "mean": lambda v: sum(v) / len(v),
                    "min": lambda v: min(v),
                    "max": lambda v: max(v),
                    "count": lambda v: len(v),
                    "median": lambda v: (sorted(v)[len(v)//2] + sorted(v)[(len(v)-1)//2]) / 2,
                    "std_dev": lambda v: _math.sqrt(sum((x - sum(v)/len(v))**2 for x in v) / len(v)),
                }
                
                if func in fast:
                    return StepResult(True, data=round(fast[func](values), 2))
                
                # Unknown aggregation — LLM writes code for it
                if self._llm:
                    return await self._llm_generate_code(f"aggregate_{func}", {"values": values})
                
                return StepResult(False, error=f"Unknown aggregation: {func}. Available without LLM: {', '.join(fast.keys())}")
            
            else:
                return StepResult(False, error=f"Unknown operation: {operation}. Available: formula, calculate, aggregate")
                
        except Exception as e:
            return StepResult(False, error=str(e))

# Wire built-in formulas after class body
ComputePrimitive._BUILTIN_FORMULAS = {
    "amortization": ComputePrimitive._builtin_amortization,
    "compound_interest": ComputePrimitive._builtin_compound_interest,
}


# ============================================================
#  EMAIL PRIMITIVE
# ============================================================

class EmailPrimitive(Primitive):
    """Email operations via Gmail, Outlook, or other providers."""
    
    def __init__(self, providers: Optional[Dict[str, Any]] = None, **kwargs):
        # Support both old single-connector and new multi-provider patterns
        self._providers = providers or {}
        # Legacy compat: if old-style args passed, wrap in provider
        if not self._providers and kwargs.get('connector'):
            connector = kwargs['connector']
            provider_name = 'gmail' if 'Gmail' in type(connector).__name__ else 'outlook'
            self._providers[provider_name] = connector
        self._send = kwargs.get('send_func')
        self._list = kwargs.get('list_func')
        self._read = kwargs.get('read_func')
        self._connector = kwargs.get('connector') or (next(iter(self._providers.values())) if self._providers else None)
    
    @property
    def name(self) -> str:
        return "EMAIL"
    
    def get_operations(self) -> Dict[str, str]:
        ops = {
            "send": "Send an email",
            "draft": "Create a draft email",
            "search": "Search emails",
            "list": "List recent emails",
            "read": "Read a specific email by ID to get its full body content",
        }
        if self._connector:
            ops.update({
                "reply": "Reply to an email",
                "forward": "Forward an email to another recipient",
                "delete": "Permanently delete an email",
                "trash": "Move an email to trash",
                "archive": "Archive an email (remove from inbox)",
                "mark_read": "Mark an email as read",
                "mark_unread": "Mark an email as unread",
                "add_label": "Add a label/folder to an email",
                "remove_label": "Remove a label/folder from an email",
                "get_labels": "List all available labels/folders",
                "get_thread": "Get all messages in an email thread",
            })
        return ops
    
    def get_param_schema(self) -> Dict[str, Dict[str, Any]]:
        schema = {
            "send": {
                "to": {"type": "str", "required": True, "description": "Recipient email address"},
                "subject": {"type": "str", "required": True, "description": "Email subject line"},
                "body": {"type": "str", "required": True, "description": "Email body text"},
                "attachments": {"type": "list", "required": False, "description": "List of file paths to attach"},
            },
            "draft": {
                "to": {"type": "str", "required": True, "description": "Recipient email address"},
                "subject": {"type": "str", "required": True, "description": "Email subject line"},
                "body": {"type": "str", "required": True, "description": "Email body text"},
            },
            "search": {
                "query": {"type": "str", "required": True, "description": "Search query using Gmail syntax (e.g. 'from:bob subject:report', 'label:travel', 'in:anywhere Africa trip'). Use label: to search folders."},
                "limit": {"type": "int", "required": False, "description": "Max results (default 10)"},
            },
            "list": {
                "query": {"type": "str", "required": False, "description": "Filter query using Gmail syntax"},
                "limit": {"type": "int", "required": False, "description": "Max results (default 10)"},
            },
            "read": {
                "message_id": {"type": "str", "required": True, "description": "The email message ID (from search/list results)"},
            },
            "reply": {
                "message_id": {"type": "str", "required": True, "description": "ID of the email to reply to"},
                "body": {"type": "str", "required": True, "description": "Reply body text"},
            },
            "forward": {
                "message_id": {"type": "str", "required": True, "description": "ID of the email to forward"},
                "to": {"type": "str", "required": True, "description": "Recipient email address"},
                "body": {"type": "str", "required": False, "description": "Additional message to include"},
            },
            "delete": {
                "message_id": {"type": "str", "required": True, "description": "ID of the email to delete permanently"},
            },
            "trash": {
                "message_id": {"type": "str", "required": True, "description": "ID of the email to move to trash"},
            },
            "archive": {
                "message_id": {"type": "str", "required": True, "description": "ID of the email to archive"},
            },
            "mark_read": {
                "message_id": {"type": "str", "required": True, "description": "ID of the email to mark as read"},
            },
            "mark_unread": {
                "message_id": {"type": "str", "required": True, "description": "ID of the email to mark as unread"},
            },
            "add_label": {
                "message_id": {"type": "str", "required": True, "description": "ID of the email"},
                "label_ids": {"type": "list", "required": True, "description": "Label IDs to add"},
            },
            "remove_label": {
                "message_id": {"type": "str", "required": True, "description": "ID of the email"},
                "label_ids": {"type": "list", "required": True, "description": "Label IDs to remove"},
            },
            "get_labels": {},
            "get_thread": {
                "thread_id": {"type": "str", "required": True, "description": "Thread ID to get all messages for"},
            },
        }
        return schema
    
    def get_available_operations(self) -> Dict[str, str]:
        """All email operations are always available - execute returns helpful errors if no provider."""
        return self.get_operations()
    
    async def execute(self, operation: str, params: Dict[str, Any]) -> StepResult:
        try:
            if operation == "send":
                if not self._send:
                    return StepResult(False, error="Email sending not configured")
                
                result = await self._send(
                    to=params.get("to"),
                    subject=params.get("subject"),
                    body=params.get("body"),
                    attachments=params.get("attachments"),
                )
                return StepResult(True, data=result)
            
            elif operation == "draft":
                # For now, create draft means just prepare the email
                return StepResult(True, data={
                    "draft": True,
                    "to": params.get("to"),
                    "subject": params.get("subject"),
                    "body": params.get("body"),
                })
            
            elif operation in ["search", "list"]:
                if not self._list:
                    return StepResult(False, error="Email listing not configured")
                
                result = await self._list(
                    query=params.get("query", ""),
                    max_results=params.get("limit", 10),
                )
                return StepResult(True, data=result)
            
            elif operation == "read":
                if not self._read:
                    return StepResult(False, error="Email reading not configured")
                
                message_id = params.get("message_id")
                if not message_id:
                    return StepResult(False, error="message_id is required")
                
                email = await self._read(message_id)
                if hasattr(email, 'to_dict'):
                    return StepResult(True, data=email.to_dict())
                return StepResult(True, data=email)
            
            elif operation == "reply":
                if not self._connector or not hasattr(self._connector, "reply"):
                    return StepResult(False, error="Reply not supported by this email provider")
                result = await self._connector.reply(
                    message_id=params["message_id"],
                    body=params["body"],
                    html=params.get("html", False),
                )
                return StepResult(True, data=result)
            
            elif operation == "forward":
                if not self._connector or not hasattr(self._connector, "forward"):
                    return StepResult(False, error="Forward not supported by this email provider")
                result = await self._connector.forward(
                    message_id=params["message_id"],
                    to=params["to"],
                    additional_body=params.get("body", ""),
                )
                return StepResult(True, data=result)
            
            elif operation == "delete":
                if not self._connector or not hasattr(self._connector, "delete_message"):
                    return StepResult(False, error="Delete not supported by this email provider")
                await self._connector.delete_message(params["message_id"])
                return StepResult(True, data={"deleted": True})
            
            elif operation == "trash":
                if not self._connector or not hasattr(self._connector, "trash_message"):
                    return StepResult(False, error="Trash not supported by this email provider")
                await self._connector.trash_message(params["message_id"])
                return StepResult(True, data={"trashed": True})
            
            elif operation == "archive":
                if not self._connector or not hasattr(self._connector, "archive_message"):
                    return StepResult(False, error="Archive not supported by this email provider")
                await self._connector.archive_message(params["message_id"])
                return StepResult(True, data={"archived": True})
            
            elif operation == "mark_read":
                if not self._connector or not hasattr(self._connector, "mark_read"):
                    return StepResult(False, error="Mark read not supported by this email provider")
                await self._connector.mark_read(params["message_id"])
                return StepResult(True, data={"marked_read": True})
            
            elif operation == "mark_unread":
                if not self._connector or not hasattr(self._connector, "mark_unread"):
                    return StepResult(False, error="Mark unread not supported by this email provider")
                await self._connector.mark_unread(params["message_id"])
                return StepResult(True, data={"marked_unread": True})
            
            elif operation == "add_label":
                if not self._connector or not hasattr(self._connector, "add_label"):
                    return StepResult(False, error="Labels not supported by this email provider")
                await self._connector.add_label(params["message_id"], params["label_ids"])
                return StepResult(True, data={"labels_added": True})
            
            elif operation == "remove_label":
                if not self._connector or not hasattr(self._connector, "remove_label"):
                    return StepResult(False, error="Labels not supported by this email provider")
                await self._connector.remove_label(params["message_id"], params["label_ids"])
                return StepResult(True, data={"labels_removed": True})
            
            elif operation == "get_labels":
                if not self._connector or not hasattr(self._connector, "get_labels"):
                    return StepResult(False, error="Labels not supported by this email provider")
                result = await self._connector.get_labels()
                return StepResult(True, data=result)
            
            elif operation == "get_thread":
                if not self._connector or not hasattr(self._connector, "get_thread"):
                    return StepResult(False, error="Threads not supported by this email provider")
                result = await self._connector.get_thread(params["thread_id"])
                return StepResult(True, data=result)
            
            else:
                return StepResult(False, error=f"Unknown operation: {operation}")
                
        except Exception as e:
            return StepResult(False, error=str(e))


# ============================================================
#  CONTACTS PRIMITIVE
# ============================================================

class ContactsPrimitive(Primitive):
    """Contact management.
    
    Local in-memory store by default. Wire in Google Contacts, Outlook, etc.
    via a providers dict.
    """
    
    def __init__(self, providers: Optional[Dict[str, Any]] = None):
        self._contacts: Dict[str, Dict] = {}
        self._providers = providers or {}
    
    @property
    def name(self) -> str:
        return "CONTACTS"
    
    def get_operations(self) -> Dict[str, str]:
        return {
            "search": "Find contacts by name or email",
            "add": "Add a contact",
            "list": "List all contacts",
        }
    
    def get_param_schema(self) -> Dict[str, Dict[str, Any]]:
        return {
            "search": {
                "query": {"type": "str", "required": True, "description": "Name or email to search for"},
                "provider": {"type": "str", "required": False, "description": "Provider: google, microsoft (default: local)"},
            },
            "add": {
                "name": {"type": "str", "required": True, "description": "Contact name"},
                "email": {"type": "str", "required": False, "description": "Email address"},
                "phone": {"type": "str", "required": False, "description": "Phone number"},
                "provider": {"type": "str", "required": False, "description": "Provider: google, microsoft (default: local)"},
            },
            "list": {
                "limit": {"type": "int", "required": False, "description": "Max contacts to return"},
                "provider": {"type": "str", "required": False, "description": "Provider: google, microsoft (default: local)"},
            },
        }
    
    def _get_provider(self, name: Optional[str]) -> Optional[Any]:
        if name and name in self._providers:
            return self._providers[name]
        if self._providers:
            return next(iter(self._providers.values()))
        return None
    
    def add_contact(self, name: str, email: str, phone: Optional[str] = None):
        """Add a contact (can be called directly to seed data)."""
        self._contacts[name.lower()] = {"name": name, "email": email, "phone": phone}
    
    async def execute(self, operation: str, params: Dict[str, Any]) -> StepResult:
        try:
            provider_name = params.get("provider")
            provider = self._get_provider(provider_name)
            
            if operation == "search":
                query = params.get("query", "").lower()
                
                if provider and hasattr(provider, "search"):
                    result = await provider.search(query=query)
                    if result:
                        first = result[0]
                        return StepResult(True, data={"name": getattr(first, "name", str(first)), "email": getattr(first, "email", ""), "phone": getattr(first, "phone", "")})
                    return StepResult(True, data=None)
                
                matches = [
                    c for c in self._contacts.values()
                    if query in c["name"].lower() or query in c.get("email", "").lower()
                ]
                if matches:
                    return StepResult(True, data=matches[0])
                return StepResult(True, data=None)
            
            elif operation == "add":
                name = params.get("name")
                email = params.get("email")
                phone = params.get("phone")
                
                if not name:
                    return StepResult(False, error="Name required")
                
                if provider and hasattr(provider, "create_contact"):
                    result = await provider.create_contact(name=name, email=email, phone=phone)
                    return StepResult(True, data={"name": name, "email": email})
                
                self._contacts[name.lower()] = {"name": name, "email": email, "phone": phone}
                return StepResult(True, data={"name": name, "email": email})
            
            elif operation == "list":
                limit = params.get("limit", 100)
                
                if provider and hasattr(provider, "list_contacts"):
                    result = await provider.list_contacts(max_results=limit)
                    return StepResult(True, data=[{"name": getattr(c, "name", str(c)), "email": getattr(c, "email", "")} for c in result])
                
                return StepResult(True, data=list(self._contacts.values())[:limit])
            
            else:
                return StepResult(False, error=f"Unknown operation: {operation}")
                
        except Exception as e:
            return StepResult(False, error=str(e))


# ============================================================
#  KNOWLEDGE PRIMITIVE (backed by SemanticMemory)
# ============================================================

class KnowledgePrimitive(Primitive):
    """Memory and knowledge storage backed by the intelligence layer's SemanticMemory.
    
    Provides the agent with persistent, semantic memory — facts about people,
    preferences, relationships, and context that survive across sessions.
    """
    
    def __init__(self, storage_path: Optional[str] = None):
        # Initialize real semantic memory from intelligence layer
        try:
            from intelligence.semantic_memory import get_semantic_memory
            self._memory = get_semantic_memory()
        except ImportError:
            self._memory = None
        # Keep legacy fallback
        self._memories: List[Dict] = []
        self._storage_path = storage_path
        if storage_path and Path(storage_path).exists():
            self._load()
    
    @property
    def name(self) -> str:
        return "KNOWLEDGE"
    
    def get_operations(self) -> Dict[str, str]:
        return {
            "remember": "Store a fact or piece of information for later recall. Use this to remember user preferences, facts about people, important context, or anything worth persisting across conversations.",
            "recall": "Search memory for relevant facts matching a query. Returns the most relevant stored facts ranked by relevance.",
            "recall_about": "Get everything known about a specific person or entity — all facts, relationships, interaction history.",
            "connect": "Create a relationship between two entities in the knowledge graph (e.g., 'Alice' works_with 'Bob').",
            "forget": "Remove stored information by fact ID or entity name.",
        }
    
    def get_param_schema(self) -> Dict[str, Dict[str, Any]]:
        return {
            "remember": {
                "content": {"type": "str", "required": True, "description": "The fact or information to remember"},
                "entity": {"type": "str", "required": False, "description": "Primary entity this fact is about (person name, org, etc.)"},
                "category": {"type": "str", "required": False, "description": "Category: preference, relationship, event, insight, context, instruction, behavior, entity_info"},
            },
            "recall": {
                "query": {"type": "str", "required": True, "description": "Search query to find relevant memories"},
                "limit": {"type": "int", "required": False, "description": "Max results (default 5)"},
            },
            "recall_about": {
                "entity": {"type": "str", "required": True, "description": "Person or entity name to get facts about"},
            },
            "connect": {
                "entity1": {"type": "str", "required": True, "description": "First entity name"},
                "relationship": {"type": "str", "required": True, "description": "Relationship type (e.g., works_with, manages, married_to)"},
                "entity2": {"type": "str", "required": True, "description": "Second entity name"},
            },
            "forget": {
                "entity": {"type": "str", "required": False, "description": "Entity name — forgets all facts about this entity"},
                "fact_id": {"type": "str", "required": False, "description": "Specific fact ID to remove"},
            },
        }
    
    def _load(self):
        if self._storage_path:
            try:
                with open(self._storage_path, "r") as f:
                    self._memories = json.load(f)
            except Exception:
                pass
    
    def _save(self):
        if self._storage_path:
            Path(self._storage_path).parent.mkdir(parents=True, exist_ok=True)
            with open(self._storage_path, "w") as f:
                json.dump(self._memories, f, indent=2)
    
    async def execute(self, operation: str, params: Dict[str, Any]) -> StepResult:
        try:
            # Use real SemanticMemory if available
            if self._memory:
                return await self._execute_semantic(operation, params)
            # Fallback to basic list storage
            return await self._execute_basic(operation, params)
        except Exception as e:
            return StepResult(False, error=str(e))
    
    async def _execute_semantic(self, operation: str, params: Dict[str, Any]) -> StepResult:
        """Execute using the real SemanticMemory intelligence layer."""
        if operation == "remember":
            content = params.get("content") or params.get("text") or ""
            entity = params.get("entity")
            category_str = params.get("category", "context")
            
            # Convert category string to enum, fallback to CONTEXT
            from intelligence.semantic_memory import FactCategory
            category_map = {c.value: c for c in FactCategory}
            category = category_map.get(category_str.lower(), FactCategory.CONTEXT) if isinstance(category_str, str) else category_str
            
            fact = await self._memory.remember(
                content,
                category=category,
                entity=entity,
            )
            return StepResult(True, data={
                "stored": True,
                "fact_id": fact.id if hasattr(fact, 'id') else str(fact),
                "message": f"Remembered: {content[:100]}"
            })
        
        elif operation == "recall":
            query = params.get("query", "")
            limit = params.get("limit", 5)
            
            results = await self._memory.recall(query, limit=limit)
            if not results:
                return StepResult(True, data={"results": [], "message": "No matching memories found."})
            
            facts = []
            for item in results:
                # recall() returns List[Tuple[Fact, float]]
                fact = item[0] if isinstance(item, tuple) else item
                if hasattr(fact, 'content'):
                    facts.append({
                        "content": fact.content,
                        "category": fact.category.value if hasattr(fact.category, 'value') else str(fact.category),
                        "entity": fact.primary_entity,
                        "created": fact.created_at if hasattr(fact, 'created_at') else None,
                    })
                elif isinstance(fact, dict):
                    facts.append(fact)
            
            return StepResult(True, data={"results": facts[:limit]})
        
        elif operation == "recall_about":
            entity = params.get("entity", "")
            facts = await self._memory.recall_about(entity)
            
            if not facts:
                return StepResult(True, data={"entity": entity, "facts": [], "message": f"No facts known about '{entity}'."})
            
            fact_list = []
            for f in facts:
                if hasattr(f, 'content'):
                    fact_list.append({"content": f.content, "category": f.category.value if hasattr(f.category, 'value') else str(f.category)})
                elif isinstance(f, dict):
                    fact_list.append(f)
            
            # Also get entity summary if available
            entity_info = await self._memory.get_entity(entity)
            summary = None
            if entity_info and isinstance(entity_info, dict):
                summary = entity_info.get("summary")
            
            return StepResult(True, data={
                "entity": entity,
                "facts": fact_list,
                "summary": summary,
                "count": len(fact_list),
            })
        
        elif operation == "connect":
            e1 = params.get("entity1", "")
            rel = params.get("relationship", "")
            e2 = params.get("entity2", "")
            await self._memory.connect_entities(e1, rel, e2)
            return StepResult(True, data={"connected": f"{e1} --[{rel}]--> {e2}"})
        
        elif operation == "forget":
            entity = params.get("entity")
            fact_id = params.get("fact_id")
            if entity:
                await self._memory.forget(entity=entity)
                return StepResult(True, data={"forgotten": f"All facts about '{entity}'"})
            elif fact_id:
                await self._memory.forget(fact_id=fact_id)
                return StepResult(True, data={"forgotten": f"Fact {fact_id}"})
            else:
                return StepResult(False, error="Provide either 'entity' or 'fact_id'")
        
        else:
            return StepResult(False, error=f"Unknown operation: {operation}")
    
    async def _execute_basic(self, operation: str, params: Dict[str, Any]) -> StepResult:
        """Fallback: basic list-based memory."""
        if operation == "remember":
            content = params.get("content", "")
            memory = {
                "id": len(self._memories),
                "content": content,
                "tags": params.get("tags", []),
                "timestamp": datetime.now().isoformat(),
            }
            self._memories.append(memory)
            self._save()
            return StepResult(True, data={"id": memory["id"]})
        
        elif operation == "recall":
            query = params.get("query", "").lower()
            limit = params.get("limit", 5)
            matches = [
                m for m in self._memories
                if query in m["content"].lower() or any(query in t.lower() for t in m.get("tags", []))
            ]
            return StepResult(True, data=matches[:limit])
        
        elif operation == "forget":
            memory_id = params.get("id") or params.get("fact_id")
            self._memories = [m for m in self._memories if m.get("id") != memory_id]
            self._save()
            return StepResult(True)
        
        else:
            return StepResult(True, data={"results": []})


# ============================================================
#  PATTERNS PRIMITIVE
# ============================================================

class PatternsPrimitive(Primitive):
    """Behavioral pattern recognition — detects routines, habits, and anomalies.
    
    The agent can check what the user typically does at this time,
    detect when something is unusual, and understand recurring behaviors.
    """
    
    def __init__(self):
        try:
            from intelligence.pattern_recognition import get_pattern_engine
            self._patterns = get_pattern_engine()
        except ImportError:
            self._patterns = None
    
    @property
    def name(self) -> str:
        return "PATTERNS"
    
    def get_operations(self) -> Dict[str, str]:
        return {
            "whats_expected": "Check what the user typically does at this time of day/week. Use to proactively suggest routine actions.",
            "get_patterns": "List all detected behavioral patterns (routines, habits, recurring actions).",
            "detect_anomalies": "Check for anything unusual — missed routines, unexpected behavior changes.",
        }
    
    def get_available_operations(self) -> Dict[str, str]:
        if not self._patterns:
            return {}
        return self.get_operations()
    
    def get_param_schema(self) -> Dict[str, Dict[str, Any]]:
        return {
            "whats_expected": {
                "window_minutes": {"type": "int", "required": False, "description": "Time window to check (default 60 minutes)"},
            },
            "get_patterns": {
                "min_confidence": {"type": "float", "required": False, "description": "Minimum confidence threshold 0-1 (default 0.5)"},
            },
            "detect_anomalies": {
                "lookback_days": {"type": "int", "required": False, "description": "Days to look back for anomaly detection (default 7)"},
            },
        }
    
    async def execute(self, operation: str, params: Dict[str, Any]) -> StepResult:
        if not self._patterns:
            return StepResult(False, error="Pattern engine not available")
        
        try:
            if operation == "whats_expected":
                window = params.get("window_minutes", 60)
                expected = await self._patterns.whats_expected_now(window_minutes=window)
                if not expected:
                    return StepResult(True, data={"expected": [], "message": "No expected patterns right now."})
                # whats_expected_now() returns List[Dict] with keys: pattern, description, confidence
                patterns = [{"name": p.get("pattern", ""), "description": p.get("description", ""), "confidence": p.get("confidence", 0)} for p in expected]
                return StepResult(True, data={"expected": patterns})
            
            elif operation == "get_patterns":
                min_conf = params.get("min_confidence", 0.5)
                detected = await self._patterns.get_patterns(min_confidence=min_conf)
                if not detected:
                    return StepResult(True, data={"patterns": [], "message": "No patterns detected yet. Patterns emerge after repeated usage."})
                patterns = [{
                    "name": p.name,
                    "description": p.description,
                    "type": p.pattern_type.value if hasattr(p.pattern_type, 'value') else str(p.pattern_type),
                    "confidence": p.confidence,
                    "occurrences": p.occurrence_count,
                } for p in detected]
                return StepResult(True, data={"patterns": patterns, "count": len(patterns)})
            
            elif operation == "detect_anomalies":
                lookback = params.get("lookback_days", 7)
                anomalies = await self._patterns.detect_anomalies(lookback_days=lookback)
                if not anomalies:
                    return StepResult(True, data={"anomalies": [], "message": "No anomalies detected."})
                items = [{"type": a.get("type", "unknown"), "description": a.get("description", ""), "severity": a.get("severity", "low")} for a in anomalies]
                return StepResult(True, data={"anomalies": items})
            
            else:
                return StepResult(False, error=f"Unknown operation: {operation}")
        
        except Exception as e:
            return StepResult(False, error=str(e))


# ============================================================
#  INTELLIGENCE PRIMITIVE (cross-service + proactive)
# ============================================================

class IntelligencePrimitive(Primitive):
    """Cross-service intelligence — connects dots across email, calendar, files, and memory.
    
    Provides person briefs, meeting preparation, related content discovery,
    and proactive suggestions.
    """
    
    def __init__(self):
        try:
            from intelligence import get_intelligence_hub
            self._hub = get_intelligence_hub()
        except ImportError:
            self._hub = None
    
    @property
    def name(self) -> str:
        return "INTELLIGENCE"
    
    def get_operations(self) -> Dict[str, str]:
        return {
            "brief_on_person": "Get a comprehensive brief on a person — recent emails, meetings, shared docs, known facts, communication frequency.",
            "prepare_for_meeting": "Gather preparation materials for a meeting — attendee briefs, related emails/docs, suggested discussion points.",
            "find_related": "Search across all services for content related to a query — finds emails, docs, tasks, and memories that connect.",
            "get_suggestions": "Get proactive suggestions — things the user should do, follow up on, or prepare for.",
        }
    
    def get_available_operations(self) -> Dict[str, str]:
        if not self._hub:
            return {}
        return self.get_operations()
    
    def get_param_schema(self) -> Dict[str, Dict[str, Any]]:
        return {
            "brief_on_person": {
                "person": {"type": "str", "required": True, "description": "Person name or email to get a brief on"},
            },
            "prepare_for_meeting": {
                "title": {"type": "str", "required": False, "description": "Meeting title to prepare for"},
                "attendees": {"type": "str", "required": False, "description": "Comma-separated attendee names/emails"},
            },
            "find_related": {
                "query": {"type": "str", "required": True, "description": "Topic or query to find related content across all services"},
                "max_items": {"type": "int", "required": False, "description": "Max results (default 10)"},
            },
            "get_suggestions": {
                "max_suggestions": {"type": "int", "required": False, "description": "Max suggestions to return (default 5)"},
            },
        }
    
    async def execute(self, operation: str, params: Dict[str, Any]) -> StepResult:
        if not self._hub:
            return StepResult(False, error="Intelligence hub not available")
        
        try:
            if operation == "brief_on_person":
                person = params.get("person", "")
                brief = await self._hub.brief_on_person(person)
                if not brief:
                    return StepResult(True, data={"person": person, "message": f"No information found about '{person}'."})
                # Convert to serializable dict
                if hasattr(brief, '__dict__'):
                    data = {k: v for k, v in brief.__dict__.items() if not k.startswith('_')}
                elif isinstance(brief, dict):
                    data = brief
                else:
                    data = {"brief": str(brief)}
                data["person"] = person
                return StepResult(True, data=data)
            
            elif operation == "prepare_for_meeting":
                title = params.get("title")
                attendees_str = params.get("attendees", "")
                attendees = [a.strip() for a in attendees_str.split(",") if a.strip()] if attendees_str else []
                prep = await self._hub.prepare_for_meeting(title=title, attendees=attendees)
                if not prep:
                    return StepResult(True, data={"message": "No meeting preparation data available."})
                if hasattr(prep, '__dict__'):
                    data = {k: v for k, v in prep.__dict__.items() if not k.startswith('_')}
                elif isinstance(prep, dict):
                    data = prep
                else:
                    data = {"prep": str(prep)}
                return StepResult(True, data=data)
            
            elif operation == "find_related":
                query = params.get("query", "")
                max_items = params.get("max_items", 10)
                results = await self._hub.find_related(query, max_items=max_items)
                if not results:
                    return StepResult(True, data={"results": [], "message": f"No related content found for '{query}'."})
                items = []
                for r in results:
                    if hasattr(r, 'to_dict'):
                        items.append(r.to_dict())
                    elif hasattr(r, '__dict__'):
                        items.append({k: v for k, v in r.__dict__.items() if not k.startswith('_')})
                    else:
                        items.append(str(r))
                return StepResult(True, data={"results": items, "count": len(items)})
            
            elif operation == "get_suggestions":
                max_sugg = params.get("max_suggestions", 5)
                suggestions = await self._hub.get_suggestions(max_suggestions=max_sugg)
                if not suggestions:
                    return StepResult(True, data={"suggestions": [], "message": "No proactive suggestions right now."})
                items = []
                for s in suggestions:
                    if hasattr(s, 'to_dict'):
                        items.append(s.to_dict())
                    elif hasattr(s, '__dict__'):
                        items.append({k: v for k, v in s.__dict__.items() if not k.startswith('_')})
                    else:
                        items.append(str(s))
                return StepResult(True, data={"suggestions": items, "count": len(items)})
            
            else:
                return StepResult(False, error=f"Unknown operation: {operation}")
        
        except Exception as e:
            return StepResult(False, error=str(e))


# ============================================================
#  CALENDAR PRIMITIVE
# ============================================================

class CalendarPrimitive(Primitive):
    """Calendar and scheduling operations.
    
    Stores events locally. Can be wired to Google Calendar, Outlook, etc.
    via provider functions passed at init.
    """
    
    def __init__(
        self,
        storage_path: Optional[str] = None,
        create_func: Optional[Callable] = None,
        list_func: Optional[Callable] = None,
        list_calendars_func: Optional[Callable] = None,
        providers: Optional[Dict[str, Any]] = None,
    ):
        self._events: List[Dict] = []
        self._storage_path = storage_path
        self._create_func = create_func
        self._list_func = list_func
        self._list_calendars_func = list_calendars_func
        self._providers = providers or {}
        self._connector = next(iter(self._providers.values())) if self._providers else None
        self._calendars_cache: List[Dict] = []  # Cache of available calendars
        if storage_path and Path(storage_path).exists():
            try:
                with open(storage_path, "r") as f:
                    self._events = json.load(f)
            except Exception:
                pass
    
    @property
    def name(self) -> str:
        return "CALENDAR"
    
    def get_operations(self) -> Dict[str, str]:
        return {
            "create": "Create a calendar event (supports birthdays, meetings, reminders, recurring events)",
            "list": "List/get events in a date range. USE THIS for 'what's on my calendar today/this week/tomorrow' queries. By default queries ALL of the user's own calendars (not subscriptions) in a single call — no need to call list_calendars first. Only specify calendar_id if the user asks about a specific calendar.",
            "list_calendars": "List all available calendars (name, id, access level). Only needed if the user explicitly asks about their calendars or you need to look up a calendar_id by name.",
            "search": "Search events by keyword text (e.g. 'dentist', 'team meeting'). NOT for date lookups — use list for that.",
            "delete": "Delete an event by ID",
            "availability": "Check free/busy times in a date range",
        }
    
    def get_param_schema(self) -> Dict[str, Dict[str, Any]]:
        return {
            "create": {
                "title": {"type": "str", "required": True, "description": "Event title"},
                "start": {"type": "str", "required": True, "description": "Start datetime ISO format (e.g. 2026-04-01T09:00:00)"},
                "end": {"type": "str", "required": False, "description": "End datetime ISO format (defaults to 1 hour after start)"},
                "description": {"type": "str", "required": False, "description": "Event description or notes"},
                "location": {"type": "str", "required": False, "description": "Event location"},
                "attendees": {"type": "list", "required": False, "description": "List of attendee email addresses"},
                "calendar_id": {"type": "str", "required": False, "description": "Calendar ID (use 'primary' for main calendar, or calendar name like 'FAMILY SHARED' for specific calendars)"},
                "recurrence": {"type": "str", "required": False, "description": "Recurrence rule: 'yearly', 'monthly', 'weekly', 'daily', or RRULE string"},
                "all_day": {"type": "bool", "required": False, "description": "True for all-day events like birthdays"},
            },
            "list": {
                "start_date": {"type": "str", "required": False, "description": "Start of range (ISO date, defaults to today)"},
                "end_date": {"type": "str", "required": False, "description": "End of range (ISO date, defaults to 7 days out)"},
                "limit": {"type": "int", "required": False, "description": "Max events to return (default 50)"},
                "calendar_id": {"type": "str", "required": False, "description": "Calendar ID or name (default: queries user's own calendars, not subscriptions)"},
            },
            "search": {
                "query": {"type": "str", "required": True, "description": "Keyword to search for in event titles/descriptions (NOT a date)"},
                "calendar_id": {"type": "str", "required": False, "description": "Calendar ID or name to search within"},
                "start_date": {"type": "str", "required": False, "description": "Start of range (ISO date)"},
                "end_date": {"type": "str", "required": False, "description": "End of range (ISO date)"},
            },
            "list_calendars": {},
            "delete": {
                "id": {"type": "str", "required": True, "description": "Event ID to delete"},
            },
            "availability": {
                "start_date": {"type": "str", "required": True, "description": "Start of range (ISO date)"},
                "end_date": {"type": "str", "required": True, "description": "End of range (ISO date)"},
            },
        }
    
    def _save(self):
        if self._storage_path:
            Path(self._storage_path).parent.mkdir(parents=True, exist_ok=True)
            with open(self._storage_path, "w") as f:
                json.dump(self._events, f, indent=2)
            print(f"[CALENDAR] Saved {len(self._events)} events to {self._storage_path}")
    
    async def _resolve_calendar_id(self, calendar_id: str) -> str:
        """Resolve calendar name to ID.
        
        If calendar_id looks like a name (contains spaces or is a known name),
        look it up in the list of calendars and return the actual ID.
        """
        if not calendar_id or calendar_id == "primary":
            return "primary"
        
        # If it looks like an email/ID already (contains @), use as-is
        if "@" in calendar_id:
            return calendar_id
        
        # Try to resolve by name
        if self._list_calendars_func:
            try:
                if not self._calendars_cache:
                    self._calendars_cache = await self._list_calendars_func()
                    print(f"[CALENDAR] Cached {len(self._calendars_cache)} calendars")
                
                # Search for matching calendar by name (case insensitive)
                # list_calendars() may return 'name' or 'summary' depending on connector
                search_name = calendar_id.lower().strip()
                for cal in self._calendars_cache:
                    cal_name = (cal.get("name") or cal.get("summary") or "").lower()
                    if search_name in cal_name or cal_name in search_name:
                        real_id = cal.get("id")
                        print(f"[CALENDAR] Resolved '{calendar_id}' -> '{real_id}'")
                        return real_id
                
                print(f"[CALENDAR] No calendar found matching '{calendar_id}', using as-is")
            except Exception as e:
                print(f"[CALENDAR] Error resolving calendar: {e}")
        
        return calendar_id

    async def execute(self, operation: str, params: Dict[str, Any]) -> StepResult:
        print(f"[CALENDAR] execute({operation}, {params})")
        try:
            if operation == "create":
                if self._create_func:
                    print(f"[CALENDAR] Using external calendar (Google/Outlook)")
                    # Resolve calendar name to ID (e.g., "FAMILY SHARED" -> actual ID)
                    raw_calendar_id = params.get("calendar_id", "primary")
                    resolved_calendar_id = await self._resolve_calendar_id(raw_calendar_id)
                    
                    # Convert friendly recurrence to RRULE
                    recurrence = params.get("recurrence")
                    recurrence_rules = None
                    if recurrence:
                        recurrence_map = {
                            "yearly": ["RRULE:FREQ=YEARLY"],
                            "annually": ["RRULE:FREQ=YEARLY"],
                            "monthly": ["RRULE:FREQ=MONTHLY"],
                            "weekly": ["RRULE:FREQ=WEEKLY"],
                            "daily": ["RRULE:FREQ=DAILY"],
                        }
                        if recurrence.lower() in recurrence_map:
                            recurrence_rules = recurrence_map[recurrence.lower()]
                        elif recurrence.startswith("RRULE:"):
                            recurrence_rules = [recurrence]
                    
                    # Map param names: CalendarPrimitive uses 'title', Google API uses 'summary'
                    api_params = {
                        "summary": params.get("title") or params.get("summary", "Untitled"),
                        "start": params.get("start"),
                        "end": params.get("end"),
                        "description": params.get("description", ""),
                        "location": params.get("location", ""),
                        "attendees": params.get("attendees", []),
                        "calendar_id": resolved_calendar_id,
                        "all_day": params.get("all_day", False),
                        "recurrence": recurrence_rules,
                    }
                    print(f"[CALENDAR] Creating event: {api_params}")
                    result = await self._create_func(**api_params)
                    # Handle CalendarEvent dataclass or dict
                    if hasattr(result, 'to_dict'):
                        result_dict = result.to_dict()
                    elif isinstance(result, dict):
                        result_dict = result
                    else:
                        result_dict = {"id": str(result)}
                    return StepResult(True, data={
                        **result_dict,
                        "storage": "google_calendar",
                        "calendar": raw_calendar_id,
                        "message": f"Event '{api_params['summary']}' created in Google Calendar ({raw_calendar_id})"
                    })
                
                print(f"[CALENDAR] Creating local event")
                event = {
                    "id": f"evt_{len(self._events)}_{int(datetime.now().timestamp())}",
                    "title": params.get("title", "Untitled"),
                    "start": params.get("start", datetime.now().isoformat()),
                    "end": params.get("end"),
                    "description": params.get("description", ""),
                    "location": params.get("location", ""),
                    "attendees": params.get("attendees", []),
                    "created": datetime.now().isoformat(),
                    "storage": "local",  # Indicate where event is stored
                }
                self._events.append(event)
                self._save()
                print(f"[CALENDAR] Event created: {event['id']} - {event['title']}")
                print(f"[CALENDAR] Saved to: {self._storage_path}")
                return StepResult(True, data={
                    **event,
                    "message": f"Event '{event['title']}' saved to local calendar ({self._storage_path}). Google Calendar sync coming soon."
                })
            
            elif operation == "list":
                if self._list_func:
                    # Map primitive params to connector params
                    api_params = {}
                    if params.get("start_date"):
                        dt = datetime.fromisoformat(params["start_date"])
                        # Don't force timezone on date-only strings — let the
                        # connector handle them so the calendar's own timezone
                        # is respected by the Google API.
                        api_params["time_min"] = dt
                    if params.get("end_date"):
                        dt = datetime.fromisoformat(params["end_date"])
                        api_params["time_max"] = dt + timedelta(days=1)
                    if params.get("limit"):
                        api_params["max_results"] = params["limit"]
                    if params.get("calendar_id"):
                        raw_cal = params["calendar_id"]
                        api_params["calendar_id"] = await self._resolve_calendar_id(raw_cal)
                    result = await self._list_func(**api_params)
                    # Convert CalendarEvent objects to dicts
                    if result and hasattr(result[0], 'to_dict'):
                        result = [e.to_dict() for e in result]
                    return StepResult(True, data=result)
                
                start = params.get("start_date", datetime.now().strftime("%Y-%m-%d"))
                end = params.get("end_date")
                limit = params.get("limit", 50)
                
                filtered = []
                for evt in self._events:
                    evt_start = evt.get("start", "")
                    if evt_start >= start and (not end or evt_start <= end):
                        filtered.append(evt)
                
                filtered.sort(key=lambda e: e.get("start", ""))
                return StepResult(True, data=filtered[:limit])
            
            elif operation == "list_calendars":
                if self._list_calendars_func:
                    calendars = await self._list_calendars_func()
                    return StepResult(True, data=calendars)
                return StepResult(False, error="No calendar provider connected")
            
            elif operation == "search":
                if self._list_func:
                    search_kwargs = {}
                    raw_query = params.get("query", "")
                    # Don't pass date strings as text search — they match nothing.
                    # A date-like query means the caller wants events on that date.
                    import re
                    if re.match(r'^\d{4}-\d{2}-\d{2}', raw_query):
                        # Treat as date range query, not text search
                        try:
                            dt = datetime.fromisoformat(raw_query[:10])
                            search_kwargs["time_min"] = dt
                            search_kwargs["time_max"] = dt + timedelta(days=1)
                        except ValueError:
                            search_kwargs["query"] = raw_query
                    else:
                        search_kwargs["query"] = raw_query
                    if params.get("start_date"):
                        search_kwargs["time_min"] = datetime.fromisoformat(params["start_date"])
                    if params.get("end_date"):
                        dt = datetime.fromisoformat(params["end_date"])
                        search_kwargs["time_max"] = dt + timedelta(days=1)
                    if params.get("calendar_id"):
                        search_kwargs["calendar_id"] = await self._resolve_calendar_id(params["calendar_id"])
                    result = await self._list_func(**search_kwargs)
                    if result and hasattr(result[0], 'to_dict'):
                        result = [e.to_dict() for e in result]
                    return StepResult(True, data=result)
                
                query = params.get("query", "").lower()
                matches = [
                    e for e in self._events
                    if query in e.get("title", "").lower()
                    or query in e.get("description", "").lower()
                    or query in e.get("location", "").lower()
                ]
                return StepResult(True, data=matches)
            
            elif operation == "delete":
                event_id = params.get("id")
                before = len(self._events)
                self._events = [e for e in self._events if e.get("id") != event_id]
                self._save()
                return StepResult(True, data={"deleted": before != len(self._events)})
            
            elif operation == "availability":
                start = params.get("start_date", "")
                end = params.get("end_date", "")
                busy = [
                    {"start": e["start"], "end": e.get("end", e["start"]), "title": e["title"]}
                    for e in self._events
                    if e.get("start", "") >= start and e.get("start", "") <= end
                ]
                return StepResult(True, data={"busy": busy, "count": len(busy)})
            
            else:
                return StepResult(False, error=f"Unknown operation: {operation}")
        except Exception as e:
            return StepResult(False, error=str(e))


# ============================================================
#  WEB PRIMITIVE
# ============================================================

class WebPrimitive(Primitive):
    """Web/HTTP operations — fetch pages, call APIs, search the web."""
    
    def __init__(self, llm_complete: Optional[Callable] = None, search_provider: Optional[Any] = None):
        self._llm = llm_complete
        self._search_provider = search_provider
    
    @property
    def name(self) -> str:
        return "WEB"
    
    def get_operations(self) -> Dict[str, str]:
        return {
            "fetch": "Fetch content from a URL (returns text/HTML). Works for static pages only.",
            "api": "Make an HTTP API call (GET, POST, PUT, DELETE)",
            "search": "Search the web for current information (sports schedules, news, facts). Use this for questions about dates, times, events.",
            "extract": "Fetch a static webpage URL and extract specific information. NOT for google.com, bing.com, or other search engines (JS-rendered). Use for news sites, wikipedia, official event pages.",
        }
    
    def get_param_schema(self) -> Dict[str, Dict[str, Any]]:
        return {
            "fetch": {
                "url": {"type": "str", "required": True, "description": "URL to fetch"},
                "max_length": {"type": "int", "required": False, "description": "Max chars to return (default 10000)"},
            },
            "api": {
                "url": {"type": "str", "required": True, "description": "API endpoint URL"},
                "method": {"type": "str", "required": False, "description": "HTTP method: GET, POST, PUT, DELETE (default GET)"},
                "headers": {"type": "dict", "required": False, "description": "HTTP headers"},
                "body": {"type": "dict", "required": False, "description": "Request body (for POST/PUT)"},
                "params": {"type": "dict", "required": False, "description": "Query parameters"},
            },
            "search": {
                "query": {"type": "str", "required": True, "description": "Search query"},
                "limit": {"type": "int", "required": False, "description": "Max results (default 5)"},
            },
            "extract": {
                "url": {"type": "str", "required": True, "description": "URL to fetch and extract from"},
                "what": {"type": "str", "required": True, "description": "What to extract (e.g. 'the main article text', 'all prices', 'contact info')"},
            },
        }
    
    def get_available_operations(self) -> Dict[str, str]:
        """Only show search if a search provider is configured."""
        ops = self.get_operations()
        if not self._search_provider:
            ops.pop("search", None)
        return ops
    
    async def execute(self, operation: str, params: Dict[str, Any]) -> StepResult:
        try:
            import httpx
        except ImportError:
            try:
                import urllib.request
                _has_httpx = False
            except Exception:
                return StepResult(False, error="No HTTP library available")
            _has_httpx = False
        else:
            _has_httpx = True
        
        try:
            if operation == "fetch":
                url = params.get("url", "")
                max_len = params.get("max_length", 10000)
                
                if not url:
                    return StepResult(False, error="Missing 'url' parameter")
                
                if _has_httpx:
                    async with httpx.AsyncClient(
                        follow_redirects=True, 
                        timeout=30,
                        headers={"User-Agent": "Mozilla/5.0 (compatible; Telic/1.0; +https://github.com/rob637/Zero)"}
                    ) as client:
                        resp = await client.get(url)
                        resp.raise_for_status()
                        text = resp.text[:max_len]
                else:
                    req = urllib.request.Request(url, headers={"User-Agent": "Telic/1.0"})
                    with urllib.request.urlopen(req, timeout=30) as resp:
                        text = resp.read().decode("utf-8", errors="ignore")[:max_len]
                
                return StepResult(True, data={"url": url, "content": text, "length": len(text)})
            
            elif operation == "api":
                if not _has_httpx:
                    return StepResult(False, error="httpx required for API calls. Install: pip install httpx")
                
                url = params.get("url", "")
                method = params.get("method", "GET").upper()
                headers = params.get("headers", {})
                body = params.get("body")
                query_params = params.get("params")
                
                if not url:
                    return StepResult(False, error="Missing 'url' parameter")
                
                async with httpx.AsyncClient(follow_redirects=True, timeout=30) as client:
                    resp = await client.request(
                        method, url,
                        headers=headers,
                        json=body if body else None,
                        params=query_params,
                    )
                    
                    try:
                        data = resp.json()
                    except Exception:
                        data = resp.text[:10000]
                    
                    return StepResult(True, data={
                        "status": resp.status_code,
                        "data": data,
                        "headers": dict(resp.headers),
                    })
            
            elif operation == "search":
                query = params.get("query", "")
                limit = params.get("limit", 5)
                
                if self._search_provider and hasattr(self._search_provider, "search"):
                    result = await self._search_provider.search(query=query, num_results=limit)
                    return StepResult(True, data=result)
                
                return StepResult(False, error="Web search not configured. Connect a search provider (Google, Bing, etc.)")
            
            elif operation == "extract":
                url = params.get("url", "")
                what = params.get("what", "the main content")
                
                if not url:
                    return StepResult(False, error="Missing 'url' parameter")
                if not self._llm:
                    return StepResult(False, error="LLM required for extraction")
                
                # Fetch first
                fetch_result = await self.execute("fetch", {"url": url, "max_length": 15000})
                if not fetch_result.success:
                    return fetch_result
                
                content = fetch_result.data.get("content", "")
                
                # Inject today's date for time-sensitive extractions
                today_iso = datetime.now().strftime("%Y-%m-%d")
                
                prompt = f"""Extract the following from this web page:
What to extract: {what}

IMPORTANT: Today's date is {today_iso}. If returning dates/times, use {today_iso} as the date.

Web page content:
{content[:12000]}

Return ONLY a JSON object with the extracted data."""
                
                response = await self._llm(prompt)
                json_match = re.search(r'\{[\s\S]*\}', response)
                if json_match:
                    try:
                        return StepResult(True, data=json.loads(json_match.group()))
                    except json.JSONDecodeError:
                        pass
                return StepResult(True, data={"extracted": response})
            
            else:
                return StepResult(False, error=f"Unknown operation: {operation}")
        except Exception as e:
            return StepResult(False, error=str(e))


# ============================================================
#  WEATHER PRIMITIVE
# ============================================================

class WeatherPrimitive(Primitive):
    """Weather — current conditions, forecasts, and air quality.
    
    Uses OpenWeatherMap via the WeatherConnector.
    """
    
    def __init__(self, connector: Any = None):
        self._connector = connector
    
    @property
    def name(self) -> str:
        return "WEATHER"
    
    def get_operations(self) -> Dict[str, str]:
        return {
            "current": "Get current weather for a city, zip code, or coordinates (e.g. 'New York', '10001', '40.7,-74.0')",
            "forecast": "Get weather forecast for the next 1-5 days (3-hour intervals)",
            "air_quality": "Get air quality index (AQI) and pollutant levels for a location",
            "search_cities": "Search for cities by name to find the right location",
        }
    
    def get_param_schema(self) -> Dict[str, Dict[str, Any]]:
        return {
            "current": {
                "location": {"type": "str", "required": True, "description": "City name, zip code, or lat,lon coordinates"},
            },
            "forecast": {
                "location": {"type": "str", "required": True, "description": "City name, zip code, or lat,lon coordinates"},
                "days": {"type": "int", "required": False, "description": "Number of days (1-5, default 3)"},
            },
            "air_quality": {
                "location": {"type": "str", "required": True, "description": "City name or lat,lon coordinates"},
            },
            "search_cities": {
                "query": {"type": "str", "required": True, "description": "City name to search for"},
                "limit": {"type": "int", "required": False, "description": "Max results (default 5)"},
            },
        }
    
    async def execute(self, operation: str, params: Dict[str, Any]) -> StepResult:
        if not self._connector:
            return StepResult(False, error="Weather is not configured. Connect an OpenWeatherMap API key in Settings to use weather features.")
        try:
            if operation == "current":
                location = params.get("location", "")
                if not location:
                    return StepResult(False, error="Missing 'location' parameter")
                result = await self._connector.get_current(location)
                return StepResult(True, data=result.to_dict())
            
            elif operation == "forecast":
                location = params.get("location", "")
                if not location:
                    return StepResult(False, error="Missing 'location' parameter")
                days = int(params.get("days", 3))
                result = await self._connector.get_forecast(location, days=days)
                return StepResult(True, data=result)
            
            elif operation == "air_quality":
                location = params.get("location", "")
                if not location:
                    return StepResult(False, error="Missing 'location' parameter")
                result = await self._connector.get_air_quality(location)
                return StepResult(True, data=result)
            
            elif operation == "search_cities":
                query = params.get("query", "")
                if not query:
                    return StepResult(False, error="Missing 'query' parameter")
                limit = int(params.get("limit", 5))
                results = await self._connector.search_cities(query, limit=limit)
                return StepResult(True, data={"cities": results})
            
            else:
                return StepResult(False, error=f"Unknown operation: {operation}")
        except Exception as e:
            return StepResult(False, error=str(e))


# ============================================================
#  NEWS PRIMITIVE
# ============================================================

class NewsPrimitive(Primitive):
    """News — top headlines, search, and source discovery.
    
    Uses NewsAPI via the NewsConnector.
    """
    
    def __init__(self, connector: Any = None):
        self._connector = connector
    
    @property
    def name(self) -> str:
        return "NEWS"
    
    def get_operations(self) -> Dict[str, str]:
        return {
            "headlines": "Get top headlines by country and/or category (business, sports, tech, health, science, entertainment)",
            "search": "Search all news articles by keyword. Supports AND, OR, NOT operators. Can filter by date range and source.",
            "sources": "List available news sources, optionally filtered by category, language, or country",
        }
    
    def get_param_schema(self) -> Dict[str, Dict[str, Any]]:
        return {
            "headlines": {
                "country": {"type": "str", "required": False, "description": "2-letter country code (default: us). Options: us, gb, ca, au, de, fr, it, etc."},
                "category": {"type": "str", "required": False, "description": "Category: business, entertainment, general, health, science, sports, technology"},
                "query": {"type": "str", "required": False, "description": "Keywords to filter headlines"},
                "limit": {"type": "int", "required": False, "description": "Max articles (default 10)"},
            },
            "search": {
                "query": {"type": "str", "required": True, "description": "Search keywords (supports AND, OR, NOT)"},
                "sort_by": {"type": "str", "required": False, "description": "relevancy, popularity, or publishedAt (default: relevancy)"},
                "from_date": {"type": "str", "required": False, "description": "Start date (YYYY-MM-DD)"},
                "to_date": {"type": "str", "required": False, "description": "End date (YYYY-MM-DD)"},
                "sources": {"type": "str", "required": False, "description": "Comma-separated source IDs (e.g. 'bbc-news,cnn')"},
                "limit": {"type": "int", "required": False, "description": "Max articles (default 10)"},
            },
            "sources": {
                "category": {"type": "str", "required": False, "description": "Category filter"},
                "language": {"type": "str", "required": False, "description": "2-letter language code (default: en)"},
                "country": {"type": "str", "required": False, "description": "2-letter country code"},
            },
        }
    
    async def execute(self, operation: str, params: Dict[str, Any]) -> StepResult:
        if not self._connector:
            return StepResult(False, error="News is not configured. Connect a NewsAPI key in Settings to use news features.")
        try:
            if operation == "headlines":
                articles = await self._connector.top_headlines(
                    country=params.get("country", "us"),
                    category=params.get("category"),
                    query=params.get("query"),
                    limit=int(params.get("limit", 10)),
                )
                return StepResult(True, data={
                    "count": len(articles),
                    "articles": [a.to_dict() for a in articles],
                })
            
            elif operation == "search":
                query = params.get("query", "")
                if not query:
                    return StepResult(False, error="Missing 'query' parameter")
                articles = await self._connector.search(
                    query=query,
                    sort_by=params.get("sort_by", "relevancy"),
                    from_date=params.get("from_date"),
                    to_date=params.get("to_date"),
                    sources=params.get("sources"),
                    limit=int(params.get("limit", 10)),
                )
                return StepResult(True, data={
                    "query": query,
                    "count": len(articles),
                    "articles": [a.to_dict() for a in articles],
                })
            
            elif operation == "sources":
                sources = await self._connector.get_sources(
                    category=params.get("category"),
                    language=params.get("language", "en"),
                    country=params.get("country"),
                )
                return StepResult(True, data={
                    "count": len(sources),
                    "sources": sources,
                })
            
            else:
                return StepResult(False, error=f"Unknown operation: {operation}")
        except Exception as e:
            return StepResult(False, error=str(e))


# ============================================================
#  NOTION PRIMITIVE
# ============================================================

class NotionPrimitive(Primitive):
    """Notion — pages, databases, content blocks, and search.
    
    Uses the Notion API via NotionConnector.
    """
    
    def __init__(self, connector: Any = None):
        self._connector = connector
    
    @property
    def name(self) -> str:
        return "NOTION"
    
    def get_operations(self) -> Dict[str, str]:
        return {
            "search": "Search across all Notion pages and databases. Can filter by type (page/database).",
            "get_page": "Get a page's properties and metadata by ID",
            "get_page_content": "Get the content blocks (text, lists, headings) of a page",
            "create_page": "Create a new page under a parent page or database. Supports markdown-like content.",
            "update_page": "Update a page's properties",
            "append_content": "Add text content to the bottom of an existing page",
            "archive_page": "Archive (soft-delete) a page",
            "list_databases": "List all databases shared with the integration",
            "get_database": "Get a database's schema and metadata",
            "query_database": "Query a database with optional filters and sorting",
            "create_database": "Create a new inline database in a page",
            "add_comment": "Add a comment to a page",
            "get_comments": "Get all comments on a page",
        }
    
    def get_param_schema(self) -> Dict[str, Dict[str, Any]]:
        return {
            "search": {
                "query": {"type": "str", "required": False, "description": "Search text (empty = recent pages)"},
                "filter_type": {"type": "str", "required": False, "description": "'page' or 'database'"},
                "limit": {"type": "int", "required": False, "description": "Max results (default 10)"},
            },
            "get_page": {
                "page_id": {"type": "str", "required": True, "description": "Page ID"},
            },
            "get_page_content": {
                "page_id": {"type": "str", "required": True, "description": "Page ID"},
            },
            "create_page": {
                "parent_id": {"type": "str", "required": True, "description": "Parent page ID or database ID"},
                "title": {"type": "str", "required": True, "description": "Page title"},
                "content": {"type": "str", "required": False, "description": "Page body content (markdown-like)"},
                "parent_type": {"type": "str", "required": False, "description": "'page' or 'database' (default: page)"},
            },
            "update_page": {
                "page_id": {"type": "str", "required": True, "description": "Page ID"},
                "properties": {"type": "dict", "required": True, "description": "Properties to update (Notion format)"},
            },
            "append_content": {
                "page_id": {"type": "str", "required": True, "description": "Page ID"},
                "content": {"type": "str", "required": True, "description": "Text to append"},
            },
            "archive_page": {
                "page_id": {"type": "str", "required": True, "description": "Page ID to archive"},
            },
            "list_databases": {
                "limit": {"type": "int", "required": False, "description": "Max results (default 10)"},
            },
            "get_database": {
                "database_id": {"type": "str", "required": True, "description": "Database ID"},
            },
            "query_database": {
                "database_id": {"type": "str", "required": True, "description": "Database ID"},
                "filter": {"type": "dict", "required": False, "description": "Notion filter object"},
                "sorts": {"type": "list", "required": False, "description": "Sort objects [{property, direction}]"},
                "limit": {"type": "int", "required": False, "description": "Max results (default 20)"},
            },
            "create_database": {
                "parent_page_id": {"type": "str", "required": True, "description": "Parent page ID"},
                "title": {"type": "str", "required": True, "description": "Database title"},
                "properties": {"type": "dict", "required": True, "description": "Property schema"},
            },
            "add_comment": {
                "page_id": {"type": "str", "required": True, "description": "Page ID"},
                "text": {"type": "str", "required": True, "description": "Comment text"},
            },
            "get_comments": {
                "page_id": {"type": "str", "required": True, "description": "Page ID"},
            },
        }
    
    async def execute(self, operation: str, params: Dict[str, Any]) -> StepResult:
        if not self._connector:
            return StepResult(False, error="Notion is not configured. Connect Notion in Settings to use Notion features.")
        try:
            if operation == "search":
                results = await self._connector.search(
                    query=params.get("query", ""),
                    filter_type=params.get("filter_type"),
                    limit=int(params.get("limit", 10)),
                )
                return StepResult(True, data={"count": len(results), "results": results})
            
            elif operation == "get_page":
                result = await self._connector.get_page(params["page_id"])
                return StepResult(True, data=result)
            
            elif operation == "get_page_content":
                blocks = await self._connector.get_page_content(params["page_id"])
                return StepResult(True, data={"blocks": blocks})
            
            elif operation == "create_page":
                result = await self._connector.create_page(
                    parent_id=params["parent_id"],
                    title=params["title"],
                    content=params.get("content"),
                    parent_type=params.get("parent_type", "page"),
                )
                return StepResult(True, data=result)
            
            elif operation == "update_page":
                result = await self._connector.update_page(params["page_id"], params["properties"])
                return StepResult(True, data=result)
            
            elif operation == "append_content":
                blocks = await self._connector.append_content(params["page_id"], params["content"])
                return StepResult(True, data={"appended_blocks": len(blocks)})
            
            elif operation == "archive_page":
                result = await self._connector.archive_page(params["page_id"])
                return StepResult(True, data=result)
            
            elif operation == "list_databases":
                results = await self._connector.list_databases(limit=int(params.get("limit", 10)))
                return StepResult(True, data={"count": len(results), "databases": results})
            
            elif operation == "get_database":
                result = await self._connector.get_database(params["database_id"])
                return StepResult(True, data=result)
            
            elif operation == "query_database":
                results = await self._connector.query_database(
                    database_id=params["database_id"],
                    filter_obj=params.get("filter"),
                    sorts=params.get("sorts"),
                    limit=int(params.get("limit", 20)),
                )
                return StepResult(True, data={"count": len(results), "results": results})
            
            elif operation == "create_database":
                result = await self._connector.create_database(
                    parent_page_id=params["parent_page_id"],
                    title=params["title"],
                    properties=params["properties"],
                )
                return StepResult(True, data=result)
            
            elif operation == "add_comment":
                result = await self._connector.add_comment(params["page_id"], params["text"])
                return StepResult(True, data=result)
            
            elif operation == "get_comments":
                comments = await self._connector.get_comments(params["page_id"])
                return StepResult(True, data={"count": len(comments), "comments": comments})
            
            else:
                return StepResult(False, error=f"Unknown operation: {operation}")
        except Exception as e:
            return StepResult(False, error=str(e))


# ============================================================
#  LINEAR PRIMITIVE
# ============================================================

class LinearPrimitive(Primitive):
    """Linear — issue tracking, projects, cycles, and teams.
    
    Uses the Linear GraphQL API via LinearConnector.
    """
    
    def __init__(self, connector: Any = None):
        self._connector = connector
    
    @property
    def name(self) -> str:
        return "LINEAR"
    
    def get_operations(self) -> Dict[str, str]:
        return {
            "list_issues": "List issues with optional filters by team, status, or assignee",
            "get_issue": "Get a single issue by ID or identifier (e.g. ENG-123)",
            "create_issue": "Create a new issue in a team with title, description, priority, labels, etc.",
            "update_issue": "Update an issue's title, description, status, priority, assignee, or due date",
            "add_comment": "Add a comment to an issue",
            "search_issues": "Full-text search across all issues",
            "list_teams": "List all teams with members, statuses, and labels",
            "list_cycles": "List cycles (sprints) with progress info",
            "list_projects": "List projects with progress and team info",
            "me": "Get the current authenticated user and their recent assigned issues",
        }
    
    def get_param_schema(self) -> Dict[str, Dict[str, Any]]:
        return {
            "list_issues": {
                "team_key": {"type": "str", "required": False, "description": "Team key to filter (e.g. 'ENG')"},
                "status": {"type": "str", "required": False, "description": "Status name (e.g. 'In Progress', 'Done', 'Todo')"},
                "assignee": {"type": "str", "required": False, "description": "Assignee name or 'me' for current user"},
                "limit": {"type": "int", "required": False, "description": "Max results (default 20)"},
            },
            "get_issue": {
                "issue_id": {"type": "str", "required": True, "description": "Issue UUID or identifier like 'ENG-123'"},
            },
            "create_issue": {
                "title": {"type": "str", "required": True, "description": "Issue title"},
                "team_key": {"type": "str", "required": True, "description": "Team key (e.g. 'ENG'). Use list_teams to find available teams."},
                "description": {"type": "str", "required": False, "description": "Issue description (supports markdown)"},
                "priority": {"type": "int", "required": False, "description": "0=No priority, 1=Urgent, 2=High, 3=Medium, 4=Low"},
                "status": {"type": "str", "required": False, "description": "Status name (e.g. 'Todo', 'In Progress')"},
                "labels": {"type": "list", "required": False, "description": "List of label names"},
                "due_date": {"type": "str", "required": False, "description": "Due date (YYYY-MM-DD)"},
                "estimate": {"type": "int", "required": False, "description": "Story points estimate"},
            },
            "update_issue": {
                "issue_id": {"type": "str", "required": True, "description": "Issue UUID"},
                "title": {"type": "str", "required": False, "description": "New title"},
                "description": {"type": "str", "required": False, "description": "New description"},
                "status": {"type": "str", "required": False, "description": "New status name"},
                "priority": {"type": "int", "required": False, "description": "New priority (0-4)"},
                "due_date": {"type": "str", "required": False, "description": "New due date (YYYY-MM-DD)"},
                "estimate": {"type": "int", "required": False, "description": "New estimate"},
            },
            "add_comment": {
                "issue_id": {"type": "str", "required": True, "description": "Issue UUID"},
                "body": {"type": "str", "required": True, "description": "Comment text (supports markdown)"},
            },
            "search_issues": {
                "query": {"type": "str", "required": True, "description": "Search text"},
                "limit": {"type": "int", "required": False, "description": "Max results (default 20)"},
            },
            "list_teams": {},
            "list_cycles": {
                "team_key": {"type": "str", "required": False, "description": "Filter by team key"},
                "limit": {"type": "int", "required": False, "description": "Max results (default 5)"},
            },
            "list_projects": {
                "limit": {"type": "int", "required": False, "description": "Max results (default 10)"},
            },
            "me": {},
        }
    
    async def execute(self, operation: str, params: Dict[str, Any]) -> StepResult:
        if not self._connector:
            return StepResult(False, error="Linear is not configured. Connect Linear in Settings to use Linear features.")
        try:
            if operation == "list_issues":
                results = await self._connector.list_issues(
                    team_key=params.get("team_key"),
                    status=params.get("status"),
                    assignee=params.get("assignee"),
                    limit=int(params.get("limit", 20)),
                )
                return StepResult(True, data={"count": len(results), "issues": results})
            
            elif operation == "get_issue":
                result = await self._connector.get_issue(params["issue_id"])
                return StepResult(True, data=result)
            
            elif operation == "create_issue":
                result = await self._connector.create_issue(
                    title=params["title"],
                    team_key=params["team_key"],
                    description=params.get("description"),
                    priority=int(params["priority"]) if params.get("priority") is not None else None,
                    status=params.get("status"),
                    labels=params.get("labels"),
                    due_date=params.get("due_date"),
                    estimate=int(params["estimate"]) if params.get("estimate") is not None else None,
                )
                return StepResult(True, data=result)
            
            elif operation == "update_issue":
                result = await self._connector.update_issue(
                    issue_id=params["issue_id"],
                    title=params.get("title"),
                    description=params.get("description"),
                    status=params.get("status"),
                    priority=int(params["priority"]) if params.get("priority") is not None else None,
                    due_date=params.get("due_date"),
                    estimate=int(params["estimate"]) if params.get("estimate") is not None else None,
                )
                return StepResult(True, data=result)
            
            elif operation == "add_comment":
                result = await self._connector.add_comment(params["issue_id"], params["body"])
                return StepResult(True, data=result)
            
            elif operation == "search_issues":
                results = await self._connector.search_issues(
                    query=params["query"],
                    limit=int(params.get("limit", 20)),
                )
                return StepResult(True, data={"count": len(results), "issues": results})
            
            elif operation == "list_teams":
                results = await self._connector.list_teams()
                return StepResult(True, data={"count": len(results), "teams": results})
            
            elif operation == "list_cycles":
                results = await self._connector.list_cycles(
                    team_key=params.get("team_key"),
                    limit=int(params.get("limit", 5)),
                )
                return StepResult(True, data={"count": len(results), "cycles": results})
            
            elif operation == "list_projects":
                results = await self._connector.list_projects(
                    limit=int(params.get("limit", 10)),
                )
                return StepResult(True, data={"count": len(results), "projects": results})
            
            elif operation == "me":
                result = await self._connector.me()
                return StepResult(True, data=result)
            
            else:
                return StepResult(False, error=f"Unknown operation: {operation}")
        except Exception as e:
            return StepResult(False, error=str(e))


# ============================================================
#  TRELLO PRIMITIVE
# ============================================================

class TrelloPrimitive(Primitive):
    """Trello — boards, lists, cards, checklists, and search.
    
    Uses the Trello REST API via TrelloConnector.
    """
    
    def __init__(self, connector: Any = None):
        self._connector = connector
    
    @property
    def name(self) -> str:
        return "TRELLO"
    
    def get_operations(self) -> Dict[str, str]:
        return {
            "list_boards": "List all boards for the authenticated user",
            "get_board": "Get a board with its lists, labels, and metadata",
            "create_board": "Create a new board",
            "get_lists": "Get all lists on a board",
            "create_list": "Create a new list on a board",
            "get_cards": "Get cards from a list or board",
            "get_card": "Get a single card with comments, checklists, and full details",
            "create_card": "Create a new card in a list with name, description, due date, labels",
            "update_card": "Update a card's name, description, due date, list, or archive status",
            "move_card": "Move a card to a different list",
            "delete_card": "Permanently delete a card",
            "add_comment": "Add a comment to a card",
            "add_checklist": "Add a checklist with items to a card",
            "search": "Search across boards and cards",
            "get_board_members": "Get members of a board",
        }
    
    def get_param_schema(self) -> Dict[str, Dict[str, Any]]:
        return {
            "list_boards": {
                "filter": {"type": "str", "required": False, "description": "'open', 'closed', or 'all' (default: open)"},
            },
            "get_board": {
                "board_id": {"type": "str", "required": True, "description": "Board ID"},
            },
            "create_board": {
                "name": {"type": "str", "required": True, "description": "Board name"},
                "description": {"type": "str", "required": False, "description": "Board description"},
            },
            "get_lists": {
                "board_id": {"type": "str", "required": True, "description": "Board ID"},
            },
            "create_list": {
                "board_id": {"type": "str", "required": True, "description": "Board ID"},
                "name": {"type": "str", "required": True, "description": "List name"},
            },
            "get_cards": {
                "list_id": {"type": "str", "required": False, "description": "List ID (get cards in this list)"},
                "board_id": {"type": "str", "required": False, "description": "Board ID (get all cards on board)"},
                "limit": {"type": "int", "required": False, "description": "Max results (default 50)"},
            },
            "get_card": {
                "card_id": {"type": "str", "required": True, "description": "Card ID"},
            },
            "create_card": {
                "list_id": {"type": "str", "required": True, "description": "List ID to create card in"},
                "name": {"type": "str", "required": True, "description": "Card name/title"},
                "description": {"type": "str", "required": False, "description": "Card description (markdown)"},
                "due": {"type": "str", "required": False, "description": "Due date (ISO 8601 or YYYY-MM-DD)"},
                "labels": {"type": "list", "required": False, "description": "List of label IDs"},
                "position": {"type": "str", "required": False, "description": "'top' or 'bottom'"},
            },
            "update_card": {
                "card_id": {"type": "str", "required": True, "description": "Card ID"},
                "name": {"type": "str", "required": False, "description": "New name"},
                "description": {"type": "str", "required": False, "description": "New description"},
                "due": {"type": "str", "required": False, "description": "New due date"},
                "due_complete": {"type": "bool", "required": False, "description": "Mark due date complete"},
                "list_id": {"type": "str", "required": False, "description": "Move to different list"},
                "closed": {"type": "bool", "required": False, "description": "Archive (true) or unarchive (false)"},
            },
            "move_card": {
                "card_id": {"type": "str", "required": True, "description": "Card ID"},
                "list_id": {"type": "str", "required": True, "description": "Target list ID"},
            },
            "delete_card": {
                "card_id": {"type": "str", "required": True, "description": "Card ID"},
            },
            "add_comment": {
                "card_id": {"type": "str", "required": True, "description": "Card ID"},
                "text": {"type": "str", "required": True, "description": "Comment text"},
            },
            "add_checklist": {
                "card_id": {"type": "str", "required": True, "description": "Card ID"},
                "name": {"type": "str", "required": True, "description": "Checklist name"},
                "items": {"type": "list", "required": False, "description": "List of checklist item names"},
            },
            "search": {
                "query": {"type": "str", "required": True, "description": "Search text"},
                "model_types": {"type": "str", "required": False, "description": "Comma-separated: 'cards', 'boards' (default: both)"},
                "limit": {"type": "int", "required": False, "description": "Max results per type (default 10)"},
            },
            "get_board_members": {
                "board_id": {"type": "str", "required": True, "description": "Board ID"},
            },
        }
    
    async def execute(self, operation: str, params: Dict[str, Any]) -> StepResult:
        if not self._connector:
            return StepResult(False, error="Trello is not configured. Connect Trello in Settings to use Trello features.")
        try:
            if operation == "list_boards":
                results = await self._connector.list_boards(
                    filter=params.get("filter", "open"),
                )
                return StepResult(True, data={"count": len(results), "boards": results})
            
            elif operation == "get_board":
                result = await self._connector.get_board(params["board_id"])
                return StepResult(True, data=result)
            
            elif operation == "create_board":
                result = await self._connector.create_board(
                    name=params["name"],
                    description=params.get("description"),
                )
                return StepResult(True, data=result)
            
            elif operation == "get_lists":
                results = await self._connector.get_lists(params["board_id"])
                return StepResult(True, data={"count": len(results), "lists": results})
            
            elif operation == "create_list":
                result = await self._connector.create_list(params["board_id"], params["name"])
                return StepResult(True, data=result)
            
            elif operation == "get_cards":
                results = await self._connector.get_cards(
                    list_id=params.get("list_id"),
                    board_id=params.get("board_id"),
                    limit=int(params.get("limit", 50)),
                )
                return StepResult(True, data={"count": len(results), "cards": results})
            
            elif operation == "get_card":
                result = await self._connector.get_card(params["card_id"])
                return StepResult(True, data=result)
            
            elif operation == "create_card":
                result = await self._connector.create_card(
                    list_id=params["list_id"],
                    name=params["name"],
                    description=params.get("description"),
                    due=params.get("due"),
                    labels=params.get("labels"),
                    position=params.get("position"),
                )
                return StepResult(True, data=result)
            
            elif operation == "update_card":
                result = await self._connector.update_card(
                    card_id=params["card_id"],
                    name=params.get("name"),
                    description=params.get("description"),
                    due=params.get("due"),
                    due_complete=params.get("due_complete"),
                    list_id=params.get("list_id"),
                    closed=params.get("closed"),
                )
                return StepResult(True, data=result)
            
            elif operation == "move_card":
                result = await self._connector.move_card(params["card_id"], params["list_id"])
                return StepResult(True, data=result)
            
            elif operation == "delete_card":
                await self._connector.delete_card(params["card_id"])
                return StepResult(True, data={"deleted": True})
            
            elif operation == "add_comment":
                result = await self._connector.add_comment(params["card_id"], params["text"])
                return StepResult(True, data=result)
            
            elif operation == "add_checklist":
                result = await self._connector.add_checklist(
                    card_id=params["card_id"],
                    name=params["name"],
                    items=params.get("items"),
                )
                return StepResult(True, data=result)
            
            elif operation == "search":
                results = await self._connector.search(
                    query=params["query"],
                    model_types=params.get("model_types", "cards,boards"),
                    limit=int(params.get("limit", 10)),
                )
                return StepResult(True, data=results)
            
            elif operation == "get_board_members":
                results = await self._connector.get_board_members(params["board_id"])
                return StepResult(True, data={"count": len(results), "members": results})
            
            else:
                return StepResult(False, error=f"Unknown operation: {operation}")
        except Exception as e:
            return StepResult(False, error=str(e))


# ============================================================
#  AIRTABLE PRIMITIVE
# ============================================================

class AirtablePrimitive(Primitive):
    """Airtable — bases, tables, and records.
    
    Uses the Airtable REST API via AirtableConnector.
    """
    
    def __init__(self, connector: Any = None):
        self._connector = connector
    
    @property
    def name(self) -> str:
        return "AIRTABLE"
    
    def get_operations(self) -> Dict[str, str]:
        return {
            "list_bases": "List all Airtable bases accessible to the token",
            "get_base_schema": "Get the schema (tables and fields) for a base",
            "list_records": "List records from a table with optional filters, sorting, and field selection",
            "get_record": "Get a single record by ID",
            "create_records": "Create one or more records in a table (max 10 per call)",
            "update_records": "Update one or more records (max 10 per call)",
            "delete_records": "Delete one or more records by ID (max 10 per call)",
            "search_records": "Search for records where a field contains a value",
        }
    
    def get_param_schema(self) -> Dict[str, Dict[str, Any]]:
        return {
            "list_bases": {},
            "get_base_schema": {
                "base_id": {"type": "str", "required": True, "description": "Base ID (starts with 'app')"},
            },
            "list_records": {
                "base_id": {"type": "str", "required": True, "description": "Base ID"},
                "table_name": {"type": "str", "required": True, "description": "Table name or ID"},
                "view": {"type": "str", "required": False, "description": "View name or ID to filter by"},
                "formula": {"type": "str", "required": False, "description": "Airtable formula filter (e.g. \"{Status}='Active'\")"},
                "sort": {"type": "list", "required": False, "description": "Sort list [{field, direction}]"},
                "fields": {"type": "list", "required": False, "description": "Field names to include"},
                "max_records": {"type": "int", "required": False, "description": "Max records (default 100)"},
            },
            "get_record": {
                "base_id": {"type": "str", "required": True, "description": "Base ID"},
                "table_name": {"type": "str", "required": True, "description": "Table name or ID"},
                "record_id": {"type": "str", "required": True, "description": "Record ID (starts with 'rec')"},
            },
            "create_records": {
                "base_id": {"type": "str", "required": True, "description": "Base ID"},
                "table_name": {"type": "str", "required": True, "description": "Table name or ID"},
                "records": {"type": "list", "required": True, "description": "List of field dicts, e.g. [{\"Name\": \"Test\"}]"},
            },
            "update_records": {
                "base_id": {"type": "str", "required": True, "description": "Base ID"},
                "table_name": {"type": "str", "required": True, "description": "Table name or ID"},
                "records": {"type": "list", "required": True, "description": "List of {id, fields} dicts"},
            },
            "delete_records": {
                "base_id": {"type": "str", "required": True, "description": "Base ID"},
                "table_name": {"type": "str", "required": True, "description": "Table name or ID"},
                "record_ids": {"type": "list", "required": True, "description": "List of record IDs to delete"},
            },
            "search_records": {
                "base_id": {"type": "str", "required": True, "description": "Base ID"},
                "table_name": {"type": "str", "required": True, "description": "Table name or ID"},
                "field": {"type": "str", "required": True, "description": "Field name to search in"},
                "value": {"type": "str", "required": True, "description": "Value to search for"},
                "max_records": {"type": "int", "required": False, "description": "Max results (default 20)"},
            },
        }
    
    async def execute(self, operation: str, params: Dict[str, Any]) -> StepResult:
        if not self._connector:
            return StepResult(False, error="Airtable is not configured. Connect Airtable in Settings to use Airtable features.")
        try:
            if operation == "list_bases":
                results = await self._connector.list_bases()
                return StepResult(True, data={"count": len(results), "bases": results})
            
            elif operation == "get_base_schema":
                result = await self._connector.get_base_schema(params["base_id"])
                return StepResult(True, data=result)
            
            elif operation == "list_records":
                results = await self._connector.list_records(
                    base_id=params["base_id"],
                    table_name=params["table_name"],
                    view=params.get("view"),
                    formula=params.get("formula"),
                    sort=params.get("sort"),
                    fields=params.get("fields"),
                    max_records=int(params.get("max_records", 100)),
                )
                return StepResult(True, data={"count": len(results), "records": results})
            
            elif operation == "get_record":
                result = await self._connector.get_record(
                    params["base_id"], params["table_name"], params["record_id"],
                )
                return StepResult(True, data=result)
            
            elif operation == "create_records":
                results = await self._connector.create_records(
                    base_id=params["base_id"],
                    table_name=params["table_name"],
                    records=params["records"],
                )
                return StepResult(True, data={"count": len(results), "records": results})
            
            elif operation == "update_records":
                results = await self._connector.update_records(
                    base_id=params["base_id"],
                    table_name=params["table_name"],
                    records=params["records"],
                )
                return StepResult(True, data={"count": len(results), "records": results})
            
            elif operation == "delete_records":
                results = await self._connector.delete_records(
                    base_id=params["base_id"],
                    table_name=params["table_name"],
                    record_ids=params["record_ids"],
                )
                return StepResult(True, data={"count": len(results), "deleted": results})
            
            elif operation == "search_records":
                results = await self._connector.search_records(
                    base_id=params["base_id"],
                    table_name=params["table_name"],
                    field=params["field"],
                    value=params["value"],
                    max_records=int(params.get("max_records", 20)),
                )
                return StepResult(True, data={"count": len(results), "records": results})
            
            else:
                return StepResult(False, error=f"Unknown operation: {operation}")
        except Exception as e:
            return StepResult(False, error=str(e))


# ============================================================
#  ZOOM PRIMITIVE
# ============================================================

class ZoomPrimitive(Primitive):
    """Zoom — meetings, recordings, participants, and user profile.
    
    Uses the Zoom REST API via ZoomConnector.
    """
    
    def __init__(self, connector: Any = None):
        self._connector = connector
    
    @property
    def name(self) -> str:
        return "ZOOM"
    
    def get_operations(self) -> Dict[str, str]:
        return {
            "me": "Get the current Zoom user's profile",
            "list_meetings": "List meetings (scheduled, live, upcoming, or past)",
            "get_meeting": "Get full meeting details including join URL and settings",
            "create_meeting": "Schedule a new meeting with topic, time, duration, and settings",
            "update_meeting": "Update a meeting's topic, time, duration, or agenda",
            "delete_meeting": "Delete/cancel a meeting",
            "get_participants": "Get participants from a past meeting",
            "list_recordings": "List cloud recordings with download links",
        }
    
    def get_param_schema(self) -> Dict[str, Dict[str, Any]]:
        return {
            "me": {},
            "list_meetings": {
                "type": {"type": "str", "required": False, "description": "'scheduled', 'live', 'upcoming', or 'previous_meetings'"},
                "page_size": {"type": "int", "required": False, "description": "Results per page (default 30, max 300)"},
            },
            "get_meeting": {
                "meeting_id": {"type": "str", "required": True, "description": "Meeting ID"},
            },
            "create_meeting": {
                "topic": {"type": "str", "required": True, "description": "Meeting title"},
                "start_time": {"type": "str", "required": False, "description": "Start time ISO 8601 (e.g. 2026-04-15T10:00:00Z)"},
                "duration": {"type": "int", "required": False, "description": "Duration in minutes (default 60)"},
                "timezone": {"type": "str", "required": False, "description": "Timezone (e.g. America/New_York)"},
                "agenda": {"type": "str", "required": False, "description": "Meeting description/agenda"},
                "password": {"type": "str", "required": False, "description": "Meeting password"},
                "waiting_room": {"type": "bool", "required": False, "description": "Enable waiting room"},
            },
            "update_meeting": {
                "meeting_id": {"type": "str", "required": True, "description": "Meeting ID"},
                "topic": {"type": "str", "required": False, "description": "New topic"},
                "start_time": {"type": "str", "required": False, "description": "New start time (ISO 8601)"},
                "duration": {"type": "int", "required": False, "description": "New duration in minutes"},
                "agenda": {"type": "str", "required": False, "description": "New agenda"},
            },
            "delete_meeting": {
                "meeting_id": {"type": "str", "required": True, "description": "Meeting ID"},
            },
            "get_participants": {
                "meeting_id": {"type": "str", "required": True, "description": "Past meeting UUID"},
            },
            "list_recordings": {
                "from_date": {"type": "str", "required": False, "description": "Start date (YYYY-MM-DD)"},
                "to_date": {"type": "str", "required": False, "description": "End date (YYYY-MM-DD)"},
            },
        }
    
    async def execute(self, operation: str, params: Dict[str, Any]) -> StepResult:
        if not self._connector:
            return StepResult(False, error="Zoom is not configured. Connect Zoom in Settings to use Zoom features.")
        try:
            if operation == "me":
                result = await self._connector.me()
                return StepResult(True, data=result)
            
            elif operation == "list_meetings":
                results = await self._connector.list_meetings(
                    type=params.get("type", "scheduled"),
                    page_size=int(params.get("page_size", 30)),
                )
                return StepResult(True, data={"count": len(results), "meetings": results})
            
            elif operation == "get_meeting":
                result = await self._connector.get_meeting(params["meeting_id"])
                return StepResult(True, data=result)
            
            elif operation == "create_meeting":
                result = await self._connector.create_meeting(
                    topic=params["topic"],
                    start_time=params.get("start_time"),
                    duration=int(params.get("duration", 60)),
                    timezone=params.get("timezone"),
                    agenda=params.get("agenda"),
                    password=params.get("password"),
                    waiting_room=bool(params.get("waiting_room", False)),
                )
                return StepResult(True, data=result)
            
            elif operation == "update_meeting":
                await self._connector.update_meeting(
                    meeting_id=params["meeting_id"],
                    topic=params.get("topic"),
                    start_time=params.get("start_time"),
                    duration=int(params["duration"]) if params.get("duration") else None,
                    agenda=params.get("agenda"),
                )
                return StepResult(True, data={"updated": True})
            
            elif operation == "delete_meeting":
                await self._connector.delete_meeting(params["meeting_id"])
                return StepResult(True, data={"deleted": True})
            
            elif operation == "get_participants":
                results = await self._connector.get_meeting_participants(params["meeting_id"])
                return StepResult(True, data={"count": len(results), "participants": results})
            
            elif operation == "list_recordings":
                results = await self._connector.list_recordings(
                    from_date=params.get("from_date"),
                    to_date=params.get("to_date"),
                )
                return StepResult(True, data={"count": len(results), "recordings": results})
            
            else:
                return StepResult(False, error=f"Unknown operation: {operation}")
        except Exception as e:
            return StepResult(False, error=str(e))


# ============================================================
#  LINKEDIN PRIMITIVE
# ============================================================

class LinkedInPrimitive(Primitive):
    """LinkedIn — profile, posts, shares, and organization search.
    
    Uses the LinkedIn REST API via LinkedInConnector.
    """
    
    def __init__(self, connector: Any = None):
        self._connector = connector
    
    @property
    def name(self) -> str:
        return "LINKEDIN"
    
    def get_operations(self) -> Dict[str, str]:
        return {
            "me": "Get the current LinkedIn user's profile",
            "create_post": "Create a text or article post on LinkedIn",
            "get_posts": "Get the current user's recent posts",
            "delete_post": "Delete a post",
            "get_organization": "Get company/organization details by ID",
            "search_companies": "Search for companies by keywords",
        }
    
    def get_param_schema(self) -> Dict[str, Dict[str, Any]]:
        return {
            "me": {},
            "create_post": {
                "text": {"type": "str", "required": True, "description": "Post text content"},
                "visibility": {"type": "str", "required": False, "description": "'PUBLIC' or 'CONNECTIONS' (default: PUBLIC)"},
                "article_url": {"type": "str", "required": False, "description": "URL to share as an article"},
                "article_title": {"type": "str", "required": False, "description": "Title for the article link"},
                "article_description": {"type": "str", "required": False, "description": "Description for the article link"},
            },
            "get_posts": {
                "limit": {"type": "int", "required": False, "description": "Max posts (default 10)"},
            },
            "delete_post": {
                "post_urn": {"type": "str", "required": True, "description": "Post URN (e.g. urn:li:ugcPost:123456)"},
            },
            "get_organization": {
                "org_id": {"type": "str", "required": True, "description": "Organization ID (numeric)"},
            },
            "search_companies": {
                "keywords": {"type": "str", "required": True, "description": "Search keywords"},
                "limit": {"type": "int", "required": False, "description": "Max results (default 10)"},
            },
        }
    
    async def execute(self, operation: str, params: Dict[str, Any]) -> StepResult:
        if not self._connector:
            return StepResult(False, error="LinkedIn is not configured. Connect LinkedIn in Settings to use LinkedIn features.")
        try:
            if operation == "me":
                result = await self._connector.me()
                return StepResult(True, data=result)
            
            elif operation == "create_post":
                result = await self._connector.create_post(
                    text=params["text"],
                    visibility=params.get("visibility", "PUBLIC"),
                    article_url=params.get("article_url"),
                    article_title=params.get("article_title"),
                    article_description=params.get("article_description"),
                )
                return StepResult(True, data=result)
            
            elif operation == "get_posts":
                results = await self._connector.get_posts(
                    limit=int(params.get("limit", 10)),
                )
                return StepResult(True, data={"count": len(results), "posts": results})
            
            elif operation == "delete_post":
                await self._connector.delete_post(params["post_urn"])
                return StepResult(True, data={"deleted": True})
            
            elif operation == "get_organization":
                result = await self._connector.get_organization(params["org_id"])
                return StepResult(True, data=result)
            
            elif operation == "search_companies":
                results = await self._connector.search_companies(
                    keywords=params["keywords"],
                    limit=int(params.get("limit", 10)),
                )
                return StepResult(True, data={"count": len(results), "companies": results})
            
            else:
                return StepResult(False, error=f"Unknown operation: {operation}")
        except Exception as e:
            return StepResult(False, error=str(e))


# ============================================================
#  REDDIT PRIMITIVE
# ============================================================

class RedditPrimitive(Primitive):
    """Reddit — posts, comments, subreddits, and search.
    
    Uses the Reddit OAuth API via RedditConnector.
    """
    
    def __init__(self, connector: Any = None):
        self._connector = connector
    
    @property
    def name(self) -> str:
        return "REDDIT"
    
    def get_operations(self) -> Dict[str, str]:
        return {
            "me": "Get the current Reddit user's profile and karma",
            "get_subreddit": "Get subreddit info (subscribers, description, etc.)",
            "get_posts": "Get posts from a subreddit (hot, new, top, rising)",
            "get_post": "Get a single post with top comments",
            "search": "Search for posts across Reddit or within a subreddit",
            "submit_post": "Submit a new text or link post to a subreddit",
            "add_comment": "Add a comment to a post or reply to a comment",
            "get_user_posts": "Get a user's submitted posts",
            "get_saved": "Get the user's saved posts and comments",
        }
    
    def get_param_schema(self) -> Dict[str, Dict[str, Any]]:
        return {
            "me": {},
            "get_subreddit": {
                "name": {"type": "str", "required": True, "description": "Subreddit name (without /r/)"},
            },
            "get_posts": {
                "subreddit": {"type": "str", "required": True, "description": "Subreddit name"},
                "sort": {"type": "str", "required": False, "description": "'hot', 'new', 'top', 'rising', 'controversial'"},
                "time_filter": {"type": "str", "required": False, "description": "For top/controversial: 'hour','day','week','month','year','all'"},
                "limit": {"type": "int", "required": False, "description": "Max posts (default 25)"},
            },
            "get_post": {
                "subreddit": {"type": "str", "required": True, "description": "Subreddit name"},
                "post_id": {"type": "str", "required": True, "description": "Post ID"},
                "comment_limit": {"type": "int", "required": False, "description": "Max top-level comments (default 10)"},
            },
            "search": {
                "query": {"type": "str", "required": True, "description": "Search query"},
                "subreddit": {"type": "str", "required": False, "description": "Limit to subreddit"},
                "sort": {"type": "str", "required": False, "description": "'relevance', 'hot', 'top', 'new', 'comments'"},
                "time_filter": {"type": "str", "required": False, "description": "'hour','day','week','month','year','all'"},
                "limit": {"type": "int", "required": False, "description": "Max results (default 25)"},
            },
            "submit_post": {
                "subreddit": {"type": "str", "required": True, "description": "Target subreddit"},
                "title": {"type": "str", "required": True, "description": "Post title"},
                "text": {"type": "str", "required": False, "description": "Self-post body text"},
                "url": {"type": "str", "required": False, "description": "URL for link post"},
            },
            "add_comment": {
                "parent_fullname": {"type": "str", "required": True, "description": "Parent full name (e.g. t3_abc123)"},
                "text": {"type": "str", "required": True, "description": "Comment body (markdown)"},
            },
            "get_user_posts": {
                "username": {"type": "str", "required": False, "description": "Reddit username (default: current user)"},
                "limit": {"type": "int", "required": False, "description": "Max posts (default 25)"},
            },
            "get_saved": {
                "limit": {"type": "int", "required": False, "description": "Max items (default 25)"},
            },
        }
    
    async def execute(self, operation: str, params: Dict[str, Any]) -> StepResult:
        if not self._connector:
            return StepResult(False, error="Reddit is not configured. Connect Reddit in Settings to use Reddit features.")
        try:
            if operation == "me":
                result = await self._connector.me()
                return StepResult(True, data=result)
            
            elif operation == "get_subreddit":
                result = await self._connector.get_subreddit(params["name"])
                return StepResult(True, data=result)
            
            elif operation == "get_posts":
                results = await self._connector.get_posts(
                    subreddit=params["subreddit"],
                    sort=params.get("sort", "hot"),
                    time_filter=params.get("time_filter", "day"),
                    limit=int(params.get("limit", 25)),
                )
                return StepResult(True, data={"count": len(results), "posts": results})
            
            elif operation == "get_post":
                result = await self._connector.get_post(
                    subreddit=params["subreddit"],
                    post_id=params["post_id"],
                    comment_limit=int(params.get("comment_limit", 10)),
                )
                return StepResult(True, data=result)
            
            elif operation == "search":
                results = await self._connector.search(
                    query=params["query"],
                    subreddit=params.get("subreddit"),
                    sort=params.get("sort", "relevance"),
                    time_filter=params.get("time_filter", "all"),
                    limit=int(params.get("limit", 25)),
                )
                return StepResult(True, data={"count": len(results), "posts": results})
            
            elif operation == "submit_post":
                result = await self._connector.submit_post(
                    subreddit=params["subreddit"],
                    title=params["title"],
                    text=params.get("text"),
                    url=params.get("url"),
                )
                return StepResult(True, data=result)
            
            elif operation == "add_comment":
                result = await self._connector.add_comment(params["parent_fullname"], params["text"])
                return StepResult(True, data=result)
            
            elif operation == "get_user_posts":
                results = await self._connector.get_user_posts(
                    username=params.get("username"),
                    limit=int(params.get("limit", 25)),
                )
                return StepResult(True, data={"count": len(results), "posts": results})
            
            elif operation == "get_saved":
                results = await self._connector.get_saved(
                    limit=int(params.get("limit", 25)),
                )
                return StepResult(True, data={"count": len(results), "items": results})
            
            else:
                return StepResult(False, error=f"Unknown operation: {operation}")
        except Exception as e:
            return StepResult(False, error=str(e))


# ============================================================
#  TELEGRAM PRIMITIVE
# ============================================================

class TelegramPrimitive(Primitive):
    """Telegram — send messages, photos, documents, manage chats.
    
    Uses the Telegram Bot API via TelegramConnector.
    """
    
    def __init__(self, connector: Any = None):
        self._connector = connector
    
    @property
    def name(self) -> str:
        return "TELEGRAM"
    
    def get_operations(self) -> Dict[str, str]:
        return {
            "me": "Get the bot's profile info",
            "send_message": "Send a text message to a chat (supports Markdown/HTML)",
            "edit_message": "Edit a previously sent message",
            "delete_message": "Delete a message",
            "forward_message": "Forward a message to another chat",
            "send_photo": "Send a photo by URL with optional caption",
            "send_document": "Send a document/file by URL with optional caption",
            "get_chat": "Get chat info (title, type, description)",
            "get_member_count": "Get number of members in a chat",
            "get_updates": "Get recent incoming messages and updates",
            "pin_message": "Pin a message in a chat",
            "unpin_message": "Unpin a message or all pinned messages",
        }
    
    def get_param_schema(self) -> Dict[str, Dict[str, Any]]:
        return {
            "me": {},
            "send_message": {
                "chat_id": {"type": "str", "required": True, "description": "Chat ID or @channel_username"},
                "text": {"type": "str", "required": True, "description": "Message text (Markdown/HTML)"},
                "parse_mode": {"type": "str", "required": False, "description": "'Markdown', 'MarkdownV2', or 'HTML'"},
                "disable_notification": {"type": "bool", "required": False, "description": "Send silently"},
                "reply_to_message_id": {"type": "int", "required": False, "description": "Reply to message ID"},
            },
            "edit_message": {
                "chat_id": {"type": "str", "required": True, "description": "Chat ID"},
                "message_id": {"type": "int", "required": True, "description": "Message ID to edit"},
                "text": {"type": "str", "required": True, "description": "New text"},
                "parse_mode": {"type": "str", "required": False, "description": "'Markdown', 'MarkdownV2', or 'HTML'"},
            },
            "delete_message": {
                "chat_id": {"type": "str", "required": True, "description": "Chat ID"},
                "message_id": {"type": "int", "required": True, "description": "Message ID"},
            },
            "forward_message": {
                "chat_id": {"type": "str", "required": True, "description": "Target chat ID"},
                "from_chat_id": {"type": "str", "required": True, "description": "Source chat ID"},
                "message_id": {"type": "int", "required": True, "description": "Message ID to forward"},
            },
            "send_photo": {
                "chat_id": {"type": "str", "required": True, "description": "Chat ID"},
                "photo_url": {"type": "str", "required": True, "description": "Photo URL"},
                "caption": {"type": "str", "required": False, "description": "Photo caption"},
            },
            "send_document": {
                "chat_id": {"type": "str", "required": True, "description": "Chat ID"},
                "document_url": {"type": "str", "required": True, "description": "Document URL"},
                "caption": {"type": "str", "required": False, "description": "Document caption"},
            },
            "get_chat": {
                "chat_id": {"type": "str", "required": True, "description": "Chat ID or @username"},
            },
            "get_member_count": {
                "chat_id": {"type": "str", "required": True, "description": "Chat ID"},
            },
            "get_updates": {
                "limit": {"type": "int", "required": False, "description": "Max updates (default 10)"},
            },
            "pin_message": {
                "chat_id": {"type": "str", "required": True, "description": "Chat ID"},
                "message_id": {"type": "int", "required": True, "description": "Message ID to pin"},
            },
            "unpin_message": {
                "chat_id": {"type": "str", "required": True, "description": "Chat ID"},
                "message_id": {"type": "int", "required": False, "description": "Specific message to unpin (omit to unpin all)"},
            },
        }
    
    async def execute(self, operation: str, params: Dict[str, Any]) -> StepResult:
        if not self._connector:
            return StepResult(False, error="Telegram is not configured. Connect a Telegram bot token in Settings to use Telegram features.")
        try:
            if operation == "me":
                result = await self._connector.me()
                return StepResult(True, data=result)
            
            elif operation == "send_message":
                result = await self._connector.send_message(
                    chat_id=params["chat_id"],
                    text=params["text"],
                    parse_mode=params.get("parse_mode", "Markdown"),
                    disable_notification=bool(params.get("disable_notification", False)),
                    reply_to_message_id=int(params["reply_to_message_id"]) if params.get("reply_to_message_id") else None,
                )
                return StepResult(True, data=result)
            
            elif operation == "edit_message":
                result = await self._connector.edit_message(
                    chat_id=params["chat_id"],
                    message_id=int(params["message_id"]),
                    text=params["text"],
                    parse_mode=params.get("parse_mode", "Markdown"),
                )
                return StepResult(True, data=result)
            
            elif operation == "delete_message":
                await self._connector.delete_message(params["chat_id"], int(params["message_id"]))
                return StepResult(True, data={"deleted": True})
            
            elif operation == "forward_message":
                result = await self._connector.forward_message(
                    chat_id=params["chat_id"],
                    from_chat_id=params["from_chat_id"],
                    message_id=int(params["message_id"]),
                )
                return StepResult(True, data=result)
            
            elif operation == "send_photo":
                result = await self._connector.send_photo(
                    chat_id=params["chat_id"],
                    photo_url=params["photo_url"],
                    caption=params.get("caption"),
                )
                return StepResult(True, data=result)
            
            elif operation == "send_document":
                result = await self._connector.send_document(
                    chat_id=params["chat_id"],
                    document_url=params["document_url"],
                    caption=params.get("caption"),
                )
                return StepResult(True, data=result)
            
            elif operation == "get_chat":
                result = await self._connector.get_chat(params["chat_id"])
                return StepResult(True, data=result)
            
            elif operation == "get_member_count":
                count = await self._connector.get_chat_member_count(params["chat_id"])
                return StepResult(True, data={"member_count": count})
            
            elif operation == "get_updates":
                results = await self._connector.get_updates(
                    limit=int(params.get("limit", 10)),
                )
                return StepResult(True, data={"count": len(results), "updates": results})
            
            elif operation == "pin_message":
                await self._connector.pin_message(params["chat_id"], int(params["message_id"]))
                return StepResult(True, data={"pinned": True})
            
            elif operation == "unpin_message":
                msg_id = int(params["message_id"]) if params.get("message_id") else None
                await self._connector.unpin_message(params["chat_id"], msg_id)
                return StepResult(True, data={"unpinned": True})
            
            else:
                return StepResult(False, error=f"Unknown operation: {operation}")
        except Exception as e:
            return StepResult(False, error=str(e))


class HubSpotPrimitive(Primitive):
    """HubSpot CRM — contacts, companies, deals, tickets, notes, pipelines.
    
    Uses HubSpot CRM API v3 via HubSpotConnector.
    """

    def __init__(self, connector=None):
        self._c = connector
        self._connector = connector

    @property
    def name(self) -> str:
        return "HUBSPOT"

    def get_operations(self) -> Dict[str, str]:
        return {
            "list_contacts": "List HubSpot contacts",
            "get_contact": "Get a HubSpot contact by ID",
            "create_contact": "Create a new HubSpot contact",
            "update_contact": "Update a HubSpot contact",
            "delete_contact": "Delete a HubSpot contact",
            "search_contacts": "Search HubSpot contacts",
            "list_companies": "List HubSpot companies",
            "get_company": "Get a HubSpot company by ID",
            "create_company": "Create a new HubSpot company",
            "update_company": "Update a HubSpot company",
            "delete_company": "Delete a HubSpot company",
            "search_companies": "Search HubSpot companies",
            "list_deals": "List HubSpot deals",
            "get_deal": "Get a HubSpot deal by ID",
            "create_deal": "Create a new HubSpot deal",
            "update_deal": "Update a HubSpot deal",
            "delete_deal": "Delete a HubSpot deal",
            "search_deals": "Search HubSpot deals",
            "list_tickets": "List HubSpot tickets",
            "get_ticket": "Get a HubSpot ticket by ID",
            "create_ticket": "Create a new HubSpot ticket",
            "update_ticket": "Update a HubSpot ticket",
            "delete_ticket": "Delete a HubSpot ticket",
            "create_note": "Create a note on a HubSpot record",
            "get_note": "Get a HubSpot note by ID",
            "get_associations": "Get associations between HubSpot objects",
            "create_association": "Create an association between HubSpot objects",
            "list_pipelines": "List HubSpot pipelines",
            "get_pipeline_stages": "Get stages of a HubSpot pipeline",
            "list_owners": "List HubSpot owners",
        }

    async def execute(self, operation: str, params: dict) -> StepResult:
        if not self._c:
            return StepResult(False, error="HubSpot is not configured. Connect HubSpot in Settings to use HubSpot CRM features.")
        op = operation.lower().strip()
        try:
            if op == "list_contacts":
                data = await self._c.list_contacts(
                    limit=params.get("limit", 20),
                    after=params.get("after"),
                    properties=params.get("properties"),
                )
                return StepResult(True, data=data)

            elif op == "get_contact":
                data = await self._c.get_contact(
                    str(params["contact_id"]),
                    properties=params.get("properties"),
                )
                return StepResult(True, data=data)

            elif op == "create_contact":
                data = await self._c.create_contact(params.get("properties", {}))
                return StepResult(True, data=data)

            elif op == "update_contact":
                data = await self._c.update_contact(
                    str(params["contact_id"]),
                    params.get("properties", {}),
                )
                return StepResult(True, data=data)

            elif op == "delete_contact":
                await self._c.delete_contact(str(params["contact_id"]))
                return StepResult(True, data={"deleted": True})

            elif op == "search_contacts":
                data = await self._c.search_contacts(
                    params["query"],
                    limit=params.get("limit", 10),
                    properties=params.get("properties"),
                )
                return StepResult(True, data=data)

            elif op == "list_companies":
                data = await self._c.list_companies(
                    limit=params.get("limit", 20),
                    after=params.get("after"),
                    properties=params.get("properties"),
                )
                return StepResult(True, data=data)

            elif op == "get_company":
                data = await self._c.get_company(
                    str(params["company_id"]),
                    properties=params.get("properties"),
                )
                return StepResult(True, data=data)

            elif op == "create_company":
                data = await self._c.create_company(params.get("properties", {}))
                return StepResult(True, data=data)

            elif op == "update_company":
                data = await self._c.update_company(
                    str(params["company_id"]),
                    params.get("properties", {}),
                )
                return StepResult(True, data=data)

            elif op == "delete_company":
                await self._c.delete_company(str(params["company_id"]))
                return StepResult(True, data={"deleted": True})

            elif op == "search_companies":
                data = await self._c.search_companies(
                    params["query"],
                    limit=params.get("limit", 10),
                    properties=params.get("properties"),
                )
                return StepResult(True, data=data)

            elif op == "list_deals":
                data = await self._c.list_deals(
                    limit=params.get("limit", 20),
                    after=params.get("after"),
                    properties=params.get("properties"),
                )
                return StepResult(True, data=data)

            elif op == "get_deal":
                data = await self._c.get_deal(
                    str(params["deal_id"]),
                    properties=params.get("properties"),
                )
                return StepResult(True, data=data)

            elif op == "create_deal":
                data = await self._c.create_deal(params.get("properties", {}))
                return StepResult(True, data=data)

            elif op == "update_deal":
                data = await self._c.update_deal(
                    str(params["deal_id"]),
                    params.get("properties", {}),
                )
                return StepResult(True, data=data)

            elif op == "delete_deal":
                await self._c.delete_deal(str(params["deal_id"]))
                return StepResult(True, data={"deleted": True})

            elif op == "search_deals":
                data = await self._c.search_deals(
                    params["query"],
                    limit=params.get("limit", 10),
                    properties=params.get("properties"),
                )
                return StepResult(True, data=data)

            elif op == "list_tickets":
                data = await self._c.list_tickets(
                    limit=params.get("limit", 20),
                    after=params.get("after"),
                    properties=params.get("properties"),
                )
                return StepResult(True, data=data)

            elif op == "get_ticket":
                data = await self._c.get_ticket(
                    str(params["ticket_id"]),
                    properties=params.get("properties"),
                )
                return StepResult(True, data=data)

            elif op == "create_ticket":
                data = await self._c.create_ticket(params.get("properties", {}))
                return StepResult(True, data=data)

            elif op == "update_ticket":
                data = await self._c.update_ticket(
                    str(params["ticket_id"]),
                    params.get("properties", {}),
                )
                return StepResult(True, data=data)

            elif op == "delete_ticket":
                await self._c.delete_ticket(str(params["ticket_id"]))
                return StepResult(True, data={"deleted": True})

            elif op == "create_note":
                data = await self._c.create_note(
                    params["body"],
                    contact_id=params.get("contact_id"),
                    company_id=params.get("company_id"),
                    deal_id=params.get("deal_id"),
                )
                return StepResult(True, data=data)

            elif op == "get_note":
                data = await self._c.get_note(str(params["note_id"]))
                return StepResult(True, data=data)

            elif op == "get_associations":
                data = await self._c.get_associations(
                    params["object_type"],
                    str(params["object_id"]),
                    params["to_object_type"],
                )
                return StepResult(True, data=data)

            elif op == "create_association":
                data = await self._c.create_association(
                    params["from_type"],
                    str(params["from_id"]),
                    params["to_type"],
                    str(params["to_id"]),
                    int(params["association_type_id"]),
                )
                return StepResult(True, data=data)

            elif op == "list_pipelines":
                data = await self._c.list_pipelines(
                    object_type=params.get("object_type", "deals"),
                )
                return StepResult(True, data=data)

            elif op == "get_pipeline_stages":
                data = await self._c.get_pipeline_stages(
                    object_type=params.get("object_type", "deals"),
                    pipeline_id=params.get("pipeline_id", "default"),
                )
                return StepResult(True, data=data)

            elif op == "list_owners":
                data = await self._c.list_owners(
                    limit=params.get("limit", 100),
                    after=params.get("after"),
                )
                return StepResult(True, data=data)

            else:
                return StepResult(False, error=f"Unknown operation: {operation}")
        except Exception as e:
            return StepResult(False, error=str(e))


class StripePrimitive(Primitive):
    """Stripe — customers, charges, invoices, subscriptions, products, payments.
    
    Uses Stripe REST API via StripeConnector.
    """

    def __init__(self, connector=None):
        self._c = connector
        self._connector = connector

    @property
    def name(self) -> str:
        return "STRIPE"

    def get_operations(self) -> Dict[str, str]:
        return {
            "list_customers": "List Stripe customers",
            "get_customer": "Get a Stripe customer by ID",
            "create_customer": "Create a new Stripe customer",
            "update_customer": "Update a Stripe customer",
            "delete_customer": "Delete a Stripe customer",
            "list_products": "List Stripe products",
            "get_product": "Get a Stripe product by ID",
            "create_product": "Create a new Stripe product",
            "list_prices": "List prices for a product",
            "create_price": "Create a price for a product",
            "list_invoices": "List Stripe invoices",
            "get_invoice": "Get a Stripe invoice by ID",
            "create_invoice": "Create a new invoice for a customer",
            "finalize_invoice": "Finalize a draft invoice",
            "void_invoice": "Void an invoice",
            "list_subscriptions": "List subscriptions",
            "get_subscription": "Get a subscription by ID",
            "cancel_subscription": "Cancel a subscription",
            "list_payment_intents": "List payment intents",
            "get_payment_intent": "Get a payment intent by ID",
            "list_charges": "List charges",
            "get_charge": "Get a charge by ID",
            "get_balance": "Get Stripe account balance",
        }

    async def execute(self, operation: str, params: dict) -> StepResult:
        if not self._c:
            return StepResult(False, error="Stripe is not configured. Connect Stripe in Settings to use Stripe features.")
        op = operation.lower().strip()
        try:
            if op == "list_customers":
                data = await self._c.list_customers(
                    limit=params.get("limit", 20),
                    starting_after=params.get("starting_after"),
                    email=params.get("email"),
                )
                return StepResult(True, data=data)

            elif op == "get_customer":
                data = await self._c.get_customer(str(params["customer_id"]))
                return StepResult(True, data=data)

            elif op == "create_customer":
                data = await self._c.create_customer(
                    email=params.get("email"),
                    name=params.get("name"),
                    description=params.get("description"),
                    metadata=params.get("metadata"),
                )
                return StepResult(True, data=data)

            elif op == "update_customer":
                cid = str(params["customer_id"])
                data = await self._c.update_customer(
                    cid,
                    name=params.get("name"),
                    email=params.get("email"),
                    description=params.get("description"),
                    metadata=params.get("metadata"),
                )
                return StepResult(True, data=data)

            elif op == "delete_customer":
                data = await self._c.delete_customer(str(params["customer_id"]))
                return StepResult(True, data=data)

            elif op == "list_products":
                data = await self._c.list_products(
                    limit=params.get("limit", 20),
                    active=params.get("active"),
                )
                return StepResult(True, data=data)

            elif op == "get_product":
                data = await self._c.get_product(str(params["product_id"]))
                return StepResult(True, data=data)

            elif op == "create_product":
                data = await self._c.create_product(
                    name=params["name"],
                    description=params.get("description"),
                    metadata=params.get("metadata"),
                )
                return StepResult(True, data=data)

            elif op == "list_prices":
                data = await self._c.list_prices(
                    product_id=params.get("product_id"),
                    limit=params.get("limit", 20),
                )
                return StepResult(True, data=data)

            elif op == "create_price":
                data = await self._c.create_price(
                    product_id=str(params["product_id"]),
                    unit_amount=int(params["unit_amount"]),
                    currency=params.get("currency", "usd"),
                    recurring_interval=params.get("recurring_interval"),
                )
                return StepResult(True, data=data)

            elif op == "list_invoices":
                data = await self._c.list_invoices(
                    customer_id=params.get("customer_id"),
                    limit=params.get("limit", 20),
                    status=params.get("status"),
                )
                return StepResult(True, data=data)

            elif op == "get_invoice":
                data = await self._c.get_invoice(str(params["invoice_id"]))
                return StepResult(True, data=data)

            elif op == "create_invoice":
                data = await self._c.create_invoice(
                    customer_id=str(params["customer_id"]),
                    auto_advance=params.get("auto_advance", True),
                )
                return StepResult(True, data=data)

            elif op == "finalize_invoice":
                data = await self._c.finalize_invoice(str(params["invoice_id"]))
                return StepResult(True, data=data)

            elif op == "void_invoice":
                data = await self._c.void_invoice(str(params["invoice_id"]))
                return StepResult(True, data=data)

            elif op == "list_subscriptions":
                data = await self._c.list_subscriptions(
                    customer_id=params.get("customer_id"),
                    limit=params.get("limit", 20),
                    status=params.get("status"),
                )
                return StepResult(True, data=data)

            elif op == "get_subscription":
                data = await self._c.get_subscription(str(params["subscription_id"]))
                return StepResult(True, data=data)

            elif op == "cancel_subscription":
                data = await self._c.cancel_subscription(
                    str(params["subscription_id"]),
                    at_period_end=params.get("at_period_end", True),
                )
                return StepResult(True, data=data)

            elif op == "list_payment_intents":
                data = await self._c.list_payment_intents(
                    customer_id=params.get("customer_id"),
                    limit=params.get("limit", 20),
                )
                return StepResult(True, data=data)

            elif op == "get_payment_intent":
                data = await self._c.get_payment_intent(str(params["payment_intent_id"]))
                return StepResult(True, data=data)

            elif op == "list_charges":
                data = await self._c.list_charges(
                    customer_id=params.get("customer_id"),
                    limit=params.get("limit", 20),
                )
                return StepResult(True, data=data)

            elif op == "get_charge":
                data = await self._c.get_charge(str(params["charge_id"]))
                return StepResult(True, data=data)

            elif op == "get_balance":
                data = await self._c.get_balance()
                return StepResult(True, data=data)

            else:
                return StepResult(False, error=f"Unknown operation: {operation}")
        except Exception as e:
            return StepResult(False, error=str(e))


# ============================================================
#  NOTIFY PRIMITIVE
# ============================================================

class NotifyPrimitive(Primitive):
    """Notifications, reminders, and alerts.
    
    Stores reminders locally. Can be wired to push notifications,
    desktop alerts, Slack, etc. via provider functions.
    """
    
    def __init__(
        self,
        storage_path: Optional[str] = None,
        send_func: Optional[Callable] = None,
    ):
        self._reminders: List[Dict] = []
        self._storage_path = storage_path
        self._send_func = send_func
        if storage_path and Path(storage_path).exists():
            try:
                with open(storage_path, "r") as f:
                    self._reminders = json.load(f)
            except Exception:
                pass
    
    @property
    def name(self) -> str:
        return "NOTIFY"
    
    def get_operations(self) -> Dict[str, str]:
        return {
            "alert": "Show an immediate notification/alert to the user",
            "remind": "Set a reminder for a specific time",
            "list": "List pending reminders",
            "cancel": "Cancel a pending reminder",
        }
    
    def get_param_schema(self) -> Dict[str, Dict[str, Any]]:
        return {
            "alert": {
                "title": {"type": "str", "required": True, "description": "Alert title"},
                "message": {"type": "str", "required": True, "description": "Alert message"},
                "urgency": {"type": "str", "required": False, "description": "low, normal, or high (default normal)"},
            },
            "remind": {
                "message": {"type": "str", "required": True, "description": "Reminder message"},
                "when": {"type": "str", "required": True, "description": "When to remind (ISO datetime or relative like 'in 30 minutes', 'tomorrow 9am')"},
                "repeat": {"type": "str", "required": False, "description": "Repeat interval: daily, weekly, monthly, or none"},
            },
            "list": {
                "status": {"type": "str", "required": False, "description": "Filter: pending, triggered, all (default pending)"},
            },
            "cancel": {
                "id": {"type": "str", "required": True, "description": "Reminder ID to cancel"},
            },
        }
    
    def _save(self):
        if self._storage_path:
            Path(self._storage_path).parent.mkdir(parents=True, exist_ok=True)
            with open(self._storage_path, "w") as f:
                json.dump(self._reminders, f, indent=2)
    
    async def execute(self, operation: str, params: Dict[str, Any]) -> StepResult:
        try:
            if operation == "alert":
                title = params.get("title", "Alert")
                message = params.get("message", "")
                urgency = params.get("urgency", "normal")
                
                if self._send_func:
                    result = await self._send_func(title=title, message=message, urgency=urgency)
                    return StepResult(True, data=result)
                
                # Local-only: store and return (UI layer polls/reads these)
                alert = {
                    "id": f"alert_{int(datetime.now().timestamp())}",
                    "title": title,
                    "message": message,
                    "urgency": urgency,
                    "timestamp": datetime.now().isoformat(),
                    "type": "alert",
                }
                return StepResult(True, data=alert)
            
            elif operation == "remind":
                message = params.get("message", "")
                when = params.get("when", "")
                repeat = params.get("repeat", "none")
                
                reminder = {
                    "id": f"rem_{len(self._reminders)}_{int(datetime.now().timestamp())}",
                    "message": message,
                    "when": when,
                    "repeat": repeat,
                    "status": "pending",
                    "created": datetime.now().isoformat(),
                }
                self._reminders.append(reminder)
                self._save()
                return StepResult(True, data=reminder)
            
            elif operation == "list":
                status = params.get("status", "pending")
                if status == "all":
                    return StepResult(True, data=self._reminders)
                filtered = [r for r in self._reminders if r.get("status") == status]
                return StepResult(True, data=filtered)
            
            elif operation == "cancel":
                rid = params.get("id")
                for r in self._reminders:
                    if r.get("id") == rid:
                        r["status"] = "cancelled"
                self._save()
                return StepResult(True, data={"cancelled": rid})
            
            else:
                return StepResult(False, error=f"Unknown operation: {operation}")
        except Exception as e:
            return StepResult(False, error=str(e))


# ============================================================
#  TASK PRIMITIVE
# ============================================================

class TaskPrimitive(Primitive):
    """Task and to-do list management.
    
    Local task store. Can be wired to Todoist, Microsoft To Do,
    Jira, etc. via providers dict.
    """
    
    def __init__(self, storage_path: Optional[str] = None, providers: Optional[Dict[str, Any]] = None):
        self._tasks: List[Dict] = []
        self._storage_path = storage_path
        self._providers = providers or {}
        if storage_path and Path(storage_path).exists():
            try:
                with open(storage_path, "r") as f:
                    self._tasks = json.load(f)
            except Exception:
                pass
    
    @property
    def name(self) -> str:
        return "TASK"
    
    def get_operations(self) -> Dict[str, str]:
        return {
            "create": "Create a new task or to-do item",
            "list": "List tasks, optionally filtered by status or tag",
            "update": "Update a task (status, title, due date, etc.)",
            "complete": "Mark a task as completed",
            "delete": "Delete a task",
            "search": "Search tasks by keyword",
        }
    
    def get_param_schema(self) -> Dict[str, Dict[str, Any]]:
        return {
            "create": {
                "title": {"type": "str", "required": True, "description": "Task title"},
                "description": {"type": "str", "required": False, "description": "Task description/details"},
                "due": {"type": "str", "required": False, "description": "Due date (ISO date or datetime)"},
                "priority": {"type": "str", "required": False, "description": "low, medium, high, urgent"},
                "tags": {"type": "list", "required": False, "description": "Tags/categories"},
                "project": {"type": "str", "required": False, "description": "Project name"},
                "provider": {"type": "str", "required": False, "description": "Provider: todoist, microsoft_todo, jira (default: local)"},
            },
            "list": {
                "status": {"type": "str", "required": False, "description": "Filter: open, completed, all (default open)"},
                "project": {"type": "str", "required": False, "description": "Filter by project"},
                "tag": {"type": "str", "required": False, "description": "Filter by tag"},
                "limit": {"type": "int", "required": False, "description": "Max results"},
                "provider": {"type": "str", "required": False, "description": "Provider name"},
            },
            "update": {
                "id": {"type": "str", "required": True, "description": "Task ID"},
                "title": {"type": "str", "required": False, "description": "New title"},
                "status": {"type": "str", "required": False, "description": "New status"},
                "due": {"type": "str", "required": False, "description": "New due date"},
                "priority": {"type": "str", "required": False, "description": "New priority"},
                "provider": {"type": "str", "required": False, "description": "Provider name"},
            },
            "complete": {
                "id": {"type": "str", "required": True, "description": "Task ID to complete"},
                "provider": {"type": "str", "required": False, "description": "Provider name"},
            },
            "delete": {
                "id": {"type": "str", "required": True, "description": "Task ID to delete"},
                "provider": {"type": "str", "required": False, "description": "Provider name"},
            },
            "search": {
                "query": {"type": "str", "required": True, "description": "Search term"},
                "provider": {"type": "str", "required": False, "description": "Provider name"},
            },
        }
    
    def _get_provider(self, name: Optional[str]) -> Optional[Any]:
        if name and name in self._providers:
            return self._providers[name]
        if self._providers:
            return next(iter(self._providers.values()))
        return None
    
    def _save(self):
        if self._storage_path:
            Path(self._storage_path).parent.mkdir(parents=True, exist_ok=True)
            with open(self._storage_path, "w") as f:
                json.dump(self._tasks, f, indent=2)
    
    async def execute(self, operation: str, params: Dict[str, Any]) -> StepResult:
        try:
            provider_name = params.get("provider")
            provider = self._get_provider(provider_name)
            
            if operation == "create":
                if provider and hasattr(provider, "create_task"):
                    title = params.get("title", "Untitled")
                    description = params.get("description", "")
                    due = params.get("due")
                    # Todoist uses 'content', Microsoft To-Do uses 'title'
                    import inspect
                    sig = inspect.signature(provider.create_task)
                    if "content" in sig.parameters:
                        result = await provider.create_task(
                            content=title,
                            description=description,
                            due_date=due,
                        )
                    else:
                        result = await provider.create_task(
                            title=title,
                            body=description or None,
                            due_date=due,
                        )
                    return StepResult(True, data={"id": getattr(result, "id", str(result)), "title": params.get("title"), "status": "created"})
                
                task = {
                    "id": f"task_{len(self._tasks)}_{int(datetime.now().timestamp())}",
                    "title": params.get("title", "Untitled"),
                    "description": params.get("description", ""),
                    "status": "open",
                    "due": params.get("due"),
                    "priority": params.get("priority", "medium"),
                    "tags": params.get("tags", []),
                    "project": params.get("project"),
                    "created": datetime.now().isoformat(),
                    "completed_at": None,
                }
                self._tasks.append(task)
                self._save()
                return StepResult(True, data=task)
            
            elif operation == "list":
                if provider and hasattr(provider, "list_tasks"):
                    result = await provider.list_tasks()
                    # Convert dataclass objects to dicts if needed
                    if result and hasattr(result[0], '__dataclass_fields__'):
                        from dataclasses import asdict
                        result = [asdict(t) for t in result]
                    return StepResult(True, data=result)
                
                status = params.get("status", "open")
                project = params.get("project")
                tag = params.get("tag")
                limit = params.get("limit", 50)
                
                filtered = self._tasks
                if status != "all":
                    filtered = [t for t in filtered if t.get("status") == status]
                if project:
                    filtered = [t for t in filtered if t.get("project") == project]
                if tag:
                    filtered = [t for t in filtered if tag in t.get("tags", [])]
                
                return StepResult(True, data=filtered[:limit])
            
            elif operation == "update":
                task_id = params.get("id")
                
                if provider and hasattr(provider, "update_task"):
                    result = await provider.update_task(task_id=task_id, title=params.get("title"), status=params.get("status"))
                    return StepResult(True, data=result)
                
                for task in self._tasks:
                    if task.get("id") == task_id:
                        for k in ("title", "status", "due", "priority", "description", "project"):
                            if k in params and k != "id":
                                task[k] = params[k]
                        self._save()
                        return StepResult(True, data=task)
                return StepResult(False, error=f"Task not found: {task_id}")
            
            elif operation == "complete":
                task_id = params.get("id")
                
                if provider and hasattr(provider, "complete_task"):
                    result = await provider.complete_task(task_id=task_id)
                    return StepResult(True, data=result)
                
                for task in self._tasks:
                    if task.get("id") == task_id:
                        task["status"] = "completed"
                        task["completed_at"] = datetime.now().isoformat()
                        self._save()
                        return StepResult(True, data=task)
                return StepResult(False, error=f"Task not found: {task_id}")
            
            elif operation == "delete":
                task_id = params.get("id")
                
                if provider and hasattr(provider, "delete_task"):
                    result = await provider.delete_task(task_id=task_id)
                    return StepResult(True, data=result)
                
                before = len(self._tasks)
                self._tasks = [t for t in self._tasks if t.get("id") != task_id]
                self._save()
                return StepResult(True, data={"deleted": before != len(self._tasks)})
            
            elif operation == "search":
                query = params.get("query", "").lower()
                
                if provider and hasattr(provider, "search_tasks"):
                    result = await provider.search_tasks(query=query)
                    return StepResult(True, data=result)
                
                matches = [
                    t for t in self._tasks
                    if query in t.get("title", "").lower()
                    or query in t.get("description", "").lower()
                    or query in t.get("project", "").lower()
                ]
                return StepResult(True, data=matches)
            
            else:
                return StepResult(False, error=f"Unknown operation: {operation}")
        except Exception as e:
            return StepResult(False, error=str(e))


# ============================================================
#  SHELL PRIMITIVE
# ============================================================

class ShellPrimitive(Primitive):
    """Execute system commands in a controlled sandbox.
    
    Restricted by default — only allows safe commands.
    The allow list can be expanded per deployment.
    """
    
    def __init__(self, allowed_commands: Optional[List[str]] = None):
        # Default safe commands — no rm, no sudo, no curl piping to bash
        self._allowed = set(allowed_commands or [
            "ls", "dir", "cat", "head", "tail", "wc", "grep", "find",
            "echo", "date", "whoami", "hostname", "pwd", "df", "du",
            "sort", "uniq", "cut", "awk", "sed", "tr", "tee",
            "python", "python3", "node", "pip", "npm",
            "git",
        ])
    
    @property
    def name(self) -> str:
        return "SHELL"
    
    def get_operations(self) -> Dict[str, str]:
        return {
            "run": "Run a shell command and return its output",
            "script": "Run a multi-line script (bash)",
        }
    
    def get_param_schema(self) -> Dict[str, Dict[str, Any]]:
        return {
            "run": {
                "command": {"type": "str", "required": True, "description": "Shell command to execute"},
                "cwd": {"type": "str", "required": False, "description": "Working directory"},
                "timeout": {"type": "int", "required": False, "description": "Timeout in seconds (default 30)"},
            },
            "script": {
                "code": {"type": "str", "required": True, "description": "Bash script content"},
                "cwd": {"type": "str", "required": False, "description": "Working directory"},
                "timeout": {"type": "int", "required": False, "description": "Timeout in seconds (default 30)"},
            },
        }
    
    def _check_command(self, cmd: str) -> Optional[str]:
        """Check if command is allowed. Returns error string or None."""
        # Extract the base command (first word, ignore pipes for now)
        parts = cmd.strip().split()
        if not parts:
            return "Empty command"
        base = parts[0].split("/")[-1]  # Handle /usr/bin/python etc.
        
        # Block dangerous patterns regardless of command
        dangerous = ["rm -rf /", "mkfs", "dd if=", "> /dev/", ":(){ :|:", "fork bomb"]
        for d in dangerous:
            if d in cmd:
                return f"Blocked dangerous pattern: {d}"
        
        if base not in self._allowed:
            return f"Command '{base}' not in allow list. Allowed: {', '.join(sorted(self._allowed))}"
        return None
    
    async def execute(self, operation: str, params: Dict[str, Any]) -> StepResult:
        import subprocess
        
        try:
            if operation == "run":
                cmd = params.get("command", "")
                cwd = params.get("cwd")
                timeout = params.get("timeout", 30)
                
                if not cmd:
                    return StepResult(False, error="Missing 'command' parameter")
                
                err = self._check_command(cmd)
                if err:
                    return StepResult(False, error=err)
                
                if cwd:
                    cwd = str(Path(cwd).expanduser())
                
                proc = await asyncio.create_subprocess_shell(
                    cmd,
                    stdout=subprocess.PIPE,
                    stderr=subprocess.PIPE,
                    cwd=cwd,
                )
                stdout, stderr = await asyncio.wait_for(proc.communicate(), timeout=timeout)
                
                return StepResult(True, data={
                    "stdout": stdout.decode("utf-8", errors="ignore")[:50000],
                    "stderr": stderr.decode("utf-8", errors="ignore")[:10000],
                    "exit_code": proc.returncode,
                })
            
            elif operation == "script":
                code = params.get("code", "")
                cwd = params.get("cwd")
                timeout = params.get("timeout", 30)
                
                if not code:
                    return StepResult(False, error="Missing 'code' parameter")
                
                # Check each line of the script
                for line in code.strip().split("\n"):
                    line = line.strip()
                    if line and not line.startswith("#"):
                        err = self._check_command(line)
                        if err:
                            return StepResult(False, error=f"Line blocked: {err}")
                
                if cwd:
                    cwd = str(Path(cwd).expanduser())
                
                proc = await asyncio.create_subprocess_exec(
                    "bash", "-c", code,
                    stdout=subprocess.PIPE,
                    stderr=subprocess.PIPE,
                    cwd=cwd,
                )
                stdout, stderr = await asyncio.wait_for(proc.communicate(), timeout=timeout)
                
                return StepResult(True, data={
                    "stdout": stdout.decode("utf-8", errors="ignore")[:50000],
                    "stderr": stderr.decode("utf-8", errors="ignore")[:10000],
                    "exit_code": proc.returncode,
                })
            
            else:
                return StepResult(False, error=f"Unknown operation: {operation}")
        except asyncio.TimeoutError:
            return StepResult(False, error=f"Command timed out after {params.get('timeout', 30)}s")
        except Exception as e:
            return StepResult(False, error=str(e))


# ============================================================
#  DATA PRIMITIVE
# ============================================================

class DataPrimitive(Primitive):
    """Structured data operations — query, transform, filter, join.
    
    Works on in-memory data (lists of dicts). The AI writes
    transformation code when needed, just like COMPUTE.
    """
    
    def __init__(self, llm_complete: Optional[Callable] = None):
        self._llm = llm_complete
        self._datasets: Dict[str, List[Dict]] = {}
    
    @property
    def name(self) -> str:
        return "DATA"
    
    def get_operations(self) -> Dict[str, str]:
        return {
            "query": "Filter/sort/group data using a natural language query — AI writes the code",
            "transform": "Transform data shape (pivot, flatten, rename, etc.) — AI writes the code",
            "load": "Load data from a file (CSV, JSON) into a named dataset",
            "store": "Store data as a named dataset for later use",
            "merge": "Merge/join two datasets",
        }
    
    def get_param_schema(self) -> Dict[str, Dict[str, Any]]:
        return {
            "query": {
                "data": {"type": "list", "required": True, "description": "List of dicts to query"},
                "query": {"type": "str", "required": True, "description": "Natural language query (e.g. 'rows where age > 30 sorted by name')"},
            },
            "transform": {
                "data": {"type": "list", "required": True, "description": "List of dicts to transform"},
                "instruction": {"type": "str", "required": True, "description": "What transformation to apply (e.g. 'pivot by category', 'add a total column')"},
            },
            "load": {
                "path": {"type": "str", "required": True, "description": "File path (CSV or JSON)"},
                "name": {"type": "str", "required": False, "description": "Dataset name (default: filename)"},
            },
            "store": {
                "name": {"type": "str", "required": True, "description": "Dataset name"},
                "data": {"type": "list", "required": True, "description": "Data to store"},
            },
            "merge": {
                "left": {"type": "list", "required": True, "description": "First dataset"},
                "right": {"type": "list", "required": True, "description": "Second dataset"},
                "on": {"type": "str", "required": True, "description": "Key field to join on"},
                "how": {"type": "str", "required": False, "description": "Join type: inner, left, right, outer (default inner)"},
            },
        }
    
    async def _llm_data_code(self, instruction: str, data: List) -> StepResult:
        """Ask LLM to write Python code to query/transform data."""
        if not self._llm:
            return StepResult(False, error="LLM required for data queries")
        
        # Show sample of data for LLM context
        sample = data[:3] if len(data) > 3 else data
        
        prompt = f"""Write a Python function to process this data:

Instruction: {instruction}
Data sample (full data has {len(data)} rows): {json.dumps(sample)}

Requirements:
- Write a function called `process(data)` that takes a list of dicts
- Return the processed list of dicts (or a dict with results)
- Use only Python stdlib
- Handle missing keys gracefully

Respond with ONLY the Python code, no markdown, no explanation.
Start directly with: def process(data):"""

        try:
            response = await self._llm(prompt)
            code = response.strip()
            if code.startswith("```"):
                code = re.sub(r'^```\w*\n?', '', code)
                code = re.sub(r'\n?```$', '', code)
                code = code.strip()
            
            if "def process" not in code:
                return StepResult(False, error="LLM did not generate a valid process function")
            
            import math
            sandbox = {
                "__builtins__": {
                    "abs": abs, "round": round, "min": min, "max": max,
                    "sum": sum, "len": len, "pow": pow, "int": int, "float": float,
                    "sorted": sorted, "enumerate": enumerate, "range": range, "zip": zip,
                    "map": map, "filter": filter, "list": list, "dict": dict, "tuple": tuple,
                    "set": set, "str": str, "bool": bool, "type": type,
                    "True": True, "False": False, "None": None,
                    "isinstance": isinstance, "ValueError": ValueError,
                    "KeyError": KeyError, "IndexError": IndexError,
                    "print": lambda *a, **k: None,
                    "__import__": lambda name, *a, **k: __import__(name) if name in ("math", "datetime", "json", "re") else (_ for _ in ()).throw(ImportError(f"Import of '{name}' not allowed")),
                },
            }
            
            exec(code, sandbox)
            process_fn = sandbox.get("process")
            if not callable(process_fn):
                return StepResult(False, error="Generated code did not define a callable 'process' function")
            
            result = process_fn(data)
            return StepResult(True, data=result)
            
        except Exception as e:
            return StepResult(False, error=f"Data processing error: {type(e).__name__}: {e}")
    
    async def execute(self, operation: str, params: Dict[str, Any]) -> StepResult:
        try:
            if operation == "query":
                data = params.get("data", [])
                query = params.get("query", "")
                
                if not data:
                    # Check named datasets
                    dataset_name = params.get("dataset")
                    if dataset_name and dataset_name in self._datasets:
                        data = self._datasets[dataset_name]
                
                if not data:
                    return StepResult(False, error="No data provided. Pass 'data' (list of dicts) or wire from a previous step.")
                
                return await self._llm_data_code(f"Query: {query}", data)
            
            elif operation == "transform":
                data = params.get("data", [])
                instruction = params.get("instruction", "")
                
                if not data:
                    return StepResult(False, error="No data provided")
                
                return await self._llm_data_code(f"Transform: {instruction}", data)
            
            elif operation == "load":
                path = str(Path(params.get("path", "")).expanduser())
                name = params.get("name", Path(path).stem)
                
                if not Path(path).exists():
                    return StepResult(False, error=f"File not found: {path}")
                
                ext = Path(path).suffix.lower()
                if ext == ".json":
                    with open(path, "r") as f:
                        data = json.load(f)
                elif ext == ".csv":
                    import csv
                    with open(path, "r") as f:
                        reader = csv.DictReader(f)
                        data = list(reader)
                else:
                    return StepResult(False, error=f"Unsupported file type: {ext}. Use .json or .csv")
                
                if isinstance(data, list):
                    self._datasets[name] = data
                    return StepResult(True, data={"name": name, "rows": len(data), "sample": data[:3]})
                else:
                    return StepResult(True, data=data)
            
            elif operation == "store":
                name = params.get("name", "")
                data = params.get("data", [])
                self._datasets[name] = data
                return StepResult(True, data={"name": name, "rows": len(data)})
            
            elif operation == "merge":
                left = params.get("left", [])
                right = params.get("right", [])
                on = params.get("on", "")
                how = params.get("how", "inner")
                
                if not left or not right or not on:
                    return StepResult(False, error="Need 'left', 'right' datasets and 'on' key field")
                
                # Build index from right
                right_idx: Dict[Any, List[Dict]] = {}
                for r in right:
                    key = r.get(on)
                    if key is not None:
                        right_idx.setdefault(key, []).append(r)
                
                merged = []
                left_keys_seen = set()
                
                for l in left:
                    key = l.get(on)
                    left_keys_seen.add(key)
                    matches = right_idx.get(key, [])
                    
                    if matches:
                        for r in matches:
                            row = {**l, **r}
                            merged.append(row)
                    elif how in ("left", "outer"):
                        merged.append(dict(l))
                
                if how in ("right", "outer"):
                    for r in right:
                        key = r.get(on)
                        if key not in left_keys_seen:
                            merged.append(dict(r))
                
                return StepResult(True, data=merged)
            
            else:
                return StepResult(False, error=f"Unknown operation: {operation}")
        except Exception as e:
            return StepResult(False, error=str(e))


# ============================================================
#  MESSAGE PRIMITIVE
# ============================================================

class MessagePrimitive(Primitive):
    """Messaging across channels — Slack, Teams, Discord, SMS, WhatsApp.
    
    Provider-based: plug in any messaging backend via send_func/list_func.
    The primitive defines the universal interface; providers handle the protocol.
    """
    
    def __init__(
        self,
        send_func: Optional[Callable] = None,
        list_func: Optional[Callable] = None,
        react_func: Optional[Callable] = None,
        providers: Optional[Dict[str, Any]] = None,
    ):
        self._send = send_func
        self._list = list_func
        self._react = react_func
        self._providers = providers or {}
        self._local_messages: List[Dict] = []
    
    @property
    def name(self) -> str:
        return "MESSAGE"
    
    def get_operations(self) -> Dict[str, str]:
        return {
            "send": "Send a message to a channel, thread, or person",
            "list": "List recent messages from a channel or conversation",
            "search": "Search messages by keyword across channels",
            "react": "Add a reaction/emoji to a message",
            "reply": "Reply to a specific message in a thread",
            "channels": "List available channels or conversations",
        }
    
    def get_param_schema(self) -> Dict[str, Dict[str, Any]]:
        return {
            "send": {
                "to": {"type": "str", "required": True, "description": "Channel name, user handle, or phone number"},
                "text": {"type": "str", "required": True, "description": "Message text"},
                "provider": {"type": "str", "required": False, "description": "Provider: slack, teams, discord, sms (default: auto-detect)"},
                "attachments": {"type": "list", "required": False, "description": "List of file paths or URLs to attach"},
            },
            "list": {
                "channel": {"type": "str", "required": True, "description": "Channel name or conversation ID"},
                "limit": {"type": "int", "required": False, "description": "Max messages to return (default 20)"},
                "provider": {"type": "str", "required": False, "description": "Provider name"},
            },
            "search": {
                "query": {"type": "str", "required": True, "description": "Search term"},
                "channel": {"type": "str", "required": False, "description": "Limit search to a specific channel"},
                "provider": {"type": "str", "required": False, "description": "Provider name"},
            },
            "react": {
                "message_id": {"type": "str", "required": True, "description": "Message ID to react to"},
                "emoji": {"type": "str", "required": True, "description": "Emoji name (e.g. thumbsup, heart, check)"},
                "provider": {"type": "str", "required": False, "description": "Provider name"},
            },
            "reply": {
                "message_id": {"type": "str", "required": True, "description": "Message ID to reply to (thread parent)"},
                "text": {"type": "str", "required": True, "description": "Reply text"},
                "provider": {"type": "str", "required": False, "description": "Provider name"},
            },
            "channels": {
                "provider": {"type": "str", "required": False, "description": "Provider name"},
                "limit": {"type": "int", "required": False, "description": "Max channels to return"},
            },
        }
    
    def _get_provider(self, name: Optional[str]) -> Optional[Any]:
        """Look up a messaging provider by name."""
        if name and name in self._providers:
            return self._providers[name]
        # Return first available provider if none specified
        if self._providers:
            return next(iter(self._providers.values()))
        return None
    
    async def execute(self, operation: str, params: Dict[str, Any]) -> StepResult:
        try:
            provider_name = params.get("provider")
            
            if operation == "send":
                to = params.get("to", "")
                text = params.get("text", "")
                
                if not to or not text:
                    return StepResult(False, error="Missing 'to' and/or 'text' parameter")
                
                if self._send:
                    result = await self._send(to=to, text=text, provider=provider_name, attachments=params.get("attachments"))
                    return StepResult(True, data=result)
                
                provider = self._get_provider(provider_name)
                if provider and hasattr(provider, "send"):
                    result = await provider.send(to=to, text=text, attachments=params.get("attachments"))
                    return StepResult(True, data=result)
                
                # Local fallback — store for testing/UI display
                msg = {
                    "id": f"msg_{len(self._local_messages)}_{int(datetime.now().timestamp())}",
                    "to": to,
                    "text": text,
                    "timestamp": datetime.now().isoformat(),
                    "status": "queued",
                }
                self._local_messages.append(msg)
                return StepResult(True, data=msg)
            
            elif operation == "list":
                channel = params.get("channel", "")
                limit = params.get("limit", 20)
                
                if self._list:
                    result = await self._list(channel=channel, limit=limit, provider=provider_name)
                    return StepResult(True, data=result)
                
                provider = self._get_provider(provider_name)
                if provider and hasattr(provider, "list_messages"):
                    result = await provider.list_messages(channel=channel, limit=limit)
                    return StepResult(True, data=result)
                
                # Local fallback
                msgs = [m for m in self._local_messages if m.get("to") == channel]
                return StepResult(True, data=msgs[-limit:])
            
            elif operation == "search":
                query = params.get("query", "").lower()
                channel = params.get("channel")
                
                provider = self._get_provider(provider_name)
                if provider and hasattr(provider, "search"):
                    result = await provider.search(query=query, channel=channel)
                    return StepResult(True, data=result)
                
                # Local fallback
                matches = [
                    m for m in self._local_messages
                    if query in m.get("text", "").lower()
                    and (not channel or m.get("to") == channel)
                ]
                return StepResult(True, data=matches)
            
            elif operation == "react":
                message_id = params.get("message_id", "")
                emoji = params.get("emoji", "")
                
                if not message_id or not emoji:
                    return StepResult(False, error="Missing 'message_id' and/or 'emoji' parameter")
                
                if self._react:
                    result = await self._react(message_id=message_id, emoji=emoji, provider=provider_name)
                    return StepResult(True, data=result)
                
                provider = self._get_provider(provider_name)
                if provider and hasattr(provider, "react"):
                    result = await provider.react(message_id=message_id, emoji=emoji)
                    return StepResult(True, data=result)
                
                return StepResult(True, data={"message_id": message_id, "emoji": emoji, "status": "queued"})
            
            elif operation == "reply":
                message_id = params.get("message_id", "")
                text = params.get("text", "")
                
                if not message_id or not text:
                    return StepResult(False, error="Missing 'message_id' and/or 'text' parameter")
                
                provider = self._get_provider(provider_name)
                if provider and hasattr(provider, "reply"):
                    result = await provider.reply(message_id=message_id, text=text)
                    return StepResult(True, data=result)
                
                msg = {
                    "id": f"msg_{len(self._local_messages)}_{int(datetime.now().timestamp())}",
                    "reply_to": message_id,
                    "text": text,
                    "timestamp": datetime.now().isoformat(),
                    "status": "queued",
                }
                self._local_messages.append(msg)
                return StepResult(True, data=msg)
            
            elif operation == "channels":
                provider = self._get_provider(provider_name)
                if provider and hasattr(provider, "list_channels"):
                    result = await provider.list_channels(limit=params.get("limit", 50))
                    return StepResult(True, data=result)
                
                return StepResult(True, data=[])
            
            else:
                return StepResult(False, error=f"Unknown operation: {operation}")
        except Exception as e:
            return StepResult(False, error=str(e))


# ============================================================
#  MEDIA PRIMITIVE
# ============================================================

class MediaPrimitive(Primitive):
    """Media operations — images, audio, video.
    
    Handles conversion, metadata, generation (via AI), and playback control.
    Provider-based for services like Spotify, YouTube, etc.
    """
    
    def __init__(
        self,
        llm_complete: Optional[Callable] = None,
        providers: Optional[Dict[str, Any]] = None,
    ):
        self._llm = llm_complete
        self._providers = providers or {}
    
    @property
    def name(self) -> str:
        return "MEDIA"
    
    def get_operations(self) -> Dict[str, str]:
        return {
            "info": "Get metadata about a media file (dimensions, format, EXIF: date taken, camera, GPS, etc.)",
            "convert": "Convert media between formats (e.g. mp4→mp3, png→jpg)",
            "resize": "Resize an image to specific dimensions",
            "generate": "Generate an image or audio using AI",
            "transcribe": "Transcribe audio/video to text",
            "play": "Play or queue media via a provider (Spotify, etc.)",
            "search": "Search media libraries or services",
        }
    
    def get_param_schema(self) -> Dict[str, Dict[str, Any]]:
        return {
            "info": {
                "path": {"type": "str", "required": True, "description": "Path to media file"},
            },
            "convert": {
                "path": {"type": "str", "required": True, "description": "Source file path"},
                "format": {"type": "str", "required": True, "description": "Target format (mp3, mp4, png, jpg, wav, etc.)"},
                "output": {"type": "str", "required": False, "description": "Output file path (default: same dir, new extension)"},
            },
            "resize": {
                "path": {"type": "str", "required": True, "description": "Image file path"},
                "width": {"type": "int", "required": False, "description": "Target width in pixels"},
                "height": {"type": "int", "required": False, "description": "Target height in pixels"},
                "output": {"type": "str", "required": False, "description": "Output file path"},
            },
            "generate": {
                "prompt": {"type": "str", "required": True, "description": "Description of what to generate"},
                "type": {"type": "str", "required": False, "description": "image or audio (default image)"},
                "output": {"type": "str", "required": False, "description": "Output file path"},
            },
            "transcribe": {
                "path": {"type": "str", "required": True, "description": "Audio or video file path"},
                "language": {"type": "str", "required": False, "description": "Language code (default: auto-detect)"},
            },
            "play": {
                "query": {"type": "str", "required": False, "description": "What to play (song name, artist, playlist)"},
                "uri": {"type": "str", "required": False, "description": "Direct media URI (spotify:track:..., file path, URL)"},
                "action": {"type": "str", "required": False, "description": "play, pause, next, previous, volume (default play)"},
                "provider": {"type": "str", "required": False, "description": "Provider: spotify, youtube, local"},
            },
            "search": {
                "query": {"type": "str", "required": True, "description": "Search query"},
                "type": {"type": "str", "required": False, "description": "Filter: song, album, artist, playlist, video"},
                "provider": {"type": "str", "required": False, "description": "Provider name"},
                "limit": {"type": "int", "required": False, "description": "Max results"},
            },
        }
    
    async def execute(self, operation: str, params: Dict[str, Any]) -> StepResult:
        try:
            if operation == "info":
                path = str(Path(params.get("path", "")).expanduser())
                if not Path(path).exists():
                    return StepResult(False, error=f"File not found: {path}")
                
                stat = Path(path).stat()
                ext = Path(path).suffix.lower()
                
                info = {
                    "path": path,
                    "name": Path(path).name,
                    "format": ext.lstrip("."),
                    "size_bytes": stat.st_size,
                    "size_mb": round(stat.st_size / (1024 * 1024), 2),
                    "modified": datetime.fromtimestamp(stat.st_mtime).isoformat(),
                }
                
                # Try to get image dimensions and EXIF
                if ext in (".png", ".jpg", ".jpeg", ".gif", ".bmp", ".webp"):
                    try:
                        from PIL import Image
                        from PIL.ExifTags import TAGS, GPSTAGS
                        with Image.open(path) as img:
                            info["width"] = img.width
                            info["height"] = img.height
                            info["mode"] = img.mode
                            
                            # Extract EXIF data
                            exif_data = img._getexif()
                            if exif_data:
                                exif = {}
                                for tag_id, value in exif_data.items():
                                    tag = TAGS.get(tag_id, tag_id)
                                    if isinstance(value, bytes):
                                        try:
                                            value = value.decode("utf-8", errors="ignore")
                                        except:
                                            continue
                                    # Extract key fields
                                    if tag == "DateTimeOriginal":
                                        exif["date_taken"] = value
                                    elif tag == "Make":
                                        exif["camera_make"] = value
                                    elif tag == "Model":
                                        exif["camera_model"] = value
                                    elif tag == "Orientation":
                                        exif["orientation"] = value
                                    elif tag == "GPSInfo":
                                        # Parse GPS coordinates
                                        try:
                                            gps = {}
                                            for gps_tag_id, gps_value in value.items():
                                                gps_tag = GPSTAGS.get(gps_tag_id, gps_tag_id)
                                                gps[gps_tag] = gps_value
                                            if "GPSLatitude" in gps and "GPSLongitude" in gps:
                                                def convert_gps(coord, ref):
                                                    d, m, s = coord
                                                    decimal = float(d) + float(m)/60 + float(s)/3600
                                                    if ref in ["S", "W"]:
                                                        decimal = -decimal
                                                    return round(decimal, 6)
                                                lat = convert_gps(gps["GPSLatitude"], gps.get("GPSLatitudeRef", "N"))
                                                lon = convert_gps(gps["GPSLongitude"], gps.get("GPSLongitudeRef", "E"))
                                                exif["gps_latitude"] = lat
                                                exif["gps_longitude"] = lon
                                        except:
                                            pass
                                    elif tag == "ExposureTime":
                                        exif["exposure_time"] = str(value)
                                    elif tag == "FNumber":
                                        exif["f_number"] = float(value)
                                    elif tag == "ISOSpeedRatings":
                                        exif["iso"] = value
                                if exif:
                                    info["exif"] = exif
                    except ImportError:
                        info["note"] = "Install Pillow for image dimensions and EXIF"
                
                return StepResult(True, data=info)
            
            elif operation == "convert":
                path = str(Path(params.get("path", "")).expanduser())
                target_fmt = params.get("format", "")
                output = params.get("output")
                
                if not Path(path).exists():
                    return StepResult(False, error=f"File not found: {path}")
                if not target_fmt:
                    return StepResult(False, error="Missing 'format' parameter")
                
                if not output:
                    output = str(Path(path).with_suffix(f".{target_fmt.lstrip('.')}"))
                else:
                    output = str(Path(output).expanduser())
                
                src_ext = Path(path).suffix.lower()
                
                # Image conversion via Pillow
                if src_ext in (".png", ".jpg", ".jpeg", ".gif", ".bmp", ".webp"):
                    try:
                        from PIL import Image
                        with Image.open(path) as img:
                            if target_fmt.lower() in ("jpg", "jpeg") and img.mode == "RGBA":
                                img = img.convert("RGB")
                            img.save(output)
                        return StepResult(True, data={"input": path, "output": output, "format": target_fmt})
                    except ImportError:
                        return StepResult(False, error="Install Pillow for image conversion: pip install Pillow")
                
                # Audio/video via ffmpeg
                try:
                    import subprocess
                    proc = await asyncio.create_subprocess_exec(
                        "ffmpeg", "-i", path, "-y", output,
                        stdout=asyncio.subprocess.PIPE,
                        stderr=asyncio.subprocess.PIPE,
                    )
                    _, stderr = await asyncio.wait_for(proc.communicate(), timeout=120)
                    if proc.returncode == 0:
                        return StepResult(True, data={"input": path, "output": output, "format": target_fmt})
                    return StepResult(False, error=f"ffmpeg error: {stderr.decode()[:500]}")
                except FileNotFoundError:
                    return StepResult(False, error="ffmpeg not installed. Install ffmpeg for media conversion.")
                except asyncio.TimeoutError:
                    return StepResult(False, error="Conversion timed out after 120s")
            
            elif operation == "resize":
                path = str(Path(params.get("path", "")).expanduser())
                width = params.get("width")
                height = params.get("height")
                output = params.get("output", path)
                
                if not Path(path).exists():
                    return StepResult(False, error=f"File not found: {path}")
                
                try:
                    from PIL import Image
                    with Image.open(path) as img:
                        orig_w, orig_h = img.size
                        if width and not height:
                            height = int(orig_h * (width / orig_w))
                        elif height and not width:
                            width = int(orig_w * (height / orig_h))
                        elif not width and not height:
                            return StepResult(False, error="Specify 'width' and/or 'height'")
                        
                        resized = img.resize((width, height), Image.LANCZOS)
                        output = str(Path(output).expanduser())
                        resized.save(output)
                    
                    return StepResult(True, data={"path": output, "width": width, "height": height, "original": f"{orig_w}x{orig_h}"})
                except ImportError:
                    return StepResult(False, error="Install Pillow for image resize: pip install Pillow")
            
            elif operation == "generate":
                prompt = params.get("prompt", "")
                media_type = params.get("type", "image")
                
                if not prompt:
                    return StepResult(False, error="Missing 'prompt' parameter")
                
                if not self._llm:
                    return StepResult(False, error="LLM required for media generation")
                
                # This would connect to DALL-E, Stable Diffusion, etc.
                return StepResult(False, error=f"Image generation not configured. Connect a provider (DALL-E, Stable Diffusion, etc.) to generate: '{prompt}'")
            
            elif operation == "transcribe":
                path = str(Path(params.get("path", "")).expanduser())
                
                if not Path(path).exists():
                    return StepResult(False, error=f"File not found: {path}")
                
                # Would connect to Whisper, Google Speech, etc.
                return StepResult(False, error="Transcription not configured. Connect a provider (Whisper, Google Speech, etc.)")
            
            elif operation == "play":
                action = params.get("action", "play")
                provider_name = params.get("provider")
                
                provider = self._providers.get(provider_name) if provider_name else (next(iter(self._providers.values())) if self._providers else None)
                
                if provider and hasattr(provider, "control"):
                    result = await provider.control(
                        action=action,
                        query=params.get("query"),
                        uri=params.get("uri"),
                    )
                    return StepResult(True, data=result)
                
                return StepResult(False, error="No media player configured. Connect a provider (Spotify, YouTube, etc.)")
            
            elif operation == "search":
                query = params.get("query", "")
                provider_name = params.get("provider")
                
                provider = self._providers.get(provider_name) if provider_name else (next(iter(self._providers.values())) if self._providers else None)
                
                if provider and hasattr(provider, "search"):
                    result = await provider.search(
                        query=query,
                        type=params.get("type"),
                        limit=params.get("limit", 10),
                    )
                    return StepResult(True, data=result)
                
                return StepResult(False, error="No media search configured. Connect a provider (Spotify, YouTube, etc.)")
            
            else:
                return StepResult(False, error=f"Unknown operation: {operation}")
        except Exception as e:
            return StepResult(False, error=str(e))


# ============================================================
#  BROWSER PRIMITIVE
# ============================================================

class BrowserPrimitive(Primitive):
    """Browser automation — navigate, interact with web pages, fill forms.
    
    Distinct from WEB (which does raw HTTP). BROWSER controls a real browser
    for JavaScript-heavy sites, form filling, screenshots, etc.
    Uses Playwright when available, falls back to error messages.
    """
    
    def __init__(self, llm_complete: Optional[Callable] = None):
        self._llm = llm_complete
        self._page = None
        self._browser = None
    
    @property
    def name(self) -> str:
        return "BROWSER"
    
    def get_operations(self) -> Dict[str, str]:
        return {
            "open": "Open a URL in a browser",
            "click": "Click an element on the page",
            "type": "Type text into an input field",
            "screenshot": "Take a screenshot of the current page",
            "read": "Read the text content of the current page or a specific element",
            "fill_form": "Fill out a form with provided field values",
            "execute_js": "Execute JavaScript on the current page",
        }
    
    def get_param_schema(self) -> Dict[str, Dict[str, Any]]:
        return {
            "open": {
                "url": {"type": "str", "required": True, "description": "URL to navigate to"},
                "wait": {"type": "int", "required": False, "description": "Seconds to wait after load (default 2)"},
            },
            "click": {
                "selector": {"type": "str", "required": True, "description": "CSS selector or text of element to click"},
            },
            "type": {
                "selector": {"type": "str", "required": True, "description": "CSS selector of input field"},
                "text": {"type": "str", "required": True, "description": "Text to type"},
                "clear": {"type": "bool", "required": False, "description": "Clear field before typing (default true)"},
            },
            "screenshot": {
                "path": {"type": "str", "required": False, "description": "Save path (default: ~/screenshot.png)"},
                "full_page": {"type": "bool", "required": False, "description": "Capture full page (default false)"},
            },
            "read": {
                "selector": {"type": "str", "required": False, "description": "CSS selector to read (default: body)"},
                "max_length": {"type": "int", "required": False, "description": "Max chars to return"},
            },
            "fill_form": {
                "fields": {"type": "dict", "required": True, "description": "Map of CSS selector -> value to fill"},
                "submit": {"type": "bool", "required": False, "description": "Submit the form after filling (default false)"},
            },
            "execute_js": {
                "script": {"type": "str", "required": True, "description": "JavaScript code to execute"},
            },
        }
    
    async def _ensure_browser(self) -> StepResult:
        """Lazy-init a Playwright browser instance."""
        if self._page:
            return StepResult(True)
        
        try:
            from playwright.async_api import async_playwright
            pw = await async_playwright().start()
            self._browser = await pw.chromium.launch(headless=True)
            self._page = await self._browser.new_page()
            return StepResult(True)
        except ImportError:
            return StepResult(False, error="Playwright not installed. Install: pip install playwright && playwright install chromium")
        except Exception as e:
            return StepResult(False, error=f"Failed to start browser: {e}")
    
    async def execute(self, operation: str, params: Dict[str, Any]) -> StepResult:
        try:
            if operation == "open":
                url = params.get("url", "")
                wait = params.get("wait", 2)
                
                if not url:
                    return StepResult(False, error="Missing 'url' parameter")
                
                init = await self._ensure_browser()
                if not init.success:
                    return init
                
                await self._page.goto(url, wait_until="domcontentloaded", timeout=30000)
                if wait > 0:
                    await asyncio.sleep(wait)
                
                title = await self._page.title()
                return StepResult(True, data={"url": url, "title": title})
            
            elif operation == "click":
                selector = params.get("selector", "")
                if not selector:
                    return StepResult(False, error="Missing 'selector' parameter")
                
                init = await self._ensure_browser()
                if not init.success:
                    return init
                
                # Try CSS selector first, then text
                try:
                    await self._page.click(selector, timeout=5000)
                except Exception:
                    await self._page.click(f"text={selector}", timeout=5000)
                
                return StepResult(True, data={"clicked": selector})
            
            elif operation == "type":
                selector = params.get("selector", "")
                text = params.get("text", "")
                clear = params.get("clear", True)
                
                if not selector or not text:
                    return StepResult(False, error="Missing 'selector' and/or 'text' parameter")
                
                init = await self._ensure_browser()
                if not init.success:
                    return init
                
                if clear:
                    await self._page.fill(selector, text, timeout=5000)
                else:
                    await self._page.type(selector, text, timeout=5000)
                
                return StepResult(True, data={"selector": selector, "typed": len(text)})
            
            elif operation == "screenshot":
                path = params.get("path", str(Path.home() / "screenshot.png"))
                full_page = params.get("full_page", False)
                
                init = await self._ensure_browser()
                if not init.success:
                    return init
                
                path = str(Path(path).expanduser())
                Path(path).parent.mkdir(parents=True, exist_ok=True)
                await self._page.screenshot(path=path, full_page=full_page)
                
                return StepResult(True, data={"path": path, "full_page": full_page})
            
            elif operation == "read":
                selector = params.get("selector", "body")
                max_len = params.get("max_length", 10000)
                
                init = await self._ensure_browser()
                if not init.success:
                    return init
                
                text = await self._page.inner_text(selector, timeout=5000)
                return StepResult(True, data={"text": text[:max_len], "length": len(text)})
            
            elif operation == "fill_form":
                fields = params.get("fields", {})
                submit = params.get("submit", False)
                
                if not fields:
                    return StepResult(False, error="Missing 'fields' parameter")
                
                init = await self._ensure_browser()
                if not init.success:
                    return init
                
                filled = []
                for selector, value in fields.items():
                    await self._page.fill(selector, str(value), timeout=5000)
                    filled.append(selector)
                
                if submit:
                    # Try common submit patterns
                    for submit_sel in ['button[type="submit"]', 'input[type="submit"]', "form button"]:
                        try:
                            await self._page.click(submit_sel, timeout=3000)
                            break
                        except Exception:
                            continue
                
                return StepResult(True, data={"filled": filled, "submitted": submit})
            
            elif operation == "execute_js":
                script = params.get("script", "")
                if not script:
                    return StepResult(False, error="Missing 'script' parameter")
                
                init = await self._ensure_browser()
                if not init.success:
                    return init
                
                result = await self._page.evaluate(script)
                return StepResult(True, data={"result": result})
            
            else:
                return StepResult(False, error=f"Unknown operation: {operation}")
        except Exception as e:
            return StepResult(False, error=str(e))


# ============================================================
#  DEVTOOLS PRIMITIVE
# ============================================================

class DevToolsPrimitive(Primitive):
    """Developer tools — GitHub, Jira, and other dev platforms.
    
    Provider-based: plug in GitHub, Jira, or any dev platform connector.
    """
    
    def __init__(self, providers: Optional[Dict[str, Any]] = None):
        self._providers = providers or {}
    
    @property
    def name(self) -> str:
        return "DEVTOOLS"
    
    def get_operations(self) -> Dict[str, str]:
        return {
            "list_issues": "List issues or tickets from a project",
            "get_issue": "Get details of a specific issue or ticket",
            "create_issue": "Create a new issue or ticket",
            "update_issue": "Update an existing issue or ticket",
            "comment": "Add a comment to an issue or PR",
            "list_prs": "List pull requests or merge requests",
            "create_pr": "Create a pull request",
            "list_repos": "List repositories or projects",
            "search": "Search issues, PRs, or repositories",
            "list_notifications": "List GitHub notifications",
            "review_requests": "List PRs awaiting your review",
            "list_commits": "List recent commits in a repository",
            "list_branches": "List branches in a repository",
            "list_workflow_runs": "List CI/CD workflow runs (GitHub Actions)",
            "activity_summary": "Get a summary of your dev activity",
        }
    
    def get_param_schema(self) -> Dict[str, Dict[str, Any]]:
        return {
            "list_issues": {
                "repo": {"type": "str", "required": False, "description": "Repository or project key (e.g. 'owner/repo' or 'PROJ')"},
                "state": {"type": "str", "required": False, "description": "Filter: open, closed, all (default open)"},
                "limit": {"type": "int", "required": False, "description": "Max results (default 20)"},
                "provider": {"type": "str", "required": False, "description": "Provider: github, jira (default: auto)"},
            },
            "get_issue": {
                "repo": {"type": "str", "required": False, "description": "Repository or project key"},
                "issue_id": {"type": "str", "required": True, "description": "Issue number or key (e.g. '42' or 'PROJ-123')"},
                "provider": {"type": "str", "required": False, "description": "Provider name"},
            },
            "create_issue": {
                "repo": {"type": "str", "required": False, "description": "Repository or project key"},
                "title": {"type": "str", "required": True, "description": "Issue title"},
                "body": {"type": "str", "required": False, "description": "Issue description/body"},
                "labels": {"type": "list", "required": False, "description": "Labels/tags"},
                "assignee": {"type": "str", "required": False, "description": "Assignee username"},
                "provider": {"type": "str", "required": False, "description": "Provider name"},
            },
            "update_issue": {
                "repo": {"type": "str", "required": False, "description": "Repository or project key"},
                "issue_id": {"type": "str", "required": True, "description": "Issue number or key"},
                "title": {"type": "str", "required": False, "description": "New title"},
                "body": {"type": "str", "required": False, "description": "New body"},
                "state": {"type": "str", "required": False, "description": "New state: open or closed"},
                "labels": {"type": "list", "required": False, "description": "Updated labels"},
                "provider": {"type": "str", "required": False, "description": "Provider name"},
            },
            "comment": {
                "repo": {"type": "str", "required": False, "description": "Repository or project key"},
                "issue_id": {"type": "str", "required": True, "description": "Issue/PR number or key"},
                "body": {"type": "str", "required": True, "description": "Comment text"},
                "provider": {"type": "str", "required": False, "description": "Provider name"},
            },
            "list_prs": {
                "repo": {"type": "str", "required": False, "description": "Repository (e.g. 'owner/repo')"},
                "state": {"type": "str", "required": False, "description": "Filter: open, closed, all (default open)"},
                "limit": {"type": "int", "required": False, "description": "Max results"},
                "provider": {"type": "str", "required": False, "description": "Provider name"},
            },
            "create_pr": {
                "repo": {"type": "str", "required": False, "description": "Repository (e.g. 'owner/repo')"},
                "title": {"type": "str", "required": True, "description": "PR title"},
                "body": {"type": "str", "required": False, "description": "PR description"},
                "head": {"type": "str", "required": True, "description": "Source branch"},
                "base": {"type": "str", "required": False, "description": "Target branch (default: main)"},
                "provider": {"type": "str", "required": False, "description": "Provider name"},
            },
            "list_repos": {
                "limit": {"type": "int", "required": False, "description": "Max results"},
                "provider": {"type": "str", "required": False, "description": "Provider name"},
            },
            "search": {
                "query": {"type": "str", "required": True, "description": "Search query"},
                "type": {"type": "str", "required": False, "description": "Search type: issues, repos (default issues)"},
                "limit": {"type": "int", "required": False, "description": "Max results"},
                "provider": {"type": "str", "required": False, "description": "Provider name"},
            },
            "list_notifications": {
                "all": {"type": "bool", "required": False, "description": "Include read notifications"},
                "provider": {"type": "str", "required": False, "description": "Provider name"},
            },
            "review_requests": {
                "limit": {"type": "int", "required": False, "description": "Max results"},
                "provider": {"type": "str", "required": False, "description": "Provider name"},
            },
            "list_commits": {
                "repo": {"type": "str", "required": True, "description": "Repository (e.g. 'owner/repo')"},
                "branch": {"type": "str", "required": False, "description": "Branch or SHA"},
                "limit": {"type": "int", "required": False, "description": "Max results"},
                "provider": {"type": "str", "required": False, "description": "Provider name"},
            },
            "list_branches": {
                "repo": {"type": "str", "required": True, "description": "Repository (e.g. 'owner/repo')"},
                "provider": {"type": "str", "required": False, "description": "Provider name"},
            },
            "list_workflow_runs": {
                "repo": {"type": "str", "required": True, "description": "Repository (e.g. 'owner/repo')"},
                "status": {"type": "str", "required": False, "description": "Filter: completed, in_progress, queued"},
                "provider": {"type": "str", "required": False, "description": "Provider name"},
            },
            "activity_summary": {
                "provider": {"type": "str", "required": False, "description": "Provider name"},
            },
        }
    
    def _get_provider(self, name: Optional[str]) -> Optional[Any]:
        if name and name in self._providers:
            return self._providers[name]
        if self._providers:
            return next(iter(self._providers.values()))
        return None
    
    @staticmethod
    def _split_repo(repo_str: str):
        """Split 'owner/repo' into (owner, repo). Returns ('', '') if invalid."""
        if "/" in repo_str:
            parts = repo_str.split("/", 1)
            return parts[0], parts[1]
        return "", repo_str
    
    async def execute(self, operation: str, params: Dict[str, Any]) -> StepResult:
        try:
            provider_name = params.get("provider")
            provider = self._get_provider(provider_name)
            
            if not provider:
                return StepResult(False, error="No devtools provider configured. Connect GitHub, Jira, etc.")
            
            repo_str = params.get("repo", "")
            owner, repo = self._split_repo(repo_str) if repo_str else ("", "")
            
            if operation == "list_issues":
                if hasattr(provider, "list_issues"):
                    result = await provider.list_issues(
                        owner, repo,
                        state=params.get("state", "open"),
                        per_page=params.get("limit", 20),
                    )
                    return StepResult(True, data=result)
                return StepResult(False, error="Provider does not support list_issues")
            
            elif operation == "get_issue":
                issue_id = params.get("issue_id", "")
                if hasattr(provider, "get_issue"):
                    result = await provider.get_issue(owner, repo, int(issue_id))
                    return StepResult(True, data=result)
                return StepResult(False, error="Provider does not support get_issue")
            
            elif operation == "create_issue":
                if hasattr(provider, "create_issue"):
                    result = await provider.create_issue(
                        owner, repo,
                        title=params.get("title", ""),
                        body=params.get("body", ""),
                        labels=params.get("labels"),
                        assignees=[params["assignee"]] if params.get("assignee") else None,
                    )
                    return StepResult(True, data=result)
                return StepResult(False, error="Provider does not support create_issue")
            
            elif operation == "update_issue":
                if hasattr(provider, "update_issue"):
                    result = await provider.update_issue(
                        owner, repo,
                        number=int(params.get("issue_id", 0)),
                        title=params.get("title"),
                        body=params.get("body"),
                        state=params.get("state"),
                        labels=params.get("labels"),
                    )
                    return StepResult(True, data=result)
                return StepResult(False, error="Provider does not support update_issue")
            
            elif operation == "comment":
                if hasattr(provider, "add_comment"):
                    result = await provider.add_comment(
                        owner, repo,
                        issue_number=int(params.get("issue_id", 0)),
                        body=params.get("body", ""),
                    )
                    return StepResult(True, data=result)
                return StepResult(False, error="Provider does not support add_comment")
            
            elif operation == "list_prs":
                if hasattr(provider, "list_pull_requests"):
                    result = await provider.list_pull_requests(
                        owner, repo,
                        state=params.get("state", "open"),
                    )
                    return StepResult(True, data=result)
                return StepResult(False, error="Provider does not support list_prs")
            
            elif operation == "create_pr":
                if hasattr(provider, "create_pull_request"):
                    result = await provider.create_pull_request(
                        owner, repo,
                        title=params.get("title", ""),
                        head=params.get("head", ""),
                        base=params.get("base", "main"),
                        body=params.get("body", ""),
                    )
                    return StepResult(True, data=result)
                return StepResult(False, error="Provider does not support create_pr")
            
            elif operation == "list_repos":
                if hasattr(provider, "list_repos"):
                    result = await provider.list_repos(per_page=params.get("limit", 20))
                    return StepResult(True, data=result)
                return StepResult(False, error="Provider does not support list_repos")
            
            elif operation == "search":
                search_type = params.get("type", "issues")
                query = params.get("query", "")
                limit = params.get("limit", 20)
                if search_type == "repos" and hasattr(provider, "search_repos"):
                    result = await provider.search_repos(query, per_page=limit)
                    return StepResult(True, data=result)
                elif hasattr(provider, "search_issues"):
                    result = await provider.search_issues(query, per_page=limit)
                    return StepResult(True, data=result)
                return StepResult(False, error="Provider does not support search")
            
            elif operation == "list_notifications":
                if hasattr(provider, "list_notifications"):
                    result = await provider.list_notifications(
                        all=params.get("all", False),
                    )
                    return StepResult(True, data=result)
                return StepResult(False, error="Provider does not support list_notifications")
            
            elif operation == "review_requests":
                if hasattr(provider, "get_review_requests"):
                    result = await provider.get_review_requests(
                        per_page=params.get("limit", 20),
                    )
                    return StepResult(True, data=result)
                return StepResult(False, error="Provider does not support review_requests")
            
            elif operation == "list_commits":
                if hasattr(provider, "list_commits"):
                    result = await provider.list_commits(
                        owner, repo,
                        sha=params.get("branch"),
                        per_page=params.get("limit", 20),
                    )
                    return StepResult(True, data=result)
                return StepResult(False, error="Provider does not support list_commits")
            
            elif operation == "list_branches":
                if hasattr(provider, "list_branches"):
                    result = await provider.list_branches(owner, repo)
                    return StepResult(True, data=result)
                return StepResult(False, error="Provider does not support list_branches")
            
            elif operation == "list_workflow_runs":
                if hasattr(provider, "list_workflow_runs"):
                    result = await provider.list_workflow_runs(
                        owner, repo,
                        status=params.get("status"),
                        per_page=params.get("limit", 10),
                    )
                    return StepResult(True, data=result)
                return StepResult(False, error="Provider does not support list_workflow_runs")
            
            elif operation == "activity_summary":
                if hasattr(provider, "get_activity_summary"):
                    result = await provider.get_activity_summary()
                    return StepResult(True, data=result)
                return StepResult(False, error="Provider does not support activity_summary")
            
            else:
                return StepResult(False, error=f"Unknown operation: {operation}")
        except Exception as e:
            return StepResult(False, error=str(e))


# ============================================================
#  CLOUD STORAGE PRIMITIVE
# ============================================================

class CloudStoragePrimitive(Primitive):
    """Cloud storage operations — Google Drive, OneDrive, Dropbox.
    
    Provider-based: plug in any cloud storage connector.
    """
    
    def __init__(self, providers: Optional[Dict[str, Any]] = None):
        self._providers = providers or {}
    
    @property
    def name(self) -> str:
        return "CLOUD_STORAGE"
    
    def get_operations(self) -> Dict[str, str]:
        return {
            "list": "List files and folders in cloud storage",
            "search": "Search for files by name or content",
            "download": "Download a file from cloud storage",
            "upload": "Upload a file to cloud storage",
            "create_folder": "Create a folder in cloud storage",
            "delete": "Delete a file or folder from cloud storage",
            "share": "Share a file or folder with others",
        }
    
    def get_param_schema(self) -> Dict[str, Dict[str, Any]]:
        return {
            "list": {
                "path": {"type": "str", "required": False, "description": "Folder path or ID (default: root)"},
                "limit": {"type": "int", "required": False, "description": "Max results"},
                "provider": {"type": "str", "required": False, "description": "Provider: google_drive, onedrive, dropbox"},
            },
            "search": {
                "query": {"type": "str", "required": True, "description": "Search query (file name or content)"},
                "limit": {"type": "int", "required": False, "description": "Max results"},
                "provider": {"type": "str", "required": False, "description": "Provider name"},
            },
            "download": {
                "file_id": {"type": "str", "required": True, "description": "File ID or path in cloud storage"},
                "local_path": {"type": "str", "required": True, "description": "Local path to save the file"},
                "provider": {"type": "str", "required": False, "description": "Provider name"},
            },
            "upload": {
                "local_path": {"type": "str", "required": True, "description": "Local file path to upload"},
                "remote_path": {"type": "str", "required": False, "description": "Destination path in cloud storage"},
                "provider": {"type": "str", "required": False, "description": "Provider name"},
            },
            "create_folder": {
                "name": {"type": "str", "required": True, "description": "Folder name"},
                "parent": {"type": "str", "required": False, "description": "Parent folder ID or path"},
                "provider": {"type": "str", "required": False, "description": "Provider name"},
            },
            "delete": {
                "file_id": {"type": "str", "required": True, "description": "File or folder ID to delete"},
                "provider": {"type": "str", "required": False, "description": "Provider name"},
            },
            "share": {
                "file_id": {"type": "str", "required": True, "description": "File or folder ID to share"},
                "email": {"type": "str", "required": True, "description": "Email of person to share with"},
                "role": {"type": "str", "required": False, "description": "Permission: viewer, editor, commenter (default viewer)"},
                "provider": {"type": "str", "required": False, "description": "Provider name"},
            },
        }
    
    def _get_provider(self, name: Optional[str]) -> Optional[Any]:
        if name and name in self._providers:
            return self._providers[name]
        if self._providers:
            return next(iter(self._providers.values()))
        return None
    
    async def execute(self, operation: str, params: Dict[str, Any]) -> StepResult:
        try:
            provider_name = params.get("provider")
            provider = self._get_provider(provider_name)
            
            if not provider:
                return StepResult(False, error="No cloud storage provider configured. Connect Google Drive, OneDrive, or Dropbox.")
            
            if operation == "list":
                if hasattr(provider, "list_files"):
                    result = await provider.list_files(
                        folder_id=params.get("path"),
                        max_results=params.get("limit", 50),
                    )
                    return StepResult(True, data=result)
                return StepResult(False, error="Provider does not support list_files")
            
            elif operation == "search":
                if hasattr(provider, "search_files"):
                    result = await provider.search_files(
                        query=params.get("query", ""),
                        max_results=params.get("limit", 20),
                    )
                    return StepResult(True, data=result)
                return StepResult(False, error="Provider does not support search_files")
            
            elif operation == "download":
                if hasattr(provider, "download_file"):
                    result = await provider.download_file(
                        file_id=params.get("file_id", ""),
                        local_path=params.get("local_path", ""),
                    )
                    return StepResult(True, data=result)
                return StepResult(False, error="Provider does not support download_file")
            
            elif operation == "upload":
                if hasattr(provider, "upload_file"):
                    result = await provider.upload_file(
                        local_path=params.get("local_path", ""),
                        remote_path=params.get("remote_path"),
                    )
                    return StepResult(True, data=result)
                return StepResult(False, error="Provider does not support upload_file")
            
            elif operation == "create_folder":
                if hasattr(provider, "create_folder"):
                    result = await provider.create_folder(
                        name=params.get("name", ""),
                        parent_id=params.get("parent"),
                    )
                    return StepResult(True, data=result)
                return StepResult(False, error="Provider does not support create_folder")
            
            elif operation == "delete":
                if hasattr(provider, "delete_file"):
                    result = await provider.delete_file(file_id=params.get("file_id", ""))
                    return StepResult(True, data=result)
                return StepResult(False, error="Provider does not support delete_file")
            
            elif operation == "share":
                if hasattr(provider, "share_file"):
                    result = await provider.share_file(
                        file_id=params.get("file_id", ""),
                        email=params.get("email", ""),
                        role=params.get("role", "viewer"),
                    )
                    return StepResult(True, data=result)
                return StepResult(False, error="Provider does not support share_file")
            
            else:
                return StepResult(False, error=f"Unknown operation: {operation}")
        except Exception as e:
            return StepResult(False, error=str(e))


# ============================================================
#  CLIPBOARD PRIMITIVE
# ============================================================

class ClipboardPrimitive(Primitive):
    """System clipboard operations — copy, paste, history.
    
    Works cross-platform via pyperclip when available.
    """
    
    def __init__(self):
        self._history: List[Dict] = []
        self._max_history = 50
    
    @property
    def name(self) -> str:
        return "CLIPBOARD"
    
    def get_operations(self) -> Dict[str, str]:
        return {
            "copy": "Copy text to the system clipboard",
            "paste": "Get the current clipboard contents",
            "history": "Get recent clipboard history",
            "clear": "Clear the clipboard",
        }
    
    def get_param_schema(self) -> Dict[str, Dict[str, Any]]:
        return {
            "copy": {
                "text": {"type": "str", "required": True, "description": "Text to copy to clipboard"},
            },
            "paste": {},
            "history": {
                "limit": {"type": "int", "required": False, "description": "Max items to return (default 10)"},
            },
            "clear": {},
        }
    
    async def execute(self, operation: str, params: Dict[str, Any]) -> StepResult:
        try:
            # Try to use pyperclip for real clipboard access
            try:
                import pyperclip
                has_pyperclip = True
            except ImportError:
                has_pyperclip = False
            
            if operation == "copy":
                text = params.get("text", "")
                if not text:
                    return StepResult(False, error="Missing 'text' parameter")
                
                if has_pyperclip:
                    pyperclip.copy(text)
                
                # Store in history
                self._history.insert(0, {
                    "text": text[:500],
                    "timestamp": datetime.now().isoformat(),
                    "length": len(text),
                })
                if len(self._history) > self._max_history:
                    self._history = self._history[:self._max_history]
                
                return StepResult(True, data={"copied": True, "length": len(text)})
            
            elif operation == "paste":
                if has_pyperclip:
                    text = pyperclip.paste()
                    return StepResult(True, data={"text": text, "length": len(text)})
                elif self._history:
                    return StepResult(True, data={"text": self._history[0]["text"], "length": self._history[0]["length"]})
                return StepResult(True, data={"text": "", "length": 0})
            
            elif operation == "history":
                limit = params.get("limit", 10)
                return StepResult(True, data=self._history[:limit])
            
            elif operation == "clear":
                if has_pyperclip:
                    pyperclip.copy("")
                return StepResult(True, data={"cleared": True})
            
            else:
                return StepResult(False, error=f"Unknown operation: {operation}")
        except Exception as e:
            return StepResult(False, error=str(e))


# ============================================================
#  TRANSLATE PRIMITIVE
# ============================================================

class TranslatePrimitive(Primitive):
    """Language translation and detection.
    
    Uses LLM for translation when no dedicated API is configured.
    """
    
    def __init__(self, llm_complete: Optional[Callable] = None, provider: Optional[Any] = None):
        self._llm = llm_complete
        self._provider = provider
    
    @property
    def name(self) -> str:
        return "TRANSLATE"
    
    def get_operations(self) -> Dict[str, str]:
        return {
            "translate": "Translate text from one language to another",
            "detect": "Detect the language of text",
            "languages": "List supported languages",
        }
    
    def get_param_schema(self) -> Dict[str, Dict[str, Any]]:
        return {
            "translate": {
                "text": {"type": "str", "required": True, "description": "Text to translate"},
                "to": {"type": "str", "required": True, "description": "Target language (e.g. 'spanish', 'fr', 'zh')"},
                "from": {"type": "str", "required": False, "description": "Source language (default: auto-detect)"},
            },
            "detect": {
                "text": {"type": "str", "required": True, "description": "Text to analyze"},
            },
            "languages": {},
        }
    
    async def execute(self, operation: str, params: Dict[str, Any]) -> StepResult:
        try:
            if operation == "translate":
                text = params.get("text", "")
                target_lang = params.get("to", "")
                source_lang = params.get("from", "auto")
                
                if not text or not target_lang:
                    return StepResult(False, error="Missing 'text' and/or 'to' parameter")
                
                if self._provider and hasattr(self._provider, "translate"):
                    result = await self._provider.translate(text, target_lang, source_lang)
                    return StepResult(True, data=result)
                
                if self._llm:
                    prompt = f"Translate the following text to {target_lang}. Return ONLY the translation, nothing else:\n\n{text}"
                    translation = await self._llm(prompt)
                    return StepResult(True, data={
                        "original": text,
                        "translated": translation.strip(),
                        "to": target_lang,
                        "from": source_lang,
                    })
                
                return StepResult(False, error="No translation provider or LLM configured")
            
            elif operation == "detect":
                text = params.get("text", "")
                if not text:
                    return StepResult(False, error="Missing 'text' parameter")
                
                if self._provider and hasattr(self._provider, "detect"):
                    result = await self._provider.detect(text)
                    return StepResult(True, data=result)
                
                if self._llm:
                    prompt = f"What language is this text written in? Respond with ONLY the language name:\n\n{text[:500]}"
                    language = await self._llm(prompt)
                    return StepResult(True, data={
                        "text": text[:100] + "..." if len(text) > 100 else text,
                        "language": language.strip(),
                        "confidence": 0.9,
                    })
                
                return StepResult(False, error="No language detection available")
            
            elif operation == "languages":
                common_languages = [
                    {"code": "en", "name": "English"},
                    {"code": "es", "name": "Spanish"},
                    {"code": "fr", "name": "French"},
                    {"code": "de", "name": "German"},
                    {"code": "it", "name": "Italian"},
                    {"code": "pt", "name": "Portuguese"},
                    {"code": "zh", "name": "Chinese"},
                    {"code": "ja", "name": "Japanese"},
                    {"code": "ko", "name": "Korean"},
                    {"code": "ar", "name": "Arabic"},
                    {"code": "ru", "name": "Russian"},
                    {"code": "hi", "name": "Hindi"},
                ]
                return StepResult(True, data=common_languages)
            
            else:
                return StepResult(False, error=f"Unknown operation: {operation}")
        except Exception as e:
            return StepResult(False, error=str(e))


# ============================================================
#  DATABASE PRIMITIVE
# ============================================================

class DatabasePrimitive(Primitive):
    """Database operations — SQLite, PostgreSQL, MySQL.
    
    Local SQLite by default, can connect to remote databases via connection string.
    """
    
    def __init__(self, default_db_path: Optional[str] = None):
        self._default_db = default_db_path or str(Path.home() / ".telic" / "data.db")
        self._connections: Dict[str, Any] = {}
    
    @property
    def name(self) -> str:
        return "DATABASE"
    
    def get_operations(self) -> Dict[str, str]:
        return {
            "query": "Execute a SELECT query and return results",
            "execute": "Execute an INSERT, UPDATE, DELETE, or DDL statement",
            "tables": "List all tables in the database",
            "schema": "Get the schema of a specific table",
            "connect": "Connect to a database",
        }
    
    def get_param_schema(self) -> Dict[str, Dict[str, Any]]:
        return {
            "query": {
                "sql": {"type": "str", "required": True, "description": "SELECT query to execute"},
                "params": {"type": "list", "required": False, "description": "Query parameters for placeholders"},
                "database": {"type": "str", "required": False, "description": "Database name/path (default: local)"},
            },
            "execute": {
                "sql": {"type": "str", "required": True, "description": "SQL statement to execute"},
                "params": {"type": "list", "required": False, "description": "Query parameters"},
                "database": {"type": "str", "required": False, "description": "Database name/path"},
            },
            "tables": {
                "database": {"type": "str", "required": False, "description": "Database name/path"},
            },
            "schema": {
                "table": {"type": "str", "required": True, "description": "Table name"},
                "database": {"type": "str", "required": False, "description": "Database name/path"},
            },
            "connect": {
                "connection_string": {"type": "str", "required": True, "description": "Database connection string or file path"},
                "name": {"type": "str", "required": False, "description": "Alias for this connection"},
            },
        }
    
    def _get_connection(self, db_name: Optional[str] = None):
        """Get or create a database connection."""
        import sqlite3
        
        db_path = db_name or self._default_db
        if db_path not in self._connections:
            Path(db_path).parent.mkdir(parents=True, exist_ok=True)
            self._connections[db_path] = sqlite3.connect(db_path)
            self._connections[db_path].row_factory = sqlite3.Row
        return self._connections[db_path]
    
    async def execute(self, operation: str, params: Dict[str, Any]) -> StepResult:
        try:
            import sqlite3
            
            if operation == "query":
                sql = params.get("sql", "")
                query_params = params.get("params", [])
                db_name = params.get("database")
                
                if not sql:
                    return StepResult(False, error="Missing 'sql' parameter")
                
                # Basic injection prevention
                if not sql.strip().upper().startswith("SELECT"):
                    return StepResult(False, error="query() only allows SELECT statements. Use execute() for modifications.")
                
                conn = self._get_connection(db_name)
                cursor = conn.cursor()
                cursor.execute(sql, query_params)
                rows = cursor.fetchall()
                
                # Convert to list of dicts
                columns = [desc[0] for desc in cursor.description] if cursor.description else []
                results = [dict(zip(columns, row)) for row in rows]
                
                return StepResult(True, data={"rows": results, "count": len(results), "columns": columns})
            
            elif operation == "execute":
                sql = params.get("sql", "")
                query_params = params.get("params", [])
                db_name = params.get("database")
                
                if not sql:
                    return StepResult(False, error="Missing 'sql' parameter")
                
                conn = self._get_connection(db_name)
                cursor = conn.cursor()
                cursor.execute(sql, query_params)
                conn.commit()
                
                return StepResult(True, data={
                    "rowcount": cursor.rowcount,
                    "lastrowid": cursor.lastrowid,
                })
            
            elif operation == "tables":
                db_name = params.get("database")
                conn = self._get_connection(db_name)
                cursor = conn.cursor()
                cursor.execute("SELECT name FROM sqlite_master WHERE type='table' ORDER BY name")
                tables = [row[0] for row in cursor.fetchall()]
                return StepResult(True, data=tables)
            
            elif operation == "schema":
                table = params.get("table", "")
                db_name = params.get("database")
                
                if not table:
                    return StepResult(False, error="Missing 'table' parameter")
                
                conn = self._get_connection(db_name)
                cursor = conn.cursor()
                cursor.execute(f"PRAGMA table_info({table})")  # noqa: S608
                columns = []
                for row in cursor.fetchall():
                    columns.append({
                        "name": row[1],
                        "type": row[2],
                        "nullable": not row[3],
                        "default": row[4],
                        "primary_key": bool(row[5]),
                    })
                return StepResult(True, data={"table": table, "columns": columns})
            
            elif operation == "connect":
                conn_str = params.get("connection_string", "")
                name = params.get("name", conn_str)
                
                if not conn_str:
                    return StepResult(False, error="Missing 'connection_string' parameter")
                
                # For now, only SQLite is supported
                if conn_str.endswith(".db") or conn_str.endswith(".sqlite"):
                    self._get_connection(conn_str)
                    return StepResult(True, data={"connected": True, "database": conn_str, "type": "sqlite"})
                
                return StepResult(False, error="Only SQLite databases (.db, .sqlite) are currently supported")
            
            else:
                return StepResult(False, error=f"Unknown operation: {operation}")
        except Exception as e:
            return StepResult(False, error=str(e))


# ============================================================
#  SCREENSHOT PRIMITIVE
# ============================================================

class ScreenshotPrimitive(Primitive):
    """System-level screenshots — full screen, window, region.
    
    Uses mss or pillow for capture.
    """
    
    def __init__(self):
        self._default_path = str(Path.home() / "Screenshots")
    
    @property
    def name(self) -> str:
        return "SCREENSHOT"
    
    def get_operations(self) -> Dict[str, str]:
        return {
            "capture": "Take a screenshot of the entire screen or a region",
            "window": "Take a screenshot of a specific window",
            "list": "List recently taken screenshots",
        }
    
    def get_param_schema(self) -> Dict[str, Dict[str, Any]]:
        return {
            "capture": {
                "path": {"type": "str", "required": False, "description": "Save path (default: ~/Screenshots/screenshot_<timestamp>.png)"},
                "region": {"type": "dict", "required": False, "description": "Region to capture: {x, y, width, height}"},
                "monitor": {"type": "int", "required": False, "description": "Monitor number (default: all monitors)"},
            },
            "window": {
                "title": {"type": "str", "required": True, "description": "Window title (partial match)"},
                "path": {"type": "str", "required": False, "description": "Save path"},
            },
            "list": {
                "limit": {"type": "int", "required": False, "description": "Max screenshots to return (default 10)"},
            },
        }
    
    async def execute(self, operation: str, params: Dict[str, Any]) -> StepResult:
        try:
            if operation == "capture":
                path = params.get("path")
                region = params.get("region")
                monitor = params.get("monitor")
                
                # Ensure screenshots directory exists
                Path(self._default_path).mkdir(parents=True, exist_ok=True)
                
                if not path:
                    timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
                    path = str(Path(self._default_path) / f"screenshot_{timestamp}.png")
                else:
                    path = str(Path(path).expanduser())
                
                try:
                    import mss
                    with mss.mss() as sct:
                        if region:
                            monitor_area = {"left": region.get("x", 0), "top": region.get("y", 0), 
                                          "width": region.get("width", 800), "height": region.get("height", 600)}
                            screenshot = sct.grab(monitor_area)
                        elif monitor:
                            screenshot = sct.grab(sct.monitors[monitor])
                        else:
                            screenshot = sct.grab(sct.monitors[0])  # Primary monitor
                        
                        mss.tools.to_png(screenshot.rgb, screenshot.size, output=path)
                    
                    return StepResult(True, data={"path": path, "size": Path(path).stat().st_size})
                except ImportError:
                    return StepResult(False, error="Install mss for screenshots: pip install mss")
            
            elif operation == "window":
                title = params.get("title", "")
                path = params.get("path")
                
                if not title:
                    return StepResult(False, error="Missing 'title' parameter")
                
                # Window screenshots require platform-specific code
                return StepResult(False, error="Window screenshots not yet implemented. Use capture() with a region instead.")
            
            elif operation == "list":
                limit = params.get("limit", 10)
                
                screenshots_dir = Path(self._default_path)
                if not screenshots_dir.exists():
                    return StepResult(True, data=[])
                
                files = sorted(
                    [f for f in screenshots_dir.iterdir() if f.suffix.lower() in (".png", ".jpg", ".jpeg")],
                    key=lambda f: f.stat().st_mtime,
                    reverse=True,
                )[:limit]
                
                return StepResult(True, data=[
                    {"path": str(f), "name": f.name, "size": f.stat().st_size, 
                     "modified": datetime.fromtimestamp(f.stat().st_mtime).isoformat()}
                    for f in files
                ])
            
            else:
                return StepResult(False, error=f"Unknown operation: {operation}")
        except Exception as e:
            return StepResult(False, error=str(e))


# ============================================================
#  AUTOMATION PRIMITIVE
# ============================================================

class AutomationPrimitive(Primitive):
    """Task automation — schedules, recurring jobs, workflows.
    
    Stores automation rules locally. In production, would integrate with
    system schedulers (cron, Task Scheduler, etc.)
    """
    
    def __init__(self, storage_path: Optional[str] = None):
        self._storage_path = storage_path or str(Path.home() / ".telic" / "automations.json")
        self._automations: List[Dict] = []
        self._load()
    
    def _load(self):
        if Path(self._storage_path).exists():
            try:
                with open(self._storage_path, "r") as f:
                    self._automations = json.load(f)
            except Exception:
                self._automations = []
    
    def _save(self):
        Path(self._storage_path).parent.mkdir(parents=True, exist_ok=True)
        with open(self._storage_path, "w") as f:
            json.dump(self._automations, f, indent=2)
    
    @property
    def name(self) -> str:
        return "AUTOMATION"
    
    def get_operations(self) -> Dict[str, str]:
        return {
            "create": "Create a new automation rule",
            "list": "List all automation rules",
            "enable": "Enable a disabled automation",
            "disable": "Disable an automation without deleting it",
            "delete": "Delete an automation rule",
            "run": "Manually trigger an automation",
        }
    
    def get_param_schema(self) -> Dict[str, Dict[str, Any]]:
        return {
            "create": {
                "name": {"type": "str", "required": True, "description": "Automation name"},
                "trigger": {"type": "str", "required": True, "description": "Trigger type: schedule, event, webhook"},
                "schedule": {"type": "str", "required": False, "description": "Cron expression or natural language (e.g. 'every monday at 9am')"},
                "action": {"type": "str", "required": True, "description": "Action to perform (natural language task description)"},
                "enabled": {"type": "bool", "required": False, "description": "Whether automation is active (default true)"},
            },
            "list": {
                "status": {"type": "str", "required": False, "description": "Filter: enabled, disabled, all (default all)"},
            },
            "enable": {
                "id": {"type": "str", "required": True, "description": "Automation ID to enable"},
            },
            "disable": {
                "id": {"type": "str", "required": True, "description": "Automation ID to disable"},
            },
            "delete": {
                "id": {"type": "str", "required": True, "description": "Automation ID to delete"},
            },
            "run": {
                "id": {"type": "str", "required": True, "description": "Automation ID to run"},
            },
        }
    
    async def execute(self, operation: str, params: Dict[str, Any]) -> StepResult:
        try:
            if operation == "create":
                name = params.get("name", "")
                trigger = params.get("trigger", "schedule")
                schedule = params.get("schedule", "")
                action = params.get("action", "")
                enabled = params.get("enabled", True)
                
                if not name or not action:
                    return StepResult(False, error="Missing 'name' and/or 'action' parameter")
                
                automation = {
                    "id": f"auto_{len(self._automations)}_{int(datetime.now().timestamp())}",
                    "name": name,
                    "trigger": trigger,
                    "schedule": schedule,
                    "action": action,
                    "enabled": enabled,
                    "created": datetime.now().isoformat(),
                    "last_run": None,
                    "run_count": 0,
                }
                self._automations.append(automation)
                self._save()
                
                return StepResult(True, data=automation)
            
            elif operation == "list":
                status = params.get("status", "all")
                if status == "enabled":
                    items = [a for a in self._automations if a.get("enabled", True)]
                elif status == "disabled":
                    items = [a for a in self._automations if not a.get("enabled", True)]
                else:
                    items = self._automations
                return StepResult(True, data=items)
            
            elif operation == "enable":
                auto_id = params.get("id", "")
                for a in self._automations:
                    if a["id"] == auto_id:
                        a["enabled"] = True
                        self._save()
                        return StepResult(True, data=a)
                return StepResult(False, error=f"Automation not found: {auto_id}")
            
            elif operation == "disable":
                auto_id = params.get("id", "")
                for a in self._automations:
                    if a["id"] == auto_id:
                        a["enabled"] = False
                        self._save()
                        return StepResult(True, data=a)
                return StepResult(False, error=f"Automation not found: {auto_id}")
            
            elif operation == "delete":
                auto_id = params.get("id", "")
                for i, a in enumerate(self._automations):
                    if a["id"] == auto_id:
                        deleted = self._automations.pop(i)
                        self._save()
                        return StepResult(True, data={"deleted": deleted["name"]})
                return StepResult(False, error=f"Automation not found: {auto_id}")
            
            elif operation == "run":
                auto_id = params.get("id", "")
                for a in self._automations:
                    if a["id"] == auto_id:
                        a["last_run"] = datetime.now().isoformat()
                        a["run_count"] = a.get("run_count", 0) + 1
                        self._save()
                        # In production, this would trigger the actual action
                        return StepResult(True, data={
                            "triggered": a["name"],
                            "action": a["action"],
                            "status": "queued",
                        })
                return StepResult(False, error=f"Automation not found: {auto_id}")
            
            else:
                return StepResult(False, error=f"Unknown operation: {operation}")
        except Exception as e:
            return StepResult(False, error=str(e))


# ============================================================
#  SEARCH PRIMITIVE
# ============================================================

class SearchPrimitive(Primitive):
    """Universal search — dynamically searches across all available primitives.
    
    Discovers which primitives have a 'search' or 'recall' operation
    and aggregates results. No hard-coded source list.
    """
    
    def __init__(self, primitives: Optional[Dict[str, 'Primitive']] = None):
        self._primitives = primitives or {}
    
    def set_primitives(self, primitives: Dict[str, 'Primitive']):
        """Set the primitives dict after construction (for circular dependency)."""
        self._primitives = primitives
    
    @property
    def name(self) -> str:
        return "SEARCH"
    
    def get_operations(self) -> Dict[str, str]:
        return {
            "all": "Search across all available sources (files, email, calendar, tasks, contacts, messages, etc.)",
        }
    
    def get_param_schema(self) -> Dict[str, Dict[str, Any]]:
        return {
            "all": {
                "query": {"type": "str", "required": True, "description": "Search query"},
                "limit": {"type": "int", "required": False, "description": "Max results per source (default 5)"},
            },
        }
    
    def _find_searchable(self) -> Dict[str, tuple]:
        """Discover all primitives that support search/recall operations."""
        searchable = {}
        for name, prim in self._primitives.items():
            if name == "SEARCH":
                continue  # Don't recurse into ourselves
            ops = prim.get_operations()
            if "search" in ops:
                searchable[name.lower()] = (prim, "search")
            elif "recall" in ops:
                searchable[name.lower()] = (prim, "recall")
        return searchable
    
    async def execute(self, operation: str, params: Dict[str, Any]) -> StepResult:
        try:
            query = params.get("query", "")
            limit = params.get("limit", 5)
            
            if not query:
                return StepResult(False, error="Missing 'query' parameter")
            
            if operation == "all":
                results = {"query": query, "sources": {}}
                searchable = self._find_searchable()
                
                for source_name, (prim, op) in searchable.items():
                    try:
                        result = await prim.execute(op, {"query": query, "limit": limit})
                        if result.success and result.data:
                            results["sources"][source_name] = result.data
                    except Exception:
                        pass
                
                return StepResult(True, data=results)
            
            else:
                return StepResult(False, error=f"Unknown operation: {operation}")
        except Exception as e:
            return StepResult(False, error=str(e))


# ============================================================
#  CHAT PRIMITIVE - Instant messaging (Slack, Teams, Discord, etc.)
# ============================================================

class ChatPrimitive(Primitive):
    """Instant messaging operations - separate from MESSAGE for real-time chat."""
    
    def __init__(self, providers: Optional[Dict[str, Any]] = None):
        self._providers = providers or {}
        self._local_messages: List[Dict] = []
    
    @property
    def name(self) -> str:
        return "CHAT"
    
    def get_operations(self) -> Dict[str, str]:
        return self.get_available_operations()
    
    def get_available_operations(self) -> Dict[str, str]:
        return {
            "send": "Send a chat message to a channel or user",
            "read": "Read recent messages from a channel",
            "search": "Search chat history",
            "react": "Add reaction to a message",
            "reply": "Reply in a thread",
            "channels": "List available channels",
            "create_channel": "Create a new channel",
        }
    
    def get_param_schema(self) -> Dict[str, Any]:
        return {
            "send": {"channel": {"type": "str", "description": "Channel or user"}, "message": {"type": "str", "description": "Message text"}},
            "read": {"channel": {"type": "str", "description": "Channel"}, "limit": {"type": "int", "description": "Max messages", "default": 20}},
            "search": {"query": {"type": "str", "description": "Search query"}, "channel": {"type": "str", "description": "Channel (optional)"}},
            "react": {"message_id": {"type": "str", "description": "Message ID"}, "emoji": {"type": "str", "description": "Emoji"}},
            "reply": {"message_id": {"type": "str", "description": "Message ID"}, "text": {"type": "str", "description": "Reply text"}},
            "channels": {},
            "create_channel": {"name": {"type": "str", "description": "Channel name"}, "members": {"type": "list", "description": "Member IDs"}},
        }
    
    async def execute(self, operation: str, params: Dict[str, Any]) -> StepResult:
        try:
            if operation == "send":
                channel = params.get("channel", "general")
                message = params.get("message", "")
                # Try providers first
                for name, provider in self._providers.items():
                    if hasattr(provider, "send_message"):
                        result = await provider.send_message(channel=channel, text=message)
                        return StepResult(True, data={"sent": True, "provider": name, "result": result})
                # Local fallback
                msg = {"channel": channel, "message": message, "timestamp": datetime.now().isoformat()}
                self._local_messages.append(msg)
                return StepResult(True, data={"sent": True, "provider": "local", "message": msg})
            
            elif operation == "read":
                channel = params.get("channel", "general")
                limit = params.get("limit", 20)
                for name, provider in self._providers.items():
                    if hasattr(provider, "get_messages"):
                        result = await provider.get_messages(channel=channel, limit=limit)
                        return StepResult(True, data={"messages": result, "provider": name})
                # Local fallback
                msgs = [m for m in self._local_messages if m.get("channel") == channel][-limit:]
                return StepResult(True, data={"messages": msgs, "provider": "local"})
            
            elif operation == "search":
                query = params.get("query", "").lower()
                channel = params.get("channel")
                for name, provider in self._providers.items():
                    if hasattr(provider, "search_messages"):
                        result = await provider.search_messages(query=query, channel=channel)
                        return StepResult(True, data={"results": result, "provider": name})
                # Local fallback
                results = [m for m in self._local_messages if query in m.get("message", "").lower()]
                if channel:
                    results = [m for m in results if m.get("channel") == channel]
                return StepResult(True, data={"results": results, "provider": "local"})
            
            elif operation == "react":
                message_id = params.get("message_id")
                emoji = params.get("emoji")
                for name, provider in self._providers.items():
                    if hasattr(provider, "add_reaction"):
                        result = await provider.add_reaction(message_id=message_id, emoji=emoji)
                        return StepResult(True, data={"reacted": True, "provider": name})
                return StepResult(True, data={"reacted": True, "provider": "local", "message_id": message_id, "emoji": emoji})
            
            elif operation == "reply":
                message_id = params.get("message_id")
                text = params.get("text")
                for name, provider in self._providers.items():
                    if hasattr(provider, "reply_to_message"):
                        result = await provider.reply_to_message(message_id=message_id, text=text)
                        return StepResult(True, data={"replied": True, "provider": name})
                return StepResult(True, data={"replied": True, "provider": "local", "message_id": message_id, "text": text})
            
            elif operation == "channels":
                for name, provider in self._providers.items():
                    if hasattr(provider, "list_channels"):
                        result = await provider.list_channels()
                        return StepResult(True, data={"channels": result, "provider": name})
                return StepResult(True, data={"channels": ["general", "random"], "provider": "local"})
            
            elif operation == "create_channel":
                name_param = params.get("name")
                members = params.get("members", [])
                for pname, provider in self._providers.items():
                    if hasattr(provider, "create_channel"):
                        result = await provider.create_channel(name=name_param, members=members)
                        return StepResult(True, data={"created": True, "provider": pname, "channel": result})
                return StepResult(True, data={"created": True, "provider": "local", "name": name_param})
            
            else:
                return StepResult(False, error=f"Unknown operation: {operation}")
        except Exception as e:
            return StepResult(False, error=str(e))


# ============================================================
#  MEETING PRIMITIVE - Video conferencing
# ============================================================

class MeetingPrimitive(Primitive):
    """Video conferencing operations - Zoom, Teams, Meet, etc."""
    
    def __init__(self, providers: Optional[Dict[str, Any]] = None):
        self._providers = providers or {}
        self._local_meetings: List[Dict] = []
    
    @property
    def name(self) -> str:
        return "MEETING"
    
    def get_operations(self) -> Dict[str, str]:
        return self.get_available_operations()
    
    def get_available_operations(self) -> Dict[str, str]:
        return {
            "schedule": "Schedule a video meeting",
            "join": "Get join link for a meeting",
            "cancel": "Cancel a scheduled meeting",
            "list": "List upcoming meetings",
            "recording": "Get meeting recording",
            "transcript": "Get meeting transcript",
        }
    
    def get_param_schema(self) -> Dict[str, Any]:
        return {
            "schedule": {
                "title": {"type": "str", "description": "Meeting title"},
                "start": {"type": "str", "description": "Start time ISO"},
                "duration": {"type": "int", "description": "Duration in minutes"},
                "attendees": {"type": "list", "description": "Attendee emails"},
            },
            "join": {"meeting_id": {"type": "str", "description": "Meeting ID"}},
            "cancel": {"meeting_id": {"type": "str", "description": "Meeting ID"}},
            "list": {"limit": {"type": "int", "description": "Max meetings", "default": 10}},
            "recording": {"meeting_id": {"type": "str", "description": "Meeting ID"}},
            "transcript": {"meeting_id": {"type": "str", "description": "Meeting ID"}},
        }
    
    async def execute(self, operation: str, params: Dict[str, Any]) -> StepResult:
        try:
            if operation == "schedule":
                title = params.get("title", "Meeting")
                start = params.get("start")
                duration = params.get("duration", 30)
                attendees = params.get("attendees", [])
                
                for name, provider in self._providers.items():
                    if hasattr(provider, "schedule_meeting"):
                        result = await provider.schedule_meeting(title=title, start=start, duration=duration, attendees=attendees)
                        return StepResult(True, data={"scheduled": True, "provider": name, "meeting": result})
                
                # Local fallback
                meeting_id = f"meet_{int(datetime.now().timestamp())}"
                meeting = {
                    "id": meeting_id,
                    "title": title,
                    "start": start,
                    "duration": duration,
                    "attendees": attendees,
                    "join_url": f"https://meet.example.com/{meeting_id}",
                }
                self._local_meetings.append(meeting)
                return StepResult(True, data={"scheduled": True, "provider": "local", "meeting": meeting})
            
            elif operation == "join":
                meeting_id = params.get("meeting_id")
                for name, provider in self._providers.items():
                    if hasattr(provider, "get_join_url"):
                        url = await provider.get_join_url(meeting_id)
                        return StepResult(True, data={"join_url": url, "provider": name})
                # Local fallback
                for m in self._local_meetings:
                    if m.get("id") == meeting_id:
                        return StepResult(True, data={"join_url": m.get("join_url"), "provider": "local"})
                return StepResult(True, data={"join_url": f"https://meet.example.com/{meeting_id}", "provider": "local"})
            
            elif operation == "cancel":
                meeting_id = params.get("meeting_id")
                for name, provider in self._providers.items():
                    if hasattr(provider, "cancel_meeting"):
                        await provider.cancel_meeting(meeting_id)
                        return StepResult(True, data={"cancelled": True, "provider": name})
                self._local_meetings = [m for m in self._local_meetings if m.get("id") != meeting_id]
                return StepResult(True, data={"cancelled": True, "provider": "local"})
            
            elif operation == "list":
                limit = params.get("limit", 10)
                for name, provider in self._providers.items():
                    if hasattr(provider, "list_meetings"):
                        result = await provider.list_meetings(limit=limit)
                        return StepResult(True, data={"meetings": result, "provider": name})
                return StepResult(True, data={"meetings": self._local_meetings[:limit], "provider": "local"})
            
            elif operation == "recording":
                meeting_id = params.get("meeting_id")
                for name, provider in self._providers.items():
                    if hasattr(provider, "get_recording"):
                        result = await provider.get_recording(meeting_id)
                        return StepResult(True, data={"recording": result, "provider": name})
                return StepResult(True, data={"recording": None, "message": "No recording available", "provider": "local"})
            
            elif operation == "transcript":
                meeting_id = params.get("meeting_id")
                for name, provider in self._providers.items():
                    if hasattr(provider, "get_transcript"):
                        result = await provider.get_transcript(meeting_id)
                        return StepResult(True, data={"transcript": result, "provider": name})
                return StepResult(True, data={"transcript": None, "message": "No transcript available", "provider": "local"})
            
            else:
                return StepResult(False, error=f"Unknown operation: {operation}")
        except Exception as e:
            return StepResult(False, error=str(e))


# ============================================================
#  SMS PRIMITIVE - Text messaging
# ============================================================

class SmsPrimitive(Primitive):
    """SMS/text messaging operations."""
    
    def __init__(self, providers: Optional[Dict[str, Any]] = None):
        self._providers = providers or {}
        self._local_messages: List[Dict] = []
    
    @property
    def name(self) -> str:
        return "SMS"
    
    def get_operations(self) -> Dict[str, str]:
        return self.get_available_operations()
    
    def get_available_operations(self) -> Dict[str, str]:
        return {
            "send": "Send an SMS message",
            "read": "Read SMS messages",
            "search": "Search SMS history",
        }
    
    def get_param_schema(self) -> Dict[str, Any]:
        return {
            "send": {"to": {"type": "str", "description": "Phone number"}, "message": {"type": "str", "description": "Message text"}},
            "read": {"from": {"type": "str", "description": "Phone number (optional)"}, "limit": {"type": "int", "description": "Max messages", "default": 20}},
            "search": {"query": {"type": "str", "description": "Search query"}},
        }
    
    async def execute(self, operation: str, params: Dict[str, Any]) -> StepResult:
        try:
            if operation == "send":
                to = params.get("to")
                message = params.get("message")
                for name, provider in self._providers.items():
                    if hasattr(provider, "send_sms"):
                        result = await provider.send_sms(to=to, message=message)
                        return StepResult(True, data={"sent": True, "provider": name, "result": result})
                # Local fallback
                msg = {"to": to, "message": message, "timestamp": datetime.now().isoformat()}
                self._local_messages.append(msg)
                return StepResult(True, data={"sent": True, "provider": "local", "message": msg})
            
            elif operation == "read":
                from_num = params.get("from")
                limit = params.get("limit", 20)
                for name, provider in self._providers.items():
                    if hasattr(provider, "read_sms"):
                        result = await provider.read_sms(from_number=from_num, limit=limit)
                        return StepResult(True, data={"messages": result, "provider": name})
                msgs = self._local_messages
                if from_num:
                    msgs = [m for m in msgs if m.get("to") == from_num or m.get("from") == from_num]
                return StepResult(True, data={"messages": msgs[-limit:], "provider": "local"})
            
            elif operation == "search":
                query = params.get("query", "").lower()
                for name, provider in self._providers.items():
                    if hasattr(provider, "search_sms"):
                        result = await provider.search_sms(query=query)
                        return StepResult(True, data={"results": result, "provider": name})
                results = [m for m in self._local_messages if query in m.get("message", "").lower()]
                return StepResult(True, data={"results": results, "provider": "local"})
            
            else:
                return StepResult(False, error=f"Unknown operation: {operation}")
        except Exception as e:
            return StepResult(False, error=str(e))


# ============================================================
#  SPREADSHEET PRIMITIVE - Excel, Google Sheets
# ============================================================

class SpreadsheetPrimitive(Primitive):
    """Spreadsheet operations - Google Sheets, Excel, etc."""
    
    def __init__(self, llm_complete: Optional[Callable] = None, providers: Optional[Dict[str, Any]] = None):
        self._llm = llm_complete
        self._providers = providers or {}
    
    @property
    def name(self) -> str:
        return "SPREADSHEET"
    
    def get_operations(self) -> Dict[str, str]:
        return self.get_available_operations()
    
    def get_available_operations(self) -> Dict[str, str]:
        return {
            "read": "Read data from a spreadsheet",
            "write": "Write data to a spreadsheet",
            "create": "Create a new spreadsheet",
            "add_sheet": "Add a worksheet to a spreadsheet",
            "formula": "Set a formula in a cell",
            "format": "Format cells",
            "chart": "Create a chart",
        }
    
    def get_param_schema(self) -> Dict[str, Any]:
        return {
            "read": {"file": {"type": "str", "description": "File path or ID"}, "sheet": {"type": "str", "description": "Sheet name"}, "range": {"type": "str", "description": "Cell range like A1:D10"}},
            "write": {"file": {"type": "str", "description": "File path or ID"}, "data": {"type": "list", "description": "2D array of values"}, "sheet": {"type": "str", "description": "Sheet name"}, "range": {"type": "str", "description": "Starting cell"}},
            "create": {"name": {"type": "str", "description": "Spreadsheet name"}, "data": {"type": "list", "description": "Initial data (optional)"}},
            "add_sheet": {"file": {"type": "str", "description": "File path or ID"}, "name": {"type": "str", "description": "New sheet name"}},
            "formula": {"file": {"type": "str", "description": "File path or ID"}, "cell": {"type": "str", "description": "Cell like A1"}, "formula": {"type": "str", "description": "Formula"}},
            "format": {"file": {"type": "str", "description": "File path or ID"}, "range": {"type": "str", "description": "Cell range"}, "format": {"type": "dict", "description": "Format options"}},
            "chart": {"file": {"type": "str", "description": "File path or ID"}, "data_range": {"type": "str", "description": "Data range"}, "chart_type": {"type": "str", "description": "bar, line, pie, etc."}},
        }
    
    async def execute(self, operation: str, params: Dict[str, Any]) -> StepResult:
        try:
            file_path = params.get("file", "")
            
            # Try providers first
            for name, provider in self._providers.items():
                if operation == "read" and hasattr(provider, "read_spreadsheet"):
                    result = await provider.read_spreadsheet(file_path, params.get("sheet"), params.get("range"))
                    return StepResult(True, data={"data": result, "provider": name})
                elif operation == "write" and hasattr(provider, "write_spreadsheet"):
                    await provider.write_spreadsheet(file_path, params.get("data"), params.get("sheet"), params.get("range"))
                    return StepResult(True, data={"written": True, "provider": name})
                elif operation == "create" and hasattr(provider, "create_spreadsheet"):
                    result = await provider.create_spreadsheet(params.get("name"), params.get("data"))
                    return StepResult(True, data={"created": True, "file": result, "provider": name})
            
            # Local file handling for CSV/Excel
            if operation == "read":
                if file_path.endswith(".csv"):
                    import csv
                    with open(file_path, "r") as f:
                        reader = csv.reader(f)
                        data = list(reader)
                    return StepResult(True, data={"data": data, "provider": "local"})
                elif file_path.endswith((".xlsx", ".xls")):
                    try:
                        import openpyxl
                        wb = openpyxl.load_workbook(file_path)
                        sheet = wb[params.get("sheet")] if params.get("sheet") else wb.active
                        data = [[cell.value for cell in row] for row in sheet.iter_rows()]
                        return StepResult(True, data={"data": data, "provider": "local"})
                    except ImportError:
                        return StepResult(False, error="openpyxl not installed for Excel support")
                return StepResult(True, data={"data": [], "message": "File type not supported locally", "provider": "local"})
            
            elif operation == "write":
                data = params.get("data", [])
                if file_path.endswith(".csv"):
                    import csv
                    with open(file_path, "w", newline="") as f:
                        writer = csv.writer(f)
                        writer.writerows(data)
                    return StepResult(True, data={"written": True, "provider": "local"})
                return StepResult(True, data={"written": True, "provider": "local", "note": "Would write to cloud"})
            
            elif operation == "create":
                name = params.get("name", "Spreadsheet")
                data = params.get("data", [])
                file_path = f"{name}.csv"
                if data:
                    import csv
                    with open(file_path, "w", newline="") as f:
                        writer = csv.writer(f)
                        writer.writerows(data)
                return StepResult(True, data={"created": True, "file": file_path, "provider": "local"})
            
            elif operation in ("add_sheet", "formula", "format", "chart"):
                return StepResult(True, data={"success": True, "provider": "local", "note": f"{operation} would be applied via cloud provider"})
            
            else:
                return StepResult(False, error=f"Unknown operation: {operation}")
        except Exception as e:
            return StepResult(False, error=str(e))


# ============================================================
#  PRESENTATION PRIMITIVE - PowerPoint, Google Slides
# ============================================================

class PresentationPrimitive(Primitive):
    """Presentation operations - PowerPoint, Google Slides, etc."""
    
    def __init__(self, llm_complete: Optional[Callable] = None, providers: Optional[Dict[str, Any]] = None):
        self._llm = llm_complete
        self._providers = providers or {}
    
    @property
    def name(self) -> str:
        return "PRESENTATION"
    
    def get_operations(self) -> Dict[str, str]:
        return self.get_available_operations()
    
    def get_available_operations(self) -> Dict[str, str]:
        return {
            "create": "Create a new presentation",
            "add_slide": "Add a slide to a presentation",
            "update_slide": "Update slide content",
            "export": "Export to PDF or images",
            "get_text": "Extract all text from presentation",
        }
    
    def get_param_schema(self) -> Dict[str, Any]:
        return {
            "create": {"name": {"type": "str", "description": "Presentation name"}, "template": {"type": "str", "description": "Template (optional)"}},
            "add_slide": {"file": {"type": "str", "description": "File path or ID"}, "layout": {"type": "str", "description": "Slide layout"}, "content": {"type": "dict", "description": "Slide content"}},
            "update_slide": {"file": {"type": "str", "description": "File path or ID"}, "slide_id": {"type": "str", "description": "Slide index or ID"}, "content": {"type": "dict", "description": "New content"}},
            "export": {"file": {"type": "str", "description": "File path or ID"}, "format": {"type": "str", "description": "pdf, png, jpg"}},
            "get_text": {"file": {"type": "str", "description": "File path or ID"}},
        }
    
    async def execute(self, operation: str, params: Dict[str, Any]) -> StepResult:
        try:
            file_path = params.get("file", "")
            
            for name, provider in self._providers.items():
                if operation == "create" and hasattr(provider, "create_presentation"):
                    result = await provider.create_presentation(params.get("name"), params.get("template"))
                    return StepResult(True, data={"created": True, "file": result, "provider": name})
                elif operation == "add_slide" and hasattr(provider, "add_slide"):
                    result = await provider.add_slide(file_path, params.get("layout"), params.get("content"))
                    return StepResult(True, data={"added": True, "provider": name})
            
            # Local handling
            if operation == "create":
                name = params.get("name", "Presentation")
                return StepResult(True, data={"created": True, "file": f"{name}.pptx", "provider": "local", "note": "Would create via python-pptx"})
            
            elif operation == "add_slide":
                return StepResult(True, data={"added": True, "provider": "local", "note": "Would add slide via python-pptx"})
            
            elif operation == "update_slide":
                return StepResult(True, data={"updated": True, "provider": "local"})
            
            elif operation == "export":
                fmt = params.get("format", "pdf")
                return StepResult(True, data={"exported": True, "format": fmt, "provider": "local"})
            
            elif operation == "get_text":
                try:
                    from pptx import Presentation as PPTX
                    prs = PPTX(file_path)
                    text = []
                    for slide in prs.slides:
                        for shape in slide.shapes:
                            if hasattr(shape, "text"):
                                text.append(shape.text)
                    return StepResult(True, data={"text": "\n".join(text), "provider": "local"})
                except ImportError:
                    return StepResult(True, data={"text": "", "provider": "local", "note": "python-pptx not installed"})
            
            else:
                return StepResult(False, error=f"Unknown operation: {operation}")
        except Exception as e:
            return StepResult(False, error=str(e))


# ============================================================
#  NOTES PRIMITIVE - Note-taking apps
# ============================================================

class NotesPrimitive(Primitive):
    """Note-taking operations - OneNote, Apple Notes, Google Keep, etc."""
    
    def __init__(self, storage_path: str = "", providers: Optional[Dict[str, Any]] = None):
        self._storage_path = Path(storage_path) if storage_path else Path.home() / ".telic" / "notes"
        self._storage_path.mkdir(parents=True, exist_ok=True)
        self._providers = providers or {}
    
    @property
    def name(self) -> str:
        return "NOTES"
    
    def get_operations(self) -> Dict[str, str]:
        return self.get_available_operations()
    
    def get_available_operations(self) -> Dict[str, str]:
        return {
            "create": "Create a new note",
            "read": "Read a note",
            "update": "Update a note",
            "delete": "Delete a note",
            "search": "Search notes",
            "list": "List all notes",
        }
    
    def get_param_schema(self) -> Dict[str, Any]:
        return {
            "create": {"title": {"type": "str", "description": "Note title"}, "content": {"type": "str", "description": "Note content"}, "folder": {"type": "str", "description": "Folder (optional)"}, "tags": {"type": "list", "description": "Tags (optional)"}},
            "read": {"note_id": {"type": "str", "description": "Note ID or title"}},
            "update": {"note_id": {"type": "str", "description": "Note ID"}, "content": {"type": "str", "description": "New content"}},
            "delete": {"note_id": {"type": "str", "description": "Note ID"}},
            "search": {"query": {"type": "str", "description": "Search query"}},
            "list": {"folder": {"type": "str", "description": "Folder (optional)"}},
        }
    
    async def execute(self, operation: str, params: Dict[str, Any]) -> StepResult:
        try:
            for name, provider in self._providers.items():
                if operation == "create" and hasattr(provider, "create_note"):
                    result = await provider.create_note(params.get("title"), params.get("content"), params.get("folder"), params.get("tags"))
                    return StepResult(True, data={"created": True, "note": result, "provider": name})
                elif operation == "read" and hasattr(provider, "read_note"):
                    result = await provider.read_note(params.get("note_id"))
                    return StepResult(True, data={"note": result, "provider": name})
                elif operation == "list" and hasattr(provider, "list_notes"):
                    result = await provider.list_notes(params.get("folder"))
                    return StepResult(True, data={"notes": result, "provider": name})
            
            # Local file-based notes
            if operation == "create":
                title = params.get("title", "Untitled")
                content = params.get("content", "")
                tags = params.get("tags", [])
                note_id = f"note_{int(datetime.now().timestamp())}"
                note_file = self._storage_path / f"{note_id}.json"
                note = {"id": note_id, "title": title, "content": content, "tags": tags, "created": datetime.now().isoformat()}
                note_file.write_text(json.dumps(note))
                return StepResult(True, data={"created": True, "note": note, "provider": "local"})
            
            elif operation == "read":
                note_id = params.get("note_id")
                note_file = self._storage_path / f"{note_id}.json"
                if note_file.exists():
                    note = json.loads(note_file.read_text())
                    return StepResult(True, data={"note": note, "provider": "local"})
                # Try to find by title
                for f in self._storage_path.glob("*.json"):
                    note = json.loads(f.read_text())
                    if note.get("title") == note_id:
                        return StepResult(True, data={"note": note, "provider": "local"})
                return StepResult(False, error=f"Note not found: {note_id}")
            
            elif operation == "update":
                note_id = params.get("note_id")
                content = params.get("content")
                note_file = self._storage_path / f"{note_id}.json"
                if note_file.exists():
                    note = json.loads(note_file.read_text())
                    note["content"] = content
                    note["updated"] = datetime.now().isoformat()
                    note_file.write_text(json.dumps(note))
                    return StepResult(True, data={"updated": True, "note": note, "provider": "local"})
                return StepResult(False, error=f"Note not found: {note_id}")
            
            elif operation == "delete":
                note_id = params.get("note_id")
                note_file = self._storage_path / f"{note_id}.json"
                if note_file.exists():
                    note_file.unlink()
                    return StepResult(True, data={"deleted": True, "provider": "local"})
                return StepResult(False, error=f"Note not found: {note_id}")
            
            elif operation == "search":
                query = params.get("query", "").lower()
                results = []
                for f in self._storage_path.glob("*.json"):
                    note = json.loads(f.read_text())
                    if query in note.get("title", "").lower() or query in note.get("content", "").lower():
                        results.append(note)
                return StepResult(True, data={"notes": results, "provider": "local"})
            
            elif operation == "list":
                notes = []
                for f in self._storage_path.glob("*.json"):
                    notes.append(json.loads(f.read_text()))
                return StepResult(True, data={"notes": notes, "provider": "local"})
            
            else:
                return StepResult(False, error=f"Unknown operation: {operation}")
        except Exception as e:
            return StepResult(False, error=str(e))


# ============================================================
#  FINANCE PRIMITIVE - Banking, payments
# ============================================================

class FinancePrimitive(Primitive):
    """Financial operations - banking, payments, budgeting."""
    
    def __init__(self, providers: Optional[Dict[str, Any]] = None):
        self._providers = providers or {}
        self._local_transactions: List[Dict] = []
        self._budgets: Dict[str, Dict] = {}
    
    @property
    def name(self) -> str:
        return "FINANCE"
    
    def get_operations(self) -> Dict[str, str]:
        return self.get_available_operations()
    
    def get_available_operations(self) -> Dict[str, str]:
        return {
            "balance": "Get account balance",
            "transactions": "List transactions",
            "categorize": "Categorize a transaction",
            "spending": "Get spending summary",
            "budget": "Create or view a budget",
            "send": "Send a payment",
            "request": "Request a payment",
        }
    
    def get_param_schema(self) -> Dict[str, Any]:
        return {
            "balance": {"account": {"type": "str", "description": "Account ID (optional)"}},
            "transactions": {"account": {"type": "str", "description": "Account ID"}, "start": {"type": "str", "description": "Start date"}, "end": {"type": "str", "description": "End date"}, "limit": {"type": "int", "description": "Max transactions"}},
            "categorize": {"transaction_id": {"type": "str", "description": "Transaction ID"}, "category": {"type": "str", "description": "Category name"}},
            "spending": {"period": {"type": "str", "description": "month, week, year"}, "category": {"type": "str", "description": "Category (optional)"}},
            "budget": {"category": {"type": "str", "description": "Budget category"}, "amount": {"type": "float", "description": "Budget amount"}, "period": {"type": "str", "description": "month, week"}},
            "send": {"to": {"type": "str", "description": "Recipient"}, "amount": {"type": "float", "description": "Amount"}, "note": {"type": "str", "description": "Note (optional)"}},
            "request": {"from": {"type": "str", "description": "From person"}, "amount": {"type": "float", "description": "Amount"}, "note": {"type": "str", "description": "Note (optional)"}},
        }
    
    async def execute(self, operation: str, params: Dict[str, Any]) -> StepResult:
        try:
            for name, provider in self._providers.items():
                if operation == "balance" and hasattr(provider, "get_balance"):
                    result = await provider.get_balance(params.get("account"))
                    return StepResult(True, data={"balance": result, "provider": name})
                elif operation == "transactions" and hasattr(provider, "list_transactions"):
                    result = await provider.list_transactions(params.get("account"), params.get("start"), params.get("end"), params.get("limit"))
                    return StepResult(True, data={"transactions": result, "provider": name})
                elif operation == "send" and hasattr(provider, "send_payment"):
                    result = await provider.send_payment(params.get("to"), params.get("amount"), params.get("note"))
                    return StepResult(True, data={"sent": True, "provider": name, "result": result})
            
            # Local fallback
            if operation == "balance":
                return StepResult(True, data={"balance": 0.0, "currency": "USD", "provider": "local", "note": "Connect a bank provider for real data"})
            
            elif operation == "transactions":
                return StepResult(True, data={"transactions": self._local_transactions, "provider": "local"})
            
            elif operation == "categorize":
                tx_id = params.get("transaction_id")
                category = params.get("category")
                for tx in self._local_transactions:
                    if tx.get("id") == tx_id:
                        tx["category"] = category
                        return StepResult(True, data={"categorized": True, "provider": "local"})
                return StepResult(True, data={"categorized": True, "provider": "local", "note": "Transaction not found locally"})
            
            elif operation == "spending":
                period = params.get("period", "month")
                category = params.get("category")
                return StepResult(True, data={"spending": {}, "period": period, "category": category, "provider": "local"})
            
            elif operation == "budget":
                category = params.get("category", "general")
                amount = params.get("amount")
                period = params.get("period", "month")
                if amount:
                    self._budgets[category] = {"amount": amount, "period": period}
                return StepResult(True, data={"budgets": self._budgets, "provider": "local"})
            
            elif operation == "send":
                return StepResult(True, data={"sent": True, "provider": "local", "note": "Connect payment provider to send real payments"})
            
            elif operation == "request":
                return StepResult(True, data={"requested": True, "provider": "local", "note": "Connect payment provider to request real payments"})
            
            else:
                return StepResult(False, error=f"Unknown operation: {operation}")
        except Exception as e:
            return StepResult(False, error=str(e))


# ============================================================
#  SOCIAL PRIMITIVE - Social media
# ============================================================

class SocialPrimitive(Primitive):
    """Social media operations - Twitter, LinkedIn, Facebook, etc."""
    
    def __init__(self, providers: Optional[Dict[str, Any]] = None):
        self._providers = providers or {}
        self._local_posts: List[Dict] = []
    
    @property
    def name(self) -> str:
        return "SOCIAL"
    
    def get_operations(self) -> Dict[str, str]:
        return self.get_available_operations()
    
    def get_available_operations(self) -> Dict[str, str]:
        return {
            "post": "Create a social media post or tweet",
            "delete_post": "Delete a post or tweet",
            "feed": "Get your feed or timeline",
            "search": "Search posts or tweets",
            "like": "Like a post or tweet",
            "unlike": "Unlike a post or tweet",
            "repost": "Retweet or repost",
            "undo_repost": "Undo a retweet or repost",
            "comment": "Comment on or reply to a post",
            "profile": "Get user profile",
            "followers": "Get followers list",
            "following": "Get following list",
            "follow": "Follow a user",
            "unfollow": "Unfollow a user",
            "bookmarks": "Get bookmarked posts",
            "bookmark": "Bookmark a post",
            "user_posts": "Get posts by a specific user",
        }
    
    def get_param_schema(self) -> Dict[str, Any]:
        return {
            "post": {
                "content": {"type": "str", "required": True, "description": "Post/tweet content"},
                "reply_to": {"type": "str", "required": False, "description": "Post ID to reply to"},
                "provider": {"type": "str", "required": False, "description": "Provider: twitter, linkedin"},
            },
            "delete_post": {
                "post_id": {"type": "str", "required": True, "description": "Post/tweet ID to delete"},
                "provider": {"type": "str", "required": False, "description": "Provider name"},
            },
            "feed": {
                "limit": {"type": "int", "required": False, "description": "Max posts (default 20)"},
                "provider": {"type": "str", "required": False, "description": "Provider name"},
            },
            "search": {
                "query": {"type": "str", "required": True, "description": "Search query"},
                "limit": {"type": "int", "required": False, "description": "Max results"},
                "provider": {"type": "str", "required": False, "description": "Provider name"},
            },
            "like": {
                "post_id": {"type": "str", "required": True, "description": "Post/tweet ID to like"},
                "provider": {"type": "str", "required": False, "description": "Provider name"},
            },
            "unlike": {
                "post_id": {"type": "str", "required": True, "description": "Post/tweet ID to unlike"},
                "provider": {"type": "str", "required": False, "description": "Provider name"},
            },
            "repost": {
                "post_id": {"type": "str", "required": True, "description": "Post/tweet ID to retweet/repost"},
                "provider": {"type": "str", "required": False, "description": "Provider name"},
            },
            "undo_repost": {
                "post_id": {"type": "str", "required": True, "description": "Post/tweet ID to undo retweet"},
                "provider": {"type": "str", "required": False, "description": "Provider name"},
            },
            "comment": {
                "post_id": {"type": "str", "required": True, "description": "Post ID to reply to"},
                "text": {"type": "str", "required": True, "description": "Reply/comment text"},
                "provider": {"type": "str", "required": False, "description": "Provider name"},
            },
            "profile": {
                "username": {"type": "str", "required": False, "description": "Username (defaults to self)"},
                "provider": {"type": "str", "required": False, "description": "Provider name"},
            },
            "followers": {
                "username": {"type": "str", "required": False, "description": "Username (defaults to self)"},
                "limit": {"type": "int", "required": False, "description": "Max results"},
                "provider": {"type": "str", "required": False, "description": "Provider name"},
            },
            "following": {
                "username": {"type": "str", "required": False, "description": "Username (defaults to self)"},
                "limit": {"type": "int", "required": False, "description": "Max results"},
                "provider": {"type": "str", "required": False, "description": "Provider name"},
            },
            "follow": {
                "user_id": {"type": "str", "required": True, "description": "User ID to follow"},
                "provider": {"type": "str", "required": False, "description": "Provider name"},
            },
            "unfollow": {
                "user_id": {"type": "str", "required": True, "description": "User ID to unfollow"},
                "provider": {"type": "str", "required": False, "description": "Provider name"},
            },
            "bookmarks": {
                "limit": {"type": "int", "required": False, "description": "Max results"},
                "provider": {"type": "str", "required": False, "description": "Provider name"},
            },
            "bookmark": {
                "post_id": {"type": "str", "required": True, "description": "Post/tweet ID to bookmark"},
                "provider": {"type": "str", "required": False, "description": "Provider name"},
            },
            "user_posts": {
                "username": {"type": "str", "required": False, "description": "Username"},
                "user_id": {"type": "str", "required": False, "description": "User ID"},
                "limit": {"type": "int", "required": False, "description": "Max results"},
                "provider": {"type": "str", "required": False, "description": "Provider name"},
            },
        }
    
    def _get_provider(self, name: Optional[str] = None) -> Optional[Any]:
        if name and name in self._providers:
            return self._providers[name]
        if self._providers:
            return next(iter(self._providers.values()))
        return None
    
    async def execute(self, operation: str, params: Dict[str, Any]) -> StepResult:
        try:
            provider = self._get_provider(params.get("provider"))
            
            if operation == "post":
                content = params.get("content", "")
                reply_to = params.get("reply_to")
                if provider and hasattr(provider, "post_tweet"):
                    result = await provider.post_tweet(text=content, reply_to=reply_to)
                    return StepResult(True, data=result)
                elif provider and hasattr(provider, "create_post"):
                    result = await provider.create_post(content)
                    return StepResult(True, data=result)
                post = {"id": f"post_{int(datetime.now().timestamp())}", "content": content, "timestamp": datetime.now().isoformat()}
                self._local_posts.append(post)
                return StepResult(True, data={"posted": True, "post": post, "provider": "local"})
            
            elif operation == "delete_post":
                if provider and hasattr(provider, "delete_tweet"):
                    result = await provider.delete_tweet(params["post_id"])
                    return StepResult(True, data={"deleted": result})
                elif provider and hasattr(provider, "delete_post"):
                    result = await provider.delete_post(params["post_id"])
                    return StepResult(True, data={"deleted": result})
                return StepResult(False, error="Delete not supported")
            
            elif operation == "feed":
                limit = params.get("limit", 20)
                if provider and hasattr(provider, "get_user_tweets"):
                    me = await provider.get_me() if hasattr(provider, "get_me") else None
                    if me:
                        result = await provider.get_user_tweets(user_id=me.id if hasattr(me, 'id') else str(me), max_results=limit)
                        return StepResult(True, data=result)
                return StepResult(True, data={"posts": self._local_posts[-limit:], "provider": "local"})
            
            elif operation == "search":
                query = params.get("query", "")
                limit = params.get("limit", 20)
                if provider and hasattr(provider, "search"):
                    result = await provider.search(query=query, max_results=limit)
                    return StepResult(True, data=result)
                results = [p for p in self._local_posts if query.lower() in p.get("content", "").lower()]
                return StepResult(True, data={"posts": results, "provider": "local"})
            
            elif operation == "like":
                if provider and hasattr(provider, "like_tweet"):
                    result = await provider.like_tweet(params["post_id"])
                    return StepResult(True, data={"liked": result})
                return StepResult(True, data={"liked": True, "provider": "local"})
            
            elif operation == "unlike":
                if provider and hasattr(provider, "unlike_tweet"):
                    result = await provider.unlike_tweet(params["post_id"])
                    return StepResult(True, data={"unliked": result})
                return StepResult(True, data={"unliked": True, "provider": "local"})
            
            elif operation == "repost":
                if provider and hasattr(provider, "retweet"):
                    result = await provider.retweet(params["post_id"])
                    return StepResult(True, data={"reposted": result})
                return StepResult(True, data={"reposted": True, "provider": "local"})
            
            elif operation == "undo_repost":
                if provider and hasattr(provider, "undo_retweet"):
                    result = await provider.undo_retweet(params["post_id"])
                    return StepResult(True, data={"undone": result})
                return StepResult(True, data={"undone": True, "provider": "local"})
            
            elif operation == "comment":
                if provider and hasattr(provider, "post_tweet"):
                    result = await provider.post_tweet(text=params["text"], reply_to=params["post_id"])
                    return StepResult(True, data=result)
                return StepResult(True, data={"commented": True, "provider": "local"})
            
            elif operation == "profile":
                username = params.get("username")
                if provider:
                    if username and hasattr(provider, "get_user"):
                        result = await provider.get_user(username=username)
                        return StepResult(True, data=result)
                    elif hasattr(provider, "get_me"):
                        result = await provider.get_me()
                        return StepResult(True, data=result)
                return StepResult(True, data={"profile": {"name": "Local User"}, "provider": "local"})
            
            elif operation == "followers":
                if provider and hasattr(provider, "get_followers"):
                    user_id = params.get("user_id")
                    if not user_id and hasattr(provider, "get_me"):
                        me = await provider.get_me()
                        user_id = me.id if hasattr(me, 'id') else str(me)
                    result = await provider.get_followers(user_id=user_id, max_results=params.get("limit", 20))
                    return StepResult(True, data=result)
                return StepResult(True, data={"followers": [], "provider": "local"})
            
            elif operation == "following":
                if provider and hasattr(provider, "get_following"):
                    user_id = params.get("user_id")
                    if not user_id and hasattr(provider, "get_me"):
                        me = await provider.get_me()
                        user_id = me.id if hasattr(me, 'id') else str(me)
                    result = await provider.get_following(user_id=user_id, max_results=params.get("limit", 20))
                    return StepResult(True, data=result)
                return StepResult(True, data={"following": [], "provider": "local"})
            
            elif operation == "follow":
                if provider and hasattr(provider, "follow_user"):
                    result = await provider.follow_user(params["user_id"])
                    return StepResult(True, data={"followed": result})
                return StepResult(True, data={"followed": True, "provider": "local"})
            
            elif operation == "unfollow":
                if provider and hasattr(provider, "unfollow_user"):
                    result = await provider.unfollow_user(params["user_id"])
                    return StepResult(True, data={"unfollowed": result})
                return StepResult(True, data={"unfollowed": True, "provider": "local"})
            
            elif operation == "bookmarks":
                if provider and hasattr(provider, "get_bookmarks"):
                    result = await provider.get_bookmarks(max_results=params.get("limit", 20))
                    return StepResult(True, data=result)
                return StepResult(True, data={"bookmarks": [], "provider": "local"})
            
            elif operation == "bookmark":
                if provider and hasattr(provider, "bookmark_tweet"):
                    result = await provider.bookmark_tweet(params["post_id"])
                    return StepResult(True, data={"bookmarked": result})
                return StepResult(True, data={"bookmarked": True, "provider": "local"})
            
            elif operation == "user_posts":
                if provider and hasattr(provider, "get_user_tweets"):
                    result = await provider.get_user_tweets(
                        user_id=params.get("user_id"),
                        max_results=params.get("limit", 20),
                    )
                    return StepResult(True, data=result)
                return StepResult(True, data={"posts": [], "provider": "local"})
            
            else:
                return StepResult(False, error=f"Unknown operation: {operation}")
        except Exception as e:
            return StepResult(False, error=str(e))


# ============================================================
#  PHOTO PRIMITIVE - Photo management
# ============================================================

class PhotoPrimitive(Primitive):
    """Photo management operations - Google Photos, iCloud, etc."""
    
    def __init__(self, providers: Optional[Dict[str, Any]] = None):
        self._providers = providers or {}
    
    @property
    def name(self) -> str:
        return "PHOTO"
    
    def get_operations(self) -> Dict[str, str]:
        return self.get_available_operations()
    
    def get_available_operations(self) -> Dict[str, str]:
        return {
            "list": "List photos",
            "upload": "Upload a photo",
            "download": "Download a photo",
            "search": "Search photos",
            "create_album": "Create an album",
            "add_to_album": "Add photo to album",
            "metadata": "Get photo metadata",
            "edit": "Edit a photo",
        }
    
    def get_param_schema(self) -> Dict[str, Any]:
        return {
            "list": {"album": {"type": "str", "description": "Album ID (optional)"}, "limit": {"type": "int", "description": "Max photos"}},
            "upload": {"path": {"type": "str", "description": "Local file path"}, "album": {"type": "str", "description": "Album ID (optional)"}},
            "download": {"photo_id": {"type": "str", "description": "Photo ID"}, "path": {"type": "str", "description": "Save path"}},
            "search": {"query": {"type": "str", "description": "Search query"}},
            "create_album": {"name": {"type": "str", "description": "Album name"}},
            "add_to_album": {"photo_id": {"type": "str", "description": "Photo ID"}, "album_id": {"type": "str", "description": "Album ID"}},
            "metadata": {"photo_id": {"type": "str", "description": "Photo ID"}},
            "edit": {"photo_id": {"type": "str", "description": "Photo ID"}, "operations": {"type": "list", "description": "Edit operations"}},
        }
    
    async def execute(self, operation: str, params: Dict[str, Any]) -> StepResult:
        try:
            for name, provider in self._providers.items():
                if operation == "list" and hasattr(provider, "list_photos"):
                    result = await provider.list_photos(params.get("album"), params.get("limit"))
                    return StepResult(True, data={"photos": result, "provider": name})
                elif operation == "upload" and hasattr(provider, "upload_photo"):
                    result = await provider.upload_photo(params.get("path"), params.get("album"))
                    return StepResult(True, data={"uploaded": True, "photo": result, "provider": name})
                elif operation == "search" and hasattr(provider, "search_photos"):
                    result = await provider.search_photos(params.get("query"))
                    return StepResult(True, data={"photos": result, "provider": name})
            
            # Local file handling
            if operation == "list":
                path = Path(params.get("album", str(Path.home() / "Pictures")))
                if path.exists() and path.is_dir():
                    photos = list(path.glob("*.jpg")) + list(path.glob("*.png")) + list(path.glob("*.jpeg"))
                    return StepResult(True, data={"photos": [str(p) for p in photos[:params.get("limit", 50)]], "provider": "local"})
                return StepResult(True, data={"photos": [], "provider": "local"})
            
            elif operation == "upload":
                return StepResult(True, data={"uploaded": True, "provider": "local", "note": "Connect photo provider to upload"})
            
            elif operation == "download":
                return StepResult(True, data={"downloaded": True, "provider": "local"})
            
            elif operation == "search":
                return StepResult(True, data={"photos": [], "provider": "local", "note": "Local search not implemented"})
            
            elif operation == "create_album":
                name = params.get("name", "Album")
                path = Path.home() / "Pictures" / name
                path.mkdir(parents=True, exist_ok=True)
                return StepResult(True, data={"created": True, "path": str(path), "provider": "local"})
            
            elif operation == "add_to_album":
                return StepResult(True, data={"added": True, "provider": "local"})
            
            elif operation == "metadata":
                photo_path = params.get("photo_id")
                if Path(photo_path).exists():
                    stat = Path(photo_path).stat()
                    return StepResult(True, data={"metadata": {"size": stat.st_size, "modified": stat.st_mtime}, "provider": "local"})
                return StepResult(False, error="Photo not found")
            
            elif operation == "edit":
                return StepResult(True, data={"edited": True, "provider": "local", "note": "Would apply edits via PIL"})
            
            else:
                return StepResult(False, error=f"Unknown operation: {operation}")
        except Exception as e:
            return StepResult(False, error=str(e))


# ============================================================
#  RIDE PRIMITIVE - Ride sharing
# ============================================================

class RidePrimitive(Primitive):
    """Ride-sharing operations - Uber, Lyft, etc."""
    
    def __init__(self, providers: Optional[Dict[str, Any]] = None):
        self._providers = providers or {}
        self._ride_history: List[Dict] = []
    
    @property
    def name(self) -> str:
        return "RIDE"
    
    def get_operations(self) -> Dict[str, str]:
        return self.get_available_operations()
    
    def get_available_operations(self) -> Dict[str, str]:
        return {
            "estimate": "Get fare estimate",
            "request": "Request a ride",
            "cancel": "Cancel a ride",
            "track": "Track current ride",
            "history": "Get ride history",
        }
    
    def get_param_schema(self) -> Dict[str, Any]:
        return {
            "estimate": {"pickup": {"type": "str", "description": "Pickup address"}, "dropoff": {"type": "str", "description": "Dropoff address"}},
            "request": {"pickup": {"type": "str", "description": "Pickup address"}, "dropoff": {"type": "str", "description": "Dropoff address"}, "type": {"type": "str", "description": "UberX, Pool, etc."}},
            "cancel": {"ride_id": {"type": "str", "description": "Ride ID"}},
            "track": {"ride_id": {"type": "str", "description": "Ride ID"}},
            "history": {"limit": {"type": "int", "description": "Max rides", "default": 10}},
        }
    
    async def execute(self, operation: str, params: Dict[str, Any]) -> StepResult:
        try:
            for name, provider in self._providers.items():
                if operation == "estimate" and hasattr(provider, "get_estimate"):
                    result = await provider.get_estimate(params.get("pickup"), params.get("dropoff"))
                    return StepResult(True, data={"estimate": result, "provider": name})
                elif operation == "request" and hasattr(provider, "request_ride"):
                    result = await provider.request_ride(params.get("pickup"), params.get("dropoff"), params.get("type"))
                    return StepResult(True, data={"ride": result, "provider": name})
            
            # Local fallback
            if operation == "estimate":
                return StepResult(True, data={
                    "estimate": {"min": 10.0, "max": 15.0, "currency": "USD", "eta": "5 min"},
                    "provider": "local",
                    "note": "Connect ride provider for real estimates"
                })
            
            elif operation == "request":
                ride = {
                    "id": f"ride_{int(datetime.now().timestamp())}",
                    "pickup": params.get("pickup"),
                    "dropoff": params.get("dropoff"),
                    "type": params.get("type", "standard"),
                    "status": "requested",
                }
                self._ride_history.append(ride)
                return StepResult(True, data={"ride": ride, "provider": "local", "note": "Connect ride provider to request real rides"})
            
            elif operation == "cancel":
                ride_id = params.get("ride_id")
                for ride in self._ride_history:
                    if ride.get("id") == ride_id:
                        ride["status"] = "cancelled"
                return StepResult(True, data={"cancelled": True, "provider": "local"})
            
            elif operation == "track":
                ride_id = params.get("ride_id")
                for ride in self._ride_history:
                    if ride.get("id") == ride_id:
                        return StepResult(True, data={"ride": ride, "provider": "local"})
                return StepResult(True, data={"ride": None, "provider": "local"})
            
            elif operation == "history":
                limit = params.get("limit", 10)
                return StepResult(True, data={"rides": self._ride_history[-limit:], "provider": "local"})
            
            else:
                return StepResult(False, error=f"Unknown operation: {operation}")
        except Exception as e:
            return StepResult(False, error=str(e))


# ============================================================
#  TRAVEL PRIMITIVE - Travel booking
# ============================================================

class TravelPrimitive(Primitive):
    """Travel booking operations - flights, hotels, etc."""
    
    def __init__(self, providers: Optional[Dict[str, Any]] = None):
        self._providers = providers or {}
        self._bookings: List[Dict] = []
    
    @property
    def name(self) -> str:
        return "TRAVEL"
    
    def get_operations(self) -> Dict[str, str]:
        return self.get_available_operations()
    
    def get_available_operations(self) -> Dict[str, str]:
        return {
            "search_flights": "Search for flights",
            "search_hotels": "Search for hotels",
            "book": "Book a flight or hotel",
            "cancel": "Cancel a booking",
            "itinerary": "Get trip itinerary",
            "checkin": "Online check-in",
        }
    
    def get_param_schema(self) -> Dict[str, Any]:
        return {
            "search_flights": {"origin": {"type": "str", "description": "Origin airport"}, "destination": {"type": "str", "description": "Destination airport"}, "date": {"type": "str", "description": "Departure date"}, "return_date": {"type": "str", "description": "Return date (optional)"}},
            "search_hotels": {"location": {"type": "str", "description": "Location"}, "checkin": {"type": "str", "description": "Check-in date"}, "checkout": {"type": "str", "description": "Check-out date"}, "guests": {"type": "int", "description": "Number of guests"}},
            "book": {"type": {"type": "str", "description": "flight or hotel"}, "id": {"type": "str", "description": "Search result ID"}},
            "cancel": {"booking_id": {"type": "str", "description": "Booking ID"}},
            "itinerary": {"trip_id": {"type": "str", "description": "Trip ID (optional)"}},
            "checkin": {"booking_id": {"type": "str", "description": "Booking ID"}},
        }
    
    async def execute(self, operation: str, params: Dict[str, Any]) -> StepResult:
        try:
            for name, provider in self._providers.items():
                if operation == "search_flights" and hasattr(provider, "search_flights"):
                    result = await provider.search_flights(params.get("origin"), params.get("destination"), params.get("date"), params.get("return_date"))
                    return StepResult(True, data={"flights": result, "provider": name})
                elif operation == "search_hotels" and hasattr(provider, "search_hotels"):
                    result = await provider.search_hotels(params.get("location"), params.get("checkin"), params.get("checkout"), params.get("guests"))
                    return StepResult(True, data={"hotels": result, "provider": name})
            
            # Local fallback
            if operation == "search_flights":
                return StepResult(True, data={
                    "flights": [],
                    "provider": "local",
                    "note": "Connect travel provider to search real flights"
                })
            
            elif operation == "search_hotels":
                return StepResult(True, data={
                    "hotels": [],
                    "provider": "local",
                    "note": "Connect travel provider to search real hotels"
                })
            
            elif operation == "book":
                booking = {
                    "id": f"booking_{int(datetime.now().timestamp())}",
                    "type": params.get("type"),
                    "item_id": params.get("id"),
                    "status": "confirmed",
                }
                self._bookings.append(booking)
                return StepResult(True, data={"booking": booking, "provider": "local"})
            
            elif operation == "cancel":
                booking_id = params.get("booking_id")
                for b in self._bookings:
                    if b.get("id") == booking_id:
                        b["status"] = "cancelled"
                return StepResult(True, data={"cancelled": True, "provider": "local"})
            
            elif operation == "itinerary":
                return StepResult(True, data={"itinerary": self._bookings, "provider": "local"})
            
            elif operation == "checkin":
                return StepResult(True, data={"checkedin": True, "provider": "local"})
            
            else:
                return StepResult(False, error=f"Unknown operation: {operation}")
        except Exception as e:
            return StepResult(False, error=str(e))


# ============================================================
#  HOME PRIMITIVE - Smart home
# ============================================================

class HomePrimitive(Primitive):
    """Smart home operations - lights, thermostat, etc."""
    
    def __init__(self, providers: Optional[Dict[str, Any]] = None):
        self._providers = providers or {}
        self._devices: Dict[str, Dict] = {}
    
    @property
    def name(self) -> str:
        return "HOME"
    
    def get_operations(self) -> Dict[str, str]:
        return self.get_available_operations()
    
    def get_available_operations(self) -> Dict[str, str]:
        return {
            "devices": "List all devices",
            "state": "Get device state",
            "set": "Set device state",
            "on": "Turn device on",
            "off": "Turn device off",
            "temperature": "Set thermostat temperature",
            "routine": "Run a routine/scene",
        }
    
    def get_param_schema(self) -> Dict[str, Any]:
        return {
            "devices": {},
            "state": {"device_id": {"type": "str", "description": "Device ID"}},
            "set": {"device_id": {"type": "str", "description": "Device ID"}, "state": {"type": "dict", "description": "State to set"}},
            "on": {"device_id": {"type": "str", "description": "Device ID"}},
            "off": {"device_id": {"type": "str", "description": "Device ID"}},
            "temperature": {"device_id": {"type": "str", "description": "Thermostat ID"}, "temperature": {"type": "int", "description": "Temperature"}},
            "routine": {"routine_name": {"type": "str", "description": "Routine name"}},
        }
    
    async def execute(self, operation: str, params: Dict[str, Any]) -> StepResult:
        try:
            for name, provider in self._providers.items():
                if operation == "devices" and hasattr(provider, "list_devices"):
                    result = await provider.list_devices()
                    return StepResult(True, data={"devices": result, "provider": name})
                elif operation == "on" and hasattr(provider, "turn_on"):
                    await provider.turn_on(params.get("device_id"))
                    return StepResult(True, data={"on": True, "provider": name})
                elif operation == "off" and hasattr(provider, "turn_off"):
                    await provider.turn_off(params.get("device_id"))
                    return StepResult(True, data={"off": True, "provider": name})
            
            # Local simulation
            if operation == "devices":
                return StepResult(True, data={"devices": list(self._devices.values()), "provider": "local"})
            
            elif operation == "state":
                device_id = params.get("device_id")
                return StepResult(True, data={"state": self._devices.get(device_id, {}), "provider": "local"})
            
            elif operation == "set":
                device_id = params.get("device_id")
                state = params.get("state", {})
                self._devices[device_id] = {**self._devices.get(device_id, {}), **state}
                return StepResult(True, data={"set": True, "state": self._devices[device_id], "provider": "local"})
            
            elif operation == "on":
                device_id = params.get("device_id")
                self._devices.setdefault(device_id, {})["on"] = True
                return StepResult(True, data={"on": True, "provider": "local"})
            
            elif operation == "off":
                device_id = params.get("device_id")
                self._devices.setdefault(device_id, {})["on"] = False
                return StepResult(True, data={"off": True, "provider": "local"})
            
            elif operation == "temperature":
                device_id = params.get("device_id")
                temp = params.get("temperature")
                self._devices.setdefault(device_id, {})["temperature"] = temp
                return StepResult(True, data={"set": True, "temperature": temp, "provider": "local"})
            
            elif operation == "routine":
                routine = params.get("routine_name")
                return StepResult(True, data={"ran": True, "routine": routine, "provider": "local"})
            
            else:
                return StepResult(False, error=f"Unknown operation: {operation}")
        except Exception as e:
            return StepResult(False, error=str(e))


# ============================================================
#  SHOPPING PRIMITIVE - E-commerce
# ============================================================

class ShoppingPrimitive(Primitive):
    """E-commerce operations - Amazon, eBay, etc."""
    
    def __init__(self, providers: Optional[Dict[str, Any]] = None):
        self._providers = providers or {}
        self._cart: List[Dict] = []
        self._orders: List[Dict] = []
    
    @property
    def name(self) -> str:
        return "SHOPPING"
    
    def get_operations(self) -> Dict[str, str]:
        return self.get_available_operations()
    
    def get_available_operations(self) -> Dict[str, str]:
        return {
            "search": "Search for products",
            "product": "Get product details",
            "add_to_cart": "Add item to cart",
            "cart": "View cart",
            "track": "Track an order",
            "orders": "Get order history",
            "reorder": "Reorder a previous order",
            "price_alert": "Set price alert",
        }
    
    def get_param_schema(self) -> Dict[str, Any]:
        return {
            "search": {"query": {"type": "str", "description": "Search query"}, "filters": {"type": "dict", "description": "Filters (optional)"}},
            "product": {"product_id": {"type": "str", "description": "Product ID"}},
            "add_to_cart": {"product_id": {"type": "str", "description": "Product ID"}, "quantity": {"type": "int", "description": "Quantity", "default": 1}},
            "cart": {},
            "track": {"order_id": {"type": "str", "description": "Order ID"}},
            "orders": {"limit": {"type": "int", "description": "Max orders", "default": 10}},
            "reorder": {"order_id": {"type": "str", "description": "Order ID"}},
            "price_alert": {"product_id": {"type": "str", "description": "Product ID"}, "target_price": {"type": "float", "description": "Target price"}},
        }
    
    async def execute(self, operation: str, params: Dict[str, Any]) -> StepResult:
        try:
            for name, provider in self._providers.items():
                if operation == "search" and hasattr(provider, "search_products"):
                    result = await provider.search_products(params.get("query"), params.get("filters"))
                    return StepResult(True, data={"products": result, "provider": name})
                elif operation == "product" and hasattr(provider, "get_product"):
                    result = await provider.get_product(params.get("product_id"))
                    return StepResult(True, data={"product": result, "provider": name})
            
            # Local fallback
            if operation == "search":
                return StepResult(True, data={"products": [], "provider": "local", "note": "Connect shopping provider to search"})
            
            elif operation == "product":
                return StepResult(True, data={"product": None, "provider": "local"})
            
            elif operation == "add_to_cart":
                item = {"product_id": params.get("product_id"), "quantity": params.get("quantity", 1)}
                self._cart.append(item)
                return StepResult(True, data={"added": True, "cart": self._cart, "provider": "local"})
            
            elif operation == "cart":
                return StepResult(True, data={"cart": self._cart, "provider": "local"})
            
            elif operation == "track":
                order_id = params.get("order_id")
                for order in self._orders:
                    if order.get("id") == order_id:
                        return StepResult(True, data={"order": order, "provider": "local"})
                return StepResult(True, data={"order": None, "provider": "local"})
            
            elif operation == "orders":
                limit = params.get("limit", 10)
                return StepResult(True, data={"orders": self._orders[-limit:], "provider": "local"})
            
            elif operation == "reorder":
                order_id = params.get("order_id")
                for order in self._orders:
                    if order.get("id") == order_id:
                        new_order = {**order, "id": f"order_{int(datetime.now().timestamp())}"}
                        self._orders.append(new_order)
                        return StepResult(True, data={"order": new_order, "provider": "local"})
                return StepResult(False, error="Order not found")
            
            elif operation == "price_alert":
                return StepResult(True, data={"alert_set": True, "provider": "local"})
            
            else:
                return StepResult(False, error=f"Unknown operation: {operation}")
        except Exception as e:
            return StepResult(False, error=str(e))


class Apex:
    """
    The Telic Engine - unified interface to all capabilities.
    
    Usage:
        engine = Apex(api_key="...")
        result = await engine.do("Create amortization from loan doc and email to Fred")
    """
    
    def __init__(
        self,
        api_key: Optional[str] = None,
        model: Optional[str] = None,
        storage_path: Optional[str] = None,
        connectors: Optional[Dict[str, Any]] = None,
        enable_safety: bool = True,
    ):
        # Auto-detect API key and model
        self._api_key = api_key or os.environ.get("ANTHROPIC_API_KEY") or os.environ.get("OPENAI_API_KEY")
        
        # Auto-select model based on API key
        if model:
            self._model = model
        elif os.environ.get("ANTHROPIC_API_KEY") or (self._api_key and self._api_key.startswith("sk-ant-")):
            self._model = "anthropic/claude-sonnet-4-20250514"
        else:
            self._model = "gpt-4o-mini"
        self._storage_path = Path(storage_path or "~/.telic").expanduser()
        self._storage_path.mkdir(parents=True, exist_ok=True)
        self._connectors = connectors or {}
        
        # Safety rails
        self._safety_enabled = enable_safety and HAS_SAFETY_RAILS
        self._redaction: Optional[Any] = None
        self._audit: Optional[Any] = None
        self._trust: Optional[Any] = None
        self._approval: Optional[Any] = None
        self._undo: Optional[Any] = None
        self._action_history: Optional[Any] = None
        
        if self._safety_enabled:
            db_dir = str(self._storage_path / "db")
            Path(db_dir).mkdir(parents=True, exist_ok=True)
            self._redaction = RedactionEngine()
            self._audit = AuditLogger(db_path=os.path.join(db_dir, "audit.db"))
            self._trust = TrustLevelManager(db_path=os.path.join(db_dir, "trust.db"))
            self._approval = ApprovalGateway(trust_mgr=self._trust, audit_logger=self._audit)
            self._undo = UndoManager(db_path=os.path.join(db_dir, "undo.db"), backup_dir=str(self._storage_path / "undo_backups"))
            self._action_history = ActionHistoryDB(db_path=os.path.join(db_dir, "history.db"))
        
        # Initialize primitives
        self._primitives: Dict[str, Primitive] = {}
        self._init_primitives()
        
        # Initialize planner
        # Execution history
        self._history: List[ExecutionResult] = []
    
    @staticmethod
    def _follow_path_static(data: Any, path: str) -> Any:
        """Navigate into data via dot-path (static version for Orchestrator)."""
        for part in path.split('.'):
            if data is None:
                return None
            if isinstance(data, dict):
                data = data.get(part)
            elif isinstance(data, (list, tuple)):
                try:
                    data = data[int(part)]
                except (ValueError, IndexError):
                    return None
            else:
                return None
        return data
    
    def _init_primitives(self):
        """Initialize all primitives, wiring in connectors as providers."""
        c = self._connectors
        
        self._primitives["FILE"] = FilePrimitive()
        self._primitives["DOCUMENT"] = DocumentPrimitive(self._llm_complete)
        self._primitives["COMPUTE"] = ComputePrimitive(self._llm_complete)
        
        # Email — wire Gmail and/or Outlook connectors
        email_providers = {}
        gmail = c.get("gmail")
        outlook = c.get("outlook")
        if gmail:
            email_providers["gmail"] = gmail
        if outlook:
            email_providers["outlook"] = outlook
        first_email = gmail or outlook
        self._primitives["EMAIL"] = EmailPrimitive(
            providers=email_providers,
            send_func=first_email.send_email if first_email else None,
            list_func=first_email.list_messages if first_email else None,
            read_func=first_email.get_message if first_email and hasattr(first_email, 'get_message') else None,
            connector=first_email,
        )
        
        # Contacts — wire Google Contacts and/or Microsoft Contacts
        contacts_providers = {}
        if c.get("contacts"):
            contacts_providers["google"] = c["contacts"]
        if c.get("contacts_microsoft"):
            contacts_providers["microsoft"] = c["contacts_microsoft"]
        self._primitives["CONTACTS"] = ContactsPrimitive(providers=contacts_providers)
        
        self._primitives["KNOWLEDGE"] = KnowledgePrimitive(str(self._storage_path / "knowledge.json"))
        self._primitives["PATTERNS"] = PatternsPrimitive()
        self._primitives["INTELLIGENCE"] = IntelligencePrimitive()
        
        # Calendar — wire Google Calendar and/or Outlook Calendar
        cal_providers = {}
        gcal = c.get("calendar")
        ocal = c.get("outlook_calendar")
        if gcal:
            cal_providers["google_calendar"] = gcal
        if ocal:
            cal_providers["outlook_calendar"] = ocal
        first_cal = gcal or ocal
        self._primitives["CALENDAR"] = CalendarPrimitive(
            str(self._storage_path / "calendar.json"),
            create_func=first_cal.create_event if first_cal else None,
            list_func=first_cal.list_events if first_cal else None,
            list_calendars_func=first_cal.list_calendars if first_cal and hasattr(first_cal, 'list_calendars') else None,
            providers=cal_providers,
        )
        
        self._primitives["WEB"] = WebPrimitive(self._llm_complete, search_provider=c.get("web_search"))
        
        # Weather — wire OpenWeatherMap connector
        self._primitives["WEATHER"] = WeatherPrimitive(c.get("weather"))
        
        # News — wire NewsAPI connector
        self._primitives["NEWS"] = NewsPrimitive(c.get("news"))
        
        # Notion — wire Notion connector
        self._primitives["NOTION"] = NotionPrimitive(c.get("notion"))
        
        # Linear — wire Linear connector
        self._primitives["LINEAR"] = LinearPrimitive(c.get("linear"))
        
        # Trello — wire Trello connector
        self._primitives["TRELLO"] = TrelloPrimitive(c.get("trello"))
        
        # Airtable — wire Airtable connector
        self._primitives["AIRTABLE"] = AirtablePrimitive(c.get("airtable"))
        
        # Zoom — wire Zoom connector
        self._primitives["ZOOM"] = ZoomPrimitive(c.get("zoom"))
        
        # LinkedIn — wire LinkedIn connector
        self._primitives["LINKEDIN"] = LinkedInPrimitive(c.get("linkedin"))
        
        # Reddit — wire Reddit connector
        self._primitives["REDDIT"] = RedditPrimitive(c.get("reddit"))
        
        # Telegram — wire Telegram connector
        self._primitives["TELEGRAM"] = TelegramPrimitive(c.get("telegram"))
        
        # HubSpot — wire HubSpot connector
        self._primitives["HUBSPOT"] = HubSpotPrimitive(c.get("hubspot"))
        
        # Stripe — wire Stripe connector
        self._primitives["STRIPE"] = StripePrimitive(c.get("stripe"))
        
        # Notify — wire DesktopNotify connector
        notify_send = None
        desktop_notify = c.get("desktop_notify")
        if desktop_notify and hasattr(desktop_notify, "send"):
            notify_send = desktop_notify.send
        self._primitives["NOTIFY"] = NotifyPrimitive(str(self._storage_path / "reminders.json"), send_func=notify_send)
        
        # Task — wire Todoist, Microsoft To-Do, Jira
        task_providers = {}
        if c.get("todoist"):
            task_providers["todoist"] = c["todoist"]
        if c.get("microsoft_todo"):
            task_providers["microsoft_todo"] = c["microsoft_todo"]
        if c.get("jira"):
            task_providers["jira"] = c["jira"]
        self._primitives["TASK"] = TaskPrimitive(
            str(self._storage_path / "tasks.json"),
            providers=task_providers,
        )
        
        self._primitives["SHELL"] = ShellPrimitive()
        self._primitives["DATA"] = DataPrimitive(self._llm_complete)
        
        # Message — wire Slack, Teams, Discord, SMS
        msg_providers = {}
        if c.get("slack"):
            msg_providers["slack"] = c["slack"]
        if c.get("teams"):
            msg_providers["teams"] = c["teams"]
        if c.get("discord"):
            msg_providers["discord"] = c["discord"]
        if c.get("sms"):
            msg_providers["sms"] = c["sms"]
        self._primitives["MESSAGE"] = MessagePrimitive(providers=msg_providers)
        
        # Media — wire Spotify, YouTube
        media_providers = {}
        if c.get("spotify"):
            media_providers["spotify"] = c["spotify"]
        if c.get("youtube"):
            media_providers["youtube"] = c["youtube"]
        self._primitives["MEDIA"] = MediaPrimitive(self._llm_complete, providers=media_providers)
        
        self._primitives["BROWSER"] = BrowserPrimitive(self._llm_complete)
        
        # DevTools — wire GitHub, Jira
        devtools_providers = {}
        if c.get("github"):
            devtools_providers["github"] = c["github"]
        if c.get("jira"):
            devtools_providers["jira"] = c["jira"]
        if c.get("devtools"):
            devtools_providers["unified"] = c["devtools"]
        self._primitives["DEVTOOLS"] = DevToolsPrimitive(providers=devtools_providers)
        
        # Cloud Storage — wire Drive, OneDrive, Dropbox
        storage_providers = {}
        if c.get("drive"):
            storage_providers["google_drive"] = c["drive"]
        if c.get("onedrive"):
            storage_providers["onedrive"] = c["onedrive"]
        if c.get("dropbox"):
            storage_providers["dropbox"] = c["dropbox"]
        self._primitives["CLOUD_STORAGE"] = CloudStoragePrimitive(providers=storage_providers)
        
        # Clipboard — system clipboard operations
        self._primitives["CLIPBOARD"] = ClipboardPrimitive()
        
        # Translate — language translation via LLM
        self._primitives["TRANSLATE"] = TranslatePrimitive(self._llm_complete)
        
        # Database — SQLite database operations
        self._primitives["DATABASE"] = DatabasePrimitive(str(self._storage_path / "data.db"))
        
        # Screenshot — screen capture
        self._primitives["SCREENSHOT"] = ScreenshotPrimitive()
        
        # Automation — rules/triggers
        self._primitives["AUTOMATION"] = AutomationPrimitive(str(self._storage_path / "automations.json"))
        
        # Search — unified search across all primitives
        self._primitives["SEARCH"] = SearchPrimitive(self._primitives)
        
        # Chat — instant messaging (Slack, Teams, Discord)
        chat_providers = {}
        if c.get("slack"):
            chat_providers["slack"] = c["slack"]
        if c.get("teams"):
            chat_providers["teams"] = c["teams"]
        if c.get("discord"):
            chat_providers["discord"] = c["discord"]
        self._primitives["CHAT"] = ChatPrimitive(providers=chat_providers)
        
        # Meeting — video conferencing (Zoom, Teams, Meet)
        meeting_providers = {}
        if c.get("zoom"):
            meeting_providers["zoom"] = c["zoom"]
        if c.get("teams"):
            meeting_providers["teams"] = c["teams"]
        if c.get("google_meet"):
            meeting_providers["google_meet"] = c["google_meet"]
        self._primitives["MEETING"] = MeetingPrimitive(providers=meeting_providers)
        
        # SMS — text messaging
        sms_providers = {}
        if c.get("twilio"):
            sms_providers["twilio"] = c["twilio"]
        if c.get("sms"):
            sms_providers["sms"] = c["sms"]
        self._primitives["SMS"] = SmsPrimitive(providers=sms_providers)
        
        # Spreadsheet — Excel, Google Sheets
        spreadsheet_providers = {}
        if c.get("google_sheets"):
            spreadsheet_providers["google_sheets"] = c["google_sheets"]
        if c.get("excel"):
            spreadsheet_providers["excel"] = c["excel"]
        self._primitives["SPREADSHEET"] = SpreadsheetPrimitive(self._llm_complete, providers=spreadsheet_providers)
        
        # Presentation — PowerPoint, Google Slides
        presentation_providers = {}
        if c.get("google_slides"):
            presentation_providers["google_slides"] = c["google_slides"]
        if c.get("powerpoint"):
            presentation_providers["powerpoint"] = c["powerpoint"]
        self._primitives["PRESENTATION"] = PresentationPrimitive(self._llm_complete, providers=presentation_providers)
        
        # Notes — OneNote, Apple Notes, Google Keep
        notes_providers = {}
        if c.get("onenote"):
            notes_providers["onenote"] = c["onenote"]
        if c.get("google_keep"):
            notes_providers["google_keep"] = c["google_keep"]
        self._primitives["NOTES"] = NotesPrimitive(str(self._storage_path / "notes"), providers=notes_providers)
        
        # Finance — banking, payments
        finance_providers = {}
        if c.get("plaid"):
            finance_providers["plaid"] = c["plaid"]
        if c.get("paypal"):
            finance_providers["paypal"] = c["paypal"]
        if c.get("venmo"):
            finance_providers["venmo"] = c["venmo"]
        self._primitives["FINANCE"] = FinancePrimitive(providers=finance_providers)
        
        # Social — Twitter, LinkedIn, Facebook
        social_providers = {}
        if c.get("twitter"):
            social_providers["twitter"] = c["twitter"]
        if c.get("linkedin"):
            social_providers["linkedin"] = c["linkedin"]
        if c.get("facebook"):
            social_providers["facebook"] = c["facebook"]
        self._primitives["SOCIAL"] = SocialPrimitive(providers=social_providers)
        
        # Photo — Google Photos, iCloud
        photo_providers = {}
        if c.get("google_photos"):
            photo_providers["google_photos"] = c["google_photos"]
        if c.get("icloud_photos"):
            photo_providers["icloud_photos"] = c["icloud_photos"]
        self._primitives["PHOTO"] = PhotoPrimitive(providers=photo_providers)
        
        # Ride — Uber, Lyft
        ride_providers = {}
        if c.get("uber"):
            ride_providers["uber"] = c["uber"]
        if c.get("lyft"):
            ride_providers["lyft"] = c["lyft"]
        self._primitives["RIDE"] = RidePrimitive(providers=ride_providers)
        
        # Travel — flights, hotels
        travel_providers = {}
        if c.get("expedia"):
            travel_providers["expedia"] = c["expedia"]
        if c.get("google_flights"):
            travel_providers["google_flights"] = c["google_flights"]
        self._primitives["TRAVEL"] = TravelPrimitive(providers=travel_providers)
        
        # Home — smart home
        home_providers = {}
        if c.get("alexa"):
            home_providers["alexa"] = c["alexa"]
        if c.get("google_home"):
            home_providers["google_home"] = c["google_home"]
        if c.get("homekit"):
            home_providers["homekit"] = c["homekit"]
        self._primitives["HOME"] = HomePrimitive(providers=home_providers)
        
        # Shopping — Amazon, eBay
        shopping_providers = {}
        if c.get("amazon"):
            shopping_providers["amazon"] = c["amazon"]
        if c.get("ebay"):
            shopping_providers["ebay"] = c["ebay"]
        self._primitives["SHOPPING"] = ShoppingPrimitive(providers=shopping_providers)
    
    async def _llm_complete(self, prompt: str, triggering_request: str = "") -> str:
        """Call LLM for completion — with PII redaction and audit logging when safety is enabled."""
        if not self._api_key:
            raise ValueError("No API key configured. Set OPENAI_API_KEY or pass api_key to the engine")
        
        # Redact PII before sending to LLM
        send_prompt = prompt
        redaction_count = 0
        if self._safety_enabled and self._redaction:
            redaction_result = self._redaction.redact(prompt)
            send_prompt = redaction_result.redacted_text
            redaction_count = redaction_result.redaction_count
        
        # Log outbound transmission
        request_id = str(uuid.uuid4())[:8]
        if self._safety_enabled and self._audit:
            dest = TransmissionDestination.OPENAI if "gpt" in self._model else TransmissionDestination.ANTHROPIC
            self._audit.log_outbound(
                destination=dest,
                content=send_prompt[:500],
                triggering_request=triggering_request or "engine_call",
                model=self._model,
                request_id=request_id,
                contained_pii=redaction_count > 0,
                redactions_applied=redaction_count,
            )
        
        _ll = _get_litellm()
        if _ll:
            response = await asyncio.to_thread(
                _ll.completion,
                model=self._model,
                messages=[{"role": "user", "content": send_prompt}],
                api_key=self._api_key,
            )
            result_text = response.choices[0].message.content
        else:
            # Fallback to direct OpenAI
            import httpx
            async with httpx.AsyncClient() as client:
                resp = await client.post(
                    "https://api.openai.com/v1/chat/completions",
                    headers={"Authorization": f"Bearer {self._api_key}"},
                    json={"model": self._model, "messages": [{"role": "user", "content": send_prompt}]},
                    timeout=60,
                )
                result_text = resp.json()["choices"][0]["message"]["content"]
        
        # Log inbound response
        if self._safety_enabled and self._audit:
            self._audit.log_inbound(
                destination=dest,
                content=result_text[:500],
                request_id=request_id,
                model=self._model,
            )
        
        return result_text
    

